#!/usr/bin/env python3
# =============================================================================
# DISCLAIMER: This code is provided on a best-effort basis and is not in any
# way officially supported or sanctioned by Cohesity. The code is intentionally
# kept simple to retain value as example code. The code in this repository is
# provided as-is and the author accepts no liability for damages resulting
# from its use.
# =============================================================================
"""
health_check_report.py  v1.64

Multi-cluster Cohesity health check — 25-tab Excel workbook + Word document + comprehensive PowerPoint deck.
Designed for enterprise customer business reviews (EBRs) and SE trusted-advisor
engagements.  Gathers live data from Cohesity Helios (or directly from a single
cluster) and produces:

  Excel sheets
  ────────────
  1  Executive Summary       — per-cluster health score & grade
  2  Infrastructure          — versions, nodes, disks, config
  3  Node Hardware           — per-node model, serial, EOL status
  4  Disk Health             — per-disk status, SSD wear %, encryption
  5  Protection Health       — success rates, SLA, RPO gaps, DataLock per group
  6  Storage & Capacity      — utilization, data reduction, runway
  7  Policy Audit            — retention, replication, archival, DataLock per policy
  8  Policy → Groups         — every group with the policy that governs it
  9  Alerts                  — open critical/warning with age
  10 Security                — expanded checklist: encryption, FIPS, audit log, MFA,
                               NTP auth, tunnel, SSO, cert expiry, vault, FK, DataLock,
                               KMS sub-section
  11 Agent Health            — per-host agent version, status, upgradability
  12 Source Coverage         — registered sources with protected/unprotected counts
  13 Unprotected Objects     — named unprotected objects per source
  14 Recovery Audit          — recovery testing history, success rate, days since
  15 Replication/Archive     — targets, last transfer, FortKnox
  16 FortKnox Data Transfer  — per-group transfer to external targets (30d full / 1d quick)
  17 Data Services           — NAS views, quota utilization
  18 Data Exposure           — NAS share risk (SMB discovery, NFS open mount, S3)
  19 Coverage Gaps           — groups with no recent successful backup
  20 User Security           — user accounts, MFA, locked status, roles, last login,
                               custom role definitions sub-section
  21 Trends (30d)            — daily success rate + storage growth charts
  22 Recommendations         — prioritized action list
  23 Workload Risk Heatmap   — all protection groups scored by recovery risk, worst-first
  24 Audit Log               — 30-day configuration change summary and detail log
  25 DataLock Verification   — snapshot-level WORM confirmation for DataLock groups
  + Guide tab (first)        — sheet descriptions, scoring methodology, color legend

  Word document
  ─────────────
  Cover, Executive Summary, Environment, Protection, Storage,
  Security, Agent & Source Coverage, User Security, Recommendations,
  Methodology appendix.

Usage
─────
  # Helios — API key (key saved to OS keychain after first prompt)
  python3 health_check_report.py --apikey <key> [options]

  # Direct cluster — username/password + optional MFA (bypasses Helios)
  python3 health_check_report.py --cluster-host 10.1.2.3 --username admin [--domain LOCAL]

  NOTE: Helios does NOT expose a scriptable username/password login endpoint.
  Its web login uses OIDC/OAuth2 which cannot be driven from a script.
  Generate a Helios API key at:
    helios.cohesity.com → Settings → Access Management → API Keys

Options
───────
  --apikey KEY        Helios API key (stored in keychain if omitted)
  --cluster-host HOST Bypass Helios; connect directly to this cluster hostname/IP
  --username USER     Cluster username (required with --cluster-host)
  --password PASS     Password (prompted securely if omitted; saved to keychain)
  --domain DOMAIN     Auth domain — LOCAL or AD name (default: LOCAL)
  --mfa-code CODE     TOTP/OTP code for MFA-enabled cluster accounts
  --cluster NAME      (Helios mode) Limit to one cluster by partial name match
  --clear-credentials Remove stored credentials from keychain and exit
  --days N            Lookback window for alerts/runs (default 30)
  --customer NAME     Customer name printed on report cover
  --output PATH       Base path for output files (no extension)
  --quick             Skip per-group run history; use last-run only (faster)
  --excel-only        Skip Word document generation
  --word-only         Skip Excel generation
  --debug             Print raw API payloads for troubleshooting

Requirements
────────────
  Required:  pip install requests openpyxl python-docx python-pptx
  Optional:  pip install matplotlib keyring

  The script checks all prerequisites at startup and reports missing
  packages with exact install commands.

Version history
───────────────
  1.66 (2026-04-22) — fix(health_check): prereq security hardening. python-pptx
                     promoted to required (PPT is always generated). keyring
                     description updated to explain the security rationale
                     (prevents API key exposure in shell history / process
                     lists). Security note printed at startup when keyring is
                     absent so users understand the credential exposure risk.
                     Single pip install command now covers all missing packages
                     (required + optional) in one output block.
  1.65 (2026-04-21) — feat(health_check): AD/LDAP health, full cert inventory,
                     dedup/IO performance trends, audit-log 403 surfacing. New
                     fetch functions: _active_directory, _ldap_providers,
                     _all_certs, _dedup_trend, _io_stats. New sheets: AD /
                     Identity Health, Certificate Inventory, Capacity & Perf
                     Trends. _fetch_audit_log now returns "PERMISSION_DENIED"
                     sentinel on HTTP 403 so Audit Log sheet shows actionable
                     message instead of silently empty. Recommendations engine
                     extended: cert expiry CRITICAL/HIGH, AD/LDAP disconnect
                     HIGH. Guide tab updated for all new sheets.
  1.64 (2026-04-21) — feat(health_check): recovery testing, KMS, custom roles,
                     DataLock snapshots. Four new data collections (recoveries,
                     kmsConfig, roles, per-group snapshot DataLock check) added
                     to collect_cluster(). New sheets: Recovery Audit (recovery
                     task history, success rate, days-since-last-recovery) and
                     DataLock Verification (snapshot-level WORM confirmation).
                     Enhanced Security sheet with KMS sub-section (KMS type,
                     server, risk). Enhanced User Security sheet with Custom Role
                     Definitions sub-section (privilege counts, modify flags,
                     risk level). Recommendations engine extended: no-recovery-
                     test HIGH, internal-KMS MEDIUM, high-risk-custom-role HIGH,
                     DataLock-not-applied CRITICAL. Guide tab updated for all
                     new/modified sheets.
  1.63 (2026-04-21) — feat(health_check): security & recoverability analytics
                     (no new API calls). Nine new analytical features: stale
                     privileged account risk flags in User Security sheet;
                     NAS share exposure sheet (SMB discovery, NFS open mounts,
                     S3, quota); security & anomaly alert sub-section in
                     Security sheet; node software version skew column in Node
                     Hardware sheet; cluster fault tolerance margin column in
                     Infrastructure sheet; FortKnox transfer recency columns in
                     FortKnox Data Transfer sheet; replication lag RPO column in
                     Replication & Archive sheet; named unprotected objects
                     sheet; average backup duration trend column + chart in
                     Trends sheet. Recommendations engine extended with stale
                     admin, NAS exposure, node skew, and fault tolerance checks.
  1.62 (2026-04-17) — fix(security): input validation and credential hygiene.
                     --days bounded 1–3650 (previously unchecked). --cluster-host
                     validated as hostname or IPv4 before first network call.
                     --output rejected if path contains '..' traversal sequences.
                     --ca-bundle existence check added.
                     --password / --mfa-code print a process-list exposure warning
                     when supplied on the command line (still accepted — callers
                     directed toward interactive prompt or system keychain).
                     Raw API response bodies removed from authentication error
                     messages (r.text[:120] / [:200] / [:400] → HTTP status only)
                     to prevent token or credential fragments leaking into terminal
                     history and CI logs. Same changes applied to utils/cohesity_auth.py.
  1.61 (2026-04-17) — feat(pptx): Topology diagram — legend removed; FortKnox
                     column given wider separation gap (T_FK_GAP=1.60" vs
                     T_COL_GAP=0.80"); all non-source columns (R, A, F) now
                     vertically centred relative to source cluster column height
                     via _col_y0(); diagram uses full slide height (DIAG_BOT
                     pushed to H-0.30" from H-0.60").
  1.60 (2026-04-17) — fix: Excel column widths now truly content-driven —
                     min_width lowered from 12→6, max_width raised 45→60,
                     wrap_text removed so all content fits without line breaks.
                     Word tables get proportional column widths via new
                     _word_fit_widths() measured after all rows are added
                     (10 tables). PPT _table() auto-computes column widths from
                     content and centres the table horizontally; passed col_w
                     arrays are overridden.
  1.59 (2026-04-17) — feat: Excel tabs zoom 130%+autofit+left-align; Word tables
                     autofit. PPT: copyright 2026, page numbers on all data
                     slides, table headers 10pt left-aligned, contents slide
                     bars dynamically sized to widest item, Protection Summary
                     gets inline scoring methodology, topology legend is now a
                     horizontal row (active columns only) with bottom margin
                     increased to 0.60" for clear visibility.
  1.58 (2026-04-17) — feat(pptx): Topology diagram — wider node spacing for
                     elegant connector-line visibility. T_NODE_W 2.50→2.30",
                     T_NODE_H 0.80→0.85", T_V_GAP 0.18→0.32" (near-double
                     vertical breathing room), T_COL_GAP 0.40→0.80" (double
                     connector gap). Connector line weight 0.025→0.030". Max
                     column height budget raised to 4.50" before auto-scale.
  1.57 (2026-04-17) — fix(pptx): Table top-alignment on all data slides.
                     Tables now anchor to y=0.90" (0.13" below the 0.77"
                     green header bar) instead of being vertically centred.
                     Mixed slides (Scorecard, Ransomware, Security, Risk)
                     also top-align their table; methodology section follows
                     immediately below. Removes the excessive bottom whitespace
                     at the cost of centering.
  1.56 (2026-04-17) — feat(pptx): PowerPoint layout and visual polish.
                     KPI tiles: solid #D9D9D9 fill with accent colour side bar
                     and outer drop shadow (30 % opacity) — replaces gradient.
                     All data slides: tables are now vertically centred in the
                     available content area (_tbl_y helper) instead of being
                     pinned to y=0.82 with large bottom whitespace.
                     Mixed slides (Health Scorecard, Ransomware, Security
                     Posture, Workload Risk Heatmap): entire table + methodology
                     block centred as a unit (_layout_mixed helper).
                     Environment Topology: diagram vertically centred; auto-
                     scales node height/gap when many nodes would overflow;
                     HDR_Y and LEG_Y derived dynamically; all topology shapes
                     are wrapped in a single PowerPoint group ("Topology Group")
                     for easy user rearrangement.
  1.55 (2026-04-17) — feat(pptx): Centre environment topology diagram horizontally.
                     Column X positions are now computed dynamically based on
                     which column types (Source Clusters, Replication, Archival,
                     FortKnox) are actually present, so the diagram is always
                     centred on the slide regardless of how many columns exist.
  1.54 (2026-04-17) — feat(pptx): PowerPoint visual improvements — green header bar
                     height increased to 0.77" on all slides; content row Y
                     shifted down to 0.82" to clear taller bar. Contents slide
                     serial numbers now display without leading zero (1 not 01).
                     Environment Snapshot KPI tiles use a two-stop linear gradient
                     fill (light tint → accent color, 120°) with white label and
                     value text for a modern look. Added "Thank You" end slide
                     matching cover/section-divider style (green block, white rule,
                     footer).
  1.53 (2026-04-17) — fix(security): TLS certificate validation and Excel formula
                     injection hardening.
                     TLS: remove all hardcoded verify=False; add --ca-bundle PATH
                     and --insecure flags. Default is now verify=True (full
                     validation). --insecure prints a loud startup warning and
                     is the only path that suppresses urllib3 warnings.
                     session.verify is set once and inherited by all _get() calls.
                     get_auth_token() and get_helios_clusters() accept verify=
                     param threaded from main(). Same changes applied to
                     utils/cohesity_auth.py.
                     Formula injection: add _safe_cell() sanitizer that prefixes
                     strings starting with =, +, -, @, tab, CR with a single
                     quote, preventing openpyxl from writing them as formula
                     cells. Applied to _reg() and all 25 inline ws.cell writes
                     that write API-sourced data.
  1.52 (2026-04-16) — feat: Inline scoring methodology — remove 4 standalone
                     methodology slides; embed compact scoring tables directly in
                     the bottom whitespace of Health Scorecard, Ransomware Readiness,
                     Security Posture Overview, and Workload Risk Heatmap slides.
                     feat: Table body font increased from 8.5pt to 10pt across all
                     slides (header row remains 8.5pt; inline methodology tables
                     use 8pt to stay compact in the smaller space).
  1.51 (2026-04-16) — fix: Storage Capacity & Growth PPT slide now reads from the
                     correct API fields (cd["info"]["stats"]["usagePerfStats"] +
                     cd.get("capacity_runway") + cd["domains"]) instead of the
                     non-existent storage_domains key — slide was previously empty.
                     feat: DataLock red/amber/green coloring added to Protection
                     Summary and Ransomware Readiness PPT slides (was uncolored).
                     feat: Two Scoring Methodology Reference slides appended at end
                     of PowerPoint deck covering all four scoring models (Overall
                     Health, Security, Ransomware Readiness, Workload Risk) with
                     RAG colour thresholds and grade bands.
  1.50 (2026-04-15) — feat: Comprehensive PowerPoint deck (write_pptx). Replaces
                     the single-slide topology .pptx and eliminates draw.io output.
                     New deck spans ~29 slides across 4 sections: Cover, Executive
                     Summary (Snapshot, Scorecard, Capacity, Protection, Ransomware
                     Readiness, Top At-Risk Workloads, Priority Action Items),
                     Environment Topology (full connection diagram with legend),
                     Security Deep-Dive (Posture, User Security, Audit & Governance,
                     Security Recommendations), Backup Engineering (Software &
                     Hardware Lifecycle, Coverage Gaps, Agent Health, Source
                     Coverage, Workload Risk Heatmap, Engineering Recommendations).
                     All data slides use RAG colour-coding. draw.io generation
                     removed entirely. python-pptx remains optional — skipped
                     gracefully if not installed.
  1.49 (2026-04-15) — perf: Parallel API data collection and cached derived state.
                     collect_cluster() now runs all independent API calls in two
                     parallel ThreadPoolExecutor waves (wave-1: 17 independent
                     endpoints; wave-2: disks + FortKnox that depend on wave-1).
                     Per-group run fetches (non-quick mode) also parallelised.
                     Expected wall-clock speedup: ~40–60% per cluster.
                     cd["fk_status"], cd["policy_by_id"], cd["audit_enabled"]
                     pre-computed once in collect_cluster() and reused by all
                     scoring, recommendation, and sheet-generation functions
                     (previously _fk_status() was called 6×, policy_by_id dict
                     rebuilt 4×, audit config parsed 2×).
                     cd["end_usecs"] used in place of _now_usecs() throughout
                     scoring and sheet functions for a consistent reference point.
                     Dual agent loops and dual disk comprehensions in
                     _build_recommendations() merged to single passes each.
                     isinstance(x, dict) filtering for agents and sources moved to
                     collection time, removing per-item guards from hot loops.
                     _FONT_NORMAL / _FONT_BOLD openpyxl singletons defined at
                     module level; 24 per-cell _font() allocations eliminated.
                     No output changes — all modifications are internal.
  1.48 (2026-04-13) — feat: Topology diagram overhaul. Fixed missing FK vault
                     connections in _topology_data() — vaults registered via the
                     cluster vaults list but absent from policies now correctly
                     emit edges. draw.io diagram enhanced with arrow labels
                     (Replication / Archive / FortKnox / Indelible), cluster
                     node labels include software version and group/source counts,
                     and a legend block added at bottom-right. New optional
                     PowerPoint (.pptx) output via python-pptx — 16:9 widescreen
                     slide with colored rounded-rect nodes, STRAIGHT connectors,
                     arrow labels, and a legend; gracefully skipped if python-pptx
                     is not installed. python-pptx added to the optional prereq
                     check table.
  1.47 (2026-04-13) — feat: Guide tab added as the first tab in every generated
                     workbook. Covers all 22 tabs with descriptions, the overall
                     health scoring model (weights + grade table), security score
                     deductions, ransomware readiness scoring breakdown, workload
                     risk score factors, recommendation priority levels, and a
                     full color/RAG legend. Static content — no API calls required.
  1.46 (2026-04-12) — feat: Ransomware Readiness Score (0-100) computed per
                     cluster from DataLock coverage, FortKnox vault activity,
                     recovery window depth, quorum, and admin MFA. Displayed
                     in Executive Summary, Security sheet, and Word Security
                     Posture table (now 11 columns) with a pre-table environment
                     average callout.
                     feat: Predictive Capacity Runway — linear regression on
                     30-day daily usage from /v1/public/statistics/timeSeriesStats
                     projects how many days until 80% utilization. Shown in
                     Storage sheet (cols 13-15), Executive Summary (col 14),
                     and Word Storage section with a runway forecast paragraph.
                     feat: Per-Workload Risk Heatmap (sheet 20) — all protection
                     groups scored 0-100 on composite risk (last-run status,
                     SLA, RPO gap, DataLock); sorted worst-first with full-row
                     RAG coloring. Top 5 Critical/High groups also appear in the
                     Word Executive Summary as a "Top At-Risk Workloads" table.
                     feat: Audit Log / Change Summary (sheet 21) — last 30 days
                     of configuration changes from /v1/public/auditLog, grouped
                     by category (policy, group, user, vault/security, cluster
                     config) with high-risk event highlighting. Word Security
                     section adds a "Governance & Change Activity" summary with
                     high-risk event warnings.
                     Sheet count: 19 → 21.
  1.45 (2026-04-10) — fix: get_auth_token() v1 failure is no longer silently
                     swallowed. Now prints "[!] v1 endpoint failed (HTTP NNN)
                     — trying v2 ..." so the user can see both failures.
                     Final error message includes v1 response details and
                     actionable troubleshooting hints: run --clear-credentials
                     to reset a stale stored password, use --domain for AD
                     accounts, or check account status in the Cohesity UI.
                     cohesity_auth.py v1.6.
  1.44 (2026-04-10) — feat: Prereq check now always reports ALL five packages
                     (requests, openpyxl, python-docx, matplotlib, keyring)
                     in a full status table — required vs optional, installed
                     vs NOT FOUND. Previously matplotlib and keyring were only
                     reported as small NOTE lines when missing; they are now
                     first-class entries in the startup check every run.
  1.43 (2026-04-10) — fix: Standalone / single-file distribution on Windows.
                     import urllib3 failed with ModuleNotFoundError when
                     requests was not installed (urllib3 is a dep of requests,
                     not a top-level package). Now catches ImportError at the
                     requests import itself and prints a clear install command.
                     urllib3 warning suppression falls back to
                     requests.packages.urllib3 when urllib3 is not importable
                     as a standalone package.
                     fix: cohesity_auth and formatters imports wrapped in
                     try/except with full inline fallbacks — script now runs
                     as a single downloaded .py file without needing the repo's
                     utils/ directory alongside it.
  1.42 (2026-04-10) — fix: Direct cluster auth — get_auth_token() now tries
                     v1 endpoint first (POST /irisservices/api/v1/public/
                     accessTokens) then falls back to v2 (POST /v2/users/
                     sessions). The v1 endpoint has the widest compatibility
                     across Cohesity versions and account types; v2 returned
                     400 "KValidationError Access denied" on some clusters.
                     fix: collect_cluster() now stores "mode" in the cd dict
                     so downstream code (quorum N/A, Word security table)
                     correctly detects direct mode.
                     cohesity_auth.py v1.5.
  1.41 (2026-04-10) — feat: Startup prerequisite check. main() now validates
                     all required packages (requests, openpyxl, python-docx)
                     before execution and prints clear error messages with
                     exact pip install commands for any missing package.
                     Optional packages (matplotlib, keyring) produce a NOTE
                     at startup but do not block. README.md fully rewritten
                     with complete prerequisite table, all CLI options
                     (including direct-cluster auth flags), and updated
                     descriptions for Word topology (native shapes) and
                     security table (full RAG coloring).
  1.40 (2026-04-10) — feat: Word doc Security Posture table (section 5) now
                     color-codes ALL columns: red for "No"/missing, amber for
                     "Idle"/"Partial", green for "Yes"/"Active"/"Full"; score
                     column uses RAG (red <50, amber 50-74, green 75+).
                     feat: Environment Topology diagram replaced with native
                     Word DrawingML shapes — editable rounded rectangles
                     (#96C73D green) for source clusters, cloud shapes for
                     FortKnox (#1A5C3A), cylinders for archival, rounded rects
                     for replication (#0062B1). Connector arrows drawn between
                     each cluster and its targets. matplotlib PNG kept as
                     fallback only. Dark blue header bar removed from PNG.
                     Cluster color updated from #00B388 (teal) to #96C73D
                     (Cohesity green) globally. draw.io fonts 11→13 (clusters),
                     10→12 (targets); FK shape changed to cloud; box sizes
                     increased for readability.
  1.39 (2026-04-10) — fix: Executive Summary "30d Success %" now matches the
                     Trends tab. _success_stats() rebuilt to use the same data
                     source priority as _sheet_trends() (v1_runs → group_runs
                     fallback) and now returns the average of per-day success %
                     values (= average of Trends col H) instead of a flat run-
                     count ratio across the entire window.
  1.38 (2026-04-10) — fix: Environment Topology table (Word) FortKnox Vaults
                     column now uses _fk_status() (active/idle/none) instead of
                     a narrow vault-type check. FK vault names collected via the
                     same comprehensive _is_fk logic used elsewhere (krpaas,
                     kfortknox, kfort_knox, krpaasarchival, "fortknox" in name).
                     Configured-but-idle vaults now show vault name +
                     "(Configured — not sending data)" in amber. Active vaults
                     show vault name(s) on green background. "No" only when FK
                     is genuinely not configured.
  1.37 (2026-04-10) — fix: Quorum detection revised — Quorum is a Helios control-
                     plane feature; any enabled quorum group now marks ALL Helios-
                     connected clusters as quorum-enabled (per-cluster ID matching
                     was unreliable when clusterIdentifiers is unpopulated).
                     fix: DataLock label "None" → "No"; removed broken policy-only
                     fallback from _cluster_datalock() — when no group data exists
                     return "No" rather than claiming full/partial coverage.
                     feat: Word doc Security Posture table expanded from 8→10
                     columns: "Indelible (FK/DL)" split into "Indelible (FK)" and
                     "DataLock (WORM)"; "Admins w/o MFA" column added with red
                     cell highlight when count > 0. Appendix scoring text updated.
  1.36 (2026-04-10) — feat: Environment Topology diagram. New _topology_data()
                     helper extracts cluster/replication/archival/FortKnox graph
                     from all_data. _generate_topology_drawio() writes an editable
                     draw.io XML file (.drawio) alongside the Excel/Word output —
                     open with diagrams.net or the draw.io desktop app. Nodes:
                     clusters (Cohesity teal), replication targets (blue cluster
                     boxes), archival vaults (amber cylinders), FortKnox/RPaaS
                     (dark-green, dashed arrows). _render_topology_png() generates
                     a matplotlib PNG preview (requires matplotlib) embedded in the
                     Word doc after the topology table; caption references the
                     .drawio file. Cohesity brand colors used throughout (#00B388
                     teal, #0062B1 replication, #E07A00 archival, #1A5C3A FK).
  1.35 (2026-04-10) — fix: Quorum detection now fetches Helios-level quorum groups
                     via GET /v2/mcm/quorum/groups instead of probing per-cluster
                     info dict (which never had quorum fields). quorum_enabled is
                     resolved once for all clusters and stored in cd.
                     feat: FortKnox implicit DataLock — policies archiving to
                     FortKnox/RPaaS now show "FortKnox (Indelible)" (green) in
                     Policy Audit DataLock column and are not flagged as "No DataLock".
                     feat: Admin MFA elevated from MEDIUM → HIGH recommendation;
                     missing admin MFA now deducts 10 pts from security score.
  1.34 (2026-04-10) — feat: replace FIPS with split Cluster/SD Encryption columns
                     across Excel Security sheet, Infrastructure sheet, and Word doc.
                     DataLock coverage now per-group (partial compliance shown as
                     "X/Y groups"). Security score redesigned: 6 dimensions totalling
                     100 pts. Add Quorum column to Security sheet (N/A for direct
                     mode). Rename "Immutability"→"Indelible" throughout. Agent
                     upgrade recommendation uses softer "MAY cause issues" language
                     with compatibility doc URL. Word doc gains Environment Topology
                     table (clusters → replication → archival → FortKnox).
  1.33 (2026-04-09) — fix: Agent Health — hostname vs IP routing by content:
                     new _split_host_ip() helper inspects each value and routes
                     it to Host or IP Address column based on whether it matches
                     an IPv4/IPv6 pattern (not field name). Fixes empty Host
                     column when hostIp contains FQDNs.
                     fix: Infrastructure Domain (col K) — domainNames[0] returns
                     the full cluster FQDN (e.g. "mycluster.example.com"); strip
                     the cluster hostname prefix so only the domain is shown.
  1.32 (2026-04-09) — feat: DataLock (WORM) column in Protection Health sheet
                     (col 18 per group, green=Compliance/yellow=Administrative/
                     red=None); new "FortKnox Data Transfer" sheet (sheet 14)
                     showing per-protection-group transfer to every external
                     target: logical TB, physical TB, storage consumed TB,
                     snapshots. Full mode covers the full lookback window
                     (args.days); quick mode fetches a separate 1-day window.
                     Sheet count: 18 → 19.
  1.31 (2026-04-09) — polish: formatting improvements — Word doc headings and
                     cover title standardized to #70AD47 (was #00B388); Excel
                     header rows gain a thin black border on all sides and row
                     height 32; auto_fit_columns() handles multiline cell values
                     (measures longest line) with tuned defaults (min 12, max 45);
                     Trends charts moved to the right of data (col L) to prevent
                     overlap, height increased to 16, legend placed at bottom,
                     chart title rendered at 14 pt bold with overlay disabled.
                     formatters.py: auto_fit_columns defaults min=12, max=45.
  1.30 (2026-04-09) — feat: DataLock (WORM) tracking across Policy Audit,
                     Security, and Recommendations. New helper functions
                     _policy_datalock_str() and _cluster_datalock() read
                     backupPolicy.regular.retention.dataLockConfig from the
                     v2 policies API (already fetched — no new API call).
                     Policy Audit gains a "DataLock (WORM)" column (col 15,
                     green=Compliance, yellow=Administrative, red=None) with
                     "No DataLock" added to the Gaps cell when absent.
                     Security sheet gains a "DataLock (WORM)" column (col 16)
                     showing cluster-level rollup: Compliance / Administrative
                     / Mixed / None, with "No DataLock" added to the Gaps
                     column when absent. Recommendations: MEDIUM if no policies
                     have DataLock enabled; LOW if Administrative mode only
                     (urging upgrade to Compliance mode).
  1.29 (2026-04-09) — fix: removed Helios username/password auth mode —
                     Helios /irisservices/api/v1/mcm/login requires an OIDC
                     token in the Authorization header, not credentials in
                     the POST body. Helios access requires an API key
                     (generate at helios.cohesity.com → Settings → Access
                     Management → API Keys). --username/--password/--mfa-code
                     are now only supported with --cluster-host (direct
                     cluster mode). cohesity_auth.py bumped to v1.4.
  1.28 (2026-04-09) — feat: direct cluster mode (--cluster-host) bypasses
                     Helios; authenticates via POST /v2/users/sessions →
                     Bearer token. Cluster name auto-detected from cluster's
                     own /public/cluster endpoint. Added --clear-credentials
                     flag and --mfa-code for cluster MFA. _get() now reads
                     base URL from session.base_url. cohesity_auth.py v1.2.
  1.27 (2026-04-09) — fix: Disk Health sheet still empty after v1.26 — GET
                     /v2/disks is not a valid Helios-proxied endpoint.
                     Rewrote _disks() to use GET /v2/disks/local?nodeId={id}
                     per node, matching the approach used by the Cohesity ELF
                     inventory tool (cohesityInv.ps1, apiType='ClusterV2',
                     path='/disks/local?nodeId={0}'). Updated field fallbacks
                     throughout to include ELF field names: ssdUsedPercentage
                     (was ssdWearLevelPct), capacityInBytes (was capacityBytes),
                     encryptionStatus (was encryptionEnabled).
  1.26 (2026-04-09) — fix: Disk Health sheet empty — switched _disks() from the
                     private /v1/disks/local?nodeId=X endpoint (always 400 via
                     Helios proxy) to the supported GET /v2/disks endpoint;
                     private endpoint kept as fallback. Added isinstance guard
                     on usageStats sub-object in _sheet_disk_health().
  1.25 (2026-04-09) — chore: standardize on US English throughout — utilization,
                     prioritized, normalized, color-coded, behavior, organization,
                     authorized, summarizes, maximize, minimize, penalized, labeled.
                     Affects docstring, sheet titles, Word doc narrative, comments.
  1.24 (2026-04-09) — fix: Source Coverage sheet empty on older clusters —
                     /v2/data-protect/sources returns 404 on pre-7.x firmware.
                     _sources() now falls back to v1
                     /protectionSources/registrationInfo (rootNodes[].stats.
                     protectedCount / unprotectedCount) and normalizes the
                     result into the same dict shape used by the sheet and
                     recommendations engine. Also fixed isinstance guards on
                     stats / objectCount sub-objects in _sheet_source_coverage.
  1.23 (2026-04-09) — fix: suppress debug body output for best-effort endpoints
                     (_disks, _sources, _tenants) that return 400/403/404 on
                     clusters where the endpoint is unavailable or restricted.
                     Added `silent` param to _get(); these three callers pass
                     silent=True so --debug output is not cluttered with
                     expected error bodies.
  1.22 (2026-04-09) — fix: replace all unsafe `(x or {}).get()` patterns with
                     isinstance(x, dict) guards; affected fields: healthStatus
                     (could be a string in older cluster versions), auditLogConfig,
                     clusterAuditConfig, encryptionConfig, ntpSettings,
                     remoteSupportConfig, mfaConfig, sourceInfo (could be string
                     in some source types), stats (source coverage). Added
                     isinstance(d, dict) guards to disk comprehensions in
                     _build_recommendations(). Fixed _ag_st() helper in Word
                     doc section to handle str/dict/None healthStatus uniformly.
  1.21 (2026-04-08) — feat: Tier 1/2/3 ELF inventory additions. New data
                     collection: _certificates(), _agents(), _disks(),
                     _sources(), _users(), _idps(), _tenants(). New Excel
                     sheets: Disk Health (per-disk status/SSD wear/encryption),
                     Agent Health (agent version/status/upgradability), Source
                     Coverage (registered sources with protected/unprotected
                     ratio), User Security (MFA/locked/roles/last-login).
                     Security sheet expanded: audit log, NTP auth, remote
                     tunnel, cluster MFA, SSO/IDP, TLS cert expiry. New
                     recommendations: cert expiry CRITICAL/HIGH, unhealthy
                     agents HIGH, failed/missing disks CRITICAL, SSD wear HIGH,
                     unprotected sources HIGH, users without MFA MEDIUM, no SSO
                     LOW, audit log not enabled MEDIUM. Sheet count: 14 → 18.
  1.20 (2026-04-08) — fix: 6.8.1 LTS had EOS November 15 2024 (OOS) — was
                     incorrectly matched by "6.8" prefix as in-support; split
                     table so only 6.8.2+ maps to "In Support — LTS Jun 2026".
                     Hardware EOL table expanded: added C2100, C2200, C2510,
                     C4500 model variants confirmed in Cohesity EOL doc (Oct 2025).
  1.19 (2026-04-08) — feat: software and hardware lifecycle tracking. Added
                     SOFTWARE_EOS and HARDWARE_EOL reference tables from
                     Cohesity Product Life Cycle Policy (Oct 2025). Infrastructure
                     sheet gets 3 new columns: SW Lifecycle Status (color-coded
                     red/orange/green), SW EOS Date, Days to EOS. New "Node Hardware"
                     sheet (sheet 3): per-node model, serial, node type, capacity,
                     storage tiers, HW EOL date, HW EOL status color-coded.
                     Recommendations engine updated: CRITICAL for out-of-support
                     software or past-EOL hardware, HIGH/MEDIUM for upcoming dates.
                     Word document gains a Software & Hardware Lifecycle section.
                     Sheet count: 13 → 14.
  1.18 (2026-04-08) — fix: Trends tab still truncated after v1.17 — previous
                     pagination used `len(page) < PAGE_SIZE` to detect last
                     page, but the actual server cap varies per cluster version
                     and is often lower than 1000. That check fired immediately
                     on page 1 for busy clusters, stopping pagination early.
                     Fix: remove the page-size check entirely; page purely by
                     timestamp cursor until the response is empty, oldest run
                     predates the window start, or no progress is made. Safety
                     limit raised to 200 pages.
  1.17 (2026-04-08) — fix: Trends tab shows fewer than 30 days on clusters with
                     many protection groups — v1 protectionRuns caps per page;
                     added pagination by advancing endTimeUsecs cursor.
  1.16 (2026-04-08) — feat: Trends chart — visible X/Y axes with tick marks,
                     Y axis fixed 0-100 with "0.0" format, circle data-point
                     markers (size 5, Cohesity green), 2pt line weight, StrRef
                     category fix so date strings render as labels in Excel,
                     auto tick-density reduction for >14 and >60 day ranges,
                     chart title updated to note warnings are included.
  1.15 (2026-04-08) — feat: include kWarning runs in Success % calculation
                     throughout — Trends "Success % (incl. Warnings)" col H,
                     _score_protection() quick mode, _success_stats() quick
                     mode. Warning runs complete successfully (data protected)
                     but may have minor issues (skipped objects). They still
                     appear in the Warning column for visibility. Word doc
                     Methodology appendix now documents this behavior.
  1.14 (2026-04-08) — fix: Trends tab inaccurate — switched primary data source
                     from per-group v2 runs to v1 /public/protectionRuns which
                     is the same source powering the Helios reporting page.
                     Single call per cluster; backupRun.slaViolated (col I) and
                     backupRun.stats.totalLogicalBackupSizeBytes (col J) are
                     accurate. v2 per-group runs retained as fallback. Works in
                     --quick mode (v1 call not gated behind the per-group loop).
  1.13 (2026-04-08) — fix: Trends tab cols I and J always blank — col I
                     (SLA Violations) used run.get("isSlaViolated") but the
                     field lives inside localBackupInfo, not the top-level run
                     object; fixed to lb.get("isSlaViolated"). Col J (Logical GB)
                     used lb.get("stats").get("logicalSizeBytes") but the correct
                     sub-object is localSnapshotStats; fixed to
                     lb.get("localSnapshotStats").get("logicalSizeBytes").
  1.12 (2026-04-08) — fix: Replication & Archive col E and col K still empty
                     after v1.11 — policy→group chain produced no rows because
                     _arch_name() returns "" when v2 policies store archival
                     targets by vault ID only (no name/targetName field), so
                     policy_arch_targets stayed empty and all_arch_targets was
                     empty. Fix: supplement arch_groups directly from
                     fk_data.dataTransferPerProtectionJob[].protectionJobName
                     per vault (reliable even without policy resolution); build
                     all_arch_targets as union of policy chain + fk_data vault
                     names + vault list names so at least one source always
                     produces rows and E/K are populated.
  1.11 (2026-04-08) — feat: Replication & Archive col E changed from policy
                     count to protection group count; new col K lists all group
                     names using each target. Status updated: "Configured" when
                     groups are assigned, "No groups assigned" when not.
                     Rewired target discovery through policy→group chain so E
                     and K reflect actual group assignment, not policy count.
  1.10 (2026-04-08) — fix: Replication & Archive col E (Policies Using) showed
                     0 because _arch_name() returned "" when the v2 policy stores
                     archival targets by vault ID only (no name/targetName field).
                     Added vault ID → name resolution via /public/vaults lookup.
                     Columns G/H renamed to include API field names and note that
                     values are cumulative totals. "Vault only — no policy" status
                     reworded to "Vault exists — not referenced by any policy".
  1.9 (2026-04-08) — fix: Replication & Archive still blank after v1.8.
                     dataTransferToVaults API requires vaultIds filter to
                     return any rows — added vault ID extraction and pass-
                     through from _vaults() to _fortknox_data(). Vault type
                     column D now read from archivalTargetConfig.targetType
                     inside the policy target (previously name-based lookup
                     against vault list failed on every mismatch).
  1.8 (2026-04-08) — fix: Replication & Archive sheet columns D-H blank.
                     Root causes: (1) _fortknox_data used startTimeUsecs/
                     endTimeUsecs — API expects startTimeMsecs/endTimeMsecs,
                     so fk_data was always empty; (2) archival rows driven by
                     cd["vaults"] list whose name may differ from policy target
                     names — flipped to drive rows from arch_use (policy targets)
                     and enrich vault type/FK stats via name-then-ID lookup;
                     (3) G/H (Logical/Physical Transferred) hardcoded "" — now
                     populated from numLogicalBytesTransferred/numPhysical.
                     Orphaned vaults (in API but not in any policy) still shown.
  1.7 (2026-04-08) — fix: _fk_status() false-negatives — lastRun.archivalInfo
                     is only populated when the most recent primary run included
                     archival; groups with daily local backups + less-frequent
                     archival showed as idle even when active. Fix: scan full
                     group_runs history (non-quick mode) for successful FK
                     archival results matched by target type AND vault name;
                     lastRun remains the fallback for --quick mode only.
  1.6 (2026-04-08) — feat: detect FortKnox configured-but-idle state. New
                     _fk_status() returns "active"/"idle"/"none" by inspecting
                     group lastRun.archivalInfo (dataTransferToVaults storageConsumed
                     is cumulative and not reliable for recency). Security sheet
                     column 5 shows "Active" (green) / "Configured — Idle" (orange) /
                     "No" (red). Idle triggers a HIGH recommendation and scores
                     10/20 points instead of the full 20.
  1.5 (2026-04-08) — fix: drop Alerts columns K (Node ID) and L (Entity) —
                     sparse propertyList fields only present on hardware alerts;
                     description + alert code already capture the context.
  1.4 (2026-04-08) — fix: Policy Audit column N (Groups Using) now computed by
                     counting groups whose policyId matches each policy — the v2
                     API does not return numProtectionGroups.
                     feat: new "Policy → Groups" sheet (sheet 6) listing every
                     protection group with the policy that governs it, sorted by
                     policy then group name, paused/failed rows highlighted.
  1.3 (2026-04-08) — fix: infrastructure disk count now reads diskCount integer
                     per node (v1 /public/nodes field) instead of looking for
                     non-existent list fields; NTP servers shows "Not configured"
                     instead of blank when the cluster has no NTP entries; added
                     ntpServersList and ntpServerIps fallback paths.
  1.2 (2026-04-08) — fix: correct field names throughout — firstTimestampUsecs,
                     clusterSoftwareVersion, localSnapshotStats, successfulObjectsCount,
                     policyId→name lookup, v2 schedule/target/retention paths,
                     broadened FortKnox detection, N/A capacity fallback.
                     feat: output filename includes customer name + local timestamp.
  1.0 (2026-04-08) — feat: initial release. 12-sheet Excel + Word document.
                     Health scoring, recommendations engine, trend charts.
"""

__version__ = "1.66"

import argparse
import datetime
import os
import re
import sys
import time
from concurrent.futures import ThreadPoolExecutor as _ThreadPoolExecutor

# ── Core HTTP client ──────────────────────────────────────────────────────────
# requests must be present.  If missing, exit with a clear install command
# rather than propagating a cryptic "No module named 'urllib3'" traceback.
# urllib3 InsecureRequestWarning is only suppressed when --insecure is passed
# explicitly; by default TLS verification is on.
try:
    import requests
except ImportError:
    print("=" * 70)
    print("  MISSING REQUIRED PACKAGE: requests")
    print()
    print("  Install all required and optional packages at once:")
    print("    pip install requests openpyxl python-docx python-pptx matplotlib keyring")
    print("=" * 70)
    sys.exit(1)

# ── Auth + formatting helpers ─────────────────────────────────────────────────
# Try the repo's utils/ directory first (normal dev/clone usage).
# If the utils modules are not found (e.g. single-file standalone download),
# fall back to inline implementations so the script works out of the box.
_utils_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "utils")
sys.path.insert(0, _utils_dir)
try:
    from cohesity_auth import (
        get_api_key, get_helios_clusters, make_helios_headers, HELIOS_HOST,
        get_cluster_password, get_auth_token, make_headers,
        clear_stored_credentials,
    )
    from formatters import (
        bytes_to_gb, bytes_to_tb, usecs_to_datetime, auto_fit_columns,
    )
except ImportError:
    # ── Bundled fallbacks — used only when running as a standalone file ───────
    import getpass as _getpass

    HELIOS_HOST     = "helios.cohesity.com"
    _KR_SVC_HELIOS  = "cohesity_helios"
    _KR_USER_HELIOS = "apikey"
    _KR_SVC_CLUSTER = "cohesity_cluster"

    def _keyring_available():
        try:
            import keyring  # noqa: F401
            return True
        except ImportError:
            return False

    def get_api_key(cli_key=None):
        if cli_key:
            return cli_key
        if _keyring_available():
            import keyring
            stored = keyring.get_password(_KR_SVC_HELIOS, _KR_USER_HELIOS)
            if stored:
                print("[*] Using Helios API key stored in system keychain.")
                return stored
            print("[*] No stored Helios API key found.")
        else:
            print("    NOTE: 'keyring' not installed — key will not be saved.")
        key = _getpass.getpass("    Enter Helios API key: ").strip()
        if not key:
            print("ERROR: API key cannot be empty.")
            sys.exit(1)
        if _keyring_available():
            import keyring
            keyring.set_password(_KR_SVC_HELIOS, _KR_USER_HELIOS, key)
            print("[*] Helios API key saved to system keychain.")
        return key

    def get_cluster_password(cluster, username, domain, cli_password=None):
        if cli_password:
            return cli_password
        kr_user = f"{cluster}:{domain}:{username}"
        if _keyring_available():
            import keyring
            stored = keyring.get_password(_KR_SVC_CLUSTER, kr_user)
            if stored:
                print(f"[*] Using stored password for {domain}\\{username}@{cluster}")
                return stored
        pwd = _getpass.getpass(f"Password for {domain}\\{username}@{cluster}: ").strip()
        if not pwd:
            print("ERROR: Password cannot be empty.")
            sys.exit(1)
        if _keyring_available():
            import keyring
            keyring.set_password(_KR_SVC_CLUSTER, kr_user, pwd)
            print(f"[*] Password saved to keychain for {domain}\\{username}@{cluster}")
        return pwd

    def clear_stored_credentials(cluster=None, username="admin", domain="LOCAL"):
        if not _keyring_available():
            print("ERROR: 'keyring' not installed. Nothing to clear.")
            sys.exit(1)
        import keyring
        if cluster:
            kr_user = f"{cluster}:{domain}:{username}"
            if keyring.get_password(_KR_SVC_CLUSTER, kr_user):
                keyring.delete_password(_KR_SVC_CLUSTER, kr_user)
                print(f"[*] Stored password for {domain}\\{username}@{cluster} removed.")
            else:
                print(f"[*] No stored password found for {domain}\\{username}@{cluster}.")
        else:
            if keyring.get_password(_KR_SVC_HELIOS, _KR_USER_HELIOS):
                keyring.delete_password(_KR_SVC_HELIOS, _KR_USER_HELIOS)
                print("[*] Stored Helios API key removed.")
            else:
                print("[*] No stored Helios API key found — nothing to clear.")

    def get_auth_token(cluster, username, password, domain, mfa_code=None,
                       verify=True):
        _hdrs  = {"Content-Type": "application/json", "Accept": "application/json"}
        url_v1 = f"https://{cluster}/irisservices/api/v1/public/accessTokens"
        _v1_note = ""
        try:
            r = requests.post(url_v1,
                              json={"username": username, "password": password,
                                    "domain": domain},
                              headers=_hdrs, verify=verify, timeout=30)
            r.raise_for_status()
            data = r.json()
            token_val = data.get("accessToken", "")
            if token_val:
                print(f"[+] Authenticated to {cluster} as {domain}\\{username} (v1)")
                return f"{data.get('tokenType', 'Bearer')} {token_val}"
        except requests.exceptions.ConnectionError:
            print(f"ERROR: Cannot connect to '{cluster}'. Check hostname/IP and network.")
            sys.exit(1)
        except requests.exceptions.HTTPError:
            _v1_note = f"HTTP {r.status_code}"
            print(f"  [!] v1 endpoint failed ({r.status_code}) — trying v2 ...")
        url_v2 = f"https://{cluster}/v2/users/sessions"
        pl2 = {"username": username, "password": password, "domain": domain}
        if mfa_code:
            pl2["otpCode"] = mfa_code
            pl2["otpType"] = "Totp"
        try:
            r = requests.post(url_v2, json=pl2, headers=_hdrs, verify=verify, timeout=30)
            r.raise_for_status()
        except requests.exceptions.ConnectionError:
            print(f"ERROR: Cannot connect to '{cluster}'. Check hostname/IP and network.")
            sys.exit(1)
        except requests.exceptions.HTTPError as e:
            _body = r.text[:400]
            if r.status_code == 401 and "otp" in _body.lower():
                print("ERROR: Cluster requires MFA. Re-run with --mfa-code <code>.")
                sys.exit(1)
            print(f"ERROR: Authentication failed on both endpoints.")
            if _v1_note:
                print(f"  v1: {_v1_note}")
            print(f"  v2: HTTP {r.status_code}")
            print()
            print("  Troubleshooting:")
            print("    • Wrong password?   Run --clear-credentials to wipe the stored")
            print("      password, then retry (you will be prompted for the correct one).")
            print(f"    • Wrong domain?    Use --domain <AD_DOMAIN> for AD accounts")
            print(f"      (current: {domain}). Use LOCAL for local cluster accounts.")
            print("    • Account locked?  Check Access Management in the Cohesity UI.")
            sys.exit(1)
        data = r.json()
        token_val = data.get("accessToken", "")
        if not token_val:
            print("ERROR: No access token in response.")
            sys.exit(1)
        print(f"[+] Authenticated to {cluster} as {domain}\\{username} (v2)")
        return f"{data.get('tokenType', 'Bearer')} {token_val}"

    def make_headers(token):
        return {"Authorization": token, "Accept": "application/json",
                "Content-Type": "application/json"}

    def make_helios_headers(api_key, cluster_id=None):
        h = {"apiKey": api_key, "Accept": "application/json",
             "Content-Type": "application/json"}
        if cluster_id is not None:
            h["accessClusterId"] = str(cluster_id)
        return h

    def get_helios_clusters(api_key, verify=True):
        url = f"https://{HELIOS_HOST}/mcm/clusters/connectionStatus"
        try:
            r = requests.get(url, headers=make_helios_headers(api_key),
                             verify=verify, timeout=30)
            r.raise_for_status()
        except requests.exceptions.ConnectionError:
            print("ERROR: Cannot connect to Helios. Check your network/VPN.")
            sys.exit(1)
        except requests.exceptions.HTTPError:
            print(f"ERROR: Helios auth failed ({r.status_code}). Check your API key.")
            sys.exit(1)
        clusters = r.json()
        if not isinstance(clusters, list):
            clusters = clusters.get("clusterList", [])
        result = [{"name": c.get("name", ""), "clusterId": c.get("clusterId")}
                  for c in clusters]
        print(f"[+] Connected to Helios — {len(result)} cluster(s)")
        return result

    # ── Formatting helpers ────────────────────────────────────────────────────
    from datetime import datetime as _dt, timezone as _tz_utc

    def usecs_to_datetime(usecs, tz=None):
        if not usecs:
            return ""
        dt = _dt.fromtimestamp(usecs / 1_000_000, tz=_tz_utc.utc)
        if tz is not None:
            dt = dt.astimezone(tz)
        return dt.strftime("%Y-%m-%d %H:%M:%S")

    def bytes_to_gb(b):
        return round(b / 1_000_000_000, 4) if b else 0.0

    def bytes_to_tb(b):
        return round(b / 1_000_000_000_000, 4) if b else 0.0

    def auto_fit_columns(ws, min_width=6, max_width=60):
        try:
            from openpyxl.cell.cell import MergedCell
            from openpyxl.utils import get_column_letter as _gcl
            from openpyxl.styles import Alignment as _Aln2
        except ImportError:
            return
        try:
            ws.sheet_view.zoomScale = 130
        except Exception:
            pass
        for col in ws.columns:
            if not col:
                continue
            first = col[0]
            col_letter = (first.column_letter if hasattr(first, "column_letter")
                          else _gcl(first.column))
            best = min_width
            for cell in col:
                if isinstance(cell, MergedCell):
                    continue
                if cell.value is not None:
                    lines = str(cell.value).split("\n")
                    best = max(best, max(len(ln) for ln in lines) + 2)
            ws.column_dimensions[col_letter].width = min(best, max_width)
        _al = _Aln2(horizontal="left", vertical="center")
        for row_cells in ws.iter_rows(min_row=2):
            for cell in row_cells:
                if not isinstance(cell, MergedCell) and cell.alignment.horizontal is None:
                    cell.alignment = _al

try:
    from openpyxl import Workbook
    from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
    from openpyxl.chart import LineChart, Reference
    from openpyxl.chart.series import SeriesLabel
    EXCEL_OK = True
    # Pre-allocated style singletons — reused across all cells to avoid
    # constructing a new object per cell during large workbook generation.
    _FONT_NORMAL = Font(size=10, color="000000")
    _FONT_BOLD   = Font(bold=True, size=10, color="000000")
except ImportError:
    EXCEL_OK = False
    _FONT_NORMAL = None
    _FONT_BOLD   = None

try:
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    DOCX_OK = True
except ImportError:
    DOCX_OK = False

try:
    from pptx import Presentation as _PptxPresentation
    from pptx.util import Inches as _PptxInches, Pt as _PptxPt
    from pptx.dml.color import RGBColor as _PptxRGB
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as _MSO_SHAPE
    from pptx.enum.dml import MSO_THEME_COLOR as _MSO_THEME
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

# ── Thresholds / best-practice baselines ─────────────────────────────────────
CAPACITY_WARN_PCT    = 70       # cluster used % — warning
CAPACITY_CRIT_PCT    = 85       # cluster used % — critical
SUCCESS_WARN_PCT     = 95.0     # backup success rate below this = warning
SUCCESS_CRIT_PCT     = 90.0     # backup success rate below this = critical
RPO_WARN_HOURS       = 24       # hours since last successful run — warning
RPO_CRIT_HOURS       = 48       # hours since last successful run — critical
ALERT_AGE_CRIT_DAYS  = 7        # open critical alert older than N days = finding
VIEW_QUOTA_WARN_PCT  = 80       # view at this % of quota = warning
VIEW_QUOTA_CRIT_PCT  = 90       # view at this % of quota = critical
VERSION_WARN_PREFIX  = "7."     # clusters below this major version get a flag
SSD_WEAR_WARN_PCT    = 70       # SSD wear % — warning
SSD_WEAR_CRIT_PCT    = 80       # SSD wear % — critical
CERT_WARN_DAYS       = 90       # TLS cert expiry warning threshold (days)
CERT_CRIT_DAYS       = 30       # TLS cert expiry critical threshold (days)

# ── Software End-of-Support matrix (checked top-to-bottom, first prefix match wins) ──
# Source: Cohesity Product Life Cycle Policy (updated Oct 2025)
SOFTWARE_EOS = [
    # (version_prefix, status, eos_date_str_or_None, display_label)
    # Evaluated in order — first prefix match wins; put longer/more-specific prefixes first.
    ("7.3",   "current",     None,         "Current"),
    ("7.2",   "in_support",  "2026-06-08", "In Support — Feature Release"),
    ("7.1.2", "in_support",  "2026-06-08", "In Support — LTS"),
    ("7.1",   "eol",         None,         "Out of Support"),   # 7.1.0 and 7.1.1
    ("7.0",   "eol",         None,         "Out of Support"),
    ("6.8.2", "in_support",  "2026-06-08", "In Support — LTS"),
    ("6.8",   "eol",         None,         "Out of Support"),   # 6.8.0 / 6.8.1 EOS Nov 2024
    # anything older falls through to the default → eol
]

# ── Hardware End-of-Life matrix (case-insensitive substring match on model string) ──
# Source: Cohesity Platforms EOL Terms and Service Express EOL database
HARDWARE_EOL = [
    # (model_substring, eol_date_str)
    # Source: Cohesity Products End of Support and End of Life Dates (Oct 2025)
    # C2xx0 series — EOL Jan 31, 2023
    ("C2000", "2023-01-31"), ("C2100", "2023-01-31"),
    ("C2200", "2023-01-31"), ("C2300", "2023-01-31"),
    # C2xx5 series — EOL May 13, 2024
    ("C2500", "2024-05-13"), ("C2505", "2024-05-13"), ("C2510", "2024-05-13"),
    # C3500 — EOL Feb 18, 2025
    ("C3500", "2025-02-18"),
    # C4xx0 / CN series — EOL Feb 22, 2027
    ("C4000", "2027-02-22"), ("C4300", "2027-02-22"), ("C4500", "2027-02-22"),
    ("CN3",   "2027-02-22"), ("CN4",   "2027-02-22"),
    # C60x5 — no EOL date announced as of Oct 2025 (omitted intentionally)
]
HW_EOL_WARN_DAYS = 180    # flag hardware EOL within this many days
SW_EOS_WARN_DAYS = 90     # flag software EOS within this many days

# Scoring weights — must sum to 1.0
_W = {
    "protection": 0.25,
    "sla":        0.20,
    "infra":      0.15,
    "capacity":   0.15,
    "security":   0.15,
    "alerts":     0.10,
}

# ── Excel palette ─────────────────────────────────────────────────────────────
COH_GREEN   = "70AD47"
DARK_GREEN  = "4E7A2E"
LT_GREEN    = "E2EFDA"
RED         = "FF4C4C"
ORANGE      = "FF8C00"
YELLOW      = "FFD700"
LIGHT_GRAY  = "F2F2F2"
WHITE       = "FFFFFF"


# ── openpyxl style helpers (lazy — only called when EXCEL_OK is True) ─────────
_IP4_RE = re.compile(r'^\d{1,3}(\.\d{1,3}){3}$')
_IP6_RE = re.compile(r'^[0-9a-fA-F]{0,4}(:[0-9a-fA-F]{0,4}){2,7}$')

def _split_host_ip(raw_host, raw_ip):
    """Route values to (hostname, ip_address) by content, not field name.

    The Cohesity agents API sometimes puts FQDNs in hostIp and leaves hostName
    empty. Inspect each value: if it looks like an IPv4/IPv6 address put it in
    the ip slot; otherwise put it in the host slot.
    """
    host = ""
    ip   = ""
    for val in (raw_host, raw_ip):
        val = (val or "").strip()
        if not val:
            continue
        if _IP4_RE.match(val) or _IP6_RE.match(val):
            if not ip:   ip   = val
        else:
            if not host: host = val
    return host, ip


def _fill(hex_color):
    return PatternFill("solid", fgColor=hex_color)

def _font(bold=False, color="000000", size=10):
    return Font(bold=bold, color=color, size=size)

def _rag_fill(score):
    """Return fill based on 0-100 score."""
    if score >= 90: return _fill(LT_GREEN)
    if score >= 75: return _fill("FFFF99")
    if score >= 60: return _fill("FFD580")
    if score >= 40: return _fill("FFAA80")
    return _fill("FF8080")

def _hdr(ws, row, cols):
    """Write a bold Cohesity-green header row with a thin black border."""
    _thin   = Side(style="thin", color="000000")
    _border = Border(left=_thin, right=_thin, top=_thin, bottom=_thin)
    for c, val in enumerate(cols, 1):
        cell = ws.cell(row=row, column=c, value=val)
        cell.font      = Font(bold=True, color=WHITE, size=10)
        cell.fill      = _fill(COH_GREEN)
        cell.alignment = Alignment(wrap_text=True, vertical="center",
                                   horizontal="center")
        cell.border    = _border
    ws.row_dimensions[row].height = 32

def _title(ws, text, span):
    ws.merge_cells(f"A1:{span}1")
    c = ws["A1"]
    c.value     = text
    c.font      = Font(bold=True, color=WHITE, size=13)
    c.fill      = _fill(COH_GREEN)
    c.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 22

_FORMULA_PREFIXES = ("=", "+", "-", "@", "\t", "\r")

def _safe_cell(val):
    """Prefix strings starting with formula metacharacters to prevent injection."""
    if isinstance(val, str) and val and val[0] in _FORMULA_PREFIXES:
        return "'" + val
    return val

def _reg(ws, row, col, val):
    cell = ws.cell(row=row, column=col, value=_safe_cell(val))
    cell.font = _FONT_NORMAL
    return cell

# ── Grade from score ──────────────────────────────────────────────────────────
def grade(score):
    if score >= 90: return "Excellent"
    if score >= 75: return "Good"
    if score >= 60: return "Fair"
    if score >= 40: return "At Risk"
    return "Critical"


def _sw_eos(version: str):
    """Return (status, eos_date_or_None, label) for a cluster software version string.
    status: 'current' | 'in_support' | 'eol' | 'unknown'
    """
    import datetime as _dt
    if not version:
        return ("unknown", None, "Unknown")
    v = version.strip().lower()
    for prefix, status, eos_str, label in SOFTWARE_EOS:
        if v.startswith(prefix.lower()):
            eos = _dt.date.fromisoformat(eos_str) if eos_str else None
            return (status, eos, label)
    # No match — treat anything pre-6.8 as out of support
    return ("eol", None, "Out of Support")


def _hw_eol(model: str):
    """Return eol_date (datetime.date) or None if model not in EOL matrix."""
    import datetime as _dt
    if not model:
        return None
    m = model.upper()
    for substr, eol_str in HARDWARE_EOL:
        if substr.upper() in m:
            return _dt.date.fromisoformat(eol_str)
    return None


# ── Time helpers ──────────────────────────────────────────────────────────────
def _now_usecs():
    return int(time.time() * 1_000_000)

def _days_ago_usecs(n):
    return int((time.time() - n * 86400) * 1_000_000)


# ─────────────────────────────────────────────────────────────────────────────
# API HELPERS
# ─────────────────────────────────────────────────────────────────────────────

def _get(session, headers, path, params=None, debug=False, silent=False):
    """GET from Helios or a direct cluster.

    The base URL is read from session.base_url (set in main() based on auth
    mode). Defaults to Helios if not set.

    silent=True suppresses the debug body print on non-200 responses — use for
    best-effort endpoints where 400/403/404 is expected on some clusters.
    """
    base = getattr(session, "base_url", f"https://{HELIOS_HOST}")
    url  = f"{base}{path}"
    try:
        r = session.get(url, headers=headers, params=params, timeout=60)
        if debug:
            print(f"    DEBUG {path} → {r.status_code}")
        if r.status_code == 200:
            return r.json()
        if debug and not silent:
            print(f"    DEBUG body: {r.text[:300]}")
        return None
    except Exception as exc:
        print(f"    WARN: GET {path} failed: {exc}")
        return None


# ─────────────────────────────────────────────────────────────────────────────
# DATA COLLECTION  (one function per resource type)
# ─────────────────────────────────────────────────────────────────────────────

def _cluster_info(session, h, debug):
    d = _get(session, h, "/irisservices/api/v1/public/cluster",
             params={"fetchStats": "true"}, debug=debug)
    return d or {}

def _nodes(session, h, debug):
    d = _get(session, h, "/irisservices/api/v1/public/nodes", debug=debug)
    return d or []

def _alerts(session, h, days, debug):
    d = _get(session, h, "/irisservices/api/v1/public/alerts",
             params={
                 "alertStateList": "kOpen",
                 "maxAlerts": 1000,
                 "startDateUsecs": _days_ago_usecs(days),
                 "endDateUsecs":   _now_usecs(),
             }, debug=debug)
    return d or []

def _protection_groups(session, h, debug):
    d = _get(session, h, "/v2/data-protect/protection-groups",
             params={
                 "isDeleted": "false",
                 "includeTenants": "true",
                 "includeLastRunInfo": "true",
             }, debug=debug)
    return (d or {}).get("protectionGroups") or []

def _group_runs(session, h, group_id, start_usecs, end_usecs, debug):
    d = _get(session, h,
             f"/v2/data-protect/protection-groups/{group_id}/runs",
             params={
                 "startTimeUsecs": start_usecs,
                 "endTimeUsecs":   end_usecs,
                 "numRuns": 500,
                 "includeTenants": "true",
             }, debug=debug)
    return (d or {}).get("runs") or []

def _v1_protection_runs(session, h, start_usecs, end_usecs, debug):
    """Paginated fetch of ALL protection runs in the time window via v1 API.

    The API returns runs newest-first. The actual per-page cap varies by
    cluster version; we do NOT rely on comparing len(page) to a constant
    because that assumption breaks when the server cap is lower than expected.
    Instead we page purely by timestamp cursor until:
      - the response is empty, or
      - the oldest run on the page predates our window start, or
      - no timestamp progress is made (safety against infinite loop).
    Safety limit: 200 pages.
    """
    all_runs = []
    page_end  = end_usecs

    for page_num in range(200):
        page = _get(session, h, "/irisservices/api/v1/public/protectionRuns",
                    params={"startTimeUsecs": start_usecs,
                            "endTimeUsecs":   page_end,
                            "numRuns":        1000},
                    debug=debug) or []
        if not page:
            break
        all_runs.extend(page)
        if debug:
            print(f"    DEBUG v1_runs page {page_num + 1}: {len(page)} runs "
                  f"(total so far: {len(all_runs)})")
        # Find the oldest startTimeUsecs on this page to use as the next cursor
        oldest_t = None
        for r in page:
            t = ((r.get("backupRun") or {}).get("stats") or {}).get("startTimeUsecs") or 0
            if t and (oldest_t is None or t < oldest_t):
                oldest_t = t
        if not oldest_t:
            break   # no timestamps found — can't advance cursor
        if oldest_t <= start_usecs:
            break   # full window covered
        if oldest_t >= page_end:
            break   # no progress — prevent infinite loop
        page_end = oldest_t - 1  # advance cursor to just before oldest run seen

    return all_runs

def _policies(session, h, debug):
    d = _get(session, h, "/v2/data-protect/policies",
             params={"includeTenants": "true"}, debug=debug)
    return (d or {}).get("policies") or []

def _storage_domains(session, h, debug):
    d = _get(session, h, "/v2/storage-domains",
             params={"includeStats": "true"}, debug=debug)
    return (d or {}).get("storageDomains") or []

def _views(session, h, debug):
    d = _get(session, h, "/v2/file-services/views",
             params={"includeStats": "true", "maxCount": 2000}, debug=debug)
    return (d or {}).get("views") or []

def _vaults(session, h, debug):
    d = _get(session, h, "/irisservices/api/v1/public/vaults",
             params={"includeFortKnoxVault": "true"}, debug=debug)
    return d or []


def _helios_quorum_groups(session, api_key, debug):
    """Fetch Helios-level quorum groups from /v2/mcm/quorum/groups.

    This is a Helios-scoped endpoint — do NOT include accessClusterId in headers.
    Returns a list of group dicts, each with:
      isEnabled          — True/False
      clusterIdentifiers — list of {"clusterId": <int>, ...}
    """
    h = make_helios_headers(api_key)   # no cluster_id — Helios-level call
    d = _get(session, h, "/v2/mcm/quorum/groups", debug=debug, silent=True)
    if isinstance(d, dict):
        return d.get("quorumGroups") or []
    if isinstance(d, list):
        return d
    return []


def _fortknox_data(session, h, start_usecs, end_usecs, vault_ids, debug):
    # API expects milliseconds; also requires vaultIds filter to return results
    params = {"startTimeMsecs": start_usecs // 1000,
              "endTimeMsecs":   end_usecs   // 1000}
    if vault_ids:
        params["vaultIds"] = vault_ids  # requests repeats list values correctly
    elif not vault_ids:
        return []   # no vaults — nothing to query
    d = _get(session, h,
             "/irisservices/api/v1/public/reports/dataTransferToVaults",
             params=params, debug=debug)
    return (d or {}).get("dataTransferSummary") or []


def _certificates(session, h, debug):
    """Fetch web server TLS certificate(s). Returns list of cert dicts only."""
    d = _get(session, h, "/irisservices/api/v1/public/certificates/webServer",
             debug=debug)
    if isinstance(d, list):
        return [c for c in d if isinstance(c, dict)]
    if isinstance(d, dict) and d:
        return [d]
    return []


def _agents(session, h, debug):
    """Fetch Cohesity agent deployment/health status per registered host."""
    d = _get(session, h, "/irisservices/api/v1/public/reports/agents",
             debug=debug)
    if isinstance(d, list):
        return d
    return (d or {}).get("agentDeploymentStatusList") or []


def _disks(session, h, nodes, debug):
    """Fetch per-disk info using GET /v2/disks/local?nodeId={id} per node.

    This mirrors the approach used by the Cohesity ELF inventory tool
    (cohesityInv.ps1, apiType='ClusterV2', path='/disks/local?nodeId={0}').
    Each node is queried individually; results are annotated with _nodeId
    and _nodeIp for display convenience.

    Field names returned by this endpoint (ELF script reference):
      serialNumber, model, status, location, type, encryptionStatus,
      ssdUsedPercentage (raw 0-100), removalProgressList, capacityInBytes
    """
    all_disks = []
    for n in nodes:
        nid = n.get("id")
        if not nid:
            continue
        d = _get(session, h, "/v2/disks/local",
                 params={"nodeId": nid}, debug=debug, silent=True)
        if not d:
            continue
        dl = d if isinstance(d, list) else (d.get("disksList") or d.get("disks") or [])
        for disk in dl:
            if not isinstance(disk, dict):
                continue
            disk["_nodeId"] = nid
            disk["_nodeIp"] = n.get("ip", "")
        all_disks.extend(dk for dk in dl if isinstance(dk, dict))
    return all_disks


def _sources(session, h, debug):
    """Fetch registered protection sources with coverage stats.

    Tries v2 first (includeTenants).  Falls back to v1
    /protectionSources/registrationInfo on older cluster firmware where the
    v2 endpoint returns 404.  Returns a list of normalized source dicts that
    _sheet_source_coverage and _build_recommendations can consume via the
    standard protectedObjectsCount / unprotectedObjectsCount keys.
    """
    d = _get(session, h, "/v2/data-protect/sources",
             params={"includeTenants": "true"}, debug=debug, silent=True)
    if isinstance(d, list) and d:
        return d
    if isinstance(d, dict):
        v2 = d.get("sources") or []
        if v2:
            return v2

    # v1 fallback — works on all supported cluster versions
    d1 = _get(session, h,
              "/irisservices/api/v1/public/protectionSources/registrationInfo",
              debug=debug)
    rows = d1 if isinstance(d1, list) else ((d1 or {}).get("rootNodes") or [])
    out = []
    for row in rows:
        if not isinstance(row, dict):
            continue
        ps   = row.get("rootNode") or row.get("protectionSource") or {}
        ps   = ps if isinstance(ps, dict) else {}
        st   = row.get("stats") or {}
        st   = st if isinstance(st, dict) else {}
        out.append({
            "name":                    ps.get("name", ""),
            "environment":             ps.get("environment", ""),
            "protectedObjectsCount":   st.get("protectedCount")   or st.get("protectedObjectsCount")   or 0,
            "unprotectedObjectsCount": st.get("unprotectedCount") or st.get("unprotectedObjectsCount") or 0,
        })
    return out


def _users(session, h, debug):
    """Fetch local and domain user accounts."""
    d = _get(session, h, "/irisservices/api/v1/public/users",
             params={"allUnderHierarchy": "true"}, debug=debug)
    return d or []


def _idps(session, h, debug):
    """Fetch configured SSO / IDP providers."""
    d = _get(session, h, "/irisservices/api/v1/public/idps", debug=debug)
    return d or []


def _tenants(session, h, debug):
    """Fetch tenant / organization list (multi-tenancy inventory)."""
    d = _get(session, h, "/irisservices/api/v1/public/tenants", debug=debug, silent=True)
    return d or []


def _capacity_timeseries(session, h, cluster_id, debug):
    """Fetch 30-day daily capacity time series for the cluster.

    Uses /irisservices/api/v1/public/statistics/timeSeriesStats with
    metric kMorphedUsageBytes (total physical bytes stored on cluster).
    Returns a list of {"timestampMsecs": int, "data": {"int64Value": int}}
    points, or [] on any error or when cluster_id is unavailable.
    """
    if not cluster_id:
        return []
    now_ms   = int(time.time() * 1000)
    start_ms = now_ms - 30 * 86400 * 1000
    d = _get(session, h,
             "/irisservices/api/v1/public/statistics/timeSeriesStats",
             params={
                 "schemaName":         "kBridgeClusterStats",
                 "metricName":         "kMorphedUsageBytes",
                 "entityId":           str(cluster_id),
                 "rollupIntervalSecs": 86400,
                 "rollupFunction":     "latest",
                 "startTimeMsecs":     start_ms,
                 "endTimeMsecs":       now_ms,
             },
             debug=debug, silent=True)
    if not d:
        return []
    return d.get("dataPointVec") or []


def _fetch_audit_log(session, h, days, debug):
    """Fetch recent audit log events (last N days).

    Returns a list of event dicts, the sentinel string "PERMISSION_DENIED" on
    HTTP 401/403 (so the sheet can surface a helpful message), or [] on any
    other error or empty response.
    """
    base = getattr(session, "base_url", f"https://{HELIOS_HOST}")
    url  = f"{base}/irisservices/api/v1/public/auditLog"
    params = {"numLogs": 1000, "startTimeUsecs": _days_ago_usecs(days)}
    try:
        r = session.get(url, headers=h, params=params, timeout=60)
        if debug:
            print(f"    DEBUG /auditLog → {r.status_code}")
        if r.status_code == 200:
            d = r.json()
            if isinstance(d, list):
                return d
            if isinstance(d, dict):
                return d.get("auditLogs") or d.get("events") or d.get("logs") or []
            return []
        if r.status_code in (401, 403):
            return "PERMISSION_DENIED"
        return []
    except Exception as exc:
        print(f"    WARN: GET /auditLog failed: {exc}")
        return []


def _recoveries(session, h, start_usecs, end_usecs, debug):
    """Fetch recovery tasks in the lookback window."""
    d = _get(session, h, "/v2/data-protect/recoveries",
             params={"startTimeUsecs": start_usecs,
                     "endTimeUsecs":   end_usecs,
                     "numRecoveries":  1000},
             debug=debug, silent=True)
    return (d.get("recoveries") if isinstance(d, dict) else None) or []


def _kms_config(session, h, debug):
    d = _get(session, h,
             "/irisservices/api/v1/public/kmsConfig",
             debug=debug, silent=True)
    if isinstance(d, dict):
        return d
    if isinstance(d, list) and d and isinstance(d[0], dict):
        return d[0]
    return {}


def _roles(session, h, debug):
    """Fetch all role definitions (built-in and custom)."""
    d = _get(session, h,
             "/irisservices/api/v1/public/roles",
             debug=debug, silent=True)
    if isinstance(d, list):
        return d
    return (d.get("roles") if isinstance(d, dict) else None) or []


def _active_directory(session, h, debug):
    d = _get(session, h, "/irisservices/api/v1/public/activeDirectory",
             debug=debug, silent=True)
    if isinstance(d, list):
        return d
    return (d.get("activeDirectoryEntry") if isinstance(d, dict) else None) or []


def _ldap_providers(session, h, debug):
    d = _get(session, h, "/irisservices/api/v1/public/ldapProvider",
             debug=debug, silent=True)
    if isinstance(d, list):
        return d
    return (d.get("ldapProviders") if isinstance(d, dict) else None) or []


def _all_certs(session, h, debug):
    d = _get(session, h, "/irisservices/api/v1/public/certificates",
             debug=debug, silent=True)
    if isinstance(d, list):
        return [c for c in d if isinstance(c, dict)]
    return (d.get("certificates") if isinstance(d, dict) else None) or []


def _dedup_trend(session, h, cluster_id, debug):
    """Fetch 30-day daily logical-capacity time series for dedup ratio computation."""
    if not cluster_id:
        return []
    now_ms   = int(time.time() * 1000)
    start_ms = now_ms - 30 * 86400 * 1000
    d = _get(session, h,
             "/irisservices/api/v1/public/statistics/timeSeriesStats",
             params={
                 "schemaName":         "kBridgeClusterStats",
                 "metricName":         "kLogicalCapacityBytes",
                 "entityId":           str(cluster_id),
                 "rollupIntervalSecs": 86400,
                 "rollupFunction":     "latest",
                 "startTimeMsecs":     start_ms,
                 "endTimeMsecs":       now_ms,
             },
             debug=debug, silent=True)
    return (d.get("dataPointVec") or [] if d else [])


def _io_stats(session, h, cluster_id, debug):
    """Fetch 30-day daily read/write throughput time series."""
    if not cluster_id:
        return {"reads": [], "writes": []}
    now_ms   = int(time.time() * 1000)
    start_ms = now_ms - 30 * 86400 * 1000
    base_p   = {
        "schemaName":         "kBridgeClusterStats",
        "entityId":           str(cluster_id),
        "rollupIntervalSecs": 86400,
        "rollupFunction":     "average",
        "startTimeMsecs":     start_ms,
        "endTimeMsecs":       now_ms,
    }
    dr = _get(session, h,
              "/irisservices/api/v1/public/statistics/timeSeriesStats",
              params={**base_p, "metricName": "kNumBytesRead"},
              debug=debug, silent=True)
    dw = _get(session, h,
              "/irisservices/api/v1/public/statistics/timeSeriesStats",
              params={**base_p, "metricName": "kNumBytesWritten"},
              debug=debug, silent=True)
    return {
        "reads":  (dr.get("dataPointVec") or [] if dr else []),
        "writes": (dw.get("dataPointVec") or [] if dw else []),
    }


def _snapshot_datalock(session, h, groups, debug):
    """Check actual DataLock status on recent snapshots for groups with DataLock policies.

    Only queries groups that have a DataLock policy. Limited to 20 groups max
    to avoid excessive API calls and timeouts.
    """
    results = []
    queried = 0
    for g in (groups or []):
        if queried >= 20:
            break
        gid = g.get("id")
        if not gid:
            continue
        d = _get(session, h,
                 "/v2/data-protect/snapshots",
                 params={"protectionGroupIds": gid,
                         "maxCount": 5},
                 debug=debug, silent=True)
        snaps = (d.get("snapshots") if isinstance(d, dict) else None) or []
        results.append({"group_id": gid,
                        "group_name": g.get("name", ""),
                        "snapshots": snaps})
        queried += 1
    return results


def collect_cluster(session, api_key, cluster, args):
    """Collect all health-check data for one cluster. Returns a dict.

    cluster dict shapes:
      Helios mode: {"name": "...", "clusterId": 123}
      Direct mode: {"name": "...", "clusterId": None, "mode": "direct",
                    "host": "10.1.2.3", "token": "Bearer ..."}
    """
    name = cluster.get("name", "Unknown")
    cid  = cluster.get("clusterId")
    dbg  = args.debug

    if cluster.get("mode") == "direct":
        h = make_headers(cluster["token"])
    else:
        h = make_helios_headers(api_key, cluster_id=cid)

    start_usecs = _days_ago_usecs(args.days)
    end_usecs   = _now_usecs()

    # ── Wave 1: all independent API calls in parallel ─────────────────────────
    print(f"  [{name}] fetching in parallel —")
    print(f"    infrastructure, nodes, alerts, groups, policies, storage, views,")
    print(f"    vaults, run-summary, certs, agents, sources, users, IDPs, cap-trend, audit-log,")
    print(f"    recoveries, KMS config, roles, AD/LDAP, all-certs, dedup-trend, I/O...",
          end="", flush=True)
    with _ThreadPoolExecutor(max_workers=8) as _ex:
        _fut_info    = _ex.submit(_cluster_info,          session, h, dbg)
        _fut_nodes   = _ex.submit(_nodes,                 session, h, dbg)
        _fut_alerts  = _ex.submit(_alerts,                session, h, args.days, dbg)
        _fut_groups  = _ex.submit(_protection_groups,     session, h, dbg)
        _fut_policies= _ex.submit(_policies,              session, h, dbg)
        _fut_domains = _ex.submit(_storage_domains,       session, h, dbg)
        _fut_views   = _ex.submit(_views,                 session, h, dbg)
        _fut_vaults  = _ex.submit(_vaults,                session, h, dbg)
        _fut_v1runs  = _ex.submit(_v1_protection_runs,    session, h, start_usecs, end_usecs, dbg)
        _fut_certs   = _ex.submit(_certificates,          session, h, dbg)
        _fut_agents  = _ex.submit(_agents,                session, h, dbg)
        _fut_sources = _ex.submit(_sources,               session, h, dbg)
        _fut_users   = _ex.submit(_users,                 session, h, dbg)
        _fut_idps    = _ex.submit(_idps,                  session, h, dbg)
        _fut_tenants = _ex.submit(_tenants,               session, h, dbg)
        _fut_audit   = _ex.submit(_fetch_audit_log,       session, h, args.days, dbg)
        _fut_recov   = _ex.submit(_recoveries,            session, h, start_usecs, end_usecs, dbg)
        _fut_kms     = _ex.submit(_kms_config,            session, h, dbg)
        _fut_roles   = _ex.submit(_roles,                 session, h, dbg)
        _fut_ad      = _ex.submit(_active_directory,      session, h, dbg)
        _fut_ldap    = _ex.submit(_ldap_providers,        session, h, dbg)
        _fut_allcerts= _ex.submit(_all_certs,             session, h, dbg)

        info         = _fut_info.result()
        nodes        = _fut_nodes.result()
        alerts       = _fut_alerts.result()
        groups       = _fut_groups.result()
        policies     = _fut_policies.result()
        domains      = _fut_domains.result()
        views        = _fut_views.result()
        vaults       = _fut_vaults.result()
        v1_runs      = _fut_v1runs.result()
        certs        = _fut_certs.result()
        # Filter to dicts at source — removes isinstance guards from hot loops
        agents       = [a for a in _fut_agents.result()  if isinstance(a, dict)]
        sources      = [s for s in _fut_sources.result() if isinstance(s, dict)]
        users        = _fut_users.result()
        idps         = _fut_idps.result()
        tenants      = _fut_tenants.result()
        audit_events = _fut_audit.result()
        recoveries   = _fut_recov.result()
        kms          = _fut_kms.result()
        roles        = _fut_roles.result()
        active_directory  = _fut_ad.result()
        ldap_providers    = _fut_ldap.result()
        all_certs         = _fut_allcerts.result()
    print(" done.")

    # ── Wave 2: calls that depend on wave-1 results ────────────────────────────
    # In direct mode clusterId is None on the cluster dict; resolve from info.
    effective_cid = info.get("id") or cid
    vault_ids = [v["id"] for v in vaults if "id" in v]
    print(f"    disks (per node), FortKnox transfer data, stats trends...", end="", flush=True)
    with _ThreadPoolExecutor(max_workers=6) as _ex:
        _fut_disks = _ex.submit(_disks, session, h, nodes, dbg)
        _fut_fk    = _ex.submit(_fortknox_data, session, h, start_usecs, end_usecs, vault_ids, dbg)
        # For the FortKnox detail tab: quick mode shows last 1 day only
        _fut_fk1d  = (_ex.submit(_fortknox_data, session, h, _days_ago_usecs(1), end_usecs, vault_ids, dbg)
                      if args.quick else None)
        _fut_capts = _ex.submit(_capacity_timeseries, session, h, effective_cid, dbg)
        _fut_dedup = _ex.submit(_dedup_trend,         session, h, effective_cid, dbg)
        _fut_io    = _ex.submit(_io_stats,            session, h, effective_cid, dbg)

        disks      = _fut_disks.result()
        fk_data    = _fut_fk.result()
        fk_data_1d = _fut_fk1d.result() if _fut_fk1d else []
        cap_ts     = _fut_capts.result()
        dedup_ts   = _fut_dedup.result()
        io_ts      = _fut_io.result()
    print(" done.")

    # ── Per-group run history (non-quick mode): parallelised ──────────────────
    group_runs = {}
    if not args.quick:
        print(f"    run history for {len(groups)} protection group(s)...", end="", flush=True)
        with _ThreadPoolExecutor(max_workers=8) as _ex:
            _grp_futs = {
                _ex.submit(_group_runs, session, h, g["id"], start_usecs, end_usecs, dbg): g["id"]
                for g in groups if g.get("id")
            }
            for _fut, _gid in _grp_futs.items():
                _runs = _fut.result()
                if _runs:
                    group_runs[_gid] = _runs
        print(" done.")

    # Build policy ID lookups used throughout the report (computed once here)
    policy_map   = {p.get("id"): p.get("name", "") for p in policies if p.get("id")}
    policy_by_id = {p.get("id"): p              for p in policies if p.get("id")}

    # ── Wave 3: DataLock snapshot verification — per-group, serial ───────────
    # Only query groups where the governing policy has DataLock configured.
    _DL_GROUPS = [
        g for g in groups
        if g.get("id") and _policy_datalock_str(
            policy_by_id.get(
                g.get("policyId") or (g.get("policy") or {}).get("id") or "", {}
            )
        ) != "None"
    ]
    snapshot_dl = []
    if _DL_GROUPS:
        print(f"    DataLock snapshot verification for up to 20 of "
              f"{len(_DL_GROUPS)} DataLock group(s)...", end="", flush=True)
        snapshot_dl = _snapshot_datalock(session, h, _DL_GROUPS, dbg)
        print(" done.")

    cd = {
        "name":        name,
        "id":          cid,
        "mode":        cluster.get("mode"),       # "direct" or None (Helios)
        "info":        info,
        "nodes":       nodes,
        "alerts":      alerts,
        "groups":      groups,
        "group_runs":  group_runs,
        "policies":    policies,
        "policy_map":  policy_map,
        "policy_by_id": policy_by_id,
        "domains":     domains,
        "views":       views,
        "vaults":      vaults,
        "fk_data":     fk_data,
        "v1_runs":     v1_runs,
        "certs":       certs,
        "agents":      agents,
        "disks":       disks,
        "sources":     sources,
        "users":       users,
        "idps":        idps,
        "tenants":     tenants,
        "start_usecs":    start_usecs,
        "end_usecs":      end_usecs,
        "days":           args.days,
        "quick":          args.quick,
        "fk_data_1d":     fk_data_1d,
        "quorum_enabled": cluster.get("quorum_enabled", False),
        "cap_timeseries": cap_ts,
        "audit_events":   audit_events,
        "recoveries":        recoveries,
        "kms":               kms,
        "roles":             roles,
        "snapshot_dl":       snapshot_dl,
        "active_directory":  active_directory,
        "ldap_providers":    ldap_providers,
        "all_certs":         all_certs,
        "dedup_ts":          dedup_ts,
        "io_ts":             io_ts,
    }

    # Pre-compute shared derived values — avoids repeated recalculation in
    # scoring, recommendations, and every sheet generation function.
    _ac = info.get("auditLogConfig")
    _ca = info.get("clusterAuditConfig")
    cd["audit_enabled"] = bool(
        info.get("auditLogEnabled")
        or (isinstance(_ac, dict) and _ac.get("enabled"))
        or (isinstance(_ca, dict) and _ca.get("enabled"))
    )
    cd["fk_status"] = _fk_status(cd)   # used by scores, recs, and all sheet fns

    cd["scores"]           = _compute_scores(cd, args.quick)
    cd["capacity_runway"]  = _capacity_runway(cd)       # must precede recommendations
    cd["recommendations"]  = _build_recommendations(cd)
    return cd


# ─────────────────────────────────────────────────────────────────────────────
# SCORING ENGINE
# ─────────────────────────────────────────────────────────────────────────────

def _score_protection(cd, quick):
    """0-100 based on backup success rate."""
    groups = cd["groups"]
    if not groups:
        return 100
    if quick:
        total   = len(groups)
        success = sum(
            1 for g in groups
            if (g.get("lastRun") or {}).get("localBackupInfo", {}).get("status", "")
               in ("kSuccess", "Succeeded", "kWarning")
        )
    else:
        all_runs = [r for runs in cd["group_runs"].values() for r in runs]
        if not all_runs:
            return 100
        total   = len(all_runs)
        success = sum(
            1 for r in all_runs
            if (r.get("localBackupInfo") or {}).get("status", "")
               in ("kSuccess", "Succeeded", "kWarning")
        )
    pct = (success / total * 100) if total else 100
    if pct >= SUCCESS_WARN_PCT: return 100
    if pct >= SUCCESS_CRIT_PCT: return int(70 + (pct - SUCCESS_CRIT_PCT) * 3)
    return max(0, int(pct))


def _score_sla(cd, quick):
    """0-100 based on SLA pass rate."""
    groups = cd["groups"]
    if not groups:
        return 100
    if quick:
        total = len(groups)
        viols = sum(1 for g in groups
                    if (g.get("lastRun") or {}).get("isSlaViolated", False))
    else:
        all_runs = [r for runs in cd["group_runs"].values() for r in runs]
        if not all_runs:
            return 100
        total = len(all_runs)
        viols = sum(1 for r in all_runs if r.get("isSlaViolated", False))
    pct_ok = ((total - viols) / total * 100) if total else 100
    if pct_ok >= 99: return 100
    if pct_ok >= 95: return 80
    if pct_ok >= 90: return 60
    if pct_ok >= 80: return 40
    return 20


def _score_infra(cd):
    """0-100 based on node health."""
    nodes = cd["nodes"]
    if not nodes:
        return 70
    healthy = sum(
        1 for n in nodes
        if ((n.get("nodeStatus") or {}).get("overallStatus", "") == "kNormal"
            or n.get("removalState") in (None, "kDontRemove"))
    )
    return int(healthy / len(nodes) * 100)


def _score_capacity(cd):
    """0-100 based on cluster capacity used %."""
    info   = cd["info"]
    usage  = ((info.get("stats") or {}).get("usagePerfStats") or {})
    usable = usage.get("physicalCapacityBytes") or 0
    used   = usage.get("totalPhysicalUsageBytes") or 0
    if not usable:
        return 70
    pct = used / usable * 100
    if pct <= CAPACITY_WARN_PCT:                return 100
    if pct <= CAPACITY_CRIT_PCT:                return int(100 - (pct - CAPACITY_WARN_PCT) * 2)
    return max(0, int(40 - (pct - CAPACITY_CRIT_PCT) * 2))


def _has_fortknox(vaults, policies):
    """Return True if FortKnox / RPaaS immutability is configured on this cluster."""
    fk_types = {"krpaas", "kfortknox", "kfort_knox", "krpaasarchival"}
    for v in vaults:
        vt = (v.get("vaultType") or "").lower()
        vn = (v.get("name") or "").lower()
        if any(t in vt for t in fk_types) or "fortknox" in vn:
            return True
    for p in policies:
        for t in ((p.get("remoteTargetPolicy") or {}).get("archivalTargets") or []):
            cfg = (t.get("archivalTargetConfig") or {})
            tt  = (cfg.get("targetType") or t.get("targetType") or "").lower()
            if any(x in tt for x in fk_types) or "fortknox" in tt:
                return True
    return False


def _fk_status(cd):
    """
    Return one of three states for FortKnox immutability:
      "active"  — FK configured AND a successful FK archival run exists within
                  the collection window (data is flowing)
      "idle"    — FK configured in vault/policy but no successful FK archival
                  found in run history (configured but not sending data)
      "none"    — FK not configured at all

    Detection strategy
    ──────────────────
    We CANNOT rely on lastRun.archivalInfo alone: that field is only populated
    when the most recent primary run included archival. Groups that run daily
    local backups but archive on a different schedule (e.g. weekly) will show
    empty archivalInfo in lastRun even when FK is working correctly.

    Primary source: group_runs (full run history, available in non-quick mode).
    Every run in the window is scanned for a successful archival result whose
    target type or name indicates FortKnox/RPaaS.

    Fallback (--quick mode): lastRun.archivalInfo on each group (less reliable
    — may miss groups whose last run was local-only).
    """
    vaults     = cd["vaults"]
    policies   = cd["policies"]
    groups     = cd["groups"]
    group_runs = cd.get("group_runs") or {}

    if not _has_fortknox(vaults, policies):
        return "none"

    fk_types = {"krpaas", "kfortknox", "kfort_knox", "krpaasarchival"}
    success_states = {"kSuccess", "Succeeded", "kSuccessWithWarning",
                      "SucceededWithWarning", "kWarning"}

    # Build set of FK vault names for name-based matching in archival results
    fk_vault_names = set()
    for v in vaults:
        vt = (v.get("vaultType") or "").lower()
        vn = (v.get("name") or "").lower()
        if any(t in vt for t in fk_types) or "fortknox" in vn:
            fk_vault_names.add(vn)
    for p in policies:
        for t in ((p.get("remoteTargetPolicy") or {}).get("archivalTargets") or []):
            cfg  = t.get("archivalTargetConfig") or {}
            name = (cfg.get("name") or cfg.get("vaultName")
                    or t.get("targetName") or "").lower()
            tt   = (cfg.get("targetType") or t.get("targetType") or "").lower()
            if any(x in tt for x in fk_types) or "fortknox" in tt:
                if name:
                    fk_vault_names.add(name)

    def _is_fk_success(run):
        """Return True if this run dict has a successful archival to an FK target."""
        arch_info = (run.get("archivalInfo") or {})
        for r in (arch_info.get("archivalTargetResults") or []):
            if r.get("status") not in success_states:
                continue
            # Match by archival target type
            ttype = (r.get("archivalTargetType")
                     or r.get("targetType") or "").lower()
            if any(x in ttype for x in fk_types) or "fortknox" in ttype:
                return True
            # Match by target name against known FK vault names
            tname = (r.get("targetName") or "").lower()
            if tname and (tname in fk_vault_names or "fortknox" in tname):
                return True
        return False

    # Primary: scan full run history (non-quick mode)
    for runs in group_runs.values():
        for run in runs:
            if _is_fk_success(run):
                return "active"

    # Fallback: check lastRun on each group (quick mode)
    for g in groups:
        if _is_fk_success(g.get("lastRun") or {}):
            return "active"

    return "idle"


def _score_security(cd):
    """0-100 security score across 6 dimensions (no FIPS — requires cluster rebuild).

    Cluster Encryption  15 pts — cluster-wide encryption at rest enabled
    SD Encryption       15 pts — all storage domains encrypted (0 if any unencrypted)
    Vault configured    15 pts — at least one archival target exists
    Replication         15 pts — at least one replication target exists
    FortKnox/Indelible  20 pts — 20=active, 10=configured-idle, 0=absent
    DataLock (WORM)     20 pts — 20=full Compliance, 15=full any mode,
                                  8=partial ≥50%, 3=partial <50%, 0=none
    Admin MFA penalty   -10 pts — deducted if any admin account lacks MFA
    """
    vaults   = cd["vaults"]
    policies = cd["policies"]
    groups   = cd["groups"]
    score    = 0

    # Cluster encryption (15 pts)
    cluster_enc, sd_enc = _sd_enc_status(cd)
    if cluster_enc:
        score += 15

    # Storage domain encryption (15 pts)
    if sd_enc is True:
        score += 15
    elif sd_enc is None and cluster_enc:
        score += 8   # No domain data but cluster enc is on — partial credit

    # Vault (15 pts)
    if vaults:
        score += 15

    # Replication (15 pts)
    if any((p.get("remoteTargetPolicy") or {}).get("replicationTargets")
           for p in policies):
        score += 15

    # FortKnox / Indelible archival (20 pts)
    fk = cd.get("fk_status") or _fk_status(cd)
    if fk == "active":
        score += 20
    elif fk == "idle":
        score += 10

    # DataLock coverage (20 pts)
    dl_status, _, covered, total = _cluster_datalock(policies, groups)
    if dl_status == "full_compliance":
        score += 20
    elif dl_status == "full_any":
        score += 15
    elif dl_status == "partial" and total > 0:
        pct = covered / total
        score += 8 if pct >= 0.5 else 3

    # Admin MFA penalty (-10 pts if any admin lacks MFA)
    admin_roles = {"COHESITY_ADMIN", "Admin", "kAdmin", "kSuperAdmin"}
    if any(
        isinstance(u, dict)
        and not (u.get("mfaEnabled") or u.get("isMfaEnabled"))
        and bool(set(u.get("roles") or []) & admin_roles)
        for u in (cd.get("users") or [])
    ):
        score = max(0, score - 10)

    return score


def _score_alerts(cd):
    """0-100: deducted for open critical and warning alerts."""
    alerts = cd["alerts"]
    crits  = sum(1 for a in alerts if a.get("severity") == "kCritical")
    warns  = sum(1 for a in alerts if a.get("severity") == "kWarning")
    return max(0, 100 - crits * 10 - warns * 2)


def _ransomware_score(cd):
    """Return (score 0-100, label, gaps) measuring ransomware recovery readiness.

    Dimensions:
      DataLock (WORM) coverage  30 pts
      FortKnox / indelible vault 25 pts
      Clean recovery window      20 pts  (oldest successful snapshot age)
      Quorum                     10 pts
      Admin MFA                  15 pts

    Labels: Strong (>=80) / Moderate (60-79) / High Risk (40-59) / Critical (<40)
    """
    gaps     = []
    policies = cd.get("policies") or []
    groups   = cd.get("groups")   or []

    # ── DataLock coverage (30 pts) ────────────────────────────────────────
    dl_status, dl_label, dl_covered, dl_total = _cluster_datalock(policies, groups)
    if dl_status == "full_compliance":
        dl_pts = 30
    elif dl_status == "full_any":
        dl_pts = 20
    elif dl_status == "partial":
        pct    = dl_covered / dl_total if dl_total else 0
        dl_pts = 15 if pct >= 0.5 else 8
        gaps.append(f"DataLock partial ({dl_label})")
    else:
        dl_pts = 0
        gaps.append("No DataLock/WORM protection on any group")

    # ── FortKnox / indelible vault (25 pts) ──────────────────────────────
    fk = cd.get("fk_status") or _fk_status(cd)
    if fk == "active":
        fk_pts = 25
    elif fk == "idle":
        fk_pts = 10
        gaps.append("FortKnox configured but not actively receiving data")
    else:
        fk_pts = 0
        gaps.append("No indelible/FortKnox vault configured")

    # ── Recovery window (20 pts) — oldest clean successful snapshot ───────
    _success_st = {"kSuccess", "Succeeded", "kSuccessWithWarning",
                   "SucceededWithWarning", "kWarning"}
    oldest_clean = None

    for run in (cd.get("v1_runs") or []):
        br = run.get("backupRun") or {}
        st = br.get("status", "")
        su = (br.get("stats") or {}).get("startTimeUsecs") or 0
        if st in _success_st and su:
            if oldest_clean is None or su < oldest_clean:
                oldest_clean = su

    if oldest_clean is None:
        for runs in (cd.get("group_runs") or {}).values():
            for run in runs:
                lb = run.get("localBackupInfo") or {}
                st = lb.get("status", "")
                su = lb.get("startTimeUsecs") or 0
                if st in _success_st and su:
                    if oldest_clean is None or su < oldest_clean:
                        oldest_clean = su

    if oldest_clean is None:
        rw_pts = 0
        gaps.append("No successful backup history found in lookback window")
    else:
        days_old = (cd.get("end_usecs", _now_usecs()) - oldest_clean) / (86_400 * 1_000_000)
        if days_old >= 30:   rw_pts = 20
        elif days_old >= 14: rw_pts = 15
        elif days_old >= 7:  rw_pts = 10
        elif days_old >= 2:  rw_pts = 5
        else:
            rw_pts = 2
            gaps.append("Recovery window is less than 2 days — limited clean restore points")

    # ── Quorum (10 pts) ───────────────────────────────────────────────────
    if cd.get("mode") == "direct":
        q_pts = 10   # Quorum is a Helios control-plane concept; N/A in direct mode
    elif cd.get("quorum_enabled"):
        q_pts = 10
    else:
        q_pts = 0
        gaps.append("Quorum not enabled (change-approval enforcement absent)")

    # ── Admin MFA (15 pts) ────────────────────────────────────────────────
    _admin_roles = {"COHESITY_ADMIN", "Admin", "kAdmin", "kSuperAdmin"}
    all_admins = [u for u in (cd.get("users") or [])
                  if isinstance(u, dict)
                  and bool(set(u.get("roles") or []) & _admin_roles)]
    if not all_admins:
        mfa_pts = 15   # no admin accounts returned — no penalty
    elif all(u.get("mfaEnabled") or u.get("isMfaEnabled") for u in all_admins):
        mfa_pts = 15
    elif any(u.get("mfaEnabled") or u.get("isMfaEnabled") for u in all_admins):
        n_mfa = sum(1 for u in all_admins
                    if u.get("mfaEnabled") or u.get("isMfaEnabled"))
        mfa_pts = 8
        gaps.append(f"Partial admin MFA ({n_mfa}/{len(all_admins)} admins enabled)")
    else:
        mfa_pts = 0
        gaps.append("No admin accounts have MFA enabled")

    score = dl_pts + fk_pts + rw_pts + q_pts + mfa_pts
    label = ("Strong"    if score >= 80 else
             "Moderate"  if score >= 60 else
             "High Risk" if score >= 40 else "Critical")
    return score, label, gaps


def _ransomware_rag_fill(score):
    """Fill color for ransomware readiness score cells."""
    if score >= 80: return _fill(LT_GREEN)
    if score >= 60: return _fill(YELLOW)
    if score >= 40: return _fill(ORANGE)
    return _fill(RED)


def _capacity_runway(cd):
    """Fit a linear regression to 30-day daily capacity time series.

    Returns a dict:
      days_to_80       int or None  — days until cluster reaches 80% used;
                                      None if stable/shrinking or data unavailable
      projected_80     date or None — calendar date of projected 80% threshold
      daily_growth_tb  float        — TB/day growth rate (may be negative = shrinking)
      data_quality     str          — "regression" | "insufficient" | "no_data"
    """
    import datetime as _dt

    ts_points = cd.get("cap_timeseries") or []
    valid = []
    for pt in ts_points:
        t = pt.get("timestampMsecs") or 0
        dv = pt.get("data") or {}
        v  = (dv.get("int64Value") or dv.get("doubleValue")
              or pt.get("value") or 0)
        if t and v:
            valid.append((t, int(v)))
    valid.sort(key=lambda x: x[0])

    info   = cd.get("info") or {}
    usage  = ((info.get("stats") or {}).get("usagePerfStats") or {})
    total_bytes = usage.get("physicalCapacityBytes") or 0
    used_bytes  = usage.get("totalPhysicalUsageBytes") or 0
    total_tb = total_bytes / 1e12
    used_tb  = used_bytes  / 1e12

    _no_data = {
        "days_to_80": None, "projected_80": None,
        "daily_growth_tb": 0.0, "data_quality": "no_data"
    }
    if total_tb <= 0:
        return _no_data
    if len(valid) < 3:
        _no_data["data_quality"] = "no_data" if not valid else "insufficient"
        return _no_data

    # Least-squares linear regression: x=days from first point, y=TB used
    t0 = valid[0][0]
    xs = [(t - t0) / (86_400 * 1_000) for t, _ in valid]
    ys = [v / 1e12                     for _, v in valid]
    n  = len(xs)
    sx = sum(xs); sy = sum(ys)
    sxx = sum(x * x for x in xs)
    sxy = sum(x * y for x, y in zip(xs, ys))
    denom = n * sxx - sx * sx
    slope = (n * sxy - sx * sy) / denom if denom else 0.0  # TB/day

    if slope <= 0:
        return {"days_to_80": None, "projected_80": None,
                "daily_growth_tb": round(slope, 6), "data_quality": "regression"}

    target_80_tb = total_tb * 0.80
    remaining_tb = max(0.0, target_80_tb - used_tb)
    days_to_80   = int(remaining_tb / slope)
    proj_date    = _dt.date.today() + _dt.timedelta(days=days_to_80)
    return {
        "days_to_80":      days_to_80,
        "projected_80":    proj_date,
        "daily_growth_tb": round(slope, 6),
        "data_quality":    "regression",
    }


def _group_risk_score(group, cd):
    """Return (score 0-100, level) for a single protection group.

    Higher score = lower risk (consistent with health scoring convention).

    Dimensions:
      Last run status   35 pts
      SLA compliance    25 pts
      RPO gap           25 pts
      DataLock          15 pts

    Levels: Low (>=75) / Medium (50-74) / High (25-49) / Critical (<25)
    """
    lr = group.get("lastRun") or {}
    lb = lr.get("localBackupInfo") or {}

    # Last run status (35 pts)
    status = lb.get("status", "")
    if status in ("kSuccess", "Succeeded"):                   run_pts = 35
    elif status in ("kWarning", "SucceededWithWarning"):      run_pts = 25
    elif status in ("kRunning", "Running"):                   run_pts = 20
    elif status in ("kSkipped", "Skipped"):                   run_pts = 10
    elif group.get("isPaused"):                               run_pts = 5
    else:                                                      run_pts = 0

    # SLA compliance (25 pts)
    sla_violated = bool(lb.get("isSlaViolated") or lr.get("isSlaViolated"))
    sla_pts = 0 if sla_violated else 25

    # RPO gap (25 pts)
    start_u = lb.get("startTimeUsecs") or 0
    if not start_u:
        rpo_pts = 0
    else:
        rpo_hrs = (cd.get("end_usecs", _now_usecs()) - start_u) / 3_600_000_000
        if   rpo_hrs <=  4: rpo_pts = 25
        elif rpo_hrs <=  8: rpo_pts = 20
        elif rpo_hrs <= 24: rpo_pts = 15
        elif rpo_hrs <= 48: rpo_pts =  8
        else:               rpo_pts =  0

    # DataLock (15 pts)
    policy_by_id = cd.get("policy_by_id") or {p.get("id"): p for p in (cd.get("policies") or []) if p.get("id")}
    policy = policy_by_id.get(group.get("policyId"), {})
    dl_str = _policy_datalock_str(policy).lower()
    if "compliance" in dl_str:                dl_pts = 15
    elif "administrative" in dl_str:          dl_pts = 10
    elif "fortknox" in dl_str or "indelible" in dl_str: dl_pts = 15
    else:                                     dl_pts = 0

    score = run_pts + sla_pts + rpo_pts + dl_pts
    level = ("Low"      if score >= 75 else
             "Medium"   if score >= 50 else
             "High"     if score >= 25 else "Critical")
    return score, level


def _audit_category(event):
    """Map an audit log event to a display category, or None for read-only ops."""
    entity = (event.get("entityType") or "").lower()
    action = (event.get("action") or event.get("type") or "").lower()

    _read_ops = {"view", "list", "get", "read", "login", "logout", "search"}
    if any(op in action for op in _read_ops):
        return None

    if "policy" in entity:                                    return "Policy Changes"
    if "protection job" in entity or "protectiongroup" in entity: return "Group Changes"
    if "user" in entity or "role" in entity:                  return "User Changes"
    if "vault" in entity or "externaltarget" in entity:       return "Vault/Security Changes"
    if "cluster" in entity or "settings" in entity:           return "Cluster Config Changes"
    return "Other"


def _audit_high_risk(event):
    """Return True if this audit event represents a high-risk operation."""
    entity  = (event.get("entityType") or "").lower()
    action  = (event.get("action")     or "").lower()
    details = str(event.get("details") or event.get("description") or "").lower()

    if "vault" in entity          and "delete" in action:  return True
    if "user" in entity           and "delete" in action:  return True
    if "encryption" in details    and "disable" in action: return True
    if ("admin" in details or "superadmin" in details) and (
            "grant" in action or "create" in action or "modify" in action): return True
    return False


def _compute_scores(cd, quick):
    s = {
        "protection": _score_protection(cd, quick),
        "sla":        _score_sla(cd, quick),
        "infra":      _score_infra(cd),
        "capacity":   _score_capacity(cd),
        "security":   _score_security(cd),
        "alerts":     _score_alerts(cd),
    }
    s["overall"] = round(sum(s[k] * _W[k] for k in _W), 1)
    s["grade"]   = grade(s["overall"])
    rw_score, rw_label, rw_gaps = _ransomware_score(cd)
    s["ransomware"]       = rw_score
    s["ransomware_label"] = rw_label
    s["ransomware_gaps"]  = rw_gaps
    return s


# ─────────────────────────────────────────────────────────────────────────────
# RECOMMENDATIONS ENGINE
# ─────────────────────────────────────────────────────────────────────────────

_PRIORITY_ORDER = {"CRITICAL": 0, "HIGH": 1, "MEDIUM": 2, "LOW": 3}

def _r(priority, category, finding, action, impact):
    return {"priority": priority, "category": category,
            "finding": finding, "action": action, "impact": impact}


def _build_recommendations(cd):
    recs     = []
    info     = cd["info"]
    nodes    = cd["nodes"]
    alerts   = cd["alerts"]
    groups   = cd["groups"]
    policies = cd["policies"]
    views    = cd["views"]
    vaults   = cd["vaults"]
    usage    = ((info.get("stats") or {}).get("usagePerfStats") or {})

    # ── Infrastructure ────────────────────────────────────────────────────────

    # ── Software lifecycle ────────────────────────────────────────────────────
    import datetime as _dt
    today = _dt.date.today()
    sw_ver = (cd["info"].get("clusterSoftwareVersion") or "").strip()
    sw_status, sw_eos_date, sw_label = _sw_eos(sw_ver)
    if sw_status == "eol":
        recs.append(_r("CRITICAL", "Software Lifecycle",
             f"Software version {sw_ver or '(unknown)'} is out of support — "
             "security patches and bug fixes are no longer available",
             "Upgrade to a supported release (7.1.2_u2 LTS, 7.2.x, or later) immediately",
             "Running unsupported software exposes the cluster to unpatched vulnerabilities"))
    elif sw_status == "in_support" and sw_eos_date:
        days_left = (sw_eos_date - today).days
        if days_left < SW_EOS_WARN_DAYS:
            recs.append(_r("HIGH", "Software Lifecycle",
                 f"Software version {sw_ver} reaches end of support in {days_left} days ({sw_eos_date})",
                 "Plan and schedule upgrade before EOS date to maintain support eligibility",
                 "Post-EOS clusters are ineligible for Cohesity technical support"))
        elif days_left < 365:
            recs.append(_r("MEDIUM", "Software Lifecycle",
                 f"Software version {sw_ver} EOS in {days_left} days ({sw_eos_date})",
                 "Begin upgrade planning — schedule upgrade within next quarter",
                 "Proactive planning avoids rushed upgrades near the EOS deadline"))

    # ── Hardware lifecycle ────────────────────────────────────────────────────
    hw_eol_past   = []
    hw_eol_soon   = []
    for n in cd["nodes"]:
        model = n.get("hardwareModel") or n.get("productModel") or ""
        eol   = _hw_eol(model)
        if not eol:
            continue
        days_left = (eol - today).days
        if days_left < 0 and model not in hw_eol_past:
            hw_eol_past.append(model)
        elif days_left <= HW_EOL_WARN_DAYS and model not in hw_eol_soon:
            hw_eol_soon.append(model)

    if hw_eol_past:
        models = ", ".join(sorted(set(hw_eol_past)))
        recs.append(_r("CRITICAL", "Hardware Lifecycle",
             f"Hardware model(s) {models} are past EOL — no longer eligible for manufacturer support",
             "Initiate hardware refresh or extended support contract immediately",
             "EOL hardware cannot be replaced under standard support; failure risk increases"))
    if hw_eol_soon:
        models = ", ".join(sorted(set(hw_eol_soon)))
        recs.append(_r("HIGH", "Hardware Lifecycle",
             f"Hardware model(s) {models} reaching EOL within {HW_EOL_WARN_DAYS} days",
             "Begin hardware refresh planning — lead times for Cohesity hardware can be 8–16 weeks",
             "Planning ahead avoids gaps in hardware support coverage"))

    unhealthy = [n for n in nodes
                 if not ((n.get("nodeStatus") or {}).get("overallStatus") == "kNormal"
                         or n.get("removalState") in (None, "kDontRemove"))]
    if unhealthy:
        recs.append(_r("CRITICAL", "Infrastructure",
            f"{len(unhealthy)} node(s) in unhealthy state",
            "Check hardware health, disk status, and cluster logs immediately",
            "Unhealthy nodes risk data availability and backup failures"))

    # ── Capacity ──────────────────────────────────────────────────────────────
    usable = usage.get("physicalCapacityBytes") or 0
    used   = usage.get("totalPhysicalUsageBytes") or 0
    if usable:
        cap_pct = used / usable * 100
        if cap_pct >= CAPACITY_CRIT_PCT:
            recs.append(_r("CRITICAL", "Capacity",
                f"Cluster at {cap_pct:.1f}% capacity — critically high",
                "Immediately expand capacity or archive/delete cold data",
                "Backup jobs fail when capacity is exhausted"))
        elif cap_pct >= CAPACITY_WARN_PCT:
            recs.append(_r("HIGH", "Capacity",
                f"Cluster at {cap_pct:.1f}% capacity — approaching limit",
                "Plan capacity expansion or data archival within 30-60 days",
                "Continued growth risks backup failures within weeks"))

    # Predictive capacity runway based on 30-day growth trend
    runway = cd.get("capacity_runway") or {}
    d80 = runway.get("days_to_80")
    if isinstance(d80, int):
        proj = runway.get("projected_80")
        proj_str = proj.isoformat() if proj else "soon"
        growth   = runway.get("daily_growth_tb", 0)
        if d80 < 30:
            recs.append(_r("CRITICAL", "Capacity",
                f"Cluster projected to reach 80% capacity in {d80} day(s) "
                f"(est. {proj_str} at {growth:.4f} TB/day growth rate)",
                "Expand cluster capacity immediately or accelerate data archival/deletion",
                "Backup jobs fail when storage is exhausted — less than 30 days "
                "to projected 80% threshold based on current growth trend"))
        elif d80 < 90:
            recs.append(_r("HIGH", "Capacity",
                f"Cluster projected to reach 80% capacity in {d80} day(s) "
                f"(est. {proj_str} at {growth:.4f} TB/day growth rate)",
                "Plan capacity expansion or data archival within the next 90 days",
                f"At the current growth rate of {growth:.4f} TB/day, storage will "
                "reach 80% utilization within 3 months — plan ahead to avoid disruption"))

    # ── Security ─────────────────────────────────────────────────────────────
    cluster_enc, sd_enc = _sd_enc_status(cd)
    if not cluster_enc:
        recs.append(_r("HIGH", "Security",
            "Cluster-wide encryption at rest is not enabled",
            "Enable cluster-wide encryption during cluster setup (requires rebuild if "
            "not done at creation time); enable storage domain encryption as an interim step",
            "Unencrypted data on physical media is exposed if drives are removed or stolen"))
    if sd_enc is False:
        recs.append(_r("HIGH", "Security",
            "One or more storage domains are not encrypted",
            "Enable encryption on all storage domains under cluster security settings",
            "Unencrypted storage domains expose data on physical media removal"))
    elif sd_enc is None and not cluster_enc:
        recs.append(_r("MEDIUM", "Security",
            "Unable to verify storage domain encryption — cluster encryption is off",
            "Review storage domain encryption settings and enable where missing",
            "Unencrypted storage domains expose data if media is removed"))

    if not vaults:
        recs.append(_r("MEDIUM", "Security",
            "No vault (archival target) is configured",
            "Configure at least one external vault for off-site retention",
            "No off-site copy of data — single point of failure"))

    fk_st = cd.get("fk_status") or _fk_status(cd)
    if fk_st == "none":
        recs.append(_r("MEDIUM", "Security",
            "No indelible/FortKnox archival target detected",
            "Configure FortKnox or WORM archival to protect against ransomware",
            "Without an indelible copy, all backups can be deleted or encrypted by attackers"))
    elif fk_st == "idle":
        recs.append(_r("HIGH", "Security",
            "FortKnox configured but no recent successful archival detected",
            "Verify protection group policies and schedules targeting FortKnox; "
            "check for failed archival runs and resolve blocking issues",
            "Indelible copy is configured but data is not flowing — "
            "ransomware protection gap despite vault being present"))

    dl_status, dl_label, dl_covered, dl_total = _cluster_datalock(policies, groups)
    if dl_status == "none":
        recs.append(_r("MEDIUM", "Security",
            "No protection groups have DataLock (WORM) enabled",
            "Enable DataLock on protection policies — prefer Compliance mode to "
            "prevent any user or admin from deleting or shortening backup snapshots",
            "Without DataLock, local backups can be deleted or modified by ransomware "
            "or a compromised admin; FortKnox alone does not protect local snapshots"))
    elif dl_status == "partial":
        ungrouped = dl_total - dl_covered
        recs.append(_r("MEDIUM", "Security",
            f"DataLock partially configured — {dl_covered}/{dl_total} groups protected; "
            f"{ungrouped} group(s) have no DataLock on their governing policy",
            "Review the Policy Audit tab and enable DataLock on policies covering "
            "the unprotected groups; Compliance mode provides the strongest protection",
            "Partial DataLock coverage leaves unprotected groups vulnerable to "
            "admin-level or ransomware deletion of local snapshots"))
    elif dl_status == "full_any":
        recs.append(_r("LOW", "Security",
            "DataLock covers all groups but some policies use Administrative mode — "
            "Compliance mode provides stronger protection",
            "Upgrade DataLock from Administrative to Compliance mode on all policies",
            "Administrative mode DataLock can be overridden by a Cohesity admin; "
            "Compliance mode is indelible even to admins"))

    has_repl = any((p.get("remoteTargetPolicy") or {}).get("replicationTargets")
                   for p in policies)
    if not has_repl:
        recs.append(_r("MEDIUM", "Resilience",
            "No policies have a replication target configured",
            "Add replication to a secondary cluster for site-level resilience",
            "Single-site data is vulnerable to site-wide outage"))

    # ── Alerts ────────────────────────────────────────────────────────────────
    crits = [a for a in alerts if a.get("severity") == "kCritical"]
    if crits:
        _now_u = cd["end_usecs"]
        old = [a for a in crits
               if (_now_u - (a.get("firstOccurrenceTime") or _now_u))
                  > ALERT_AGE_CRIT_DAYS * 86400 * 1_000_000]
        sev = "CRITICAL" if old else "HIGH"
        msg = f"{len(crits)} open critical alert(s)"
        if old:
            msg += f", {len(old)} older than {ALERT_AGE_CRIT_DAYS} days"
        recs.append(_r(sev, "Availability",
            msg,
            "Review and resolve open critical alerts; investigate root causes",
            "Unresolved critical alerts indicate active system issues"))

    # ── Protection ───────────────────────────────────────────────────────────
    paused = [g for g in groups if g.get("isPaused")]
    if paused:
        recs.append(_r("HIGH", "Protection",
            f"{len(paused)} protection group(s) are paused",
            f"Review and resume: {', '.join(g['name'] for g in paused[:3])}",
            "Paused groups do not create backups — data is unprotected"))

    rpo_gaps = []
    for g in groups:
        lr      = g.get("lastRun") or {}
        lb      = lr.get("localBackupInfo") or {}
        start_u = lb.get("startTimeUsecs") or 0
        if not start_u:
            continue
        hrs = (cd["end_usecs"] - start_u) / 3_600_000_000
        if hrs > RPO_CRIT_HOURS:
            rpo_gaps.append((g.get("name", "?"), hrs))
    if rpo_gaps:
        examples = ", ".join(f"{n} ({h:.0f}h)" for n, h in rpo_gaps[:3])
        recs.append(_r("HIGH", "Protection",
            f"{len(rpo_gaps)} group(s) with >{RPO_CRIT_HOURS}h RPO gap",
            f"Investigate: {examples}" + (" and more" if len(rpo_gaps) > 3 else ""),
            "Objects may be unrecoverable to a recent point in time"))

    # ── Policies ─────────────────────────────────────────────────────────────
    no_arch = [p for p in policies
               if not (p.get("remoteTargetPolicy") or {}).get("archivalTargets")]
    if policies and len(no_arch) > len(policies) * 0.5:
        recs.append(_r("MEDIUM", "Governance",
            f"{len(no_arch)}/{len(policies)} policies lack an archival target",
            "Add archival targets for long-term retention and ransomware recovery",
            "Without archival, data is unrecoverable after local retention expires"))

    # ── Views ────────────────────────────────────────────────────────────────
    no_quota = [v for v in views
                if not (v.get("logicalQuota") or
                        (v.get("storagePolicy") or {}).get("storageQuotaPolicy"))]
    if no_quota:
        recs.append(_r("LOW", "Governance",
            f"{len(no_quota)} NAS view(s) have no storage quota",
            "Apply quotas to all views to prevent uncontrolled capacity growth",
            "Unquoted views can silently consume all available capacity"))

    near_quota = []
    for v in views:
        lq  = v.get("logicalQuota") or \
              (v.get("storagePolicy") or {}).get("storageQuotaPolicy") or {}
        q   = lq.get("hardLimitBytes") or 0
        u   = ((v.get("stats") or {}).get("dataUsageStats") or {}) \
                .get("totalLogicalUsageBytes") or 0
        if q and u and (u / q * 100) >= VIEW_QUOTA_CRIT_PCT:
            near_quota.append(v.get("name", "?"))
    if near_quota:
        recs.append(_r("HIGH", "Capacity",
            f"{len(near_quota)} view(s) at ≥{VIEW_QUOTA_CRIT_PCT}% of quota",
            f"Expand quota or clean up: {', '.join(near_quota[:3])}",
            "Views at quota block all new writes to that share"))

    # ── TLS Certificate expiry ────────────────────────────────────────────────
    for cert in (cd.get("certs") or []):
        if not isinstance(cert, dict):
            continue
        exp_usecs = (cert.get("expiryUsecs")
                     or cert.get("expiryTimeMsecs")
                     or cert.get("notAfter"))
        if not exp_usecs:
            continue
        # Handle usecs (>1e15) vs msecs (>1e12) vs seconds
        if exp_usecs > 1_000_000_000_000_000:
            exp_dt = _dt.date.fromtimestamp(exp_usecs / 1_000_000)
        elif exp_usecs > 1_000_000_000_000:
            exp_dt = _dt.date.fromtimestamp(exp_usecs / 1_000)
        else:
            exp_dt = _dt.date.fromtimestamp(exp_usecs)
        days_left = (exp_dt - today).days
        cn = (cert.get("commonName") or cert.get("subject")
              or cert.get("certName") or "cluster TLS certificate")
        if days_left < 0:
            recs.append(_r("CRITICAL", "Certificate",
                f"TLS certificate '{cn}' expired {abs(days_left)} days ago",
                "Replace the expired certificate immediately to avoid browser/API errors",
                "Expired certificates break HTTPS access and API connectivity"))
        elif days_left < CERT_CRIT_DAYS:
            recs.append(_r("CRITICAL", "Certificate",
                f"TLS certificate '{cn}' expires in {days_left} days ({exp_dt})",
                "Renew or replace the certificate before expiry to prevent service disruption",
                "Expired certificates break HTTPS access and API connectivity"))
        elif days_left < CERT_WARN_DAYS:
            recs.append(_r("HIGH", "Certificate",
                f"TLS certificate '{cn}' expires in {days_left} days ({exp_dt})",
                "Plan certificate renewal now — lead time for CA-signed certs can be 2–4 weeks",
                "Certificate expiry causes browser warnings and may break automated API calls"))

    # ── Agent health — single pass for unhealthy + upgradable ────────────────
    unhealthy_agents  = []
    upgradable_agents = []
    for agent in (cd.get("agents") or []):
        _hs = agent.get("healthStatus")
        st  = ((isinstance(_hs, dict) and _hs.get("status"))
               or (isinstance(_hs, str) and _hs)
               or agent.get("agentStatus") or agent.get("status") or "")
        _rh, _ri = _split_host_ip(
            agent.get("hostName") or agent.get("host") or "",
            agent.get("hostIp") or "")
        host = _rh or _ri or "unknown"
        if st.lower() in ("kunhealthy", "unhealthy", "unreachable", "kfailed", "failed"):
            unhealthy_agents.append(host)
        upg = (agent.get("upgradability") or agent.get("upgradable") or "")
        if "upgradable" in upg.lower():
            upgradable_agents.append(host)
    if unhealthy_agents:
        examples = ", ".join(unhealthy_agents[:3])
        recs.append(_r("HIGH", "Agent Health",
            f"{len(unhealthy_agents)} agent(s) unhealthy or unreachable: {examples}"
            + (" and more" if len(unhealthy_agents) > 3 else ""),
            "Check agent connectivity, firewall rules, and reinstall if necessary",
            "Unhealthy agents silently prevent backups from running on affected hosts"))

    if upgradable_agents:
        examples = ", ".join(upgradable_agents[:3])
        recs.append(_r("MEDIUM", "Agent Health",
            f"{len(upgradable_agents)} agent(s) are running older versions and may be "
            f"upgraded: {examples}" + (" and more" if len(upgradable_agents) > 3 else ""),
            "Review agent compatibility matrix and upgrade agents where appropriate — "
            "see https://docs.cohesity.com/7_4/Web/UserGuide/Content/ReleaseNotes/"
            "agent-compatibility.htm for supported back-versions",
            "Outdated agents MAY cause issues with newer cluster features; Cohesity "
            "supports a limited range of back-versions — check compatibility docs"))

    # ── Disk health and SSD wear — single pass ────────────────────────────────
    failed_disks = []
    high_wear    = []
    for d in (cd.get("disks") or []):
        if (d.get("diskStatus") or d.get("status") or "").lower() \
                in ("kfailed", "failed", "kmissing", "missing"):
            failed_disks.append(d)
        wear = (d.get("ssdUsedPercentage") or d.get("ssdWearLevelPct")
                or d.get("wearLevel") or d.get("lifeUsedPct") or 0)
        if wear >= SSD_WEAR_CRIT_PCT:
            high_wear.append((d.get("_nodeIp", "?"), wear))
    if failed_disks:
        recs.append(_r("CRITICAL", "Disk Health",
            f"{len(failed_disks)} disk(s) in failed or missing state",
            "Replace failed disks immediately and verify cluster redundancy",
            "Failed disks reduce storage redundancy and risk data loss"))

    if high_wear:
        examples = ", ".join(f"{ip} ({w}%)" for ip, w in high_wear[:3])
        recs.append(_r("HIGH", "Disk Health",
            f"{len(high_wear)} SSD(s) with ≥{SSD_WEAR_CRIT_PCT}% wear: {examples}",
            "Plan SSD replacement before wear reaches 100% to avoid sudden failure",
            "Worn SSDs have elevated failure risk and may degrade cluster performance"))

    # ── Source coverage ───────────────────────────────────────────────────────
    low_coverage = []
    for src in (cd.get("sources") or []):
        _si = src.get("sourceInfo")
        si   = _si if isinstance(_si, dict) else src
        name = si.get("name") or si.get("sourceName") or src.get("name") or ""
        _src_stats = src.get("stats")
        prot   = (src.get("protectedObjectsCount")
                  or src.get("numProtectedObjects")
                  or (isinstance(_src_stats, dict) and _src_stats.get("protectedObjectsCount")) or 0)
        unprot = (src.get("unprotectedObjectsCount")
                  or src.get("numUnprotectedObjects")
                  or (isinstance(_src_stats, dict) and _src_stats.get("unprotectedObjectsCount")) or 0)
        total  = prot + unprot
        if total > 0 and unprot > 0 and (unprot / total * 100) >= 25:
            low_coverage.append(f"{name} ({unprot}/{total} unprotected)")
    if low_coverage:
        examples = ", ".join(low_coverage[:3])
        recs.append(_r("HIGH", "Coverage",
            f"{len(low_coverage)} source(s) have ≥25% unprotected objects: {examples}"
            + (" and more" if len(low_coverage) > 3 else ""),
            "Review unprotected objects in each source and add them to a protection group",
            "Unprotected objects are unrecoverable in the event of loss or corruption"))

    # ── User security ─────────────────────────────────────────────────────────
    import datetime as _dt_rec
    _now_rec = _dt_rec.datetime.now(tz=_dt_rec.timezone.utc)
    _STALE_DAYS_REC = 30
    admin_roles = {"COHESITY_ADMIN", "Admin", "kAdmin", "kSuperAdmin"}
    users_no_mfa = [
        u.get("username", "?") for u in (cd.get("users") or [])
        if isinstance(u, dict)
        and not (u.get("mfaEnabled") or u.get("isMfaEnabled"))
        and bool(set(u.get("roles") or []) & admin_roles)
    ]
    if users_no_mfa:
        examples = ", ".join(users_no_mfa[:3])
        recs.append(_r("HIGH", "Security",
            f"{len(users_no_mfa)} admin user(s) without MFA: {examples}"
            + (" and more" if len(users_no_mfa) > 3 else ""),
            "Enable MFA for all administrator accounts in Helios → Access Management",
            "Admin accounts without MFA are the most common ransomware entry point; "
            "a stolen credential gives full cluster access without MFA as a second factor"))

    # ── Stale privileged accounts ─────────────────────────────────────────────
    stale_admins = []
    for u in (cd.get("users") or []):
        if not isinstance(u, dict):
            continue
        user_roles_lower = {r.lower() for r in (u.get("roles") or [])}
        is_admin_rec = bool(
            user_roles_lower & {"cohesity_admin", "admin", "kadmin", "ksuperadmin",
                                "super_admin", "cluster_admin"}
            or any("admin" in r for r in user_roles_lower)
        )
        if not is_admin_rec:
            continue
        last_ms = u.get("lastLoginTimeMsecs") or u.get("lastLoginTime") or 0
        if last_ms:
            if last_ms > 1_000_000_000_000_000:
                last_ms = last_ms // 1000
            last_dt = _dt_rec.datetime.fromtimestamp(last_ms / 1000, tz=_dt_rec.timezone.utc)
            stale = (_now_rec - last_dt).days > _STALE_DAYS_REC
        else:
            stale = True
        if stale:
            stale_admins.append(u.get("username", "?"))
    if stale_admins:
        examples = ", ".join(stale_admins[:3])
        recs.append(_r("HIGH", "Security",
            f"{len(stale_admins)} stale privileged account(s) — no login in >{_STALE_DAYS_REC} days: "
            f"{examples}" + (" and more" if len(stale_admins) > 3 else ""),
            "Disable or remove admin accounts that have not logged in recently; "
            "review necessity and rotate credentials for retained accounts",
            "Stale admin accounts are high-value targets — unused credentials with "
            "admin privileges provide attackers full cluster access if compromised"))

    # ── Node software version skew ────────────────────────────────────────────
    cluster_sw = (cd["info"].get("clusterSoftwareVersion") or "").strip()
    mismatched_nodes = [
        n for n in cd["nodes"]
        if (n.get("nodeSoftwareVersion") or "").strip() != cluster_sw
        and (n.get("nodeSoftwareVersion") or "").strip()
    ]
    if mismatched_nodes:
        recs.append(_r("HIGH", "Infrastructure",
            f"{len(mismatched_nodes)} node(s) running a different software version than "
            f"the cluster ({cluster_sw}) — version skew detected",
            "Complete or re-run the cluster upgrade to bring all nodes to the same version; "
            "check upgrade logs for stuck or failed nodes",
            "Version skew indicates a partial or failed upgrade — cluster stability and "
            "feature parity cannot be guaranteed until all nodes match"))

    # ── Fault tolerance margin ────────────────────────────────────────────────
    _total_nodes = len(cd["nodes"])
    _healthy_nodes = sum(
        1 for n in cd["nodes"]
        if ((n.get("nodeStatus") or {}).get("overallStatus") == "kNormal"
            or n.get("removalState") in (None, "kDontRemove"))
    )
    if _total_nodes > 0:
        _tolerance = 2 if _total_nodes >= 6 else 1
        _failed_n  = _total_nodes - _healthy_nodes
        _margin    = _tolerance - _failed_n
        if _margin <= 0:
            recs.append(_r("CRITICAL", "Infrastructure",
                f"Cluster fault tolerance {'exceeded' if _margin < 0 else 'at limit'} — "
                f"{_failed_n} node(s) failed/unhealthy out of {_total_nodes} "
                f"(tolerance: {_tolerance})",
                "Restore failed nodes immediately; do not add workloads until "
                "fault tolerance margin is restored",
                "A further node failure will cause data unavailability or cluster outage"))

    # ── NAS share exposure ────────────────────────────────────────────────────
    high_risk_views = [
        v for v in (cd.get("views") or [])
        if v.get("enableSmbViewDiscovery") or
           (v.get("nfsMountPaths") and v.get("enableNfsViewDiscovery"))
    ]
    if high_risk_views:
        recs.append(_r("HIGH", "Security",
            f"{len(high_risk_views)} NAS view(s) with open discovery or uncontrolled "
            "S3 access — data exposure risk",
            "Disable SMB view discovery and NFS open mount exposure on views that "
            "do not require broad network visibility; enable quotas on all views",
            "Views with open discovery are visible to all network users — unintended "
            "data access and exfiltration risk increases significantly"))

    # ── Audit log ─────────────────────────────────────────────────────────────
    if not cd.get("audit_enabled"):
        recs.append(_r("MEDIUM", "Security",
            "Cluster audit logging is not enabled",
            "Enable audit logging in cluster security settings to capture all admin actions",
            "Without audit logs, unauthorized changes cannot be detected or investigated"))

    # ── SSO / IDP ─────────────────────────────────────────────────────────────
    if not cd.get("idps"):
        recs.append(_r("LOW", "Security",
            "No SSO / IDP is configured",
            "Configure SAML or LDAP-based SSO for centralized identity management",
            "Without SSO, local accounts with individual passwords increase credential risk"))

    # ── Recovery testing ─────────────────────────────────────────────────────
    _recov_list = cd.get("recoveries") or []
    _succ_recov = [
        r for r in _recov_list
        if any(s in (r.get("status") or "") for s in ("Succeeded", "kSuccess"))
    ]
    if not _recov_list:
        recs.append(_r("HIGH", "Recoverability",
            f"No recovery tasks found in the last {cd.get('days', 30)} days — "
            "recoverability is unproven",
            "Schedule a recovery test drill — test restore at least one critical workload "
            "per month to verify recoverability",
            "Without regular recovery testing, backup data integrity cannot be confirmed; "
            "untested restores fail at the worst possible moment"))
    elif not _succ_recov:
        recs.append(_r("HIGH", "Recoverability",
            f"No successful recovery found in the last {cd.get('days', 30)} days — "
            "all recovery attempts failed or are in progress",
            "Investigate failed recovery tasks and schedule a successful recovery test drill",
            "Failed recoveries indicate restore procedures or backup data may be unreliable"))

    # ── KMS configuration ─────────────────────────────────────────────────────
    _kms = cd.get("kms") or {}
    if _kms:
        _kms_type = (_kms.get("serverType") or _kms.get("type") or "").lower()
        if "internal" not in _kms_type and "builtin" not in _kms_type and "local" not in _kms_type:
            pass  # External KMS configured — no recommendation needed
        else:
            recs.append(_r("MEDIUM", "Security",
                "Encryption keys stored in internal KMS — no external key management configured",
                "Consider deploying an external KMS (KMIP-compatible or AWS KMS) for "
                "separation of duties between data custodian and key custodian",
                "Internal KMS means the encryption keys reside on the same system as the "
                "encrypted data — a compromised cluster admin can access both"))
    else:
        # Empty kms dict — internal or unknown
        cluster_enc = (cd["info"].get("encryptionEnabled") or
                       bool((cd["info"].get("encryptionConfig") or {}).get("encryptionEnabled")))
        if cluster_enc:
            recs.append(_r("MEDIUM", "Security",
                "Encryption keys stored internally — consider external KMS (KMIP/AWS) "
                "for separation of duties",
                "Deploy a KMIP-compatible external key management server or AWS KMS and "
                "configure it in Cohesity Security → Key Management",
                "Internal key storage means a compromised cluster admin can access both "
                "the encrypted data and the decryption keys"))

    # ── Custom roles ──────────────────────────────────────────────────────────
    _BUILTIN_ROLES = {"Admin", "Viewer", "BackupOperator", "BackupAdmin",
                      "RestoreOperator", "ReplicationOperator",
                      "COHESITY_ADMIN", "kAdmin", "kSuperAdmin"}
    _custom_roles = [
        r for r in (cd.get("roles") or [])
        if (r.get("isCustomRole") or r.get("label") == "Custom"
            or r.get("name") not in _BUILTIN_ROLES)
        and r.get("name") not in _BUILTIN_ROLES
    ]
    _high_risk_roles = []
    for _cr in _custom_roles:
        _privs = [str(p) for p in (_cr.get("privileges") or _cr.get("permissions") or [])]
        _privs_upper = " ".join(_privs).upper()
        _has_cluster_mod = any(
            "CLUSTER" in _p and any(x in _p for x in ("MODIFY", "CREATE", "DELETE"))
            for _p in _privs_upper.split()
        ) or ("CLUSTER" in _privs_upper and
              any(x in _privs_upper for x in ("MODIFY", "CREATE", "DELETE")))
        _has_user_mod = "USER" in _privs_upper and "MODIFY" in _privs_upper
        if _has_cluster_mod or _has_user_mod:
            _high_risk_roles.append(_cr.get("name", "?"))
    if _high_risk_roles:
        examples = ", ".join(_high_risk_roles[:3])
        recs.append(_r("HIGH", "Security",
            f"Custom roles with cluster or user modify privileges detected: {examples}"
            + (" and more" if len(_high_risk_roles) > 3 else ""),
            "Review all assignments of high-privilege custom roles; apply least-privilege "
            "principle and remove unnecessary cluster or user modify permissions",
            "Custom roles with modify privileges grant broad administrative capability — "
            "if assigned to compromised accounts, full cluster control is at risk"))

    # ── DataLock snapshot verification ────────────────────────────────────────
    _dl_not_locked = []
    for _entry in (cd.get("snapshot_dl") or []):
        _snaps = _entry.get("snapshots") or []
        if not _snaps:
            continue
        _locked_count = sum(
            1 for _s in _snaps
            if (_s.get("dataLockConstraints") or
                (isinstance(_s.get("wormProperties"), dict) and
                 _s["wormProperties"].get("isLocked")) or
                _s.get("retentionType") in ("kCompliance", "kAdministrative"))
        )
        if _locked_count == 0:
            _dl_not_locked.append(_entry.get("group_name", _entry.get("group_id", "?")))
    if _dl_not_locked:
        examples = ", ".join(_dl_not_locked[:3])
        recs.append(_r("CRITICAL", "Data Protection",
            f"DataLock policy configured but snapshots are not locked for "
            f"{len(_dl_not_locked)} group(s): {examples}"
            + (" and more" if len(_dl_not_locked) > 3 else ""),
            "Investigate why DataLock is not being applied to snapshots — check policy "
            "DataLock mode, cluster DataLock enablement, and StorageDomain WORM settings",
            "Immutability protection is not in effect — backups that appear DataLocked "
            "can still be deleted or modified, leaving no true indelible copy"))

    # ── Certificate expiry ────────────────────────────────────────────────────
    _now_ms = int(time.time() * 1000)
    _cert_src = cd.get("all_certs") or cd.get("certs") or []
    for _cert in _cert_src:
        _expiry_ms = (_cert.get("expiryDateMsecs") or _cert.get("notAfterMsecs") or 0)
        if not _expiry_ms:
            _expiry_usecs = _cert.get("expiryDateUsecs") or _cert.get("expirationDateUsecs") or 0
            if _expiry_usecs:
                _expiry_ms = _expiry_usecs // 1000
        if not _expiry_ms:
            continue
        _days = (_expiry_ms - _now_ms) / (1000 * 86400)
        _cname = _cert.get("name") or _cert.get("subject") or "Certificate"
        if _days < 0:
            recs.append(_r("CRITICAL", "Security",
                f"Certificate '{_cname}' has EXPIRED",
                "Renew certificate immediately to restore secure TLS communications",
                "Expired certificates cause connection failures and browser security warnings"))
        elif _days < 30:
            recs.append(_r("CRITICAL", "Security",
                f"Certificate '{_cname}' expires in {int(_days)} day(s)",
                "Renew certificate immediately — fewer than 30 days remaining",
                "Near-expiry certificates risk service disruptions if not renewed before cutover"))
        elif _days < 90:
            recs.append(_r("HIGH", "Security",
                f"Certificate '{_cname}' expires in {int(_days)} day(s)",
                "Schedule certificate renewal — fewer than 90 days remaining",
                "Proactive renewal avoids emergency changes and potential service outages"))

    # ── Active Directory / LDAP health ───────────────────────────────────────
    for _ad in (cd.get("active_directory") or []):
        _status = (_ad.get("connectionStatus") or _ad.get("status") or "").lower()
        _domain = _ad.get("domainName") or "Unknown Domain"
        if any(w in _status for w in ("disconnect", "error", "fail", "unreachable")):
            recs.append(_r("HIGH", "Identity",
                f"Active Directory domain '{_domain}' is not connected",
                "Check DC connectivity, DNS resolution, and Kerberos/time sync configuration",
                "AD disconnection breaks user authentication and RBAC for all AD-backed accounts"))
    for _lp in (cd.get("ldap_providers") or []):
        _status = (_lp.get("connectionStatus") or _lp.get("status") or "").lower()
        _lname  = _lp.get("name") or "LDAP Provider"
        if any(w in _status for w in ("disconnect", "error", "fail", "unreachable")):
            recs.append(_r("HIGH", "Identity",
                f"LDAP provider '{_lname}' is not connected",
                "Check LDAP server availability, port, and bind credentials",
                "LDAP disconnection breaks authentication for all LDAP-backed user accounts"))

    recs.sort(key=lambda x: _PRIORITY_ORDER.get(x["priority"], 99))
    return recs


# ─────────────────────────────────────────────────────────────────────────────
# EXCEL SHEETS
# ─────────────────────────────────────────────────────────────────────────────

def _success_stats(cd):
    """Return (succ_pct, sla_pct) from run history.

    succ_pct: average of per-day Success % (incl. Warnings) — identical to
              the average of column H in the Trends sheet, so the Executive
              Summary stays consistent with the Trends tab.
    sla_pct:  percentage of total runs without an SLA violation.

    Data source priority mirrors _sheet_trends():
      1. v1_runs  — Helios protectionRuns (preferred; same as Helios reports)
      2. group_runs — v2 per-group run history (fallback when v1 unavailable)
      3. lastRun per group — quick-mode last resort
    """
    # ── Build daily buckets (same logic as _sheet_trends) ────────────────────
    daily = {}
    for run in (cd.get("v1_runs") or []):
        br  = run.get("backupRun") or {}
        su  = (br.get("stats") or {}).get("startTimeUsecs") or 0
        if not su:
            continue
        day = datetime.datetime.fromtimestamp(
            su / 1_000_000, tz=datetime.timezone.utc).strftime("%Y-%m-%d")
        d  = daily.setdefault(day, dict(total=0, succ=0, warn=0, viols=0))
        st = br.get("status", "")
        d["total"] += 1
        if st == "kSuccess":   d["succ"]  += 1
        elif st == "kWarning": d["warn"]  += 1
        if br.get("slaViolated"): d["viols"] += 1

    if not daily:
        for runs in cd["group_runs"].values():
            for run in runs:
                lb = run.get("localBackupInfo") or {}
                su = lb.get("startTimeUsecs") or 0
                if not su:
                    continue
                day = datetime.datetime.fromtimestamp(
                    su / 1_000_000, tz=datetime.timezone.utc).strftime("%Y-%m-%d")
                d  = daily.setdefault(day, dict(total=0, succ=0, warn=0, viols=0))
                st = lb.get("status", "")
                d["total"] += 1
                if st in ("kSuccess", "Succeeded"): d["succ"]  += 1
                elif st == "kWarning":              d["warn"]  += 1
                if lb.get("isSlaViolated"): d["viols"] += 1

    # ── Average of per-day success % (= average of Trends col H) ─────────────
    if daily:
        day_pcts = [(d["succ"] + d["warn"]) / d["total"] * 100
                    for d in daily.values() if d["total"]]
        succ_pct    = round(sum(day_pcts) / len(day_pcts), 1) if day_pcts else 100
        total_runs  = sum(d["total"] for d in daily.values())
        total_viols = sum(d["viols"] for d in daily.values())
        sla_pct     = round((total_runs - total_viols) / total_runs * 100, 1) if total_runs else 100
        return succ_pct, sla_pct

    # ── quick-mode last resort: last-run per group ────────────────────────────
    groups = cd["groups"]
    if not groups:
        return (100, 100)
    total = len(groups)
    succ  = sum(1 for g in groups
                if (g.get("lastRun") or {}).get("localBackupInfo", {})
                   .get("status", "") in ("kSuccess", "Succeeded", "kWarning"))
    viols = sum(1 for g in groups
                if (g.get("lastRun") or {}).get("isSlaViolated", False))
    return (
        round(succ / total * 100, 1),
        round((total - viols) / total * 100, 1),
    )


def _cap_stats(cd):
    """
    Return (usable_bytes, used_bytes, cap_pct_or_NA) for a cluster.
    Primary source: usagePerfStats from /v1/public/cluster?fetchStats=true.
    Falls back to summing storageConsumedBytes across storage domains when the
    cluster-level stats field comes back null (happens on some versions/configs).
    Returns cap_pct as a float, or the string "N/A" when capacity is unknown.
    """
    info   = cd["info"]
    usage  = ((info.get("stats") or {}).get("usagePerfStats") or {})
    usable = usage.get("physicalCapacityBytes") or 0
    used   = usage.get("totalPhysicalUsageBytes") or 0

    if not usable:
        # Fallback: sum physical used across domains (no usable total available)
        domain_used = sum(
            (d.get("stats") or {}).get("storageConsumedBytes") or 0
            for d in cd.get("domains", [])
        )
        return 0, domain_used, "N/A"

    pct = round(used / usable * 100, 1)
    return usable, used, pct


def _sheet_summary(wb, all_data):
    ws = wb.create_sheet("Executive Summary")
    ws.freeze_panes = "A3"
    _title(ws, "Cohesity Health Check — Executive Summary", "O")

    cols = ["Cluster", "Software Version", "Node Count",
            "Health Score", "Grade",
            "30d Success %", "SLA Pass %", "Capacity Used %",
            "Open Criticals", "Open Warnings",
            "Security Score /100",
            "Ransomware Readiness", "Ransomware Score /100",
            "Capacity Runway",
            "Recommendations", "Top Finding"]
    _hdr(ws, 2, cols)

    for cd in all_data:
        info   = cd["info"]
        scores = cd["scores"]
        alerts = cd["alerts"]
        _, _, cap_pct = _cap_stats(cd)
        crits  = sum(1 for a in alerts if a.get("severity") == "kCritical")
        warns  = sum(1 for a in alerts if a.get("severity") == "kWarning")
        succ_pct, sla_pct = _success_stats(cd)
        top = cd["recommendations"][0]["finding"] if cd["recommendations"] else "None identified"

        # Capacity runway display
        runway  = cd.get("capacity_runway") or {}
        dq      = runway.get("data_quality", "no_data")
        d80     = runway.get("days_to_80")
        if dq in ("no_data", "insufficient"):
            runway_val = "Insufficient data"
        elif d80 is None:
            runway_val = "N/A — stable"
        else:
            runway_val = f"{d80}d"

        row = [cd["name"],
               info.get("clusterSoftwareVersion", ""), len(cd["nodes"]),
               scores["overall"], scores["grade"],
               succ_pct, sla_pct, cap_pct,
               crits, warns,
               scores["security"],
               scores.get("ransomware_label", "—"),
               scores.get("ransomware", "—"),
               runway_val,
               len(cd["recommendations"]), top]
        r = ws.max_row + 1
        for c, val in enumerate(row, 1):
            ws.cell(row=r, column=c, value=_safe_cell(val)).font = _FONT_NORMAL

        # RAG color on overall score + grade
        for col in (4, 5):
            ws.cell(row=r, column=col).fill = _rag_fill(scores["overall"])
            ws.cell(row=r, column=col).font = _font(bold=True)

        # Red crits
        if crits:
            ws.cell(row=r, column=9).fill = _fill(RED)
            ws.cell(row=r, column=9).font = _font(bold=True, color=WHITE)

        # Ransomware score RAG (col 13)
        rw_score = scores.get("ransomware")
        if isinstance(rw_score, int):
            ws.cell(row=r, column=13).fill = _ransomware_rag_fill(rw_score)
            ws.cell(row=r, column=13).font = _font(bold=True)
            ws.cell(row=r, column=12).fill = _ransomware_rag_fill(rw_score)

        # Capacity runway RAG (col 14)
        if isinstance(d80, int):
            rc = ws.cell(row=r, column=14)
            if d80 < 30:
                rc.fill = _fill(RED);    rc.font = _font(bold=True, color=WHITE)
            elif d80 < 90:
                rc.fill = _fill(ORANGE); rc.font = _font(bold=True)
            else:
                rc.fill = _fill(LT_GREEN)

    auto_fit_columns(ws)
    ws.column_dimensions[ws.cell(row=2, column=16).column_letter].width = 55


def _sheet_infrastructure(wb, all_data):
    ws = wb.create_sheet("Infrastructure")
    ws.freeze_panes = "A3"
    _title(ws, "Infrastructure — Cluster & Node Health", "Q")

    cols = ["Cluster", "Software Version", "SW Lifecycle Status", "SW EOS Date", "Days to EOS",
            "Cluster Type", "Cluster ID", "Node Count", "Healthy Nodes", "Disk Count",
            "Domain", "DNS Servers", "NTP Servers", "Timezone",
            "Cluster Encryption", "SD Encryption", "Helios Status", "Fault Tolerance"]
    _hdr(ws, 2, cols)

    for cd in all_data:
        info  = cd["info"]
        nodes = cd["nodes"]
        healthy = sum(
            1 for n in nodes
            if ((n.get("nodeStatus") or {}).get("overallStatus") == "kNormal"
                or n.get("removalState") in (None, "kDontRemove"))
        )
        # Disk count — v1 /public/nodes returns diskCount as an integer per node
        disks = sum(n.get("diskCount") or 0 for n in nodes)

        sw_ver = info.get("clusterSoftwareVersion", "")
        import datetime as _dt_mod
        sw_status, sw_eos_date, sw_label = _sw_eos(sw_ver)
        today = _dt_mod.date.today()
        days_to_eos = (sw_eos_date - today).days if sw_eos_date else ""
        eos_date_str = sw_eos_date.isoformat() if sw_eos_date else ("—" if sw_status == "current" else "Unknown")

        cluster_enc, sd_enc = _sd_enc_status(cd)
        dns  = ", ".join(info.get("dnsServerIps") or [])
        # NTP: v1 API returns ntpSettings.ntpServers or a flat ntpServers list
        _ntp_s  = info.get("ntpSettings")
        ntp_raw = ((isinstance(_ntp_s, dict) and (_ntp_s.get("ntpServers") or _ntp_s.get("ntpServersList")))
                   or info.get("ntpServers")
                   or info.get("ntpServerIps")
                   or [])
        if isinstance(ntp_raw, list):
            ntp = ", ".join(str(s) for s in ntp_raw) if ntp_raw else "Not configured"
        else:
            ntp = str(ntp_raw) if ntp_raw else "Not configured"
        domain_raw = (info.get("domainNames") or [""])[0]
        # domainNames returns the cluster FQDN (e.g. "mycluster.example.com").
        # Strip the cluster hostname prefix so only the domain portion is shown.
        _cluster_prefix = (info.get("name") or "").lower().rstrip(".")
        if _cluster_prefix and domain_raw.lower().startswith(_cluster_prefix + "."):
            domain = domain_raw[len(_cluster_prefix) + 1:]
        else:
            domain = domain_raw

        sd_enc_str = ("Yes" if sd_enc is True else "No" if sd_enc is False else "N/A")

        # Fault tolerance margin (col 18)
        total_nodes = len(nodes)
        if total_nodes > 0:
            tolerance = 2 if total_nodes >= 6 else 1
            failed_n  = total_nodes - healthy
            margin    = tolerance - failed_n
            if margin > 0:
                ft_str = f"OK (can lose {margin} more)"
            elif margin == 0:
                ft_str = "AT LIMIT"
            else:
                ft_str = "EXCEEDED"
        else:
            ft_str  = "N/A"
            margin  = 1   # default safe for coloring

        row = [cd["name"], sw_ver,
               sw_label, eos_date_str, days_to_eos,
               info.get("clusterType", ""), str(info.get("id", "")),
               len(nodes), healthy, disks if disks else "",
               domain, dns, ntp, info.get("timezone", ""),
               "Yes" if cluster_enc else "No",
               sd_enc_str,
               "Connected", ft_str]
        rn = ws.max_row + 1
        for c, val in enumerate(row, 1):
            ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL

        # color SW Lifecycle Status cell (col 3)
        sw_cell = ws.cell(row=rn, column=3)
        if sw_status == "eol":
            sw_cell.fill = _fill(RED);    sw_cell.font = _font(bold=True, color=WHITE)
        elif sw_status == "in_support" and isinstance(days_to_eos, int) and days_to_eos < SW_EOS_WARN_DAYS:
            sw_cell.fill = _fill(ORANGE); sw_cell.font = _font(bold=True)
        elif sw_status == "current":
            sw_cell.fill = _fill(LT_GREEN)

        # color Days to EOS cell (col 5)
        eos_cell = ws.cell(row=rn, column=5)
        if isinstance(days_to_eos, int):
            if days_to_eos < 0:
                eos_cell.fill = _fill(RED);    eos_cell.font = _font(bold=True, color=WHITE)
            elif days_to_eos < SW_EOS_WARN_DAYS:
                eos_cell.fill = _fill(ORANGE); eos_cell.font = _font(bold=True)

        # color Cluster Encryption (col 15) and SD Encryption (col 16)
        enc_cell = ws.cell(row=rn, column=15)
        enc_cell.fill = _fill(LT_GREEN) if cluster_enc else _fill(RED)
        enc_cell.font = _font(bold=True, color=WHITE if not cluster_enc else "000000")
        sd_cell = ws.cell(row=rn, column=16)
        if sd_enc is True:
            sd_cell.fill = _fill(LT_GREEN); sd_cell.font = _font(bold=True)
        elif sd_enc is False:
            sd_cell.fill = _fill(RED);      sd_cell.font = _font(bold=True, color=WHITE)
        else:
            sd_cell.fill = _fill(YELLOW);   sd_cell.font = _font(bold=True)

        # Node health — now col 9
        if healthy < len(nodes):
            ws.cell(row=rn, column=9).fill = _fill(ORANGE)
            ws.cell(row=rn, column=9).font = _font(bold=True)

        # Fault Tolerance (col 18)
        ft_cell = ws.cell(row=rn, column=18)
        if ft_str == "EXCEEDED":
            ft_cell.fill = _fill(RED);    ft_cell.font = _font(bold=True, color=WHITE)
        elif ft_str == "AT LIMIT":
            ft_cell.fill = _fill(YELLOW); ft_cell.font = _font(bold=True)
        elif ft_str.startswith("OK"):
            ft_cell.fill = _fill(LT_GREEN); ft_cell.font = _font(bold=True)

    auto_fit_columns(ws)


def _sheet_hardware(wb, all_data):
    import datetime as _dt
    ws = wb.create_sheet("Node Hardware")
    ws.freeze_panes = "A3"
    _title(ws, "Node Hardware Inventory — Models, Capacity & End-of-Life Status", "L")

    cols = ["Cluster", "Node ID", "IP Address", "Node Type", "Model",
            "Serial Number", "Node SW Version",
            "Raw Capacity (TB)", "Disk Count", "Storage Tiers",
            "HW EOL Date", "HW EOL Status", "SW Version Match"]
    _hdr(ws, 2, cols)

    today = _dt.date.today()

    for cd in all_data:
        cluster_sw_ver = (cd["info"].get("clusterSoftwareVersion") or "").strip()
        for n in cd["nodes"]:
            model   = n.get("hardwareModel") or n.get("productModel") or ""
            serial  = n.get("cohesityNodeSerial") or ""
            ntype   = n.get("nodeType") or ""
            nsw     = n.get("nodeSoftwareVersion") or ""
            cap_tb  = round((n.get("maxPhysicalCapacityBytes") or 0) / 1e12, 2)
            disks   = n.get("diskCount") or 0

            # Storage tiers summary e.g. "SSD:4, HDD:8"
            tiers = []
            for t in (n.get("diskCountByTier") or []):
                tier_name = t.get("storageTier", "")
                tier_cnt  = t.get("diskCount", 0)
                if tier_name and tier_cnt:
                    tiers.append(f"{tier_name}:{tier_cnt}")
            tier_str = ", ".join(tiers) or str(disks)

            eol_date = _hw_eol(model)
            if eol_date:
                days_left = (eol_date - today).days
                if days_left < 0:
                    hw_status = f"PAST EOL ({abs(days_left)}d ago)"
                elif days_left <= HW_EOL_WARN_DAYS:
                    hw_status = f"EOL in {days_left}d"
                else:
                    hw_status = f"In Service ({days_left}d remaining)"
                eol_str = eol_date.isoformat()
            else:
                hw_status = "In Service" if model else "Unknown"
                eol_str   = "—"

            # SW version match vs cluster version (col 13)
            nsw_stripped = nsw.strip()
            if not nsw_stripped:
                sw_match = "Unknown"
            elif nsw_stripped == cluster_sw_ver:
                sw_match = "Match"
            else:
                sw_match = f"MISMATCH ({nsw_stripped})"

            row = [cd["name"],
                   n.get("id", ""), n.get("ip", ""),
                   ntype, model, serial, nsw,
                   cap_tb, disks, tier_str,
                   eol_str, hw_status, sw_match]
            rn = ws.max_row + 1
            for c, val in enumerate(row, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL

            # color HW EOL Status cell (col 12)
            eol_cell = ws.cell(row=rn, column=12)
            if eol_date:
                days_left = (eol_date - today).days
                if days_left < 0:
                    eol_cell.fill = _fill(RED);    eol_cell.font = _font(bold=True, color=WHITE)
                elif days_left <= HW_EOL_WARN_DAYS:
                    eol_cell.fill = _fill(ORANGE); eol_cell.font = _font(bold=True)
                elif days_left <= 365:
                    eol_cell.fill = _fill(YELLOW); eol_cell.font = _font(bold=True)

            # color SW Version Match cell (col 13)
            sw_cell = ws.cell(row=rn, column=13)
            if sw_match == "Match":
                sw_cell.fill = _fill(LT_GREEN); sw_cell.font = _font(bold=True)
            elif sw_match.startswith("MISMATCH"):
                sw_cell.fill = _fill(RED);      sw_cell.font = _font(bold=True, color=WHITE)

    auto_fit_columns(ws)
    ws.column_dimensions["E"].width = 30  # Model


def _sheet_protection(wb, all_data):
    ws = wb.create_sheet("Protection Health")
    ws.freeze_panes = "A3"
    _title(ws, "Protection Health — Group Status & Run Metrics", "T")

    cols = ["Cluster", "Group Name", "Policy",
            "Active", "Paused",
            "Last Run Type", "Last Run Status", "SLA Violated",
            "Last Run Start", "Last Run End", "Duration",
            "Objects OK", "Objects Failed",
            "Logical (GB)", "Physical (GB)",
            "Replication Status", "Archival Status",
            "DataLock (WORM)",
            "RPO (hrs)", "Flag"]
    _hdr(ws, 2, cols)

    for cd in all_data:
        pm         = cd.get("policy_map", {})
        policy_by_id = cd.get("policy_by_id") or {p.get("id"): p for p in cd["policies"] if p.get("id")}
        for g in cd["groups"]:
            lr  = g.get("lastRun") or {}
            lb  = lr.get("localBackupInfo") or {}

            # runType and isSlaViolated sit on localBackupInfo (confirmed from runs API)
            run_type    = lb.get("runType") or lr.get("runType", "")
            sla_violated = lb.get("isSlaViolated") or lr.get("isSlaViolated") or False
            status      = lb.get("status", "")

            start_u = lb.get("startTimeUsecs") or 0
            end_u   = lb.get("endTimeUsecs") or 0
            dur = ""
            if start_u and end_u:
                s = int((end_u - start_u) / 1_000_000)
                h, m = divmod(s // 60, 60)
                dur = f"{h}h {m:02d}m"
            rpo_hrs = round((cd["end_usecs"] - start_u) / 3_600_000_000, 1) \
                      if start_u else ""

            # Object counts are directly on localBackupInfo (not nested in stats)
            objs_ok   = lb.get("successfulObjectsCount", "")
            objs_fail = lb.get("failedObjectsCount", "")
            # Sizes are under localSnapshotStats sub-dict
            snap = lb.get("localSnapshotStats") or {}
            logical_gb  = round(bytes_to_gb(snap.get("logicalSizeBytes") or 0), 2)
            physical_gb = round(bytes_to_gb(snap.get("bytesWritten") or 0), 2)

            # Replication / archival status from lastRun (same level as localBackupInfo)
            repl_st = arch_st = ""
            for ri in ((lr.get("replicationInfo") or {})
                       .get("replicationTargetResults") or []):
                repl_st = ri.get("status", ""); break
            for ai in ((lr.get("archivalInfo") or {})
                       .get("archivalTargetResults") or []):
                arch_st = ai.get("status", ""); break

            flag = "OK"
            if g.get("isPaused"):
                flag = "PAUSED"
            elif status in ("kFailure", "Failed"):
                flag = "FAILED"
            elif sla_violated:
                flag = "SLA VIOLATED"
            elif isinstance(rpo_hrs, float) and rpo_hrs > RPO_CRIT_HOURS:
                flag = f"RPO GAP {rpo_hrs:.0f}h"
            elif isinstance(rpo_hrs, float) and rpo_hrs > RPO_WARN_HOURS:
                flag = f"RPO WARN {rpo_hrs:.0f}h"

            # Policy name: prefer policyName on group; fall back to policy_map lookup
            policy_name = (g.get("policyName")
                           or pm.get(g.get("policyId"), "")
                           or g.get("policyId", ""))

            # DataLock: look up via the full policy object
            policy = policy_by_id.get(g.get("policyId"), {})
            dl_str = _policy_datalock_str(policy)

            row = [cd["name"], g.get("name", ""),
                   policy_name,
                   "Yes" if g.get("isActive", True) else "No",
                   "Yes" if g.get("isPaused", False) else "No",
                   run_type, status,
                   "Yes" if sla_violated else "No",
                   usecs_to_datetime(start_u), usecs_to_datetime(end_u), dur,
                   objs_ok, objs_fail,
                   logical_gb, physical_gb,
                   repl_st, arch_st,
                   dl_str,
                   rpo_hrs, flag]
            rn = ws.max_row + 1
            for c, val in enumerate(row, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL

            # DataLock col 18 — color by mode
            dc = ws.cell(row=rn, column=18)
            dl_lower = dl_str.lower()
            if "compliance" in dl_lower:
                dc.fill = _fill(LT_GREEN)
            elif "administrative" in dl_lower:
                dc.fill = _fill(YELLOW); dc.font = _font(bold=True)
            else:
                dc.fill = _fill(RED); dc.font = _font(bold=True, color=WHITE)

            # Flag col 20
            fc = ws.cell(row=rn, column=20)
            if flag == "FAILED":
                fc.fill = _fill(RED);   fc.font = _font(bold=True, color=WHITE)
            elif flag in ("PAUSED",) or "RPO GAP" in flag:
                fc.fill = _fill(ORANGE); fc.font = _font(bold=True)
            elif flag == "SLA VIOLATED" or "RPO WARN" in flag:
                fc.fill = _fill(YELLOW); fc.font = _font(bold=True)
            else:
                fc.fill = _fill(LT_GREEN)

    auto_fit_columns(ws)


def _sheet_storage(wb, all_data):
    ws = wb.create_sheet("Storage & Capacity")
    ws.freeze_panes = "A3"
    _title(ws, "Storage & Capacity — Utilization & Data Reduction", "O")

    cols = ["Cluster", "Storage Domain", "Usable (TB)", "Used (TB)",
            "Free (TB)", "Used %", "Logical Data (TB)",
            "Data Reduction Ratio", "Dedup Ratio", "Compression Ratio",
            "Savings (TB)", "Status",
            "Runway (days to 80%)", "Est. 80% Full Date", "Daily Growth (TB/d)"]
    _hdr(ws, 2, cols)

    for cd in all_data:
        info   = cd["info"]
        stats  = (info.get("stats") or {})
        usage  = (stats.get("usagePerfStats") or {})
        data   = (stats.get("dataUsageStats") or {})

        usable, used, cap_pct = _cap_stats(cd)
        logical     = usage.get("dataInBytes") or 0
        physical    = usage.get("dataInBytesAfterReduction") or 0
        dedup_after = data.get("dataInBytesAfterDedup") or 0
        dr_ratio    = round((stats.get("dataReductionRatio") or 0.0), 2)
        dedup_ratio = round(logical / dedup_after, 2) if dedup_after else 0
        comp_ratio  = round(dedup_after / physical, 2) if physical else 0
        free        = usable - used
        savings     = round(bytes_to_tb(logical - physical), 2) \
                      if logical > physical else 0

        flag = ""
        if isinstance(cap_pct, float):
            flag = ("CRITICAL" if cap_pct >= CAPACITY_CRIT_PCT
                    else "WARN" if cap_pct >= CAPACITY_WARN_PCT else "OK")
        elif cap_pct == "N/A":
            flag = "N/A"

        # Capacity runway columns
        runway     = cd.get("capacity_runway") or {}
        d80        = runway.get("days_to_80")
        proj_80    = runway.get("projected_80")
        daily_gr   = runway.get("daily_growth_tb")
        dq         = runway.get("data_quality", "no_data")
        if dq == "no_data":
            runway_str = "N/A"
        elif dq == "insufficient":
            runway_str = "Insufficient data"
        elif d80 is None:
            runway_str = "N/A — stable"
        elif d80 == 0:
            runway_str = "Already >80%"
        else:
            runway_str = d80
        proj_str   = proj_80.isoformat() if proj_80 else ""
        growth_str = round(daily_gr, 4) if daily_gr is not None else ""

        row = [cd["name"], "(Cluster Total)",
               round(bytes_to_tb(usable), 2) or "",
               round(bytes_to_tb(used), 2) or "",
               round(bytes_to_tb(free), 2) or "",
               cap_pct,
               round(bytes_to_tb(logical), 2),
               dr_ratio, dedup_ratio, comp_ratio, savings, flag,
               runway_str, proj_str, growth_str]
        rn = ws.max_row + 1
        for c, val in enumerate(row, 1):
            cell = ws.cell(row=rn, column=c, value=_safe_cell(val))
            cell.font = _font(bold=(c == 2))

        sc = ws.cell(row=rn, column=12)
        if flag == "CRITICAL":
            sc.fill = _fill(RED);    sc.font = _font(bold=True, color=WHITE)
        elif flag == "WARN":
            sc.fill = _fill(ORANGE); sc.font = _font(bold=True)
        elif flag == "OK":
            sc.fill = _fill(LT_GREEN)

        # Runway col 13 coloring
        if isinstance(d80, int):
            rc = ws.cell(row=rn, column=13)
            if d80 == 0:
                rc.fill = _fill(RED);    rc.font = _font(bold=True, color=WHITE)
            elif d80 < 30:
                rc.fill = _fill(RED);    rc.font = _font(bold=True, color=WHITE)
            elif d80 < 90:
                rc.fill = _fill(ORANGE); rc.font = _font(bold=True)
            else:
                rc.fill = _fill(LT_GREEN); rc.font = _font(bold=True)

        # Per-domain rows
        for d in cd["domains"]:
            ds = d.get("stats") or {}
            ld = ds.get("dataInBytes") or 0
            da = ds.get("dataInBytesAfterDedup") or 0
            ph = ds.get("storageConsumedBytes") or 0
            drow = [cd["name"], d.get("name", ""),
                    "", "", "", "",
                    round(bytes_to_tb(ld), 2),
                    round(ld / ph, 2) if ph else 0,
                    round(ld / da, 2) if da else 0,
                    round(da / ph, 2) if ph else 0,
                    round(bytes_to_tb(ld - ph), 2) if ld > ph else 0, "",
                    "", "", ""]
            rn2 = ws.max_row + 1
            for c, val in enumerate(drow, 1):
                ws.cell(row=rn2, column=c, value=_safe_cell(val)).font = _FONT_NORMAL

    auto_fit_columns(ws)


def _ret_str(r):
    if not r: return ""
    d = r.get("duration", "")
    u = r.get("unit", "")
    return f"{d} {u}".strip() if d else u

def _sched_str(s):
    """Convert a v2 schedule dict to a human-readable string."""
    if not s: return ""
    unit = s.get("unit", "")
    if unit == "Days":
        freq = (s.get("daySchedule") or {}).get("frequency", 1)
        return f"Every {freq} day{'s' if freq != 1 else ''}"
    if unit == "Weeks":
        days = (s.get("weekSchedule") or {}).get("dayOfWeek") or []
        return f"Weekly ({', '.join(days[:2])}{'…' if len(days) > 2 else ''})" if days else "Weekly"
    if unit == "Hours":
        freq = (s.get("hourSchedule") or {}).get("frequency", 1)
        return f"Every {freq}h"
    if unit == "Minutes":
        freq = (s.get("minuteSchedule") or {}).get("frequency", 1)
        return f"Every {freq}m"
    if unit in ("Months", "Monthly"): return "Monthly"
    if unit in ("Years", "Yearly"):   return "Yearly"
    # fallback: flat frequency dict
    for key in ("minutes", "hours", "days", "weeks", "months", "years"):
        if isinstance(s, dict) and key in s:
            return f"Every {s[key]} {key}"
    return unit or ""

def _policy_datalock_str(p):
    """Return a display string for a single policy's DataLock / indelible config.

    Priority:
      1. Explicit local DataLock (backupPolicy.regular.retention.dataLockConfig).
         Returns e.g. "Compliance / 14 Days" or "Administrative / 30 Days".
      2. FortKnox / RPaaS archival target — indelible by design.
         Returns "FortKnox (Indelible)".
      3. Neither configured → "None".
    """
    bp  = p.get("backupPolicy") or {}
    reg = bp.get("regular") or {}
    ret = reg.get("retention") or bp.get("retention") or {}
    dlc = ret.get("dataLockConfig") if isinstance(ret, dict) else None
    if isinstance(dlc, dict) and dlc.get("mode"):
        mode = dlc["mode"]
        dur  = dlc.get("duration")
        unit = dlc.get("unit", "Days")
        return f"{mode} / {dur} {unit}" if dur else mode
    # No explicit WORM — check for FortKnox/RPaaS archival (implicit indelible)
    _fk = {"krpaas", "kfortknox", "kfort_knox", "krpaasarchival"}
    for t in ((p.get("remoteTargetPolicy") or {}).get("archivalTargets") or []):
        cfg   = t.get("archivalTargetConfig") or {}
        tt    = (cfg.get("targetType") or t.get("targetType") or "").lower()
        tname = (t.get("targetName") or cfg.get("name") or "").lower()
        if any(x in tt for x in _fk) or "fortknox" in tt or "fortknox" in tname:
            return "FortKnox (Indelible)"
    return "None"


def _sd_enc_status(cd):
    """Return (cluster_enc, sd_enc) booleans for encryption state.

    cluster_enc: True if cluster-wide encryption is enabled.
    sd_enc:      True if ALL storage domains have a non-empty encryptionType.
                 None if no domain data is available (unknown).
    """
    info = cd["info"]
    _enc_cfg = info.get("encryptionConfig")
    cluster_enc = bool(info.get("encryptionEnabled") or
                       (isinstance(_enc_cfg, dict) and _enc_cfg.get("encryptionEnabled")))
    domains = [d for d in (cd.get("domains") or []) if isinstance(d, dict)]
    if not domains:
        return cluster_enc, None
    sd_enc = all(
        bool((d.get("storagePolicy") or {}).get("encryptionType", ""))
        for d in domains
    )
    return cluster_enc, sd_enc


def _cluster_datalock(policies, groups=None):
    """Return (status, label, covered, total) for cluster-level DataLock coverage.

    Computes coverage at the *protection group* level: a group is covered if its
    governing policy has DataLock enabled (local retention or archival).

    status: "full_compliance" | "full_any" | "partial" | "none"
    label:  human-readable string, e.g. "Full Compliance" or "Partial — 8/15 groups"
    covered: number of groups with DataLock
    total:   total number of groups
    """
    # Build policy_id → datalock mode map
    policy_dl = {}   # policy id → "compliance" | "administrative" | "none"
    for p in policies:
        pid = p.get("id")
        if not pid:
            continue
        # Check local retention
        s = _policy_datalock_str(p).lower()
        mode = ("compliance" if "compliance" in s
                else "administrative" if "admin" in s else "none")
        # Check archival target retention (upgrade to compliance if found)
        if mode != "compliance":
            for t in ((p.get("remoteTargetPolicy") or {}).get("archivalTargets") or []):
                dlc = (t.get("retention") or {}).get("dataLockConfig")
                if isinstance(dlc, dict):
                    m = (dlc.get("mode") or "").lower()
                    if "compliance" in m:
                        mode = "compliance"; break
                    elif "admin" in m and mode == "none":
                        mode = "administrative"
        policy_dl[pid] = mode

    # Count group coverage
    total   = len(groups) if groups else 0
    covered_comp = 0
    covered_any  = 0
    if groups:
        for g in groups:
            pid  = g.get("policyId")
            mode = policy_dl.get(pid, "none")
            if mode == "compliance":
                covered_comp += 1
                covered_any  += 1
            elif mode == "administrative":
                covered_any += 1
    covered = covered_any

    if total == 0:
        # No active protection groups — cannot determine coverage.
        return ("none", "No", 0, 0)

    if covered == 0:
        return ("none", "No", 0, total)
    if covered == total and covered_comp == total:
        return ("full_compliance", "Full Compliance", covered, total)
    if covered == total:
        return ("full_any", f"Full ({covered_comp}/{total} Compliance)", covered, total)
    # Partial coverage
    return ("partial", f"Partial ({covered}/{total} groups)", covered, total)


def _sheet_policies(wb, all_data):
    ws = wb.create_sheet("Policy Audit")
    ws.freeze_panes = "A3"
    _title(ws, "Policy Audit — Retention, Replication & Archival", "P")

    cols = ["Cluster", "Policy Name", "Type",
            "Full Schedule", "Incremental Schedule",
            "Local Retention", "Log Retention",
            "Has Replication", "Replication Target", "Replication Retention",
            "Has Archival", "Archival Target", "Archival Retention",
            "Groups Using", "DataLock (WORM)", "Gaps"]
    _hdr(ws, 2, cols)

    for cd in all_data:
        # Count protection groups per policy ID so column N is populated
        policy_group_count: dict = {}
        for g in cd["groups"]:
            pid = g.get("policyId")
            if pid:
                policy_group_count[pid] = policy_group_count.get(pid, 0) + 1

        for p in cd["policies"]:
            bp  = p.get("backupPolicy") or {}
            rtp = p.get("remoteTargetPolicy") or {}
            rep_tgts  = rtp.get("replicationTargets") or []
            arch_tgts = rtp.get("archivalTargets") or []

            # v2 API replication target name: targetName or remoteTargetConfig.clusterName/name
            def _rep_name(t):
                return (t.get("targetName")
                        or (t.get("remoteTargetConfig") or {}).get("clusterName")
                        or (t.get("remoteTargetConfig") or {}).get("name")
                        or "")
            # v2 API archival target name: targetName or archivalTargetConfig.name/vaultName
            def _arch_name(t):
                return (t.get("targetName")
                        or (t.get("archivalTargetConfig") or {}).get("name")
                        or (t.get("archivalTargetConfig") or {}).get("vaultName")
                        or "")

            rep_names  = ", ".join(filter(None, (_rep_name(t)  for t in rep_tgts)))
            arch_names = ", ".join(filter(None, (_arch_name(t) for t in arch_tgts)))

            # v2 schedule: backupPolicy.regular.full.schedule / regular.incremental.schedule
            regular     = bp.get("regular") or {}
            full_sched  = (regular.get("full") or bp.get("full") or {}).get("schedule") or {}
            incr_sched  = (regular.get("incremental") or bp.get("incremental") or {}).get("schedule") or {}

            local_ret = regular.get("retention") or bp.get("retention") or {}
            log_ret   = (bp.get("log") or {}).get("retention") or {}
            rep_ret   = (rep_tgts[0].get("retention") or {}) if rep_tgts else {}
            arch_ret  = (arch_tgts[0].get("retention") or {}) if arch_tgts else {}

            gaps = []
            if not rep_tgts:  gaps.append("No replication")
            if not arch_tgts: gaps.append("No archival")

            dl_str = _policy_datalock_str(p)
            dl_low = dl_str.lower()
            if dl_str == "None":
                gaps.append("No DataLock")

            row = [cd["name"], p.get("name", ""),
                   p.get("type") or p.get("policyCategory") or "",
                   _sched_str(full_sched), _sched_str(incr_sched),
                   _ret_str(local_ret), _ret_str(log_ret),
                   "Yes" if rep_tgts else "No", rep_names, _ret_str(rep_ret),
                   "Yes" if arch_tgts else "No", arch_names, _ret_str(arch_ret),
                   policy_group_count.get(p.get("id"), 0),
                   dl_str,
                   "; ".join(gaps) if gaps else "OK"]
            rn = ws.max_row + 1
            for c, val in enumerate(row, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL

            # DataLock col 15 — color by mode
            dl_cell = ws.cell(row=rn, column=15)
            if "compliance" in dl_low:
                dl_cell.fill = _fill(LT_GREEN); dl_cell.font = _font(bold=True)
            elif "admin" in dl_low:
                dl_cell.fill = _fill(YELLOW);   dl_cell.font = _font(bold=True)
            elif "fortknox" in dl_low:
                dl_cell.fill = _fill(LT_GREEN); dl_cell.font = _font(bold=True)
            else:
                dl_cell.fill = _fill(RED);      dl_cell.font = _font(bold=True, color=WHITE)

            # Gaps col 16
            gc = ws.cell(row=rn, column=16)
            if gaps:
                gc.fill = _fill(YELLOW)
                gc.font = _font(bold=True)

    auto_fit_columns(ws)


def _sheet_policy_groups(wb, all_data):
    """Sheet: one row per protection group showing which policy governs it."""
    ws = wb.create_sheet("Policy → Groups")
    ws.freeze_panes = "A3"
    _title(ws, "Policy → Groups — Protection Groups Associated with Each Policy", "G")

    cols = ["Cluster", "Policy Name", "Group Name", "Environment",
            "Active", "Paused", "Last Run Status"]
    _hdr(ws, 2, cols)

    for cd in all_data:
        pm = cd.get("policy_map", {})
        # Sort by policy name then group name for easy reading
        sorted_groups = sorted(
            cd["groups"],
            key=lambda g: (
                g.get("policyName") or pm.get(g.get("policyId"), "") or "",
                g.get("name", "")
            )
        )
        for g in sorted_groups:
            policy_name = (g.get("policyName")
                           or pm.get(g.get("policyId"), "")
                           or g.get("policyId", ""))
            lb     = (g.get("lastRun") or {}).get("localBackupInfo") or {}
            status = lb.get("status", "")
            paused = g.get("isPaused", False)

            row = [cd["name"], policy_name, g.get("name", ""),
                   g.get("environment", ""),
                   "Yes" if g.get("isActive", True) else "No",
                   "Yes" if paused else "No",
                   status]
            rn = ws.max_row + 1
            for c, val in enumerate(row, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL

            # Highlight paused groups
            if paused:
                ws.cell(row=rn, column=6).fill = _fill(YELLOW)
                ws.cell(row=rn, column=6).font = _font(bold=True)
            # Highlight failed status
            if status in ("kFailure", "Failed"):
                ws.cell(row=rn, column=7).fill = _fill(RED)
                ws.cell(row=rn, column=7).font = _font(bold=True, color=WHITE)

    auto_fit_columns(ws)


def _sheet_alerts(wb, all_data):
    ws = wb.create_sheet("Alerts")
    ws.freeze_panes = "A3"
    _title(ws, "Open Alerts — Severity, Age & Description", "J")

    cols = ["Cluster", "Severity", "Category", "Alert Code",
            "Alert Name", "Description",
            "First Occurred", "Age (Days)", "Last Updated",
            "Acknowledged"]
    _hdr(ws, 2, cols)

    sev_order = {"kCritical": 0, "kWarning": 1, "kInfo": 2}

    for cd in all_data:
        sorted_a = sorted(cd["alerts"],
                          key=lambda a: (sev_order.get(a.get("severity"), 9),
                                         a.get("firstTimestampUsecs") or 0))
        for a in sorted_a:
            # v1 alerts API uses firstTimestampUsecs (not firstOccurrenceTime)
            first_u = a.get("firstTimestampUsecs") or 0
            last_u  = a.get("latestTimestampUsecs") or 0
            age = round((cd["end_usecs"] - first_u) / (86400 * 1_000_000), 1) \
                  if first_u else ""
            doc  = a.get("alertDocument") or {}
            desc = doc.get("alertDescription") or ""

            row = [cd["name"], a.get("severity", ""),
                   a.get("alertCategory", ""), a.get("alertCode", ""),
                   doc.get("alertName", ""),
                   desc[:120] + ("…" if len(desc) > 120 else ""),
                   usecs_to_datetime(first_u), age, usecs_to_datetime(last_u),
                   "Yes" if a.get("acknowledgeTimestampUsecs") else "No"]
            rn = ws.max_row + 1
            for c, val in enumerate(row, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL

            sc = ws.cell(row=rn, column=2)
            if a.get("severity") == "kCritical":
                sc.fill = _fill(RED);    sc.font = _font(bold=True, color=WHITE)
            elif a.get("severity") == "kWarning":
                sc.fill = _fill(ORANGE); sc.font = _font(bold=True)

    auto_fit_columns(ws)
    ws.column_dimensions["F"].width = 70


def _sheet_security(wb, all_data):
    import datetime as _dt
    ws = wb.create_sheet("Security")
    ws.freeze_panes = "A3"
    _title(ws, "Security Posture — Checklist vs Best Practices", "U")

    cols = ["Cluster", "Cluster Encryption", "SD Encryption",
            "Vault Configured", "FortKnox / Indelible",
            "Replication", "Archival",
            "Min Local Retention", "Min Archival Retention",
            "Security Score /100",
            "Audit Log", "NTP Auth", "Remote Tunnel",
            "Cluster MFA", "SSO / IDP",
            "Quorum",
            "DataLock (WORM)",
            "Soonest Cert Expiry",
            "Ransomware Score /100",
            "Ransomware Readiness",
            "Gaps"]
    _hdr(ws, 2, cols)

    _RET_MULT = {"Days": 1, "Weeks": 7, "Months": 30, "Years": 365}

    def _ret_to_days(r):
        d = (r or {}).get("duration") or 0
        u = (r or {}).get("unit", "Days")
        return d * _RET_MULT.get(u, 1) if d else 0

    def _min_local_days(policies):
        vals = []
        for p in policies:
            bp  = p.get("backupPolicy") or {}
            reg = bp.get("regular") or {}
            r   = reg.get("retention") or bp.get("retention") or {}
            v   = _ret_to_days(r)
            if v: vals.append(v)
        return f"{min(vals)} days" if vals else "N/A"

    def _min_arch_days(policies):
        vals = []
        for p in policies:
            for t in ((p.get("remoteTargetPolicy") or {}).get("archivalTargets") or []):
                v = _ret_to_days(t.get("retention") or {})
                if v: vals.append(v)
        return f"{min(vals)} days" if vals else "N/A"

    today = _dt.date.today()

    for cd in all_data:
        info     = cd["info"]
        vaults   = cd["vaults"]
        policies = cd["policies"]
        scores   = cd["scores"]

        cluster_enc, sd_enc = _sd_enc_status(cd)
        has_vault = bool(vaults)
        fk_st    = cd.get("fk_status") or _fk_status(cd)
        has_repl = any((p.get("remoteTargetPolicy") or {}).get("replicationTargets")
                       for p in policies)
        has_arch = any((p.get("remoteTargetPolicy") or {}).get("archivalTargets")
                       for p in policies)

        # ── Extended security controls ─────────────────────────────────────
        audit_on = cd.get("audit_enabled", False)
        _ntp_s2  = info.get("ntpSettings")
        ntp_auth = bool(
            isinstance(_ntp_s2, dict)
            and (_ntp_s2.get("ntpAuthenticationEnabled") or _ntp_s2.get("authEnabled"))
        )
        _remote_cfg = info.get("remoteSupportConfig")
        tunnel_on = bool(
            info.get("enableRemoteSupport")
            or (isinstance(_remote_cfg, dict) and _remote_cfg.get("enabled"))
            or info.get("remoteSupportEnabled")
        )
        _mfa_cfg = info.get("mfaConfig")
        cluster_mfa = bool(
            info.get("clusterMfaEnabled")
            or info.get("mfaEnabled")
            or (isinstance(_mfa_cfg, dict) and _mfa_cfg.get("enabled"))
        )
        has_sso = bool(cd.get("idps"))

        # Quorum — Helios feature; check info dict for known field names
        # Quorum is a Helios-level feature (Security Center → Quorum).
        # Data comes from GET /v2/mcm/quorum/groups fetched once before cluster
        # collection and stored in cd["quorum_enabled"]. Direct mode = N/A.
        if cd.get("mode") == "direct":
            quorum_str = "N/A"
        else:
            quorum_str = "Yes" if cd.get("quorum_enabled") else "No"

        # Soonest TLS cert expiry
        cert_expiry_str = ""
        soonest_days    = None
        for cert in (cd.get("certs") or []):
            exp_raw = (cert.get("expiryUsecs")
                       or cert.get("expiryTimeMsecs")
                       or cert.get("notAfter"))
            if not exp_raw:
                continue
            if exp_raw > 1_000_000_000_000_000:
                exp_dt = _dt.date.fromtimestamp(exp_raw / 1_000_000)
            elif exp_raw > 1_000_000_000_000:
                exp_dt = _dt.date.fromtimestamp(exp_raw / 1_000)
            else:
                exp_dt = _dt.date.fromtimestamp(exp_raw)
            d = (exp_dt - today).days
            if soonest_days is None or d < soonest_days:
                soonest_days    = d
                cert_expiry_str = f"{exp_dt.isoformat()} ({d}d)"

        dl_status, dl_label, dl_covered, dl_total = _cluster_datalock(policies,
                                                                      cd["groups"])

        gaps = []
        if not cluster_enc:           gaps.append("No cluster encryption")
        if sd_enc is False:           gaps.append("SD encryption missing")
        if not has_vault:             gaps.append("No vault")
        if fk_st == "none":           gaps.append("No indelible copy")
        if fk_st == "idle":           gaps.append("FK configured — not sending data")
        if not has_repl:              gaps.append("No replication")
        if not has_arch:              gaps.append("No archival")
        if dl_status == "none":       gaps.append("No DataLock")
        if dl_status == "partial":    gaps.append(f"DataLock partial ({dl_covered}/{dl_total})")
        if not audit_on:              gaps.append("Audit log off")
        if not cluster_mfa:           gaps.append("No cluster MFA")
        if not has_sso:               gaps.append("No SSO")
        if quorum_str == "No":        gaps.append("Quorum not enabled")
        if soonest_days is not None and soonest_days < CERT_WARN_DAYS:
            gaps.append(f"Cert expires {soonest_days}d")

        fk_label = {"active": "Active", "idle": "Configured — Idle", "none": "No"}[fk_st]
        yn = lambda v: "Yes" if v else "No"
        sd_enc_str = ("Yes" if sd_enc is True else "No" if sd_enc is False else "N/A")
        row = [cd["name"], yn(cluster_enc), sd_enc_str, yn(has_vault), fk_label,
               yn(has_repl), yn(has_arch),
               _min_local_days(policies), _min_arch_days(policies),
               scores["security"],
               yn(audit_on), yn(ntp_auth), yn(tunnel_on),
               yn(cluster_mfa), yn(has_sso),
               quorum_str,
               dl_label,
               cert_expiry_str or "—",
               scores.get("ransomware", ""),
               scores.get("ransomware_label", ""),
               "; ".join(gaps) if gaps else "OK"]
        rn = ws.max_row + 1
        for c, val in enumerate(row, 1):
            ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL

        # Cluster Encryption col 2, SD Encryption col 3
        ce_cell = ws.cell(row=rn, column=2)
        ce_cell.fill = _fill(LT_GREEN) if cluster_enc else _fill(RED)
        ce_cell.font = _font(bold=True, color=WHITE if not cluster_enc else "000000")
        sd_cell = ws.cell(row=rn, column=3)
        if sd_enc is True:
            sd_cell.fill = _fill(LT_GREEN); sd_cell.font = _font(bold=True)
        elif sd_enc is False:
            sd_cell.fill = _fill(RED);      sd_cell.font = _font(bold=True, color=WHITE)
        else:
            sd_cell.fill = _fill(YELLOW);   sd_cell.font = _font(bold=True)

        for col, flag in [(4, has_vault), (6, has_repl), (7, has_arch)]:
            cell = ws.cell(row=rn, column=col)
            cell.fill = _fill(LT_GREEN) if flag else _fill(RED)
            cell.font = _font(bold=True, color=WHITE if not flag else "000000")

        # FK column 5: three-state coloring
        fk_cell = ws.cell(row=rn, column=5)
        if fk_st == "active":
            fk_cell.fill = _fill(LT_GREEN); fk_cell.font = _font(bold=True)
        elif fk_st == "idle":
            fk_cell.fill = _fill(ORANGE);   fk_cell.font = _font(bold=True)
        else:
            fk_cell.fill = _fill(RED);      fk_cell.font = _font(bold=True, color=WHITE)

        # Score col 10
        sc = ws.cell(row=rn, column=10)
        sc.fill = _rag_fill(scores["security"])
        sc.font = _font(bold=True)

        # Extended controls: audit(11), NTP auth(12), tunnel(13), MFA(14), SSO(15)
        for col, flag, warn_if_off in [
            (11, audit_on,    True),
            (12, ntp_auth,    False),
            (13, tunnel_on,   False),   # tunnel: yellow if on (reduce attack surface)
            (14, cluster_mfa, True),
            (15, has_sso,     False),
        ]:
            cell = ws.cell(row=rn, column=col)
            if col == 13:
                cell.fill = _fill(YELLOW) if flag else _fill(LT_GREEN)
                cell.font = _font(bold=True)
            elif warn_if_off:
                cell.fill = _fill(LT_GREEN) if flag else _fill(RED)
                cell.font = _font(bold=True, color=WHITE if not flag else "000000")
            else:
                cell.fill = _fill(LT_GREEN) if flag else _fill(YELLOW)
                cell.font = _font(bold=True)

        # Quorum col 16 — N/A for direct mode
        qc = ws.cell(row=rn, column=16)
        if quorum_str == "Yes":
            qc.fill = _fill(LT_GREEN); qc.font = _font(bold=True)
        elif quorum_str == "No":
            qc.fill = _fill(YELLOW);   qc.font = _font(bold=True)
        # N/A: no fill

        # DataLock col 17 — four-state coloring
        dl_cell = ws.cell(row=rn, column=17)
        if dl_status == "full_compliance":
            dl_cell.fill = _fill(LT_GREEN); dl_cell.font = _font(bold=True)
        elif dl_status == "full_any":
            dl_cell.fill = _fill(LT_GREEN); dl_cell.font = _font(bold=True)
        elif dl_status == "partial":
            dl_cell.fill = _fill(YELLOW);   dl_cell.font = _font(bold=True)
        else:
            dl_cell.fill = _fill(RED);      dl_cell.font = _font(bold=True, color=WHITE)

        # Cert expiry col 18
        if soonest_days is not None:
            cc = ws.cell(row=rn, column=18)
            if soonest_days < CERT_CRIT_DAYS:
                cc.fill = _fill(RED);    cc.font = _font(bold=True, color=WHITE)
            elif soonest_days < CERT_WARN_DAYS:
                cc.fill = _fill(ORANGE); cc.font = _font(bold=True)

        # Ransomware Score col 19
        rw_score = scores.get("ransomware")
        if rw_score is not None:
            rc = ws.cell(row=rn, column=19)
            rc.fill = _ransomware_rag_fill(rw_score)
            rc.font = _font(bold=True)

    # ── Encryption Key Management (KMS) sub-section ──────────────────────────
    _kms_blank = ws.max_row + 2
    ws.cell(row=_kms_blank, column=1, value="Encryption Key Management") \
      .font = _font(bold=True, size=11)

    kms_cols = ["Cluster", "KMS Type", "KMS Server",
                "Server Status", "Key Rotation", "Risk"]
    _hdr(ws, _kms_blank + 1, kms_cols)

    for cd in all_data:
        kms       = cd.get("kms") or {}
        cluster_enc, _ = _sd_enc_status(cd)

        if not kms:
            # Empty dict — internal KMS or unknown
            kms_type   = "Internal (Built-in)"
            kms_server = "N/A"
            kms_status = "Unknown"
            risk       = "MEDIUM" if cluster_enc else "HIGH"
        else:
            kms_type_raw = (kms.get("serverType") or kms.get("type")
                            or kms.get("kmsType") or "")
            if kms_type_raw:
                kms_type = kms_type_raw
            else:
                kms_type = "Internal (Built-in)"
            kms_server   = (kms.get("serverName") or kms.get("hostname")
                            or kms.get("server") or kms.get("kmipServer") or "N/A")
            connectivity = (kms.get("connectionStatus") or kms.get("status")
                            or kms.get("serverStatus") or "")
            kms_status   = connectivity if connectivity else "Connected"
            _kms_type_l  = kms_type.lower()
            if "internal" in _kms_type_l or "builtin" in _kms_type_l or "local" in _kms_type_l:
                risk = "MEDIUM"
            else:
                risk = "LOW"

        if not cluster_enc:
            risk = "HIGH"

        key_rotation = (kms.get("keyRotationPeriodDays") or
                        kms.get("rotationPeriod") or "Unknown")

        kms_row = [cd["name"], _safe_cell(kms_type), _safe_cell(str(kms_server)),
                   _safe_cell(kms_status), _safe_cell(str(key_rotation)), risk]
        rn = ws.max_row + 1
        for c, val in enumerate(kms_row, 1):
            ws.cell(row=rn, column=c, value=val).font = _FONT_NORMAL

        # Risk coloring (col 6)
        risk_cell = ws.cell(row=rn, column=6)
        if risk == "HIGH":
            risk_cell.fill = _fill(RED);    risk_cell.font = _font(bold=True, color=WHITE)
        elif risk == "MEDIUM":
            risk_cell.fill = _fill(YELLOW); risk_cell.font = _font(bold=True)
        elif risk == "LOW":
            risk_cell.fill = _fill(LT_GREEN)

    # ── Security & Anomaly Alerts sub-section ────────────────────────────────
    import datetime as _dt_sec
    _blank_row = ws.max_row + 2
    ws.cell(row=_blank_row, column=1, value="Security & Anomaly Alerts") \
      .font = _font(bold=True, size=11)

    alert_cols = ["Cluster", "Security Alerts (Open)", "Anomaly Alerts (Open)",
                  "Oldest Unack Security Alert (days)", "Newest Security Alert"]
    _hdr(ws, _blank_row + 1, alert_cols)

    _SEC_CATS    = {"ksecurity", "kdatasecurity"}
    _ANOM_CATS   = {"kanomalousactivity", "kanomalyalert"}

    for cd in all_data:
        sec_alerts  = []
        anom_alerts = []
        for a in (cd.get("alerts") or []):
            cat = (a.get("alertCategory") or "").lower()
            if cat in _SEC_CATS:
                sec_alerts.append(a)
            elif cat in _ANOM_CATS:
                anom_alerts.append(a)

        # Oldest unacknowledged security alert
        oldest_days   = ""
        newest_str    = ""
        oldest_usecs  = None
        newest_usecs  = None
        for a in sec_alerts:
            ack = a.get("acknowledgeTimestampUsecs") or 0
            if not ack:
                first_u = a.get("firstTimestampUsecs") or a.get("firstOccurrenceTime") or 0
                if first_u and (oldest_usecs is None or first_u < oldest_usecs):
                    oldest_usecs = first_u
            latest_u = a.get("latestTimestampUsecs") or a.get("latestOccurrenceTime") or 0
            if latest_u and (newest_usecs is None or latest_u > newest_usecs):
                newest_usecs = latest_u

        if oldest_usecs:
            oldest_days = round((cd["end_usecs"] - oldest_usecs) / 86_400_000_000, 1)
        if newest_usecs:
            newest_str = _dt_sec.datetime.fromtimestamp(
                newest_usecs / 1_000_000, tz=_dt_sec.timezone.utc).strftime("%Y-%m-%d")

        alert_row = [cd["name"], len(sec_alerts), len(anom_alerts),
                     oldest_days, newest_str]
        rn = ws.max_row + 1
        for c, val in enumerate(alert_row, 1):
            ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL

        if sec_alerts or anom_alerts:
            for col in range(1, 6):
                ws.cell(row=rn, column=col).fill = _fill(RED)
                ws.cell(row=rn, column=col).font = _font(bold=True, color=WHITE)

    auto_fit_columns(ws)
    ws.column_dimensions["U"].width = 60


def _sheet_replication(wb, all_data):
    ws = wb.create_sheet("Replication & Archive")
    ws.freeze_panes = "A3"
    _title(ws, "Replication & Archival Targets", "K")

    cols = ["Cluster", "Target Type", "Target Name", "Vault Type",
            "Protection Groups",
            "Vault Storage (TB)",
            "Logical Data Transferred (TB)\n[numLogicalBytesTransferred — cumulative]",
            "Physical Data Transferred (TB)\n[numPhysicalBytesTransferred — cumulative]",
            "Status", "Notes",
            "Protection Group Names",
            "Replication Lag (hrs)"]
    _hdr(ws, 2, cols)

    for cd in all_data:
        # ── Vault ID → name / type maps ───────────────────────────────────
        vault_name_by_id   = {}   # vault id → name
        vault_type_by_name = {}   # name     → vaultType
        vault_id_by_name   = {}   # name     → id
        for v in cd["vaults"]:
            vn = v.get("name") or v.get("vaultName") or ""
            vi = v.get("id")
            vt = v.get("vaultType", "")
            if vi: vault_name_by_id[vi] = vn or f"Vault-{vi}"
            if vn:
                vault_type_by_name[vn] = vt
                if vi: vault_id_by_name[vn] = vi

        # ── Target name extractors ────────────────────────────────────────
        def _rep_name(t):
            cfg = t.get("remoteTargetConfig") or {}
            return (t.get("targetName")
                    or cfg.get("clusterName")
                    or cfg.get("name") or "")

        def _arch_name(t):
            cfg  = t.get("archivalTargetConfig") or {}
            name = (t.get("targetName")
                    or cfg.get("name")
                    or cfg.get("vaultName") or "")
            if not name:
                vid = cfg.get("vaultId") or cfg.get("id")
                if vid:
                    name = vault_name_by_id.get(vid, f"Vault-{vid}")
            return name

        # ── Policy chain: policy id → target names ───────────────────────
        policy_rep_targets  = {}   # policy id → set of replication cluster names
        policy_arch_targets = {}   # policy id → set of vault names
        arch_target_type    = {}   # vault name → targetType

        for p in cd["policies"]:
            pid = p.get("id")
            if not pid:
                continue
            rtp = p.get("remoteTargetPolicy") or {}
            for t in (rtp.get("replicationTargets") or []):
                n = _rep_name(t)
                if n:
                    policy_rep_targets.setdefault(pid, set()).add(n)
            for t in (rtp.get("archivalTargets") or []):
                n = _arch_name(t)
                if n:
                    policy_arch_targets.setdefault(pid, set()).add(n)
                    if n not in arch_target_type:
                        cfg = t.get("archivalTargetConfig") or {}
                        arch_target_type[n] = (t.get("targetType")
                                               or cfg.get("targetType") or "")

        # ── Group → target mapping (policy chain) ────────────────────────
        rep_groups  = {}   # replication cluster name → [group names]
        arch_groups = {}   # vault name               → [group names]

        for g in cd["groups"]:
            pid   = g.get("policyId")
            gname = g.get("name", "")
            if not pid or not gname:
                continue
            for tname in (policy_rep_targets.get(pid) or set()):
                rep_groups.setdefault(tname, []).append(gname)
            for tname in (policy_arch_targets.get(pid) or set()):
                arch_groups.setdefault(tname, []).append(gname)

        # ── Supplement arch_groups from fk_data (direct transfer records) ─
        # fk_data lists protectionJobName per vault — reliable even when the
        # policy chain fails to extract vault names (e.g. ID-only references).
        for fk in (cd["fk_data"] or []):
            fk_vn = fk.get("vaultName", "")
            if not fk_vn:
                continue
            for job in (fk.get("dataTransferPerProtectionJob") or []):
                jname = job.get("protectionJobName", "")
                if jname and jname not in (arch_groups.get(fk_vn) or []):
                    arch_groups.setdefault(fk_vn, []).append(jname)

        # All unique target names — union of policy chain + fk_data + vault list
        all_rep_targets  = set().union(*policy_rep_targets.values()) \
                           if policy_rep_targets else set()
        fk_vault_names   = {fk.get("vaultName") for fk in (cd["fk_data"] or [])
                            if fk.get("vaultName")}
        vault_list_names = {v.get("name") or v.get("vaultName") or ""
                            for v in cd["vaults"]} - {""}
        all_arch_targets = (
            (set().union(*policy_arch_targets.values()) if policy_arch_targets else set())
            | fk_vault_names
            | vault_list_names
        )

        # ── FK transfer stats ─────────────────────────────────────────────
        fk_by_name = {}
        fk_by_id   = {}
        for fk in (cd["fk_data"] or []):
            fk_vn  = fk.get("vaultName", "")
            fk_vid = fk.get("vaultId")
            jobs   = fk.get("dataTransferPerProtectionJob") or []
            data   = {
                "consumed": sum((j.get("storageConsumed") or 0) for j in jobs),
                "logical":  sum((j.get("numLogicalBytesTransferred") or 0) for j in jobs),
                "physical": sum((j.get("numPhysicalBytesTransferred") or 0) for j in jobs),
            }
            if fk_vn:  fk_by_name[fk_vn] = data
            if fk_vid: fk_by_id[fk_vid]   = data

        def _fk_stats(vname):
            d = fk_by_name.get(vname)
            if not d:
                vid = vault_id_by_name.get(vname)
                if vid: d = fk_by_id.get(vid)
            return d or {}

        fk_types = {"rpaas", "fortknox", "kfortknox", "krpaas"}

        def _is_fk(vname, vtype):
            return (any(x in vtype.lower() for x in fk_types)
                    or "fortknox" in vname.lower())

        # ── Replication lag per target: most recent endTimeUsecs ─────────
        # Build map: replication target name → most recent endTimeUsecs across groups
        repl_latest_end = {}   # target name → latest endTimeUsecs
        for g in cd["groups"]:
            lr = g.get("lastRun") or {}
            for ri in ((lr.get("replicationInfo") or {})
                       .get("replicationTargetResults") or []):
                tgt_cfg = ri.get("remoteTargetConfig") or {}
                tgt_n   = (ri.get("targetName")
                           or tgt_cfg.get("clusterName")
                           or tgt_cfg.get("name") or "")
                end_u   = ri.get("endTimeUsecs") or 0
                if tgt_n and end_u:
                    if tgt_n not in repl_latest_end or end_u > repl_latest_end[tgt_n]:
                        repl_latest_end[tgt_n] = end_u

        def _repl_lag(tname):
            end_u = repl_latest_end.get(tname) or 0
            if not end_u:
                return None
            return round((cd["end_usecs"] - end_u) / 3_600_000_000, 1)

        # ── Replication rows ──────────────────────────────────────────────
        for tname in sorted(all_rep_targets):
            groups = sorted(rep_groups.get(tname) or [])
            gcnt   = len(groups)
            status = "Configured" if gcnt else "No groups assigned"
            lag    = _repl_lag(tname)
            lag_str = lag if lag is not None else "No data"
            rn = ws.max_row + 1
            row = [cd["name"], "Replication", tname, "", gcnt,
                   "", "", "", status, "",
                   ", ".join(groups), lag_str]
            for c, val in enumerate(row, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL
            # Color lag cell (col 12)
            lag_cell = ws.cell(row=rn, column=12)
            if lag is None:
                lag_cell.fill = _fill(RED);    lag_cell.font = _font(bold=True, color=WHITE)
            elif lag > 72:
                lag_cell.fill = _fill(RED);    lag_cell.font = _font(bold=True, color=WHITE)
            elif lag > 24:
                lag_cell.fill = _fill(YELLOW); lag_cell.font = _font(bold=True)
            else:
                lag_cell.fill = _fill(LT_GREEN); lag_cell.font = _font(bold=True)

        # ── Archival / vault rows ─────────────────────────────────────────
        shown = set()
        for tname in sorted(all_arch_targets):
            groups = sorted(arch_groups.get(tname) or [])
            gcnt   = len(groups)
            vtype  = arch_target_type.get(tname) or vault_type_by_name.get(tname, "")
            fkd    = _fk_stats(tname)
            note   = "FortKnox/RPaaS" if _is_fk(tname, vtype) else ""
            status = "Configured" if gcnt else "No groups assigned"
            rn     = ws.max_row + 1
            row    = [cd["name"], "Archival/Vault", tname, vtype, gcnt,
                      round(bytes_to_tb(fkd.get("consumed") or 0), 3),
                      round(bytes_to_tb(fkd.get("logical")  or 0), 3),
                      round(bytes_to_tb(fkd.get("physical") or 0), 3),
                      status, note,
                      ", ".join(groups), ""]
            for c, val in enumerate(row, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL
            shown.add(tname)

        # Vaults that exist in the API but are not referenced by any policy
        for v in cd["vaults"]:
            vname = v.get("name") or v.get("vaultName") or ""
            if not vname or vname in shown:
                continue
            vtype = v.get("vaultType", "")
            fkd   = _fk_stats(vname)
            note  = "FortKnox/RPaaS" if _is_fk(vname, vtype) else ""
            rn    = ws.max_row + 1
            row   = [cd["name"], "Archival/Vault", vname, vtype, 0,
                     round(bytes_to_tb(fkd.get("consumed") or 0), 3),
                     round(bytes_to_tb(fkd.get("logical")  or 0), 3),
                     round(bytes_to_tb(fkd.get("physical") or 0), 3),
                     "Vault exists — not referenced by any policy", note, "", ""]
            for c, val in enumerate(row, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL

    auto_fit_columns(ws)
    ws.column_dimensions["K"].width = 60


def _sheet_fortknox_detail(wb, all_data):
    """Per-protection-group data transfer to FortKnox / external vault targets.

    Full mode: covers the full lookback window (args.days, default 30 days).
    Quick mode: covers the last 1 day only (separate fk_data_1d fetch).
    """
    ws = wb.create_sheet("FortKnox Data Transfer")
    ws.freeze_panes = "A3"
    _title(ws,
           "FortKnox / External Target Data Transfer — Per Protection Group", "I")

    import datetime as _dt_fk

    cols = ["Cluster", "Vault Name", "Vault Type",
            "Protection Group",
            "Logical Transferred (TB)", "Physical Transferred (TB)",
            "Storage Consumed (TB)", "Snapshots",
            "Period",
            "Last FK Archival", "Days Since FK Transfer"]
    _hdr(ws, 2, cols)

    def _vtype_label(vt):
        tl = (vt or "").lower()
        if any(x in tl for x in ("fortknox", "rpaas", "krpaas", "kfortknox")):
            return "FortKnox / RPaaS"
        if any(x in tl for x in ("s3", "glacier", "azure", "gcp", "cloud")):
            return "Cloud Tier"
        if "nas" in tl:
            return "NAS Vault"
        if "tape" in tl:
            return "Tape"
        return vt or "Unknown"

    _FK_TARGET_TYPES = {"fortknox", "rpaas", "kfortknox", "krpaas", "kfort_knox",
                        "krpaasarchival"}

    def _is_fk_target(tt):
        tl = (tt or "").lower()
        return any(x in tl for x in _FK_TARGET_TYPES) or "fortknox" in tl

    for cd in all_data:
        quick      = cd.get("quick", False)
        days       = cd.get("days", 30)
        fk_source  = cd.get("fk_data_1d") if quick else (cd.get("fk_data") or [])
        period     = "Last 1 day" if quick else f"Last {days} days"

        # Build per-group most-recent FK archival endTimeUsecs
        fk_last_end = {}   # group name → latest FK archival endTimeUsecs
        for g in cd["groups"]:
            gname = g.get("name", "")
            if not gname:
                continue
            lr = g.get("lastRun") or {}
            for ar in ((lr.get("archivalInfo") or {})
                       .get("archivalTargetResults") or []):
                ttype = (ar.get("targetType") or
                         (ar.get("archivalTargetConfig") or {}).get("targetType") or "")
                if _is_fk_target(ttype):
                    end_u = ar.get("endTimeUsecs") or 0
                    if end_u and end_u > fk_last_end.get(gname, 0):
                        fk_last_end[gname] = end_u

        has_rows = False
        for fk in (fk_source or []):
            vname = fk.get("vaultName", "")
            vtype = _vtype_label(fk.get("vaultType", ""))
            for job in (fk.get("dataTransferPerProtectionJob") or []):
                jname    = job.get("protectionJobName", "")
                logical  = round(bytes_to_tb(job.get("numLogicalBytesTransferred")  or 0), 4)
                physical = round(bytes_to_tb(job.get("numPhysicalBytesTransferred") or 0), 4)
                consumed = round(bytes_to_tb(job.get("storageConsumed")             or 0), 4)
                snaps    = job.get("numSnapshots") or ""

                # Last FK archival date and days since
                last_end_u  = fk_last_end.get(jname) or 0
                if last_end_u:
                    last_fk_str = _dt_fk.datetime.fromtimestamp(
                        last_end_u / 1_000_000,
                        tz=_dt_fk.timezone.utc).strftime("%Y-%m-%d")
                    days_since  = round((cd["end_usecs"] - last_end_u) / 86_400_000_000, 1)
                else:
                    last_fk_str = "Unknown"
                    days_since  = None

                rn  = ws.max_row + 1
                row = [cd["name"], vname, vtype, jname,
                       logical, physical, consumed, snaps, period,
                       last_fk_str, days_since if days_since is not None else "Unknown"]
                for c, val in enumerate(row, 1):
                    ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL

                # Color "Days Since FK Transfer" cell (col 11)
                dc = ws.cell(row=rn, column=11)
                if days_since is None:
                    dc.fill = _fill(YELLOW); dc.font = _font(bold=True)
                elif days_since > 14:
                    dc.fill = _fill(RED);    dc.font = _font(bold=True, color=WHITE)
                elif days_since > 7:
                    dc.fill = _fill(YELLOW); dc.font = _font(bold=True)
                else:
                    dc.fill = _fill(LT_GREEN); dc.font = _font(bold=True)

                has_rows = True

        if not has_rows:
            rn = ws.max_row + 1
            ws.cell(row=rn, column=1, value=cd["name"]).font = _FONT_NORMAL
            ws.cell(row=rn, column=4,
                    value="No data transfer records for this cluster in the period.") \
               .font = _font(color="595959")

    auto_fit_columns(ws)
    ws.column_dimensions["B"].width = 30   # Vault Name
    ws.column_dimensions["D"].width = 35   # Protection Group


def _sheet_views(wb, all_data):
    ws = wb.create_sheet("Data Services")
    ws.freeze_panes = "A3"
    _title(ws, "NAS Views — Quota Utilization & Protocols", "L")

    cols = ["Cluster", "View Name", "Storage Domain", "Protocols",
            "Quota (GB)", "Used (GB)", "Usage %",
            "Case Insensitive", "Created",
            "SMB Mount Paths", "NFS Mount Paths", "Status"]
    _hdr(ws, 2, cols)

    for cd in all_data:
        for v in cd["views"]:
            lq  = v.get("logicalQuota") or \
                  (v.get("storagePolicy") or {}).get("storageQuotaPolicy") or {}
            quota_b = lq.get("hardLimitBytes") or 0
            vstats  = ((v.get("stats") or {}).get("dataUsageStats") or {})
            used_b  = vstats.get("totalLogicalUsageBytes") or 0
            quota_gb = round(bytes_to_gb(quota_b), 2) if quota_b else ""
            used_gb  = round(bytes_to_gb(used_b), 2)
            pct      = round(used_b / quota_b * 100, 1) if quota_b else ""

            protos = []
            if v.get("enableSmbViewDiscovery") or v.get("smbMountPaths"): protos.append("SMB")
            if v.get("nfsMountPaths") or v.get("enableNfsViewDiscovery"):  protos.append("NFS")
            if v.get("s3AccessEnabled"):                                   protos.append("S3")

            smb = ", ".join(v.get("smbMountPaths") or [])
            nfs = ", ".join(v.get("nfsMountPaths") or [])
            created_ms = v.get("createTimeMsecs") or 0
            created = (datetime.datetime.fromtimestamp(created_ms / 1000,
                                                        tz=datetime.timezone.utc)
                       .strftime("%Y-%m-%d") if created_ms else "")

            flag = "OK"
            if not quota_b:
                flag = "NO QUOTA"
            elif isinstance(pct, float) and pct >= VIEW_QUOTA_CRIT_PCT:
                flag = f"CRITICAL {pct:.0f}%"
            elif isinstance(pct, float) and pct >= VIEW_QUOTA_WARN_PCT:
                flag = f"WARN {pct:.0f}%"

            row = [cd["name"], v.get("name", ""), v.get("storageDomainName", ""),
                   "/".join(protos), quota_gb, used_gb, pct,
                   "Yes" if v.get("caseInsensitiveNamesEnabled") else "No",
                   created, smb, nfs, flag]
            rn = ws.max_row + 1
            for c, val in enumerate(row, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL

            fc = ws.cell(row=rn, column=12)
            if "CRITICAL" in flag:
                fc.fill = _fill(RED);    fc.font = _font(bold=True, color=WHITE)
            elif "WARN" in flag:
                fc.fill = _fill(ORANGE); fc.font = _font(bold=True)
            elif "NO QUOTA" in flag:
                fc.fill = _fill(YELLOW); fc.font = _font(bold=True)
            else:
                fc.fill = _fill(LT_GREEN)

    auto_fit_columns(ws)


def _sheet_data_exposure(wb, all_data):
    """NAS Share Exposure — per-view risk assessment for SMB discovery, NFS, S3, quota."""
    ws = wb.create_sheet("Data Exposure")
    ws.freeze_panes = "A3"
    _title(ws, "Data Exposure — NAS View Access & Risk Assessment", "I")

    cols = ["Cluster", "View Name", "Storage Domain",
            "SMB Discovery", "NFS Open Mount", "S3 Enabled",
            "Quota Set", "Risk Level", "Notes"]
    _hdr(ws, 2, cols)

    for cd in all_data:
        for v in cd.get("views") or []:
            # ── SMB Discovery ────────────────────────────────────────────
            smb_disc = bool(v.get("enableSmbViewDiscovery"))
            smb_str  = "ENABLED" if smb_disc else "Off"

            # ── NFS Open Mount ───────────────────────────────────────────
            nfs_paths = v.get("nfsMountPaths") or []
            nfs_disc  = bool(v.get("enableNfsViewDiscovery"))
            if nfs_paths and nfs_disc:
                nfs_str = "OPEN"
            elif nfs_paths:
                nfs_str = "Restricted"
            else:
                nfs_str = "None"

            # ── S3 ───────────────────────────────────────────────────────
            s3_on   = bool(v.get("s3AccessEnabled"))
            s3_str  = "YES" if s3_on else "No"

            # ── Quota ────────────────────────────────────────────────────
            lq = v.get("logicalQuota") or \
                 (v.get("storagePolicy") or {}).get("storageQuotaPolicy") or {}
            quota_b = lq.get("hardLimitBytes") or 0
            if quota_b:
                quota_str = f"{round(bytes_to_gb(quota_b), 1)} GB"
            else:
                quota_str = "No Quota"

            # ── Risk Level ───────────────────────────────────────────────
            notes = []
            if smb_disc:
                notes.append("SMB discovery on")
            if nfs_str == "OPEN":
                notes.append("NFS open mount")
            if s3_on:
                notes.append("S3 enabled")
            if not quota_b:
                notes.append("No quota")

            if smb_disc or nfs_str == "OPEN":
                risk = "HIGH"
            elif s3_on or not quota_b:
                risk = "MEDIUM"
            else:
                risk = "LOW"

            row = [cd["name"], v.get("name", ""), v.get("storageDomainName", ""),
                   smb_str, nfs_str, s3_str, quota_str, risk,
                   "; ".join(notes) if notes else "OK"]
            rn = ws.max_row + 1
            for c, val in enumerate(row, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL

            # SMB Discovery (col 4)
            smb_cell = ws.cell(row=rn, column=4)
            if smb_disc:
                smb_cell.fill = _fill(RED);     smb_cell.font = _font(bold=True, color=WHITE)
            else:
                smb_cell.fill = _fill(LT_GREEN)

            # NFS Open Mount (col 5)
            nfs_cell = ws.cell(row=rn, column=5)
            if nfs_str == "OPEN":
                nfs_cell.fill = _fill(RED);     nfs_cell.font = _font(bold=True, color=WHITE)
            elif nfs_str == "Restricted":
                nfs_cell.fill = _fill(YELLOW);  nfs_cell.font = _font(bold=True)
            else:
                nfs_cell.fill = _fill(LT_GREEN)

            # S3 (col 6)
            s3_cell = ws.cell(row=rn, column=6)
            if s3_on:
                s3_cell.fill = _fill(YELLOW);   s3_cell.font = _font(bold=True)
            else:
                s3_cell.fill = _fill(LT_GREEN)

            # Quota (col 7)
            q_cell = ws.cell(row=rn, column=7)
            if not quota_b:
                q_cell.fill = _fill(YELLOW);    q_cell.font = _font(bold=True)

            # Risk Level (col 8)
            rc = ws.cell(row=rn, column=8)
            if risk == "HIGH":
                rc.fill = _fill(RED);    rc.font = _font(bold=True, color=WHITE)
            elif risk == "MEDIUM":
                rc.fill = _fill(YELLOW); rc.font = _font(bold=True)
            else:
                rc.fill = _fill(LT_GREEN); rc.font = _font(bold=True)

    auto_fit_columns(ws)
    ws.column_dimensions["I"].width = 40   # Notes


def _sheet_coverage(wb, all_data):
    ws = wb.create_sheet("Coverage Gaps")
    ws.freeze_panes = "A3"
    _title(ws, "Coverage Gaps — Groups with No Recent Successful Backup", "I")

    cols = ["Cluster", "Group Name", "Policy",
            "Active", "Paused",
            "Last Successful Run", "Hours Since Success",
            "Last Run Status", "Gap Severity"]
    _hdr(ws, 2, cols)

    for cd in all_data:
        pm = cd.get("policy_map", {})
        for g in cd["groups"]:
            lr  = g.get("lastRun") or {}
            lb  = lr.get("localBackupInfo") or {}
            st  = lb.get("status", "")
            sla_viol = lb.get("isSlaViolated") or lr.get("isSlaViolated") or False
            su  = lb.get("startTimeUsecs") or 0
            hrs = round((cd["end_usecs"] - su) / 3_600_000_000, 1) if su else 9999

            is_gap = (g.get("isPaused")
                      or st in ("kFailure", "Failed")
                      or hrs > RPO_WARN_HOURS
                      or su == 0)
            if not is_gap:
                continue

            if g.get("isPaused"):
                sev = "PAUSED"
            elif st in ("kFailure", "Failed"):
                sev = "FAILED"
            elif hrs > RPO_CRIT_HOURS:
                sev = "CRITICAL RPO GAP"
            elif hrs > RPO_WARN_HOURS:
                sev = "RPO WARNING"
            else:
                sev = "UNKNOWN"

            policy_name = (g.get("policyName")
                           or pm.get(g.get("policyId"), "")
                           or g.get("policyId", ""))

            row = [cd["name"], g.get("name", ""),
                   policy_name,
                   "Yes" if g.get("isActive", True) else "No",
                   "Yes" if g.get("isPaused") else "No",
                   usecs_to_datetime(su) if su else "Never",
                   hrs if hrs < 9999 else "",
                   st, sev]
            rn = ws.max_row + 1
            for c, val in enumerate(row, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL

            fc = ws.cell(row=rn, column=9)
            if sev in ("FAILED", "CRITICAL RPO GAP"):
                fc.fill = _fill(RED);    fc.font = _font(bold=True, color=WHITE)
            elif sev in ("PAUSED", "RPO WARNING"):
                fc.fill = _fill(ORANGE); fc.font = _font(bold=True)
            else:
                fc.fill = _fill(YELLOW); fc.font = _font(bold=True)

    auto_fit_columns(ws)


def _sheet_trends(wb, all_data):
    ws = wb.create_sheet("Trends (30d)")
    ws.freeze_panes = "A3"
    _title(ws, "30-Day Trends — Daily Backup Success Rate & Volume", "J")

    has_data = any(cd.get("v1_runs") or cd["group_runs"] for cd in all_data)
    if not has_data:
        ws["A3"] = "No run history available."
        ws["A3"].font = _font(bold=False, color="595959")
        return

    cols = ["Cluster", "Date", "Total Runs", "Successful", "Failed",
            "Warning", "Canceled", "Success % (incl. Warnings)", "SLA Violations",
            "Logical (GB)", "Avg Duration (min)"]
    _hdr(ws, 2, cols)

    for cd in all_data:
        daily = {}
        daily_dur = {}   # day → list of completed run durations in minutes

        # ── Primary: v1 protectionRuns — mirrors the Helios reporting page source ──
        # Single call per cluster; backupRun.slaViolated and
        # backupRun.stats.totalLogicalBackupSizeBytes are accurate.
        for run in (cd.get("v1_runs") or []):
            br    = run.get("backupRun") or {}
            stats = br.get("stats") or {}
            su    = stats.get("startTimeUsecs") or 0
            if not su:
                continue
            day = datetime.datetime.fromtimestamp(su / 1_000_000,
                                                   tz=datetime.timezone.utc) \
                          .strftime("%Y-%m-%d")
            if day not in daily:
                daily[day] = dict(total=0, succ=0, fail=0,
                                  warn=0, cancel=0, viols=0, logical=0)
            d  = daily[day]
            st = br.get("status", "")
            d["total"] += 1
            if st == "kSuccess":                d["succ"]   += 1
            elif st == "kFailure":              d["fail"]   += 1
            elif st == "kWarning":              d["warn"]   += 1
            elif "Cancel" in st:                d["cancel"] += 1
            if br.get("slaViolated"):           d["viols"]  += 1
            d["logical"] += stats.get("totalLogicalBackupSizeBytes") or 0
            # Duration
            eu = stats.get("endTimeUsecs") or 0
            if eu and su and eu > su:
                daily_dur.setdefault(day, []).append((eu - su) / 60_000_000)

        # ── Fallback: v2 per-group runs (used when v1 unavailable) ──────────────
        if not daily:
            for runs in cd["group_runs"].values():
                for run in runs:
                    lb = run.get("localBackupInfo") or {}
                    su = lb.get("startTimeUsecs") or 0
                    if not su:
                        continue
                    day = datetime.datetime.fromtimestamp(su / 1_000_000,
                                                           tz=datetime.timezone.utc) \
                                  .strftime("%Y-%m-%d")
                    if day not in daily:
                        daily[day] = dict(total=0, succ=0, fail=0,
                                          warn=0, cancel=0, viols=0, logical=0)
                    d  = daily[day]
                    st = lb.get("status", "")
                    d["total"] += 1
                    if st in ("kSuccess", "Succeeded"): d["succ"]   += 1
                    elif st in ("kFailure", "Failed"):  d["fail"]   += 1
                    elif st == "kWarning":              d["warn"]   += 1
                    elif "Cancel" in st:                d["cancel"] += 1
                    if lb.get("isSlaViolated"):         d["viols"]  += 1
                    d["logical"] += (lb.get("localSnapshotStats") or {}).get("logicalSizeBytes") or 0
                    # Duration
                    eu = lb.get("endTimeUsecs") or 0
                    if eu and su and eu > su:
                        daily_dur.setdefault(day, []).append((eu - su) / 60_000_000)

        if not daily:
            continue

        trend_start = ws.max_row + 1
        for day in sorted(daily):
            d   = daily[day]
            pct = round((d["succ"] + d["warn"]) / d["total"] * 100, 1) if d["total"] else 0
            durs = daily_dur.get(day) or []
            avg_dur = round(sum(durs) / len(durs), 1) if durs else ""
            rn  = ws.max_row + 1
            row = [cd["name"], day, d["total"], d["succ"], d["fail"],
                   d["warn"], d["cancel"], pct, d["viols"],
                   round(bytes_to_gb(d["logical"]), 1), avg_dur]
            for c, val in enumerate(row, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL
            sc = ws.cell(row=rn, column=8)
            if pct < SUCCESS_CRIT_PCT:
                sc.fill = _fill(RED);    sc.font = _font(bold=True, color=WHITE)
            elif pct < SUCCESS_WARN_PCT:
                sc.fill = _fill(ORANGE); sc.font = _font(bold=True)

        trend_end = ws.max_row
        n_days    = trend_end - trend_start + 1
        if n_days > 1:
            chart = LineChart()
            chart.style  = 10
            chart.height = 16   # taller for readability
            chart.width  = 26   # wide but fits to the right of the data cols

            # ── Chart title: larger font, no overlay ──────────────────────────
            _title_text = f"{cd['name']} — Daily Success Rate (incl. Warnings)"
            try:
                from openpyxl.chart.title import Title as _CTitle
                from openpyxl.drawing.text import (
                    RichTextBodyProperties as _BodyPr,
                    Paragraph as _Para,
                    RegularTextRun as _Run,
                    RunProperties as _RunPr,
                )
                from openpyxl.chart.text import RichText as _CRichText, Text as _CText
                _ctitle = _CTitle()
                _ctitle.overlay = False
                _rich = _CRichText()
                _rich.bodyPr = _BodyPr()
                _para = _Para()
                _run  = _Run()
                _run.t  = _title_text
                _run.rPr = _RunPr(sz=1400, b=True)   # 14 pt bold
                _para.r.append(_run)
                _rich.p.append(_para)
                _ctitle.tx = _CText(rich=_rich)
                chart.title = _ctitle
            except Exception:
                chart.title = _title_text   # fallback: plain string

            # ── Y axis: Success %, fixed 0–100, tick marks visible ───────────
            chart.y_axis.title         = "Success %"
            chart.y_axis.numFmt        = "0.0"
            chart.y_axis.majorTickMark = "out"
            chart.y_axis.tickLblPos    = "nextTo"
            chart.y_axis.delete        = False
            chart.y_axis.scaling.min   = 0
            chart.y_axis.scaling.max   = 100

            # ── X axis: date strings, tick marks visible ──────────────────────
            chart.x_axis.title         = "Date"
            chart.x_axis.majorTickMark = "out"
            chart.x_axis.tickLblPos    = "low"
            chart.x_axis.delete        = False
            # Reduce label density so dates don't overlap on long ranges
            if n_days > 60:
                chart.x_axis.tickLblSkip = 7
            elif n_days > 14:
                chart.x_axis.tickLblSkip = 2

            data_ref = Reference(ws, min_col=8, min_row=trend_start,
                                 max_row=trend_end)
            chart.add_data(data_ref)

            # ── StrRef forces <c:strRef> so Excel renders date strings ────────
            try:
                from openpyxl.chart.data_source import DataSource, StrRef
                from openpyxl.utils import get_column_letter as _gcl
                cat_formula = (f"'{ws.title}'!$B${trend_start}:$B${trend_end}")
                chart.series[0].cat = DataSource(strRef=StrRef(f=cat_formula))
            except ImportError:
                dates_ref = Reference(ws, min_col=2, min_row=trend_start,
                                      max_row=trend_end)
                chart.set_categories(dates_ref)

            # ── Series styling: Cohesity green line + circle markers ──────────
            s = chart.series[0]
            s.title = SeriesLabel(v="Success %")
            s.graphicalProperties.line.solidFill = COH_GREEN
            s.graphicalProperties.line.width     = 18000   # 2 pt in EMUs
            s.marker.symbol  = "circle"
            s.marker.size    = 5
            s.marker.graphicalProperties.solidFill   = COH_GREEN
            s.marker.graphicalProperties.line.solidFill = COH_GREEN

            # ── Legend at the bottom ──────────────────────────────────────────
            try:
                from openpyxl.chart.legend import Legend as _ChartLegend
                chart.legend = _ChartLegend()
                chart.legend.position = "b"
            except Exception:
                pass

            # Place success-rate chart to the right of the data columns (col L = 12)
            ws.add_chart(chart, f"L{trend_start}")

            # ── Duration chart: average backup duration per day ───────────────
            dur_chart = LineChart()
            dur_chart.style  = 10
            dur_chart.height = 16
            dur_chart.width  = 26

            _dur_title = f"{cd['name']} — Avg Daily Backup Duration (min)"
            try:
                from openpyxl.chart.title import Title as _CTitle2
                from openpyxl.drawing.text import (
                    RichTextBodyProperties as _BodyPr2,
                    Paragraph as _Para2,
                    RegularTextRun as _Run2,
                    RunProperties as _RunPr2,
                )
                from openpyxl.chart.text import RichText as _CRichText2, Text as _CText2
                _ct2 = _CTitle2()
                _ct2.overlay = False
                _rich2 = _CRichText2()
                _rich2.bodyPr = _BodyPr2()
                _para2 = _Para2()
                _run2  = _Run2()
                _run2.t  = _dur_title
                _run2.rPr = _RunPr2(sz=1400, b=True)
                _para2.r.append(_run2)
                _rich2.p.append(_para2)
                _ct2.tx = _CText2(rich=_rich2)
                dur_chart.title = _ct2
            except Exception:
                dur_chart.title = _dur_title

            dur_chart.y_axis.title         = "Avg Duration (min)"
            dur_chart.y_axis.numFmt        = "0.0"
            dur_chart.y_axis.majorTickMark = "out"
            dur_chart.y_axis.tickLblPos    = "nextTo"
            dur_chart.y_axis.delete        = False

            dur_chart.x_axis.title         = "Date"
            dur_chart.x_axis.majorTickMark = "out"
            dur_chart.x_axis.tickLblPos    = "low"
            dur_chart.x_axis.delete        = False
            if n_days > 60:
                dur_chart.x_axis.tickLblSkip = 7
            elif n_days > 14:
                dur_chart.x_axis.tickLblSkip = 2

            dur_ref = Reference(ws, min_col=11, min_row=trend_start,
                                max_row=trend_end)
            dur_chart.add_data(dur_ref)

            try:
                from openpyxl.chart.data_source import DataSource as _DS2, StrRef as _SR2
                dur_chart.series[0].cat = _DS2(strRef=_SR2(
                    f=f"'{ws.title}'!$B${trend_start}:$B${trend_end}"))
            except (ImportError, IndexError):
                pass

            try:
                ds2 = dur_chart.series[0]
                ds2.title = SeriesLabel(v="Avg Duration (min)")
                ds2.graphicalProperties.line.solidFill = "FF8C00"   # orange
                ds2.graphicalProperties.line.width     = 18000
                ds2.marker.symbol  = "diamond"
                ds2.marker.size    = 5
                ds2.marker.graphicalProperties.solidFill        = "FF8C00"
                ds2.marker.graphicalProperties.line.solidFill   = "FF8C00"
            except (IndexError, AttributeError):
                pass

            try:
                from openpyxl.chart.legend import Legend as _ChartLegend2
                dur_chart.legend = _ChartLegend2()
                dur_chart.legend.position = "b"
            except Exception:
                pass

            # Place duration chart below the success-rate chart (offset by 30 rows)
            ws.add_chart(dur_chart, f"L{trend_start + 30}")

    auto_fit_columns(ws)


def _sheet_recommendations(wb, all_data):
    ws = wb.create_sheet("Recommendations")
    ws.freeze_panes = "A3"
    _title(ws, "Prioritized Recommendations", "G")

    cols = ["#", "Priority", "Cluster", "Category",
            "Finding", "Recommended Action", "Business Impact"]
    _hdr(ws, 2, cols)

    all_recs = sorted(
        [(cd["name"], rec) for cd in all_data for rec in cd["recommendations"]],
        key=lambda x: _PRIORITY_ORDER.get(x[1]["priority"], 99)
    )
    for i, (cname, rec) in enumerate(all_recs, 1):
        row = [i, rec["priority"], cname, rec["category"],
               rec["finding"], rec["action"], rec["impact"]]
        rn = ws.max_row + 1
        for c, val in enumerate(row, 1):
            ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL

        pc = ws.cell(row=rn, column=2)
        if rec["priority"] == "CRITICAL":
            pc.fill = _fill(RED);        pc.font = _font(bold=True, color=WHITE)
        elif rec["priority"] == "HIGH":
            pc.fill = _fill(ORANGE);     pc.font = _font(bold=True)
        elif rec["priority"] == "MEDIUM":
            pc.fill = _fill(YELLOW);     pc.font = _font(bold=True)
        else:
            pc.fill = _fill(LIGHT_GRAY); pc.font = _font(bold=True)

    auto_fit_columns(ws)
    for col_letter in ("E", "F", "G"):
        ws.column_dimensions[col_letter].width = 55


def _sheet_disk_health(wb, all_data):
    ws = wb.create_sheet("Disk Health")
    ws.freeze_panes = "A3"
    _title(ws, "Disk Health — Per-Disk Status, SSD Wear & Encryption", "L")

    cols = ["Cluster", "Node ID", "Node IP", "Disk ID", "Disk Type",
            "Model", "Serial", "Status", "SSD Wear %",
            "Capacity (GB)", "Encrypted", "Storage Tier"]
    _hdr(ws, 2, cols)

    any_data = False
    for cd in all_data:
        disks = cd.get("disks") or []
        if not disks:
            continue
        any_data = True
        for disk in disks:
            status  = (disk.get("diskStatus") or disk.get("status") or "Unknown")
            # ssdUsedPercentage is the ELF field name (raw 0-100 integer)
            wear    = (disk.get("ssdUsedPercentage") or disk.get("ssdWearLevelPct")
                       or disk.get("wearLevel") or disk.get("lifeUsedPct"))
            _us     = disk.get("usageStats")
            # capacityInBytes is the ELF field name
            cap_b   = ((isinstance(_us, dict) and _us.get("totalPhysicalCapacityBytes"))
                       or disk.get("capacityInBytes") or disk.get("capacityBytes") or 0)
            # encryptionStatus is the ELF field name (string like "kEncrypted")
            enc_st  = disk.get("encryptionStatus") or ""
            enc     = bool(disk.get("encryptionEnabled") or disk.get("encrypted")
                           or (isinstance(enc_st, str) and "encrypt" in enc_st.lower()))
            tier    = disk.get("storageTier") or disk.get("diskTier") or ""
            dtype   = disk.get("diskType") or disk.get("type") or ""

            row = [cd["name"],
                   disk.get("_nodeId", ""), disk.get("_nodeIp", ""),
                   disk.get("diskId") or disk.get("id") or "",
                   dtype,
                   disk.get("diskModel") or disk.get("model") or "",
                   disk.get("diskSerial") or disk.get("serialNumber") or "",
                   status,
                   wear if wear is not None else "",
                   round(bytes_to_gb(cap_b), 2) if cap_b else "",
                   "Yes" if enc else "No",
                   tier]
            rn = ws.max_row + 1
            for c, val in enumerate(row, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL

            # color Status cell (col 8)
            sc = ws.cell(row=rn, column=8)
            if status.lower() in ("kfailed", "failed", "kmissing", "missing"):
                sc.fill = _fill(RED);    sc.font = _font(bold=True, color=WHITE)
            elif status.lower() in ("kdegraded", "degraded", "kmarkedforremoval",
                                    "marked_for_removal"):
                sc.fill = _fill(ORANGE); sc.font = _font(bold=True)
            elif status.lower() in ("khealthy", "healthy", "kgood", "good"):
                sc.fill = _fill(LT_GREEN)

            # color SSD Wear cell (col 9)
            if wear is not None:
                wc = ws.cell(row=rn, column=9)
                if wear >= SSD_WEAR_CRIT_PCT:
                    wc.fill = _fill(RED);    wc.font = _font(bold=True, color=WHITE)
                elif wear >= SSD_WEAR_WARN_PCT:
                    wc.fill = _fill(ORANGE); wc.font = _font(bold=True)

    if not any_data:
        ws.cell(row=3, column=1,
                value="Disk detail not available via this API path (private endpoint).") \
           .font = _font(color="595959")

    auto_fit_columns(ws)


def _sheet_agent_health(wb, all_data):
    import datetime as _dt
    ws = wb.create_sheet("Agent Health")
    ws.freeze_panes = "A3"
    _title(ws, "Agent Health — Registered Host Agent Status & Versions", "I")

    cols = ["Cluster", "Host", "IP Address", "OS Type",
            "Agent Version", "Status", "Upgradable", "Cert Expiry", "Notes"]
    _hdr(ws, 2, cols)

    today    = _dt.date.today()
    any_data = False
    for cd in all_data:
        agents = cd.get("agents") or []
        if not agents:
            continue
        any_data = True
        for agent in agents:
            _hs2   = agent.get("healthStatus")
            status = ((isinstance(_hs2, dict) and _hs2.get("status"))
                      or (isinstance(_hs2, str) and _hs2)
                      or agent.get("agentStatus") or agent.get("status") or "")
            upg    = (agent.get("upgradability") or agent.get("upgradable") or "")

            # Certificate expiry (try usecs, then msecs)
            cert_exp = ""
            exp_raw  = (agent.get("certExpiryUsecs")
                        or agent.get("certExpiryTimeMsecs"))
            if exp_raw:
                if exp_raw > 1_000_000_000_000_000:
                    exp_dt = _dt.date.fromtimestamp(exp_raw / 1_000_000)
                elif exp_raw > 1_000_000_000_000:
                    exp_dt = _dt.date.fromtimestamp(exp_raw / 1_000)
                else:
                    exp_dt = _dt.date.fromtimestamp(exp_raw)
                days_left = (exp_dt - today).days
                cert_exp  = f"{exp_dt.isoformat()} ({days_left}d)"

            notes = ""
            err   = agent.get("lastUpgradeError")
            if err:
                notes = f"Upgrade error: {str(err)[:80]}"

            # Route hostname vs IP by content — the API sometimes puts FQDNs in hostIp
            _raw_host = agent.get("hostName") or agent.get("host") or ""
            _raw_ip   = agent.get("hostIp") or ""
            _host, _ip = _split_host_ip(_raw_host, _raw_ip)

            row = [cd["name"],
                   _host,
                   _ip,
                   agent.get("hostOsType") or agent.get("osType") or "",
                   agent.get("agentVersion") or "",
                   status, upg, cert_exp, notes]
            rn = ws.max_row + 1
            for c, val in enumerate(row, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL

            sc = ws.cell(row=rn, column=6)
            if status.lower() in ("kunhealthy", "unhealthy", "unreachable",
                                  "kfailed", "failed"):
                sc.fill = _fill(RED);    sc.font = _font(bold=True, color=WHITE)
            elif status.lower() in ("kwarning", "warning", "kdegraded",
                                    "degraded"):
                sc.fill = _fill(ORANGE); sc.font = _font(bold=True)
            elif status.lower() in ("khealthy", "healthy", "krunning",
                                    "running"):
                sc.fill = _fill(LT_GREEN)

            # Upgradable flag
            uc = ws.cell(row=rn, column=7)
            if "upgradable" in upg.lower():
                uc.fill = _fill(YELLOW); uc.font = _font(bold=True)

    if not any_data:
        ws.cell(row=3, column=1,
                value="No agent data returned from /public/reports/agents.") \
           .font = _font(color="595959")

    auto_fit_columns(ws)
    ws.column_dimensions["B"].width = 35


def _sheet_source_coverage(wb, all_data):
    ws = wb.create_sheet("Source Coverage")
    ws.freeze_panes = "A3"
    _title(ws, "Source Coverage — Registered Sources & Protection Coverage", "H")

    cols = ["Cluster", "Source Name", "Environment", "Total Objects",
            "Protected", "Unprotected", "Coverage %", "Status"]
    _hdr(ws, 2, cols)

    for cd in all_data:
        sources = cd.get("sources") or []
        if not sources:
            continue
        for src in sources:
            if not isinstance(src, dict):
                continue
            _si  = src.get("sourceInfo")
            si   = _si if isinstance(_si, dict) else src
            name = (si.get("name") or si.get("sourceName")
                    or src.get("name") or "")
            env  = (si.get("environment") or si.get("sourceEnvironment")
                    or src.get("environment") or "")

            _ss  = src.get("stats");       _ss  = _ss  if isinstance(_ss,  dict) else {}
            _oc  = src.get("objectCount"); _oc  = _oc  if isinstance(_oc,  dict) else {}
            prot   = (src.get("protectedObjectsCount")
                      or src.get("numProtectedObjects")
                      or _ss.get("protectedObjectsCount")
                      or _oc.get("protectedCount") or 0)
            unprot = (src.get("unprotectedObjectsCount")
                      or src.get("numUnprotectedObjects")
                      or _ss.get("unprotectedObjectsCount")
                      or _oc.get("unprotectedCount") or 0)
            total  = prot + unprot
            pct    = round(prot / total * 100, 1) if total else ""

            if isinstance(pct, float):
                if pct >= 95:    flag = "Good"
                elif pct >= 75:  flag = f"Partial ({pct:.0f}%)"
                else:            flag = f"Low ({pct:.0f}%)"
            else:
                flag = "Unknown"

            row = [cd["name"], name, env,
                   total if total else "",
                   prot or "", unprot or "", pct, flag]
            rn = ws.max_row + 1
            for c, val in enumerate(row, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL

            fc = ws.cell(row=rn, column=8)
            if isinstance(pct, float) and pct < 75:
                fc.fill = _fill(RED);    fc.font = _font(bold=True, color=WHITE)
            elif isinstance(pct, float) and pct < 95:
                fc.fill = _fill(ORANGE); fc.font = _font(bold=True)
            elif isinstance(pct, float):
                fc.fill = _fill(LT_GREEN)

    auto_fit_columns(ws)


def _sheet_unprotected_objects(wb, all_data):
    """Named Unprotected Objects — lists sources with unprotected object details."""
    ws = wb.create_sheet("Unprotected Objects")
    ws.freeze_panes = "A3"
    _title(ws, "Unprotected Objects — Sources with Unprotected Object Details", "E")

    cols = ["Cluster", "Source Name", "Environment",
            "Unprotected Count", "Object Names"]
    _hdr(ws, 2, cols)

    any_data = False
    for cd in all_data:
        sources = cd.get("sources") or []
        for src in sources:
            if not isinstance(src, dict):
                continue
            _si   = src.get("sourceInfo")
            si    = _si if isinstance(_si, dict) else src
            name  = (si.get("name") or si.get("sourceName")
                     or src.get("name") or "")
            env   = (si.get("environment") or si.get("sourceEnvironment")
                     or src.get("environment") or "")
            _ss   = src.get("stats")
            _ss   = _ss if isinstance(_ss, dict) else {}
            unprot = (src.get("unprotectedObjectsCount")
                      or src.get("numUnprotectedObjects")
                      or _ss.get("unprotectedObjectsCount") or 0)
            if not unprot:
                continue

            # Try to extract object names from various possible structures
            obj_names = []
            for key in ("unprotectedObjects", "unprotectedObjectNames",
                        "objects", "leafObjects"):
                items = src.get(key) or si.get(key) or []
                if not isinstance(items, list):
                    continue
                for item in items:
                    if isinstance(item, str):
                        obj_names.append(item)
                    elif isinstance(item, dict):
                        n = (item.get("name") or item.get("objectName")
                             or item.get("displayName") or "")
                        if n:
                            obj_names.append(n)
                if obj_names:
                    break

            obj_str = (", ".join(obj_names[:20])
                       + (" and more" if len(obj_names) > 20 else "")
                       if obj_names else "See source details")

            row = [cd["name"], name, env, unprot, obj_str]
            rn  = ws.max_row + 1
            for c, val in enumerate(row, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL

            cnt_cell = ws.cell(row=rn, column=4)
            if unprot > 10:
                cnt_cell.fill = _fill(RED);    cnt_cell.font = _font(bold=True, color=WHITE)
                for col in range(1, 6):
                    ws.cell(row=rn, column=col).fill = _fill(RED)
                    ws.cell(row=rn, column=col).font = _font(bold=True, color=WHITE)
            else:
                for col in range(1, 6):
                    ws.cell(row=rn, column=col).fill = _fill(YELLOW)
                    ws.cell(row=rn, column=col).font = _font(bold=True)

            any_data = True

    if not any_data:
        ws.cell(row=3, column=1,
                value="No unprotected objects found across all sources.") \
           .font = _font(color="595959")

    auto_fit_columns(ws)
    ws.column_dimensions["E"].width = 60   # Object Names


def _sheet_user_security(wb, all_data):
    import datetime as _dt
    ws = wb.create_sheet("User Security")
    ws.freeze_panes = "A3"
    _title(ws, "User Security — Accounts, MFA Status & Access Control", "K")

    cols = ["Cluster", "Username", "Domain", "Type",
            "Roles", "Active", "Locked", "MFA Enabled",
            "Last Login", "Notes", "Risk Flag"]
    _hdr(ws, 2, cols)

    _ADMIN_ROLES = {"COHESITY_ADMIN", "Admin", "kAdmin", "kSuperAdmin",
                    "super_admin", "admin", "cluster_admin", "cohesity_admin"}
    _STALE_DAYS  = 30
    any_data     = False

    for cd in all_data:
        users = cd.get("users") or []
        if not users:
            continue
        any_data = True
        now_dt = _dt.datetime.now(tz=_dt.timezone.utc)

        for user in users:
            roles  = ", ".join(user.get("roles") or [])
            locked = bool(user.get("locked") or user.get("isLocked"))
            active = not bool(user.get("isDeleted") or user.get("isSuspended"))
            mfa    = bool(user.get("mfaEnabled") or user.get("isMfaEnabled"))

            last_ms = (user.get("lastLoginTimeMsecs")
                       or user.get("lastLoginTime") or 0)
            if last_ms:
                if last_ms > 1_000_000_000_000_000:  # usecs
                    last_ms = last_ms // 1000
                last_login_dt = _dt.datetime.fromtimestamp(
                    last_ms / 1000, tz=_dt.timezone.utc)
                last_login = last_login_dt.strftime("%Y-%m-%d")
            else:
                last_login_dt = None
                last_login = "Never"

            # Determine if admin: role name contains admin/Admin or matches known admin roles
            user_roles_lower = {r.lower() for r in (user.get("roles") or [])}
            is_admin = bool(
                user_roles_lower & _ADMIN_ROLES
                or any("admin" in r for r in user_roles_lower)
            )
            domain = user.get("domain") or ""
            utype  = (user.get("userType")
                      or ("Local" if domain.lower() in ("local", "")
                          else "AD/LDAP"))

            notes = []
            if locked:      notes.append("LOCKED")
            if not mfa:     notes.append("No MFA")

            # Risk Flag logic
            risk_flags = []
            if is_admin:
                # Stale Admin: last login > 30 days ago or never
                if last_login_dt is None or (now_dt - last_login_dt).days > _STALE_DAYS:
                    risk_flags.append("Stale Admin")
                if not mfa:
                    risk_flags.append("Admin No MFA")
                if locked:
                    risk_flags.append("Locked Admin")
            risk_flag_str = " | ".join(risk_flags) if risk_flags else ""

            row = [cd["name"], user.get("username", ""), domain, utype,
                   roles,
                   "Yes" if active else "No",
                   "Yes" if locked else "No",
                   "Yes" if mfa else "No",
                   last_login,
                   "; ".join(notes) if notes else "OK",
                   risk_flag_str]
            rn = ws.max_row + 1
            for c, val in enumerate(row, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL

            if locked:
                ws.cell(row=rn, column=7).fill = _fill(RED)
                ws.cell(row=rn, column=7).font = _font(bold=True, color=WHITE)

            # MFA color: red for admins, yellow for non-admins without MFA
            if not mfa and is_admin:
                ws.cell(row=rn, column=8).fill = _fill(RED)
                ws.cell(row=rn, column=8).font = _font(bold=True, color=WHITE)
            elif not mfa:
                ws.cell(row=rn, column=8).fill = _fill(YELLOW)
                ws.cell(row=rn, column=8).font = _font(bold=True)

            # Risk Flag column coloring (col 11)
            if risk_flag_str:
                rf_cell = ws.cell(row=rn, column=11)
                if "Admin No MFA" in risk_flag_str:
                    rf_cell.fill = _fill("FF0000")
                    rf_cell.font = _font(bold=True, color=WHITE)
                elif "Stale Admin" in risk_flag_str:
                    rf_cell.fill = _fill("FFD966")
                    rf_cell.font = _font(bold=True)

    if not any_data:
        ws.cell(row=3, column=1,
                value="No user data returned from /public/users.") \
           .font = _font(color="595959")

    # ── Custom Role Definitions sub-section ──────────────────────────────────
    _BUILTIN_ROLE_NAMES = {"Admin", "Viewer", "BackupOperator", "BackupAdmin",
                           "RestoreOperator", "ReplicationOperator",
                           "COHESITY_ADMIN", "kAdmin", "kSuperAdmin",
                           "ReadOnlyUser", "ManageProtectionGroupsOperator"}

    _roles_blank = ws.max_row + 2
    ws.cell(row=_roles_blank, column=1,
            value="Custom Role Definitions").font = _font(bold=True, size=11)

    role_cols = ["Cluster", "Role Name", "Description", "Privileges (count)",
                 "Has Cluster Modify", "Has User Modify", "Has Security Modify", "Risk Level"]
    _hdr(ws, _roles_blank + 1, role_cols)

    any_role_data = False
    for cd in all_data:
        for role in (cd.get("roles") or []):
            rname = role.get("name") or ""
            # Skip built-in roles
            is_builtin = (role.get("isBuiltIn") or not role.get("isCustomRole", True))
            if rname in _BUILTIN_ROLE_NAMES or is_builtin:
                continue
            # Also skip if role name matches a built-in (case-insensitive)
            if rname.lower() in {b.lower() for b in _BUILTIN_ROLE_NAMES}:
                continue

            any_role_data = True
            privs   = role.get("privileges") or role.get("permissions") or []
            priv_count = len(privs) if isinstance(privs, list) else 0
            privs_upper = " ".join(str(p) for p in (privs if isinstance(privs, list) else [])).upper()

            # Check for specific privilege types
            priv_words = privs_upper.split()
            has_cluster_mod = (
                any("CLUSTER" in pw for pw in priv_words) and
                any(x in privs_upper for x in ("MODIFY", "CREATE", "DELETE"))
            ) or ("CLUSTER" in privs_upper and
                  any(x in privs_upper for x in ("MODIFY", "CREATE", "DELETE")))
            has_user_mod    = "USER" in privs_upper and "MODIFY" in privs_upper
            has_sec_mod     = any(x in privs_upper for x in ("SECURITY", "PRINCIPAL"))

            if has_cluster_mod or has_user_mod:
                risk = "HIGH"
            elif has_sec_mod:
                risk = "MEDIUM"
            else:
                risk = "LOW"

            desc = role.get("description") or ""
            role_row = [cd["name"], _safe_cell(rname), _safe_cell(desc),
                        priv_count,
                        "Yes" if has_cluster_mod else "No",
                        "Yes" if has_user_mod    else "No",
                        "Yes" if has_sec_mod     else "No",
                        risk]
            rn = ws.max_row + 1
            for c, val in enumerate(role_row, 1):
                ws.cell(row=rn, column=c, value=val).font = _FONT_NORMAL

            # Risk coloring (col 8)
            risk_cell = ws.cell(row=rn, column=8)
            if risk == "HIGH":
                risk_cell.fill = _fill(RED);    risk_cell.font = _font(bold=True, color=WHITE)
            elif risk == "MEDIUM":
                risk_cell.fill = _fill(YELLOW); risk_cell.font = _font(bold=True)
            elif risk == "LOW":
                risk_cell.fill = _fill(LT_GREEN)

    if not any_role_data:
        rn = ws.max_row + 1
        ws.cell(row=rn, column=1,
                value="No custom roles found — all roles are built-in defaults.") \
           .font = _font(color="595959")

    auto_fit_columns(ws)
    ws.column_dimensions["E"].width = 30


# ─────────────────────────────────────────────────────────────────────────────
# TOPOLOGY DIAGRAM  (draw.io editable + matplotlib preview)
# ─────────────────────────────────────────────────────────────────────────────
#
# Cohesity brand palette used throughout:
_DG_TEAL    = "#96C73D"   # Cohesity green          — source clusters
_DG_TEAL_DK = "#6B8F2C"   # Darker green            — cluster stroke
_DG_NAVY    = "#1A3045"   # Dark navy               — header / background
_DG_REPL    = "#0062B1"   # Blue                   — replication targets
_DG_ARCH    = "#E07A00"   # Amber                  — archival vaults
_DG_FK      = "#1A5C3A"   # Dark green             — FortKnox vaults
_DG_FG      = "#FFFFFF"   # White text on dark boxes


def _topology_data(all_data):
    """Extract the topology graph from all_data.

    Returns a dict:
      clusters   — list of cluster dicts (id, name, groups, sources, workloads, cap_tb)
      repl_nodes — list of unique replication-target dicts (id, name)
      arch_nodes — list of unique archival-vault dicts (id, name)
      fk_nodes   — list of unique FortKnox-vault dicts (id, name)
      edges      — list of edge dicts (src, tgt, type="repl"|"arch"|"fk")
    """
    _fk_types = {"krpaas", "kfortknox", "kfort_knox", "krpaasarchival"}

    def _is_fk(tt, tname=""):
        tl = tt.lower(); nl = tname.lower()
        return any(x in tl for x in _fk_types) or "fortknox" in tl or "fortknox" in nl

    clusters = []
    all_repl  = {}   # name → node dict
    all_arch  = {}
    all_fk    = {}
    edges     = []

    for i, cd in enumerate(all_data):
        cid = f"c{i}"

        # Workload types (de-prefix the leading "k", title-case)
        workloads = sorted({
            g.get("environment", "").lstrip("k").title()
            for g in cd["groups"] if g.get("environment")
        })[:4]

        # Approximate protected capacity (sum logical bytes across groups)
        raw_cap = 0
        for g in cd["groups"]:
            st = g.get("latestSuccessfulRunStats") or g.get("lastSuccessfulRunStats") or {}
            raw_cap += st.get("logicalSizeBytes") or 0
        if raw_cap == 0:
            raw_cap = cd["info"].get("usedCapacityBytes") or 0
        cap_tb = raw_cap / (1024 ** 4) if raw_cap else 0

        for p in cd["policies"]:
            rtp = p.get("remoteTargetPolicy") or {}
            for t in (rtp.get("replicationTargets") or []):
                name = (t.get("targetName")
                        or (t.get("remoteTargetConfig") or {}).get("clusterName")
                        or (t.get("remoteTargetConfig") or {}).get("name") or "")
                if not name:
                    continue
                if name not in all_repl:
                    all_repl[name] = {"id": f"r{len(all_repl)}", "name": name}
                edges.append({"src": cid, "tgt": all_repl[name]["id"], "type": "repl"})

            for t in (rtp.get("archivalTargets") or []):
                cfg   = t.get("archivalTargetConfig") or {}
                tt    = cfg.get("targetType") or t.get("targetType") or ""
                tname = t.get("targetName") or cfg.get("name") or cfg.get("vaultName") or ""
                if not tname:
                    continue
                if _is_fk(tt, tname):
                    if tname not in all_fk:
                        all_fk[tname] = {"id": f"f{len(all_fk)}", "name": tname}
                    edges.append({"src": cid, "tgt": all_fk[tname]["id"], "type": "fk"})
                else:
                    if tname not in all_arch:
                        all_arch[tname] = {"id": f"a{len(all_arch)}", "name": tname}
                    edges.append({"src": cid, "tgt": all_arch[tname]["id"], "type": "arch"})

        # Pick up FK vaults that appear in the vaults list but not in policies.
        # Previously these only got a node entry but never an edge — fixed.
        for v in (cd.get("vaults") or []):
            vtype = v.get("vaultType") or ""
            vname = v.get("name") or ""
            if vname and _is_fk(vtype, vname):
                if vname not in all_fk:
                    all_fk[vname] = {"id": f"f{len(all_fk)}", "name": vname}
                edges.append({"src": cid, "tgt": all_fk[vname]["id"], "type": "fk"})

        # Pull software version for node labels
        cluster_version = (
            cd.get("version") or
            (cd.get("info") or {}).get("clusterSoftwareVersion") or
            (cd.get("info") or {}).get("softwareVersion") or ""
        )

        clusters.append({
            "id":       cid,
            "name":     cd["name"],
            "version":  cluster_version,
            "groups":   len(cd["groups"]),
            "sources":  len(cd.get("sources") or []),
            "workloads": workloads,
            "cap_tb":   cap_tb,
        })

    # Deduplicate edges (multiple policies can produce identical edges)
    seen = set()
    unique_edges = []
    for e in edges:
        k = (e["src"], e["tgt"], e["type"])
        if k not in seen:
            seen.add(k)
            unique_edges.append(e)

    return {
        "clusters":   clusters,
        "repl_nodes": list(all_repl.values()),
        "arch_nodes": list(all_arch.values()),
        "fk_nodes":   list(all_fk.values()),
        "edges":      unique_edges,
    }


# ── Comprehensive PowerPoint deck ────────────────────────────────────────────


def _generate_topology_drawio(all_data, out_path):
    """Removed in v1.50 — draw.io output superseded by write_pptx() deck."""
    return None



def write_pptx(all_data, args, out_path):
    """Write a comprehensive multi-section .pptx health check deck.

    Sections: Cover · Executive Summary (Snapshot, Scorecard, Capacity,
    Protection, Ransomware, Top Risks, Priority Actions) · Environment
    Topology · Security Deep-Dive (Posture, User Security, Audit,
    Recommendations) · Backup Engineering (Lifecycle, Coverage Gaps,
    Agents, Source Coverage, Workload Risk, Recommendations).

    Returns the output path on success, None if python-pptx is unavailable.
    """
    if not PPTX_OK:
        return None

    print("  Writing PowerPoint deck...")
    import datetime as _dt
    import lxml.etree as _letree
    from pptx.util import Pt as _Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn as _qn

    customer = getattr(args, "customer", "") or ""
    quick    = getattr(args, "quick", False)
    today    = _dt.date.today().strftime("%B %d, %Y")

    # ── Colour palette ────────────────────────────────────────────────────────
    _C_HDR    = "#67BF1B"   # Cohesity vibrant green — header bars & accents
    _C_DIV    = "#1A5C3A"   # Dark green (topology/misc legacy use)
    _C_WHT    = "#FFFFFF"
    _C_BLK    = "#0C0C0C"
    _C_DGRY   = "#404040"
    _C_LGRY   = "#F5F5F5"   # alternating table row tint
    _C_GRN_BG = "#E2EFDA"; _C_GRN_TX = "#375623"
    _C_AMB_BG = "#FFF2CC"; _C_AMB_TX = "#7F6000"
    _C_RED_BG = "#FFCCCC"; _C_RED_TX = "#9C0006"
    # Cover / divider / TOC palette — Cohesity 2025 template
    _C_BG     = "#FFFFFF"   # White background
    _C_VACC   = "#67BF1B"   # Vibrant Cohesity green accent
    _C_BGROW  = "#F4F6F4"   # Very light green-tinted white for TOC cards
    _C_FLBG   = "#F8FAF8"   # Near-white slide tint (light slides)

    # ── Presentation setup ────────────────────────────────────────────────────
    prs = _PptxPresentation()
    prs.slide_width  = _PptxInches(13.33)
    prs.slide_height = _PptxInches(7.5)
    W, H = 13.33, 7.5

    # ── Low-level helpers ─────────────────────────────────────────────────────
    def _rgb(hex6):
        h = hex6.lstrip("#")
        return _PptxRGB(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def _blank():
        return prs.slides.add_slide(prs.slide_layouts[6])

    def _rect(slide, x, y, w, h, fill, line=None, lw=0.01, rounded=False):
        st = _MSO_SHAPE.ROUNDED_RECTANGLE if rounded else _MSO_SHAPE.RECTANGLE
        sp = slide.shapes.add_shape(
            st, _PptxInches(x), _PptxInches(y),
            _PptxInches(w), _PptxInches(h))
        sp.fill.solid()
        sp.fill.fore_color.rgb = _rgb(fill)
        if line:
            sp.line.color.rgb = _rgb(line)
            sp.line.width = _PptxInches(lw)
        else:
            sp.line.fill.background()
        return sp

    def _txt(slide, text, x, y, w, h, size=10, bold=False,
             color="#000000", align="left", wrap=True):
        tb = slide.shapes.add_textbox(
            _PptxInches(x), _PptxInches(y),
            _PptxInches(w), _PptxInches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = str(text) if text is not None else ""
        run.font.size = _Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)
        return tb

    def _header(slide, title, subtitle=""):
        _rect(slide, 0, 0, W, 0.77, _C_HDR)
        _txt(slide, title, 0.15, 0.06, W - 0.3, 0.38,
             size=20, bold=True, color=_C_WHT)
        if subtitle:
            _txt(slide, subtitle, 0.15, 0.50, W - 0.3, 0.24,
                 size=9, color="#DDEEDD")

    _page_n = [0]  # mutable page counter for data slides

    def _footer(slide, show_page=True):
        # COHESITY wordmark — left
        _txt(slide, "COHESITY", 0.15, H - 0.30, 1.20, 0.24,
             size=9, bold=True, color=_C_BLK)
        # Copyright — centre
        _txt(slide, "\u00a9 2026 Cohesity Inc. All rights reserved.",
             1.40, H - 0.30, W - 4.0, 0.24,
             size=7, color="#888888")
        # Version + date — right
        _txt(slide, f"v{__version__}  \u00b7  {today}",
             W - 2.6, H - 0.30, 2.4, 0.24,
             size=7, color="#AAAAAA", align="right")
        # Page number — bottom centre (data slides only)
        if show_page:
            _page_n[0] += 1
            _txt(slide, str(_page_n[0]),
                 W / 2 - 0.20, H - 0.30, 0.40, 0.24,
                 size=7, color="#888888", align="center")

    def _section_div(title, subtitle=""):
        slide = _blank()
        _rect(slide, 0, 0, W, H, _C_BG)              # white background
        _rect(slide, 0, 0, 0.06, H, _C_VACC)         # thin green left bar
        # Green block — left ~62%, vertically centred (matches Cohesity chapter slide)
        _rect(slide, 0.06, 1.85, 8.20, 2.95, _C_VACC)
        _txt(slide, title, 0.32, 2.10, 7.70, 2.20,
             size=36, bold=True, color=_C_WHT, align="left")
        if subtitle:
            _txt(slide, subtitle, 0.15, 4.95, 9.00, 0.55,
                 size=11, color="#555555", align="left")
        _footer(slide, show_page=False)
        return slide

    def _rag(score, hi=75, lo=50):
        if score >= hi: return (_C_GRN_BG, _C_GRN_TX)
        if score >= lo: return (_C_AMB_BG, _C_AMB_TX)
        return (_C_RED_BG, _C_RED_TX)

    def _grade_color(grade):
        g = str(grade).upper()
        if g == "A": return (_C_GRN_BG, _C_GRN_TX)
        if g == "B": return ("#DFF0D8", "#2E7D32")
        if g == "C": return (_C_AMB_BG, _C_AMB_TX)
        return (_C_RED_BG, _C_RED_TX)

    def _sg(s):
        return ("A" if s >= 90 else "B" if s >= 80 else
                "C" if s >= 70 else "D" if s >= 60 else "F")

    def _ov(cd):
        return int(_score_protection(cd, quick) * 0.25
                   + _score_sla(cd, quick) * 0.15
                   + _score_infra(cd) * 0.15
                   + _score_capacity(cd) * 0.20
                   + _score_security(cd) * 0.15
                   + _score_alerts(cd) * 0.10)

    def _yn(v):
        return "Yes" if v else "No"

    def _cx(col_w):
        """Return X offset to horizontally centre a table on the slide."""
        return max(0.1, (W - sum(col_w)) / 2)

    def _notes(slide, text):
        """Set speaker-notes text on a slide."""
        try:
            slide.notes_slide.notes_text_frame.text = text
        except Exception:
            pass

    # Shape-type constants for topology nodes
    _ST_ROUND  = _MSO_SHAPE.ROUNDED_RECTANGLE
    _ST_CLOUD  = getattr(_MSO_SHAPE, "CLOUD", _MSO_SHAPE.ROUNDED_RECTANGLE)
    _ST_CAN    = getattr(_MSO_SHAPE, "CAN",   _MSO_SHAPE.ROUNDED_RECTANGLE)

    def _table(slide, headers, rows, col_w, x, y, rh=0.27,
               hdr_bg=None, fills=None, body_sz=10):
        """Styled table. fills[row][col] = (bg_hex, tx_hex) or None."""
        # Auto-fit column widths to content
        _cw_ch = max(body_sz, 10) * 0.58 / 72  # proportional char-width estimate (inches)
        _pad_w = 0.22
        auto_w = []
        for ci in range(len(headers)):
            best = len(str(headers[ci]))
            for row in rows:
                if ci < len(row) and row[ci] is not None:
                    best = max(best, max(len(s) for s in str(row[ci]).split('\n')))
            auto_w.append(max(0.50, round(best * _cw_ch + _pad_w, 3)))
        _avail = W - 0.30
        if sum(auto_w) > _avail:
            _s = _avail / sum(auto_w)
            auto_w = [round(w * _s, 3) for w in auto_w]
        col_w = auto_w
        x = round((W - sum(col_w)) / 2, 3)

        nr, nc = len(rows) + 1, len(headers)
        tw = sum(col_w)
        tbl = slide.shapes.add_table(
            nr, nc,
            _PptxInches(x), _PptxInches(y),
            _PptxInches(tw), _PptxInches(rh * nr)).table
        for ci, cw in enumerate(col_w):
            tbl.columns[ci].width = _PptxInches(cw)
        for ri in range(nr):
            tbl.rows[ri].height = _PptxInches(rh)
        hbg = hdr_bg or _C_HDR

        def _cell(cell, text, bg, fg, bold=False, sz=8.5, align="left"):
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = {"left": PP_ALIGN.LEFT,
                           "center": PP_ALIGN.CENTER,
                           "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
            run = p.add_run()
            run.text = str(text) if text is not None else ""
            run.font.size = _Pt(sz)
            run.font.bold = bold
            run.font.color.rgb = _rgb(fg)
            tc = cell._tc
            tcPr = tc.find(_qn("a:tcPr"))
            if tcPr is None:
                tcPr = _letree.SubElement(tc, _qn("a:tcPr"))
            for old in tcPr.findall(_qn("a:solidFill")):
                tcPr.remove(old)
            sf = _letree.SubElement(tcPr, _qn("a:solidFill"))
            sr = _letree.SubElement(sf, _qn("a:srgbClr"))
            sr.set("val", bg.lstrip("#"))

        for ci, hdr in enumerate(headers):
            _cell(tbl.cell(0, ci), hdr, hbg, _C_WHT,
                  bold=True, sz=10, align="left")
        for ri, row in enumerate(rows):
            alt = _C_LGRY if ri % 2 == 0 else _C_WHT
            for ci, val in enumerate(row):
                fill = (fills[ri][ci]
                        if fills and ri < len(fills)
                        and ci < len(fills[ri]) else None)
                if fill and fill[0] is not None:
                    bg, fg = fill
                else:
                    bg, fg = alt, _C_BLK
                _cell(tbl.cell(ri + 1, ci), val, bg, fg, sz=body_sz)

    def _tbl_y(n_data_rows, rh=0.27, top=0.90, bot=None):
        """Return table top Y — anchored just below the header bar."""
        return top

    def _layout_mixed(n_main_rows, meth_total_h, rh=0.27):
        """Return (main_tbl_y, meth_y) — table top-aligned, methodology below."""
        main_h = (n_main_rows + 1) * rh
        gap    = 0.22
        main_y = round(0.90, 3)
        meth_y = round(main_y + main_h + gap, 3)
        return main_y, meth_y

    def _kpi(slide, label, value, x, y, w=1.90, h=1.25, bg=None, vc=None):
        accent = bg or _C_HDR
        # Card: White, Background 1, Darker 15% = #D9D9D9; thin border
        card = _rect(slide, x, y, w, h, "#D9D9D9", line="#BBBBBB", lw=0.012, rounded=True)
        # Outer shadow — "intense" effect
        spPr = card._element.find(_qn("p:spPr"))
        eff  = _letree.SubElement(spPr, _qn("a:effectLst"))
        shdw = _letree.SubElement(eff,  _qn("a:outerShdw"))
        shdw.set("blurRad", "50800"); shdw.set("dist", "25400")
        shdw.set("dir",     "2700000"); shdw.set("rotWithShape", "0")
        sclr = _letree.SubElement(shdw, _qn("a:srgbClr")); sclr.set("val", "000000")
        _letree.SubElement(sclr, _qn("a:alpha")).set("val", "30000")
        # Thin left accent bar (inset slightly to respect rounded corners)
        _rect(slide, x + 0.015, y + 0.04, 0.07, h - 0.08, accent)
        # Label — small, uppercase, muted
        _txt(slide, label.upper(), x + 0.14, y + 0.10, w - 0.18, 0.26,
             size=7.5, color="#595959")
        # Value — large, bold, accent colour for visual coding
        _txt(slide, str(value), x + 0.14, y + 0.38, w - 0.18, 0.66,
             size=26, bold=True, color=vc if vc else accent)


    _P_COLORS = {"CRITICAL": (_C_RED_BG, _C_RED_TX),
                 "HIGH":     ("#FFE0CC", "#7F3000"),
                 "MEDIUM":   (_C_AMB_BG, _C_AMB_TX),
                 "LOW":      (_C_GRN_BG, _C_GRN_TX)}
    _ADMIN_ROLES = {"COHESITY_ADMIN", "Admin", "kAdmin", "kSuperAdmin"}

    # ── Native PowerPoint section markers ─────────────────────────────────────
    # Tracked as (section_name, first_slide_index); populated at each section
    # boundary below and applied via _register_sections() before save().
    _sec_starts = []

    def _register_sections(prs, sec_starts):
        """Inject <p14:sectionLst> into the presentation XML so PowerPoint's
        slide panel groups slides under named section headers.
        """
        import uuid as _uuid
        _NS14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
        sld_id_lst = prs.element.find(_qn("p:sldIdLst"))
        if sld_id_lst is None or not sec_starts:
            return
        slide_ids = [int(el.get("id"))
                     for el in sld_id_lst.findall(_qn("p:sldId"))]
        total = len(slide_ids)
        sec_lst = _letree.Element(f"{{{_NS14}}}sectionLst")
        for i, (name, start_idx) in enumerate(sec_starts):
            end_idx = (sec_starts[i + 1][1]
                       if i + 1 < len(sec_starts) else total)
            sec_el = _letree.SubElement(sec_lst, f"{{{_NS14}}}section")
            sec_el.set("name", name)
            sec_el.set("id", "{" + str(_uuid.uuid4()).upper() + "}")
            ids_el = _letree.SubElement(sec_el, f"{{{_NS14}}}sldIdLst")
            for idx in range(start_idx, end_idx):
                if idx < total:
                    sld_el = _letree.SubElement(ids_el, f"{{{_NS14}}}sldId")
                    sld_el.set("id", str(slide_ids[idx]))
        prs.element.append(sec_lst)

    # ── Cover slide ───────────────────────────────────────────────────────────
    slide = _blank()
    _rect(slide, 0, 0, W, H, _C_BG)          # white background
    _rect(slide, 0, 0, 0.06, H, _C_VACC)     # thin green left bar
    # Green title block (matches Cohesity cover template)
    _rect(slide, 0.06, 1.30, 8.20, 3.90, _C_VACC)
    # Title text inside green block — white bold
    _txt(slide, "Cohesity Environment",
         0.28, 1.55, 7.80, 0.80,
         size=36, bold=True, color=_C_WHT, align="left")
    _txt(slide, "Health Check Report",
         0.28, 2.38, 7.80, 0.80,
         size=36, bold=True, color=_C_WHT, align="left")
    # Thin white rule inside block
    _rect(slide, 0.28, 3.27, 3.80, 0.04, _C_WHT)
    # Customer / date inside block
    if customer:
        _txt(slide, customer, 0.28, 3.40, 7.80, 0.55,
             size=20, bold=True, color=_C_WHT, align="left")
    _txt(slide, today, 0.28, 3.97, 7.80, 0.38,
         size=13, color="#DDFFD0", align="left")
    # Below block
    _txt(slide, "Prepared using Cohesity Helios",
         0.15, 5.38, 8.00, 0.32,
         size=9, color="#888888", align="left")
    _footer(slide, show_page=False)
    _notes(slide,
           "Cohesity Environment Health Check Report. "
           "This deck was generated live from the Cohesity Helios API and covers all "
           "clusters connected to Helios at the time of the report. "
           "Sections: Executive Summary (cluster snapshot, scorecard, capacity, "
           "protection, ransomware readiness, priority actions), "
           "Environment Topology (connection diagram), "
           "Security Deep-Dive (posture, user security, audit, recommendations), "
           "Backup Engineering (lifecycle, coverage gaps, agents, source coverage, "
           "workload risk, engineering recommendations), "
           "Scoring Methodology (reference for all scoring models). "
           "Use alongside the Excel workbook (detailed data) and Word document "
           "(narrative report) for a complete customer business review."
           + (f" Customer: {customer}." if customer else ""))
    _sec_starts.append(("Cover", 0))

    # ── Contents slide (placeholder — populated after all slides are built) ───
    toc_slide = _blank()

    # ── Section: Executive Summary ────────────────────────────────────────────
    _sec_starts.append(("Executive Summary", len(prs.slides)))
    _section_div("Executive Summary",
                 "Snapshot  \u00b7  Scorecard  \u00b7  Capacity  \u00b7  "
                 "Protection  \u00b7  Ransomware  \u00b7  Actions")

    # ── Environment Snapshot ──────────────────────────────────────────────────
    slide = _blank()
    _header(slide, "Environment Snapshot",
            f"{len(all_data)} cluster(s)  \u00b7  {today}")
    _footer(slide)
    n_nodes    = sum(len(cd["nodes"])    for cd in all_data)
    n_groups   = sum(len(cd["groups"])   for cd in all_data)
    n_policies = sum(len(cd["policies"]) for cd in all_data)
    n_vaults   = sum(len(cd["vaults"])   for cd in all_data)
    n_crits    = sum(
        sum(1 for a in cd["alerts"] if a.get("severity") == "kCritical")
        for cd in all_data)
    kpi_data = [
        ("Clusters",        len(all_data), "#2E5A8E",  None),
        ("Nodes",           n_nodes,       _C_VACC,   None),
        ("Prot. Groups",    n_groups,      _C_HDR,    None),
        ("Policies",        n_policies,    "#5B6E8A", None),
        ("Vaults",          n_vaults,      _DG_ARCH,  None),
        ("Critical Alerts", n_crits,
         "#C0392B" if n_crits else "#27AE60",
         "#C0392B" if n_crits else "#27AE60"),
    ]
    kpi_w, kpi_gap = 1.90, 0.19
    kpi_x0 = (W - (kpi_w * 6 + kpi_gap * 5)) / 2
    for i, (lbl, val, bg, vc) in enumerate(kpi_data):
        _kpi(slide, lbl, val,
             kpi_x0 + i * (kpi_w + kpi_gap), 0.82, kpi_w, 1.25, bg, vc)
    hdr_snap = ["Cluster", "Version", "Nodes", "Groups",
                "Vaults", "Crit Alerts", "Overall Score", "Grade"]
    rows_snap, fills_snap = [], []
    for cd in all_data:
        ver   = (cd["info"].get("clusterSoftwareVersion") or "—").strip()
        crits = sum(1 for a in cd["alerts"] if a.get("severity") == "kCritical")
        ov    = _ov(cd)
        grd   = _sg(ov)
        rows_snap.append([cd["name"], ver, len(cd["nodes"]),
                          len(cd["groups"]), len(cd["vaults"]), crits, ov, grd])
        cr_f = (_C_RED_BG, _C_RED_TX) if crits else (_C_GRN_BG, _C_GRN_TX)
        fills_snap.append(
            [None, None, None, None, None, cr_f, _rag(ov), _grade_color(grd)])
    _cw_snap = [2.5, 1.4, 0.8, 0.8, 0.8, 1.0, 1.3, 0.8]
    _kpi_bot = 0.82 + 1.25 + 0.20          # tiles bottom + gap
    _table(slide, hdr_snap, rows_snap, _cw_snap,
           _cx(_cw_snap), _tbl_y(len(rows_snap), top=_kpi_bot),
           fills=fills_snap)
    _notes(slide,
           "Summary of all Cohesity clusters in scope. "
           "KPI tiles at the top show total clusters, nodes, protection groups, "
           "policies, vaults, and open critical alerts. "
           "The cluster table shows per-cluster software version, group and vault "
           "counts, overall health score (0-100), and letter grade (A>=90, B>=80, "
           "C>=70, D>=60, F<60). "
           "Red overall score = below 50 (At Risk), Amber = 50-74 (Fair), "
           "Green = 75+ (Good/Excellent). "
           "Critical alerts in red indicate open issues requiring immediate action.")

    # ── Health Scorecard ──────────────────────────────────────────────────────
    slide = _blank()
    _header(slide, "Health Scorecard",
            "Per-cluster scores (0\u2013100)  \u00b7  "
            "Protection\u00d725%  SLA\u00d715%  "
            "Infra\u00d715%  Capacity\u00d720%  Security\u00d715%  Alerts\u00d710%")
    _footer(slide)
    hdr_sc = ["Cluster", "Protection", "SLA", "Infra",
              "Capacity", "Security", "Alerts", "Overall", "Grade"]
    rows_sc, fills_sc = [], []
    for cd in all_data:
        p   = _score_protection(cd, quick)
        sl  = _score_sla(cd, quick)
        inf = _score_infra(cd)
        cap = _score_capacity(cd)
        sec = _score_security(cd)
        al  = _score_alerts(cd)
        ov  = int(p*0.25 + sl*0.15 + inf*0.15 + cap*0.20 + sec*0.15 + al*0.10)
        grd = _sg(ov)
        rows_sc.append([cd["name"], p, sl, inf, cap, sec, al, ov, grd])
        fills_sc.append(
            [None] + [_rag(s) for s in (p, sl, inf, cap, sec, al, ov)]
            + [_grade_color(grd)])
    _cw_sc = [2.6, 1.1, 1.0, 1.1, 1.1, 1.1, 1.0, 1.1, 0.8]
    rows_sc  = rows_sc[:10];  fills_sc  = fills_sc[:10]
    # meth block: rule(0.015) + gap(0.05) + label(0.22) + gap(0.27) + 7-row table + gap(0.08) + grade(0.22)
    _sc_meth_h = 0.015 + 0.05 + 0.22 + 0.27 + 7 * 0.22 + 0.08 + 0.22  # ≈ 2.42"
    _sc_y, _sc_my = _layout_mixed(len(rows_sc), _sc_meth_h)
    _table(slide, hdr_sc, rows_sc, _cw_sc,
           _cx(_cw_sc), _sc_y, fills=fills_sc)
    # ── Inline scoring methodology ──────────────────────────────────────────
    _rect(slide, 0.15, _sc_my,         W - 0.30, 0.015, "#CCCCCC")
    _txt(slide, "SCORING METHODOLOGY", 0.15, _sc_my + 0.05, W - 0.30, 0.22,
         size=7.5, bold=True, color="#888888")
    _hs_cols = [2.6, 1.0, 1.5, 1.5, 1.5]
    _table(slide,
           ["Dimension", "Weight", "Green \u226575", "Amber 50\u201374", "Red <50"],
           [
               ["Protection Success", "25%", "\u226595%",      "90\u201394%",   "<90%"],
               ["SLA Compliance",     "15%", "\u226599%",      "90\u201398%",   "<90%"],
               ["Infrastructure",     "15%", "100% nodes",     "\u226590%",     "<90%"],
               ["Storage Capacity",   "20%", "\u226475% full", "75\u201385%",   ">85%"],
               ["Security Score",     "15%", "\u226575 pts",   "50\u201374",    "<50"],
               ["Alerts",             "10%", "No criticals",   "\u22641 crit",  ">1 crit"],
           ],
           _hs_cols, _cx(_hs_cols), round(_sc_my + 0.32, 3), rh=0.22, body_sz=8,
           fills=[[None, None,
                   (_C_GRN_BG, _C_GRN_TX),
                   (_C_AMB_BG, _C_AMB_TX),
                   (_C_RED_BG, _C_RED_TX)]] * 6)
    _txt(slide,
         "Grade: A\u226590  B\u226580  C\u226570  D\u226560  F<60"
         "  \u00b7  Green\u226575  Amber 50\u201374  Red <50",
         0.15, round(_sc_my + 0.32 + 7 * 0.22 + 0.10, 3),
         W - 0.30, 0.22, size=7.5, color=_C_DGRY, align="center")
    _notes(slide,
           "Per-cluster health scorecard with six weighted dimensions. "
           "Protection Success (25%) — backup run success rate in the lookback window. "
           "SLA Compliance (15%) — percentage of runs completing within their SLA window. "
           "Infrastructure (15%) — ratio of healthy nodes to total nodes. "
           "Storage Capacity (20%) — score based on used % of usable capacity. "
           "Security (15%) — encryption, vaults, replication, DataLock, MFA posture. "
           "Alerts (10%) — deducted for open critical and warning alerts. "
           "Overall score is the weighted sum; Grade: A>=90, B>=80, C>=70, D>=60, F<60. "
           "Red = below 50, Amber = 50-74, Green = 75 or above.")

    # ── Capacity & Growth ─────────────────────────────────────────────────────
    slide = _blank()
    _header(slide, "Storage Capacity & Growth",
            "Per-cluster utilisation and predictive runway to 80% full")
    _footer(slide)
    hdr_cap = ["Cluster", "Storage Domain", "Used (TB)", "Usable (TB)",
               "Used %", "Data Reduction", "Runway (days)", "80% Full By"]
    rows_cap, fills_cap = [], []
    for cd in all_data:
        stats_c  = (cd["info"].get("stats") or {})
        usage_c  = (stats_c.get("usagePerfStats") or {})
        usable_c = usage_c.get("physicalCapacityBytes") or 0
        used_c   = usage_c.get("totalPhysicalUsageBytes") or 0
        dr_c     = stats_c.get("dataReductionRatio") or 0
        pct_c    = (used_c / usable_c * 100) if usable_c else 0
        runway_d = cd.get("capacity_runway") or {}
        dq       = runway_d.get("data_quality", "no_data")
        d80      = runway_d.get("days_to_80")
        proj_80  = runway_d.get("projected_80")
        _tb = lambda b: f"{b/1e12:.1f}" if b else "\u2014"
        if dq in ("no_data", "insufficient"):
            runway_str = "Insufficient data"
            date_str   = "\u2014"
        elif d80 is None:
            runway_str = "Stable"
            date_str   = "\u2014"
        elif d80 == 0:
            runway_str = ">80% now"
            date_str   = str(proj_80) if proj_80 else "\u2014"
        else:
            runway_str = str(d80)
            date_str   = str(proj_80) if proj_80 else "\u2014"
        rows_cap.append([
            cd["name"], "(Cluster Total)",
            _tb(used_c), _tb(usable_c),
            f"{pct_c:.1f}%" if usable_c else "\u2014",
            f"{dr_c:.1f}x" if dr_c else "\u2014",
            runway_str, date_str,
        ])
        rw_val = d80 if isinstance(d80, int) else 999
        fills_cap.append([
            None, None, None, None,
            _rag(100 - pct_c, hi=70, lo=30) if usable_c else None,
            None,
            _rag(rw_val, hi=90, lo=30) if isinstance(d80, int) else None,
            None,
        ])
    _cw_cap = [2.2, 1.8, 1.0, 1.0, 0.9, 1.2, 1.2, 1.8]
    _table(slide, hdr_cap, rows_cap, _cw_cap,
           _cx(_cw_cap), _tbl_y(len(rows_cap)), fills=fills_cap)
    _notes(slide,
           "Storage utilisation and capacity runway per cluster. "
           "Used and Usable columns show cluster-total capacity in terabytes. "
           "Used % is coloured: Green = below 75%, Amber = 75-85%, Red = above 85%. "
           "Data Reduction reflects the combined deduplication and compression ratio. "
           "Runway (days) is a predictive forecast based on 30-day linear regression "
           "of daily capacity growth — how many days until the cluster reaches 80% full. "
           "Runway coloring: Green = more than 90 days, Amber = 30-90 days, "
           "Red = fewer than 30 days. 'Stable' means growth is flat or shrinking. "
           "'Insufficient data' means fewer than 3 data points were available for regression.")

    # ── Protection Summary ────────────────────────────────────────────────────
    slide = _blank()
    _header(slide, "Protection Summary",
            "Backup success rates, SLA compliance, and coverage gaps per cluster")
    _footer(slide)
    hdr_prot = ["Cluster", "Groups", "Success Rate", "SLA Pass %",
                "DataLock", "Coverage Gaps", "Prot. Score"]
    rows_prot, fills_prot = [], []
    _SUC = {"kSuccess", "Succeeded", "kWarning"}
    for cd in all_data:
        groups   = cd["groups"]
        policies = cd["policies"]
        succ_pct, sla_pct = _success_stats(cd)
        _, dl_lbl, _, _ = _cluster_datalock(policies, groups)
        gaps_n    = sum(1 for g in groups
                        if (g.get("lastRun") or {})
                           .get("localBackupInfo", {})
                           .get("status", "") not in _SUC)
        p_score   = _score_protection(cd, quick)
        rows_prot.append([cd["name"], len(groups),
                          f"{succ_pct:.1f}%", f"{sla_pct:.1f}%",
                          dl_lbl or "None", gaps_n, p_score])
        su_f = _rag(succ_pct, hi=95, lo=90)
        sl_f = _rag(sla_pct, hi=99, lo=95)
        gp_f = ((_C_GRN_BG, _C_GRN_TX) if gaps_n == 0 else
                (_C_AMB_BG, _C_AMB_TX) if gaps_n <= 5 else
                (_C_RED_BG, _C_RED_TX))
        dl_f = ((_C_RED_BG, _C_RED_TX)
                if not dl_lbl or dl_lbl.lower() in ("no", "none", "no datalock") else
                (_C_AMB_BG, _C_AMB_TX)
                if "partial" in (dl_lbl or "").lower() else
                (_C_GRN_BG, _C_GRN_TX))
        fills_prot.append([None, None, su_f, sl_f, dl_f, gp_f, _rag(p_score)])
    rows_prot = rows_prot[:10]; fills_prot = fills_prot[:10]
    _cw_prot = [2.6, 1.0, 1.5, 1.5, 2.0, 1.5, 1.7]
    # meth block: rule(0.015)+gap(0.05)+label(0.22)+gap(0.27)+2-row tbl(0.44)+gap(0.08)+note(0.22)
    _prot_meth_h = 0.015 + 0.05 + 0.22 + 0.27 + 2 * 0.22 + 0.08 + 0.22  # ≈ 1.30"
    _prot_y, _prot_my = _layout_mixed(len(rows_prot), _prot_meth_h)
    _table(slide, hdr_prot, rows_prot, _cw_prot,
           _cx(_cw_prot), _prot_y, fills=fills_prot)
    # ── Inline scoring methodology ──────────────────────────────────────────
    _rect(slide, 0.15, _prot_my,         W - 0.30, 0.015, "#CCCCCC")
    _txt(slide, "SCORING METHODOLOGY", 0.15, _prot_my + 0.05, W - 0.30, 0.22,
         size=7.5, bold=True, color="#888888")
    _pm_cols = [2.6, 1.0, 1.5, 1.5, 1.5]
    _table(slide,
           ["Dimension", "Max Pts", "Green \u226575", "Amber 50\u201374", "Red <50"],
           [["Backup Success Rate", "100 pts",
             "\u226595% runs pass", "90\u201394%", "<90%"]],
           _pm_cols, _cx(_pm_cols), round(_prot_my + 0.32, 3), rh=0.22, body_sz=8,
           fills=[[None, None,
                   (_C_GRN_BG, _C_GRN_TX),
                   (_C_AMB_BG, _C_AMB_TX),
                   (_C_RED_BG, _C_RED_TX)]])
    _txt(slide,
         "Quick mode: last-run status per group  \u00b7  "
         "Full mode: average daily success % across the lookback window  \u00b7  "
         "kWarning counted as success",
         0.15, round(_prot_my + 0.32 + 2 * 0.22 + 0.10, 3),
         W - 0.30, 0.22, size=7.5, color=_C_DGRY, align="center")
    _notes(slide,
           "Backup health summary per cluster across all protection groups. "
           "Success Rate = average of per-day success % across the lookback window "
           "(kSuccess + kWarning counted as success) — identical to the Trends tab. "
           "SLA Pass % = percentage of runs completed within their SLA window. "
           "DataLock (WORM) status: Red = no DataLock on any group (immutable backups "
           "absent — ransomware risk), Amber = partial coverage, Green = all groups covered. "
           "Coverage Gaps = groups whose last run was not successful. "
           "Protection Score reflects overall backup reliability for the cluster. "
           "See Excel Protection Health sheet for per-group detail.")

    # ── Ransomware Readiness ──────────────────────────────────────────────────
    slide = _blank()
    _header(slide, "Ransomware Readiness",
            "Composite score: DataLock \u00b7 FortKnox \u00b7 "
            "Recovery Window \u00b7 Quorum \u00b7 Admin MFA")
    _footer(slide)
    hdr_rw = ["Cluster", "RW Score", "Label",
              "DataLock", "FortKnox", "Admin MFA", "Quorum"]
    rows_rw, fills_rw = [], []
    for cd in all_data:
        rw_score, rw_label, _ = _ransomware_score(cd)
        fk = cd.get("fk_status", "none")
        _, dl_lbl, _, _ = _cluster_datalock(cd["policies"], cd["groups"])
        no_mfa = sum(
            1 for u in (cd.get("users") or [])
            if isinstance(u, dict)
            and not (u.get("mfaEnabled") or u.get("isMfaEnabled"))
            and bool(set(u.get("roles") or []) & _ADMIN_ROLES))
        mfa_str = "All enabled" if no_mfa == 0 else f"{no_mfa} missing"
        quorum  = (cd.get("mode") == "direct" or bool(cd.get("quorum_enabled")))
        fk_lbl  = {"active": "Active", "idle": "Idle (no data)",
                   "none": "Not configured"}.get(fk, "\u2014")
        rows_rw.append([cd["name"], rw_score, rw_label,
                        dl_lbl or "None", fk_lbl, mfa_str,
                        "Enabled" if quorum else "Not configured"])
        lv_f  = ((_C_GRN_BG, _C_GRN_TX) if rw_label == "Strong" else
                 (_C_AMB_BG, _C_AMB_TX) if rw_label == "Moderate" else
                 (_C_RED_BG, _C_RED_TX))
        fk_f  = ((_C_GRN_BG, _C_GRN_TX) if fk == "active" else
                 (_C_AMB_BG, _C_AMB_TX) if fk == "idle" else
                 (_C_RED_BG, _C_RED_TX))
        mfa_f = (_C_GRN_BG, _C_GRN_TX) if no_mfa == 0 else (_C_RED_BG, _C_RED_TX)
        qr_f  = (_C_GRN_BG, _C_GRN_TX) if quorum else (_C_AMB_BG, _C_AMB_TX)
        dl_f  = ((_C_RED_BG, _C_RED_TX)
                 if not dl_lbl or dl_lbl.lower() in ("no", "none", "no datalock") else
                 (_C_AMB_BG, _C_AMB_TX)
                 if "partial" in (dl_lbl or "").lower() else
                 (_C_GRN_BG, _C_GRN_TX))
        fills_rw.append([None, _rag(rw_score), lv_f, dl_f, fk_f, mfa_f, qr_f])
    _cw_rw = [2.4, 1.2, 1.5, 2.2, 1.8, 1.8, 1.5]
    rows_rw  = rows_rw[:10];  fills_rw  = fills_rw[:10]
    # meth: rule + label + 6-row table (rh=0.22) + legend text
    _rw_meth_h = 0.015 + 0.05 + 0.22 + 0.27 + 6 * 0.22 + 0.08 + 0.22  # ≈ 2.20"
    _rw_y, _rw_my = _layout_mixed(len(rows_rw), _rw_meth_h)
    _table(slide, hdr_rw, rows_rw, _cw_rw,
           _cx(_cw_rw), _rw_y, fills=fills_rw)
    # ── Inline scoring methodology ──────────────────────────────────────────
    _rect(slide, 0.15, _rw_my,         W - 0.30, 0.015, "#CCCCCC")
    _txt(slide, "SCORING METHODOLOGY", 0.15, _rw_my + 0.05, W - 0.30, 0.22,
         size=7.5, bold=True, color="#888888")
    _rw_cols = [2.4, 1.1, 8.5]
    _table(slide,
           ["Dimension", "Max Pts", "Detail"],
           [
               ["DataLock (WORM) coverage",   "30 pts",
                "30=Compliance mode, 20=any mode, 15=partial\u226550%, 8=partial<50%, 0=none"],
               ["FortKnox / indelible vault",  "25 pts",
                "25=Active, 10=Idle, 0=Not configured"],
               ["Clean recovery window",       "20 pts",
                "20=\u226530d, 15=14\u201329d, 10=7\u201313d, 5=2\u20136d, 2=<2d, 0=none"],
               ["Quorum",                      "10 pts",
                "10=enabled (or N/A for direct clusters), 0=disabled"],
               ["Admin MFA",                   "15 pts",
                "15=all admins enabled, 8=partial, 0=none enabled"],
           ],
           _rw_cols, _cx(_rw_cols), round(_rw_my + 0.32, 3), rh=0.22, body_sz=8)
    _txt(slide,
         "Strong \u226580  \u00b7  Moderate 60\u201379  \u00b7  "
         "High Risk 40\u201359  \u00b7  Critical <40",
         0.15, round(_rw_my + 0.32 + 6 * 0.22 + 0.10, 3),
         W - 0.30, 0.22, size=7.5, color=_C_DGRY, align="center")
    _notes(slide,
           "Composite ransomware recovery readiness score per cluster (0-100). "
           "DataLock (WORM) coverage (30 pts): Compliance mode = 30, any mode = 20, "
           "partial >=50% = 15, partial <50% = 8, none = 0 (RED — critical gap). "
           "FortKnox / indelible vault (25 pts): Active = 25, Idle = 10, None = 0. "
           "Admin MFA (15 pts): all admins enabled = 15, partial = 8, none = 0. "
           "Clean recovery window (20 pts): based on oldest successful snapshot age. "
           "Quorum (10 pts): enabled = 10. "
           "Labels: Strong (>=80), Moderate (60-79), High Risk (40-59), Critical (<40). "
           "Clusters scoring Critical or High Risk should be prioritised for immediate remediation.")

    # ── Top At-Risk Workloads ─────────────────────────────────────────────────
    slide = _blank()
    _header(slide, "Top At-Risk Workloads",
            "Protection groups sorted by recovery risk score  "
            "\u00b7  Lower score = higher risk")
    _footer(slide)
    risk_all = []
    for cd in all_data:
        for grp in cd["groups"]:
            score, level = _group_risk_score(grp, cd)
            lr  = grp.get("lastRun") or {}
            lb  = (lr.get("localBackupInfo") or {})
            risk_all.append((score, level, cd["name"],
                             grp.get("name", ""),
                             lb.get("status", "\u2014"),
                             _yn(bool(lr.get("isSlaViolated")))))
    risk_all.sort(key=lambda x: x[0])
    hdr_risk = ["Cluster", "Protection Group", "Risk Score",
                "Risk Level", "Last Status", "SLA Violated"]
    rows_risk, fills_risk = [], []
    for score, level, cname, gname, status, sla in risk_all[:16]:
        rows_risk.append([cname, gname, score, level, status, sla])
        lv_f  = ((_C_GRN_BG, _C_GRN_TX) if level == "Low" else
                 (_C_AMB_BG, _C_AMB_TX) if level == "Medium" else
                 (_C_RED_BG, _C_RED_TX))
        sla_f = ((_C_RED_BG, _C_RED_TX) if sla == "Yes"
                 else (_C_GRN_BG, _C_GRN_TX))
        fills_risk.append([None, None, _rag(score), lv_f, None, sla_f])
    if not rows_risk:
        rows_risk  = [["\u2014", "No workloads", "\u2014", "\u2014", "\u2014", "\u2014"]]
        fills_risk = [[None] * 6]
    _cw_risk = [2.2, 3.5, 1.2, 1.4, 2.1, 1.4]
    _table(slide, hdr_risk, rows_risk, _cw_risk,
           _cx(_cw_risk), _tbl_y(len(rows_risk)), fills=fills_risk)
    _notes(slide,
           "Top 16 protection groups most at risk of recovery failure, sorted worst-first. "
           "Risk Score factors: last run status (35 pts — Success=35, Warning=25, "
           "Running=20, Skipped=10, Paused=5, Failed=0), SLA compliance (25 pts), "
           "RPO gap from last successful run (25 pts — <=4h=25, <=8h=20, <=24h=15, "
           "<=48h=8, >48h=0), DataLock coverage (15 pts). "
           "Levels: Low (>=75), Medium (50-74), High (25-49), Critical (<25). "
           "Groups with Critical or High risk represent the highest recovery exposure — "
           "address SLA violations, failed runs, and missing DataLock coverage first. "
           "See Excel Workload Risk Heatmap sheet for all groups.")

    # ── Priority Action Items ─────────────────────────────────────────────────
    slide = _blank()
    _header(slide, "Priority Action Items",
            "CRITICAL and HIGH findings requiring immediate attention")
    _footer(slide)
    pri_rows, pri_fills = [], []
    for cd in all_data:
        for rec in _build_recommendations(cd):
            if rec["priority"] in ("CRITICAL", "HIGH"):
                pri_rows.append(
                    (rec["priority"], rec["category"], cd["name"],
                     rec["finding"][:80], rec["action"][:80]))
    pri_rows.sort(key=lambda x: (0 if x[0] == "CRITICAL" else 1, x[1]))
    if not pri_rows:
        pri_rows  = [("\u2014", "No CRITICAL or HIGH findings", "", "", "")]
        pri_fills = [[(_C_GRN_BG, _C_GRN_TX)] + [None] * 4]
    else:
        for prio, *_ in pri_rows[:15]:
            p_f = (_C_RED_BG, _C_RED_TX) if prio == "CRITICAL" \
                  else ("#FFE0CC", "#7F3000")
            pri_fills.append([p_f] + [None] * 4)
        pri_rows = pri_rows[:15]
    _cw_pri = [1.2, 1.8, 2.0, 4.1, 3.8]
    _table(slide,
           ["Priority", "Category", "Cluster", "Finding", "Recommended Action"],
           [(p, c, cl, f, a) for p, c, cl, f, a in pri_rows],
           _cw_pri, _cx(_cw_pri), _tbl_y(len(pri_rows)), fills=pri_fills)
    _notes(slide,
           "CRITICAL and HIGH priority findings requiring immediate attention, "
           "sorted by priority then category. "
           "CRITICAL items represent active risks with potential data loss or compliance "
           "breach — address within 24-48 hours. "
           "HIGH items are significant gaps that should be remediated within days to weeks. "
           "Each row shows the affected cluster, finding description (truncated — see Excel "
           "for full text), and recommended action. "
           "See the Excel Recommendations sheet for the complete list including MEDIUM and "
           "LOW priority items with business impact descriptions.")

    # ── Section: Topology ─────────────────────────────────────────────────────
    _sec_starts.append(("Environment Topology", len(prs.slides)))
    _section_div("Environment Topology",
                 "Cluster connections  \u00b7  Replication  "
                 "\u00b7  Archival  \u00b7  FortKnox / RPaaS")

    # ── Topology slide ────────────────────────────────────────────────────────
    topo       = _topology_data(all_data)
    t_clusters = topo["clusters"]
    repl_nodes = topo["repl_nodes"]
    arch_nodes = topo["arch_nodes"]
    fk_nodes   = topo["fk_nodes"]
    edges      = topo["edges"]
    slide  = _blank()
    shapes = slide.shapes
    _header(slide, "Environment Topology",
            "Source Clusters  \u00b7  Replication  \u00b7  Archival  "
            "\u00b7  FortKnox / RPaaS"
            + (f"  \u2014  {customer}" if customer else ""))
    _topo_pre = len(list(slide.shapes._spTree))

    T_NODE_W = 2.30; T_NODE_H = 0.85; T_V_GAP = 0.32
    T_COL_GAP = 0.80   # gap between C / R / A
    T_FK_GAP  = 1.60   # wider visual separation before FortKnox column
    # Horizontal centering: only for active (non-empty) columns
    _t_active = [(nlist, key) for nlist, key in [
        (t_clusters, "C"), (repl_nodes, "R"), (arch_nodes, "A"), (fk_nodes, "F")
    ] if nlist]
    _has_fk = any(k == "F" for _, k in _t_active)
    _t_active_no_fk = [(n, k) for n, k in _t_active if k != "F"]
    _t_total_w = (len(_t_active_no_fk) * T_NODE_W
                  + max(0, len(_t_active_no_fk) - 1) * T_COL_GAP
                  + (bool(_has_fk and _t_active_no_fk)) * (T_FK_GAP + T_NODE_W)
                  + (bool(_has_fk and not _t_active_no_fk)) * T_NODE_W)
    _t_x0 = (W - _t_total_w) / 2
    _t_xs = {}
    _xc = _t_x0
    for i, (_, key) in enumerate(_t_active):
        _t_xs[key] = _xc
        if i < len(_t_active) - 1:
            _next_key = _t_active[i + 1][1]
            _xc += T_NODE_W + (T_FK_GAP if _next_key == "F" else T_COL_GAP)
    T_COL_C = _t_xs.get("C", 0.20)
    T_COL_R = _t_xs.get("R", T_COL_C + T_NODE_W + T_COL_GAP)
    T_COL_A = _t_xs.get("A", T_COL_R + T_NODE_W + T_COL_GAP)
    T_COL_F = _t_xs.get("F", T_COL_A + T_NODE_W + T_FK_GAP)
    # Vertical centering: fit nodes + column headers in available space
    _n_max = max(len(t_clusters), len(repl_nodes) if repl_nodes else 0,
                 len(arch_nodes) if arch_nodes else 0,
                 len(fk_nodes)   if fk_nodes   else 0, 1)
    # Auto-scale if too many rows would overflow
    _max_col_h = 4.50
    _raw_col_h = _n_max * T_NODE_H + max(0, _n_max - 1) * T_V_GAP
    if _raw_col_h > _max_col_h:
        _s = _max_col_h / _raw_col_h
        T_NODE_H = round(T_NODE_H * _s, 3); T_V_GAP = round(T_V_GAP * _s, 3)
    _col_h   = _n_max * T_NODE_H + max(0, _n_max - 1) * T_V_GAP
    _hdr_h   = 0.26; _hdr_gap = 0.14
    _total_h = _hdr_h + _hdr_gap + _col_h
    _DIAG_TOP = 0.85; _DIAG_BOT = H - 0.30
    _top_margin = max(0.0, (_DIAG_BOT - _DIAG_TOP - _total_h) / 2)
    HDR_Y     = round(_DIAG_TOP + _top_margin, 3)
    T_START_Y = round(HDR_Y + _hdr_h + _hdr_gap, 3)

    # Source-column height used to vertically centre other columns
    _n_src   = max(len(t_clusters), 1)
    _src_h   = _n_src * T_NODE_H + max(0, _n_src - 1) * T_V_GAP

    def _col_y0(n_nodes):
        """Top Y for a column centred against the source-cluster column."""
        col_h = n_nodes * T_NODE_H + max(0, n_nodes - 1) * T_V_GAP
        return T_START_Y + max(0.0, (_src_h - col_h) / 2)

    def _t_ny(idx, y0=None):
        return (y0 if y0 is not None else T_START_Y) + idx * (T_NODE_H + T_V_GAP)

    def _t_box(label, sub, x, y, fill, stroke, shape_type=None):
        sp = shapes.add_shape(
            shape_type or _ST_ROUND,
            _PptxInches(x), _PptxInches(y),
            _PptxInches(T_NODE_W), _PptxInches(T_NODE_H))
        sp.fill.solid()
        sp.fill.fore_color.rgb = _rgb(fill)
        sp.line.color.rgb = _rgb(stroke)
        sp.line.width = _PptxInches(0.015)
        tf = sp.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = label; r1.font.bold = True
        r1.font.size = _Pt(9)
        r1.font.color.rgb = _rgb("#FFFFFF")
        if sub:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = sub
            r2.font.size = _Pt(7)
            r2.font.color.rgb = _rgb("#FFFFFF")

    def _t_conn(x1, y1, x2, y2, color):
        try:
            from pptx.enum.shapes import MSO_CONNECTOR_TYPE as _MCT
            conn = shapes.add_connector(
                _MCT.STRAIGHT,
                _PptxInches(x1), _PptxInches(y1),
                _PptxInches(x2), _PptxInches(y2))
            conn.line.color.rgb = _rgb(color)
            conn.line.width = _PptxInches(0.030)
        except Exception:
            pass
        return (x1 + x2) / 2, (y1 + y2) / 2

    def _t_lbl(text, cx, cy):
        tb2 = shapes.add_textbox(
            _PptxInches(cx - 0.42), _PptxInches(cy - 0.11),
            _PptxInches(0.84), _PptxInches(0.22))
        p = tb2.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text; r.font.size = _Pt(7)
        r.font.color.rgb = _rgb("#888888")

    t_anchor = {}
    for idx, cl in enumerate(t_clusters):
        y = _t_ny(idx)
        ver = cl.get("version", "")
        sub = f"v{ver}" if ver else ""
        if cl["groups"]:  sub += (" | " if sub else "") + f"{cl['groups']} grp"
        if cl["sources"]: sub += f" | {cl['sources']} src"
        _t_box(cl["name"], sub, T_COL_C, y, _DG_TEAL, _DG_TEAL_DK)
        t_anchor[cl["id"]] = (T_COL_C + T_NODE_W, y + T_NODE_H / 2, T_COL_C)
    _r_y0 = _col_y0(len(repl_nodes)) if repl_nodes else T_START_Y
    for idx, n in enumerate(repl_nodes):
        y = _t_ny(idx, _r_y0)
        _t_box(n["name"], "Cluster", T_COL_R, y, _DG_REPL, "#004A8A")
        t_anchor[n["id"]] = (T_COL_R + T_NODE_W, y + T_NODE_H / 2, T_COL_R)
    _a_y0 = _col_y0(len(arch_nodes)) if arch_nodes else T_START_Y
    for idx, n in enumerate(arch_nodes):
        y = _t_ny(idx, _a_y0)
        _t_box(n["name"], "Archival Vault", T_COL_A, y, _DG_ARCH, "#B05A00",
               shape_type=_ST_CAN)
        t_anchor[n["id"]] = (T_COL_A + T_NODE_W, y + T_NODE_H / 2, T_COL_A)
    _f_y0 = _col_y0(len(fk_nodes)) if fk_nodes else T_START_Y
    for idx, n in enumerate(fk_nodes):
        y = _t_ny(idx, _f_y0)
        _t_box(n["name"], "FortKnox / RPaaS", T_COL_F, y, _DG_FK, "#0E3B25",
               shape_type=_ST_CLOUD)
        t_anchor[n["id"]] = (T_COL_F + T_NODE_W, y + T_NODE_H / 2, T_COL_F)

    for nodelist, cx, lbl, color in [
        (t_clusters, T_COL_C, "Source Clusters",    _DG_TEAL),
        (repl_nodes, T_COL_R, "Replication Targets", _DG_REPL),
        (arch_nodes, T_COL_A, "Archival Vaults",     _DG_ARCH),
        (fk_nodes,   T_COL_F, "FortKnox / RPaaS",   _DG_FK),
    ]:
        if not nodelist: continue
        tb3 = shapes.add_textbox(
            _PptxInches(cx), _PptxInches(HDR_Y),
            _PptxInches(T_NODE_W), _PptxInches(0.26))
        p3 = tb3.text_frame.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = lbl; r3.font.bold = True
        r3.font.size = _Pt(9)
        r3.font.color.rgb = _rgb(color)

    _E_COLORS = {"repl": _DG_REPL, "arch": _DG_ARCH, "fk": _DG_FK}
    _E_LBLS   = {"repl": "Replication", "arch": "Archive",
                 "fk": "FortKnox/\nIndelible"}
    _seen_lbl = set()
    for e in edges:
        sid, tid, etype = e["src"], e["tgt"], e["type"]
        if sid not in t_anchor or tid not in t_anchor: continue
        sx, sy, _ = t_anchor[sid]
        _, ty, tx  = t_anchor[tid]
        cx, cy = _t_conn(sx, sy, tx, ty, _E_COLORS.get(etype, "#888888"))
        if (sid, etype) not in _seen_lbl:
            _seen_lbl.add((sid, etype))
            _t_lbl(_E_LBLS.get(etype, ""), cx, cy)

    # Wrap all topology shapes into a single group for easy user rearrangement
    _topo_els = list(slide.shapes._spTree)[_topo_pre:]
    if _topo_els:
        _EMU_W = int(W * 914400); _EMU_H = int(H * 914400)
        _all_ids = [int(el.get("id", 0)) for el in slide.shapes._spTree.iter() if el.get("id")]
        _grp_id  = (max(_all_ids) + 1) if _all_ids else 100
        grp = _letree.Element(_qn("p:grpSp"))
        nvGrpSpPr = _letree.SubElement(grp, _qn("p:nvGrpSpPr"))
        cNvPr = _letree.SubElement(nvGrpSpPr, _qn("p:cNvPr"))
        cNvPr.set("id", str(_grp_id)); cNvPr.set("name", "Topology Group")
        _letree.SubElement(nvGrpSpPr, _qn("p:cNvGrpSpPr"))
        _letree.SubElement(nvGrpSpPr, _qn("p:nvPr"))
        grpSpPr = _letree.SubElement(grp, _qn("p:grpSpPr"))
        xfrm = _letree.SubElement(grpSpPr, _qn("a:xfrm"))
        off = _letree.SubElement(xfrm, _qn("a:off"))
        off.set("x", "0"); off.set("y", "0")
        ext = _letree.SubElement(xfrm, _qn("a:ext"))
        ext.set("cx", str(_EMU_W)); ext.set("cy", str(_EMU_H))
        chOff = _letree.SubElement(xfrm, _qn("a:chOff"))
        chOff.set("x", "0"); chOff.set("y", "0")
        chExt = _letree.SubElement(xfrm, _qn("a:chExt"))
        chExt.set("cx", str(_EMU_W)); chExt.set("cy", str(_EMU_H))
        for el in _topo_els:
            slide.shapes._spTree.remove(el)
            grp.append(el)
        slide.shapes._spTree.append(grp)

    _notes(slide,
           "Environment topology map. Green rounded rectangles = source clusters "
           "(showing software version, protection group count, source count). "
           "Blue rounded rectangles = replication targets. "
           "Amber cylinders = archival vaults. "
           "Dark-green clouds = FortKnox / RPaaS indelible vaults. "
           "Arrows show data-flow direction with type labels (Replication / "
           "Archive / FortKnox-Indelible). Use this diagram to verify full "
           "resilience coverage — every cluster should have at least one "
           "replication target and one indelible vault.")
    _footer(slide)

    # ── Section: Security Deep-Dive ───────────────────────────────────────────
    _sec_starts.append(("Security Deep-Dive", len(prs.slides)))
    _section_div("Security Deep-Dive",
                 "Posture  \u00b7  User Security  "
                 "\u00b7  Audit & Governance  \u00b7  Recommendations")

    # ── Security Posture Overview ─────────────────────────────────────────────
    slide = _blank()
    _header(slide, "Security Posture Overview",
            "Encryption, vaults, DataLock, MFA, and quorum per cluster")
    _footer(slide)
    hdr_sec = ["Cluster", "Cluster Enc", "SD Enc", "Replication",
               "FortKnox", "DataLock", "Admin MFA", "Sec Score"]
    rows_sec, fills_sec = [], []
    for cd in all_data:
        c_enc, sd_enc = _sd_enc_status(cd)
        fk = cd.get("fk_status", "none")
        _, dl_lbl, _, _ = _cluster_datalock(cd["policies"], cd["groups"])
        has_repl = any(
            (p.get("remoteTargetPolicy") or {}).get("replicationTargets")
            for p in cd["policies"])
        no_mfa = sum(
            1 for u in (cd.get("users") or [])
            if isinstance(u, dict)
            and not (u.get("mfaEnabled") or u.get("isMfaEnabled"))
            and bool(set(u.get("roles") or []) & _ADMIN_ROLES))
        mfa_str = "All OK" if no_mfa == 0 else f"{no_mfa} missing"
        sec_sc  = _score_security(cd)
        sd_str  = ("Yes" if sd_enc is True else
                   "No" if sd_enc is False else "N/A")
        fk_str  = {"active": "Active", "idle": "Idle",
                   "none": "None"}.get(fk, "\u2014")
        rows_sec.append([cd["name"], _yn(c_enc), sd_str, _yn(has_repl),
                         fk_str, dl_lbl or "None", mfa_str, sec_sc])
        enc_f  = (_C_GRN_BG, _C_GRN_TX) if c_enc else (_C_RED_BG, _C_RED_TX)
        sd_f   = (_C_GRN_BG, _C_GRN_TX) if sd_enc else (_C_AMB_BG, _C_AMB_TX)
        rpl_f  = (_C_GRN_BG, _C_GRN_TX) if has_repl else (_C_RED_BG, _C_RED_TX)
        fk_f   = ((_C_GRN_BG, _C_GRN_TX) if fk == "active" else
                  (_C_AMB_BG, _C_AMB_TX) if fk == "idle" else
                  (_C_RED_BG, _C_RED_TX))
        dl_f   = ((_C_RED_BG, _C_RED_TX)
                  if not dl_lbl or dl_lbl.lower() in ("no", "none", "no datalock") else
                  (_C_AMB_BG, _C_AMB_TX)
                  if "partial" in (dl_lbl or "").lower() else
                  (_C_GRN_BG, _C_GRN_TX))
        mfa_f  = (_C_GRN_BG, _C_GRN_TX) if no_mfa == 0 else (_C_RED_BG, _C_RED_TX)
        fills_sec.append(
            [None, enc_f, sd_f, rpl_f, fk_f, dl_f, mfa_f, _rag(sec_sc)])
    _cw_sec = [2.3, 1.1, 0.9, 1.2, 1.6, 2.0, 1.3, 1.2]
    rows_sec  = rows_sec[:10];  fills_sec  = fills_sec[:10]
    # meth: rule + label + 12-row table (rh=0.20) + caption
    _sec_meth_h = 0.015 + 0.05 + 0.22 + 0.27 + 12 * 0.20 + 0.08 + 0.22  # ≈ 3.27"
    _sec_y, _sec_my = _layout_mixed(len(rows_sec), _sec_meth_h)
    _table(slide, hdr_sec, rows_sec, _cw_sec,
           _cx(_cw_sec), _sec_y, fills=fills_sec)
    # ── Inline scoring methodology ──────────────────────────────────────────
    _rect(slide, 0.15, _sec_my,        W - 0.30, 0.015, "#CCCCCC")
    _txt(slide, "SCORING METHODOLOGY", 0.15, _sec_my + 0.05, W - 0.30, 0.22,
         size=7.5, bold=True, color="#888888")
    _sec_cols = [3.8, 1.8]
    _table(slide,
           ["Control", "Points"],
           [
               ["Cluster encryption at rest",         "+15 pts"],
               ["Storage domain encryption",          "+15 pts"],
               ["Vault (archival target) configured", "+15 pts"],
               ["Replication target configured",      "+15 pts"],
               ["FortKnox active",                    "+20 pts"],
               ["FortKnox idle",                      "+10 pts"],
               ["DataLock Compliance mode",           "+20 pts"],
               ["DataLock any mode (no compliance)",  "+15 pts"],
               ["DataLock partial \u226550%",         " +8 pts"],
               ["DataLock partial <50%",              " +3 pts"],
               ["Admin accounts without MFA",         "\u221210 pts  (penalty)"],
           ],
           _sec_cols, _cx(_sec_cols), round(_sec_my + 0.32, 3), rh=0.20, body_sz=8,
           fills=[
               [None, (_C_GRN_BG, _C_GRN_TX)],
               [None, (_C_GRN_BG, _C_GRN_TX)],
               [None, (_C_GRN_BG, _C_GRN_TX)],
               [None, (_C_GRN_BG, _C_GRN_TX)],
               [None, (_C_GRN_BG, _C_GRN_TX)],
               [None, (_C_AMB_BG, _C_AMB_TX)],
               [None, (_C_GRN_BG, _C_GRN_TX)],
               [None, (_C_GRN_BG, _C_GRN_TX)],
               [None, (_C_AMB_BG, _C_AMB_TX)],
               [None, (_C_AMB_BG, _C_AMB_TX)],
               [None, (_C_RED_BG, _C_RED_TX)],
           ])
    _txt(slide,
         "Score capped at 100 \u00b7 Green\u226575  Amber 50\u201374  Red <50",
         0.15, round(_sec_my + 0.32 + 12 * 0.20 + 0.10, 3),
         W - 0.30, 0.22, size=7.5, color=_C_DGRY, align="center")
    _notes(slide,
           "Security posture scorecard across key protection dimensions per cluster. "
           "Cluster Enc: cluster-wide encryption at rest (Red = off, Green = on). "
           "SD Enc: storage domain encryption (Green = all encrypted, Amber = partial/unknown). "
           "Replication: at least one replication target configured (Red = none). "
           "FortKnox: indelible/air-gapped vault status — Active (Green), Idle (Amber), "
           "None (Red). "
           "DataLock: WORM immutable backup coverage — Full Compliance (Green), "
           "Partial (Amber), No DataLock (Red — ransomware risk). "
           "Admin MFA: multi-factor authentication on admin accounts (Red = missing). "
           "Security Score: overall 0-100 security posture score. "
           "Additive model: Cluster Enc +15, SD Enc +15, Vault +15, Replication +15, "
           "FortKnox Active +20 / Idle +10, DataLock Compliance +20 / any +15 / "
           "partial>=50% +8 / partial<50% +3. Penalty: -10 per admin lacking MFA. "
           "Thresholds: Green>=75, Amber 50-74, Red<50.")

    # ── User Security & MFA ───────────────────────────────────────────────────
    slide = _blank()
    _header(slide, "User Security & MFA",
            "User accounts \u2014 MFA status, locked state, and roles")
    _footer(slide)
    hdr_usr = ["Cluster", "Username", "Domain", "Roles",
               "MFA Enabled", "Locked"]
    rows_usr, fills_usr = [], []
    for cd in all_data:
        for u in (cd.get("users") or []):
            if not isinstance(u, dict): continue
            roles  = u.get("roles") or []
            mfa_on = bool(u.get("mfaEnabled") or u.get("isMfaEnabled"))
            locked = bool(u.get("locked") or u.get("isLocked"))
            uname  = u.get("username") or u.get("name") or "\u2014"
            domain = u.get("domain") or "LOCAL"
            rows_usr.append(
                [cd["name"], uname, domain,
                 ", ".join(roles[:3]) or "\u2014",
                 _yn(mfa_on), _yn(locked)])
            mfa_f = (_C_GRN_BG, _C_GRN_TX) if mfa_on else (_C_RED_BG, _C_RED_TX)
            lk_f  = (_C_RED_BG, _C_RED_TX) if locked else (_C_GRN_BG, _C_GRN_TX)
            fills_usr.append([None, None, None, None, mfa_f, lk_f])
    if not rows_usr:
        rows_usr  = [["\u2014", "No user data", "", "", "", ""]]
        fills_usr = [[None] * 6]
    _cw_usr = [2.0, 2.2, 1.5, 3.0, 1.5, 1.5]
    _table(slide, hdr_usr, rows_usr[:20], _cw_usr,
           _cx(_cw_usr), _tbl_y(min(len(rows_usr), 20)), fills=fills_usr[:20])
    _notes(slide,
           "User account security status across all clusters. "
           "MFA Enabled: Green = multi-factor authentication is active for this user, "
           "Red = MFA is off. Admin accounts without MFA are a HIGH severity security finding "
           "— they represent the primary attack vector for credential-based intrusions. "
           "Locked: Red = account is currently locked (possible brute-force attempt or "
           "deliberate lockout). "
           "Roles: shows up to 3 roles per user. COHESITY_ADMIN / kAdmin = full admin. "
           "Table shows first 20 users — see Excel User Security sheet for the full list "
           "with last-login timestamps.")
    if len(rows_usr) > 20:
        _txt(slide,
             f"  \u2026 and {len(rows_usr) - 20} more \u2014 "
             "see Excel User Security sheet",
             0.22, H - 0.55, W - 0.4, 0.28, size=8, color=_C_DGRY)

    # ── Governance & Audit Activity ───────────────────────────────────────────
    slide = _blank()
    _header(slide, "Governance & Audit Activity",
            "30-day configuration changes and high-risk events per cluster")
    _footer(slide)
    hdr_aud = ["Cluster", "Audit Logging", "Total Events",
               "Policy Changes", "Group Changes",
               "User Changes", "Vault/Security", "High-Risk Events"]
    rows_aud, fills_aud = [], []
    _HR_KW = {"delete", "destroy", "wipe", "disable", "remove",
              "revoke", "fortknox", "datalock", "mfa", "encryption"}
    for cd in all_data:
        audit_on = cd.get("audit_enabled", False)
        events   = cd.get("audit_events") or []
        cats     = {"Policy Changes": 0, "Group Changes": 0,
                    "User Changes": 0, "Vault/Security Changes": 0}
        high_risk = 0
        for ev in events:
            cat = _audit_category(ev)
            if cat and cat in cats: cats[cat] += 1
            act = (ev.get("action") or ev.get("type") or "").lower()
            if any(kw in act for kw in _HR_KW): high_risk += 1
        rows_aud.append([
            cd["name"], _yn(audit_on), len(events),
            cats["Policy Changes"], cats["Group Changes"],
            cats["User Changes"], cats["Vault/Security Changes"], high_risk])
        al_f = (_C_GRN_BG, _C_GRN_TX) if audit_on else (_C_RED_BG, _C_RED_TX)
        hr_f = ((_C_RED_BG, _C_RED_TX) if high_risk > 5 else
                (_C_AMB_BG, _C_AMB_TX) if high_risk > 0 else
                (_C_GRN_BG, _C_GRN_TX))
        fills_aud.append([None, al_f, None, None, None, None, None, hr_f])
    _cw_aud = [2.0, 1.3, 1.2, 1.5, 1.4, 1.4, 1.4, 1.6]
    _table(slide, hdr_aud, rows_aud, _cw_aud,
           _cx(_cw_aud), _tbl_y(len(rows_aud)), fills=fills_aud)
    _notes(slide,
           "30-day configuration change audit trail per cluster. "
           "Audit Logging must be enabled (Green) to collect events — if Red, "
           "events cannot be monitored (a HIGH security gap). "
           "Total Events = all change-type events in the lookback window (reads excluded). "
           "Categories: Policy Changes, Group Changes, User Changes, Vault/Security Changes. "
           "High-Risk Events = actions matching keywords: delete, destroy, disable, remove, "
           "revoke, datalock, fortknox, mfa, encryption. "
           "High-risk count > 5 = Red (investigate), > 0 = Amber (review), 0 = Green. "
           "See Excel Audit Log sheet for the full event list with timestamps and details.")

    # ── Security Recommendations ──────────────────────────────────────────────
    slide = _blank()
    _header(slide, "Security Recommendations",
            "Security, Ransomware, Encryption, MFA, Audit, and Vault findings")
    _footer(slide)
    _SEC_CATS = {"security", "ransomware", "fortknox", "datalock",
                 "encryption", "mfa", "quorum", "audit", "tls", "user"}
    sec_recs = []
    for cd in all_data:
        for rec in _build_recommendations(cd):
            if any(c in rec.get("category", "").lower() for c in _SEC_CATS):
                sec_recs.append(
                    (rec["priority"], rec["category"], cd["name"],
                     rec["finding"][:72], rec["action"][:72]))
    sec_recs.sort(key=lambda x:
                  ({"CRITICAL": 0, "HIGH": 1, "MEDIUM": 2,
                    "LOW": 3}.get(x[0], 4), x[1]))
    rows_sr, fills_sr = [], []
    for prio, cat, cname, finding, action in sec_recs[:15]:
        rows_sr.append([prio, cat, cname, finding, action])
        fills_sr.append([_P_COLORS.get(prio, (None, None))] + [None] * 4)
    if not rows_sr:
        rows_sr  = [["\u2014", "No security findings", "", "", ""]]
        fills_sr = [[(_C_GRN_BG, _C_GRN_TX)] + [None] * 4]
    _cw_sr = [1.2, 1.8, 2.0, 4.1, 3.8]
    _table(slide,
           ["Priority", "Category", "Cluster", "Finding", "Recommended Action"],
           rows_sr, _cw_sr, _cx(_cw_sr), _tbl_y(len(rows_sr)), fills=fills_sr)
    _notes(slide,
           "Security-focused findings from the automated analysis, sorted by priority. "
           "Covers: encryption (cluster and storage domain), DataLock/WORM configuration, "
           "FortKnox vault activity, admin MFA enforcement, quorum settings, "
           "audit logging status, TLS certificate expiry, and user security posture. "
           "CRITICAL = active exploitable risk. HIGH = significant gap. "
           "MEDIUM = best-practice improvement. LOW = optimisation. "
           "Remediate CRITICAL and HIGH items before the next customer review. "
           "See Excel Recommendations sheet for the complete prioritised list.")

    # ── Section: Backup Engineering ───────────────────────────────────────────
    _sec_starts.append(("Backup Engineering", len(prs.slides)))
    _section_div("Backup Engineering",
                 "Lifecycle  \u00b7  Coverage Gaps  \u00b7  Agents  "
                 "\u00b7  Source Coverage  \u00b7  Risk  \u00b7  Recommendations")

    # ── Software & Hardware Lifecycle ─────────────────────────────────────────
    slide = _blank()
    _header(slide, "Software & Hardware Lifecycle",
            "Cluster software versions and hardware end-of-life status")
    _footer(slide)
    hdr_sw = ["Cluster", "SW Version", "SW Status",
              "EOS Date", "Nodes", "EOL Nodes", "EOL Models"]
    rows_sw, fills_sw = [], []
    _today_d = _dt.date.today()
    for cd in all_data:
        sw_ver = (cd["info"].get("clusterSoftwareVersion") or "\u2014").strip()
        sw_status, sw_eos_date, sw_label = _sw_eos(sw_ver)
        eos_str = str(sw_eos_date) if sw_eos_date else "\u2014"
        eol_models, n_eol = [], 0
        for n in cd["nodes"]:
            model = n.get("hardwareModel") or n.get("productModel") or ""
            eol   = _hw_eol(model)
            if eol and eol < _today_d:
                n_eol += 1
                if model and model not in eol_models:
                    eol_models.append(model)
        rows_sw.append([
            cd["name"], sw_ver, sw_label or sw_status, eos_str,
            len(cd["nodes"]), n_eol,
            ", ".join(eol_models) or "None"])
        sw_f  = (_C_GRN_BG, _C_GRN_TX) if sw_status == "in_support" \
                else (_C_RED_BG, _C_RED_TX)
        eol_f = (_C_RED_BG, _C_RED_TX) if n_eol > 0 else (_C_GRN_BG, _C_GRN_TX)
        fills_sw.append([None, None, sw_f, None, None, eol_f, None])
    _cw_sw = [2.2, 1.8, 1.5, 1.4, 0.8, 1.0, 4.0]
    _table(slide, hdr_sw, rows_sw, _cw_sw,
           _cx(_cw_sw), _tbl_y(len(rows_sw)), fills=fills_sw)
    _notes(slide,
           "Software version lifecycle and hardware end-of-life status per cluster. "
           "SW Status: In Support (Green) = current release still receiving updates. "
           "End of Support (Red) = version past its end-of-software-support date — "
           "upgrade required for continued security patches and bug fixes. "
           "EOS Date: the date this software version reaches end-of-support. "
           "EOL Nodes: nodes with hardware models past their hardware end-of-life date — "
           "these are no longer eligible for hardware replacement from Cohesity. "
           "EOL Models: the specific hardware model strings flagged as EOL. "
           "See Excel Infrastructure and Node Hardware sheets for per-node detail.")

    # ── Coverage Gaps ─────────────────────────────────────────────────────────
    slide = _blank()
    _header(slide, "Coverage Gaps",
            "Protection groups with no recent successful backup")
    _footer(slide)
    hdr_gap = ["Cluster", "Protection Group", "Last Status",
               "Last Run End", "SLA Violated"]
    rows_gap, fills_gap = [], []
    for cd in all_data:
        for g in cd["groups"]:
            lr = g.get("lastRun") or {}
            lb = lr.get("localBackupInfo") or {}
            st = lb.get("status", "\u2014")
            if st in _SUC: continue
            end_t = lb.get("endTimeUsecs") or lb.get("startTimeUsecs") or 0
            sla_v = bool(lb.get("isSlaViolated") or lr.get("isSlaViolated"))
            end_s = (_dt.datetime.fromtimestamp(end_t / 1_000_000, tz=_dt.timezone.utc)
                     .strftime("%Y-%m-%d %H:%M")
                     if end_t else "Never")
            rows_gap.append(
                [cd["name"], g.get("name", "\u2014"), st, end_s, _yn(sla_v)])
            st_f  = ((_C_RED_BG, _C_RED_TX)
                     if st in ("kFailed", "Failed") else
                     (_C_AMB_BG, _C_AMB_TX))
            sla_f = ((_C_RED_BG, _C_RED_TX) if sla_v
                     else (_C_GRN_BG, _C_GRN_TX))
            fills_gap.append([None, None, st_f, None, sla_f])
    if not rows_gap:
        rows_gap  = [["\u2014",
                      "No coverage gaps \u2014 all groups have "
                      "recent successful backups",
                      "", "", ""]]
        fills_gap = [[(_C_GRN_BG, _C_GRN_TX)] + [None] * 4]
    _cw_gap = [2.0, 3.2, 1.8, 2.5, 1.7]
    _table(slide, hdr_gap, rows_gap[:18], _cw_gap,
           _cx(_cw_gap), _tbl_y(min(len(rows_gap), 18)), fills=fills_gap[:18])
    if len(rows_gap) > 18:
        _txt(slide,
             f"\u2026 and {len(rows_gap) - 18} more \u2014 "
             "see Excel Coverage Gaps sheet",
             0.1, H - 0.55, W - 0.2, 0.28, size=8, color=_C_DGRY, align="center")
    _notes(slide,
           "Protection groups that do not have a recent successful backup. "
           "These represent data protection gaps — objects in these groups may not "
           "be recoverable to their expected RPO. "
           "Last Status: kFailed/Failed (Red) = hard failure. Other statuses (Amber) "
           "= running, skipped, or unknown. "
           "Last Run End: timestamp of the most recent run attempt (Never = never run). "
           "SLA Violated: Yes (Red) = the run exceeded the SLA window defined in the policy. "
           "Priority: investigate Failed groups with SLA violations first. "
           "Common causes: network timeout, source agent offline, insufficient storage, "
           "or a paused protection group. "
           "See Excel Coverage Gaps sheet for full list.")

    # ── Agent Health ──────────────────────────────────────────────────────────
    slide = _blank()
    _header(slide, "Agent Health",
            "Per-cluster agent version summary and upgrade readiness")
    _footer(slide)
    hdr_ag = ["Cluster", "Total Agents", "Healthy",
              "Needs Upgrade", "Degraded / Error", "Agent Health %"]
    rows_ag, fills_ag = [], []
    _H_ST = {"kHealthy", "Healthy", "kGood"}
    _D_ST = {"kDegraded", "Degraded", "kError", "Error", "kUnreachable"}
    _U_ST = {"kUpgradable", "Upgradable", "kSomeUpgrades"}
    for cd in all_data:
        agents   = cd.get("agents") or []
        total    = len(agents)
        healthy  = sum(1 for a in agents
                       if (a.get("healthStatus") or
                           a.get("status") or "") in _H_ST)
        degraded = sum(1 for a in agents
                       if (a.get("healthStatus") or
                           a.get("status") or "") in _D_ST)
        upgrades = sum(1 for a in agents
                       if a.get("upgradability") in _U_ST)
        ag_score = int(healthy / total * 100) if total else 100
        rows_ag.append(
            [cd["name"], total, healthy, upgrades, degraded, ag_score])
        dg_f = (_C_RED_BG, _C_RED_TX) if degraded > 0 else (_C_GRN_BG, _C_GRN_TX)
        up_f = (_C_AMB_BG, _C_AMB_TX) if upgrades > 0 else (_C_GRN_BG, _C_GRN_TX)
        fills_ag.append([None, None, None, up_f, dg_f, _rag(ag_score)])
    _cw_ag = [2.5, 1.5, 1.3, 1.5, 1.6, 1.8]
    _table(slide, hdr_ag, rows_ag, _cw_ag,
           _cx(_cw_ag), _tbl_y(len(rows_ag)), fills=fills_ag)
    _notes(slide,
           "Cohesity agent health summary per cluster. "
           "Agents are software components installed on protected hosts (Windows, Linux, "
           "NAS, SQL, Oracle, etc.) that facilitate backup and recovery. "
           "Needs Upgrade (Amber): agents running an older version — should be upgraded "
           "to match the cluster version for full compatibility and security patches. "
           "Degraded/Error (Red): agents that are unreachable, in error state, or degraded "
           "— backups for those hosts may be failing silently. "
           "Agent Health %: percentage of agents in healthy state (Green >= 90%). "
           "See Excel Agent Health sheet for per-host agent version and status detail.")

    # ── Source Coverage ───────────────────────────────────────────────────────
    slide = _blank()
    _header(slide, "Source Coverage",
            "Registered sources \u2014 protected vs. unprotected objects")
    _footer(slide)
    hdr_src = ["Cluster", "Source Name", "Environment",
               "Total Objects", "Protected", "Unprotected", "Coverage %"]
    rows_src, fills_src = [], []
    for cd in all_data:
        for src in (cd.get("sources") or []):
            if not isinstance(src, dict): continue
            name  = src.get("name") or src.get("sourceName") or "\u2014"
            env   = src.get("environment") or src.get("envType") or "\u2014"
            total = (src.get("totalLeafCount")
                     or src.get("totalObjectCount") or 0)
            prot  = (src.get("protectedLeafCount")
                     or src.get("protectedObjectCount") or 0)
            unprot  = max(0, total - prot)
            cov_val = (prot / total * 100) if total else 100
            rows_src.append(
                [cd["name"], name, env, total, prot, unprot,
                 f"{cov_val:.0f}%" if total else "\u2014"])
            up_f = (_C_RED_BG, _C_RED_TX) if unprot > 0 \
                   else (_C_GRN_BG, _C_GRN_TX)
            fills_src.append(
                [None, None, None, None, None, up_f, _rag(cov_val)])
    if not rows_src:
        rows_src  = [["\u2014", "No source data", "", "", "", "", ""]]
        fills_src = [[None] * 7]
    _cw_src = [2.0, 2.5, 1.5, 1.3, 1.2, 1.3, 1.5]
    _table(slide, hdr_src, rows_src[:18], _cw_src,
           _cx(_cw_src), _tbl_y(min(len(rows_src), 18)), fills=fills_src[:18])
    if len(rows_src) > 18:
        _txt(slide,
             f"\u2026 and {len(rows_src) - 18} more \u2014 "
             "see Excel Source Coverage sheet",
             0.1, H - 0.55, W - 0.2, 0.28, size=8, color=_C_DGRY, align="center")
    _notes(slide,
           "Registered protection sources and their object coverage per cluster. "
           "A source is a connected system (VMware vCenter, SQL Server, NAS share, "
           "Oracle DB, physical host, etc.) registered in Cohesity. "
           "Total Objects = all leaf objects discovered within the source (VMs, DBs, etc.). "
           "Protected = objects included in at least one active protection group. "
           "Unprotected (Red if > 0) = objects discovered but not yet protected — "
           "these represent data that has no backup. "
           "Coverage % (RAG) = Protected / Total. "
           "See Excel Source Coverage sheet for the full list. "
           "Unprotected objects should be reviewed and added to protection groups.")

    # ── Workload Risk Heatmap ─────────────────────────────────────────────────
    slide = _blank()
    _header(slide, "Workload Risk Heatmap",
            "Critical and High risk protection groups sorted worst-first")
    _footer(slide)
    risk_eng = []
    for cd in all_data:
        for grp in cd["groups"]:
            score, level = _group_risk_score(grp, cd)
            if level in ("Critical", "High"):
                risk_eng.append(
                    (score, level, cd["name"], grp.get("name", "\u2014")))
    risk_eng.sort(key=lambda x: x[0])
    hdr_hr = ["Cluster", "Protection Group", "Risk Score", "Risk Level"]
    rows_hr, fills_hr = [], []
    for score, level, cname, gname in risk_eng[:20]:
        rows_hr.append([cname, gname, score, level])
        lv_f = ((_C_RED_BG, _C_RED_TX) if level == "Critical"
                else (_C_AMB_BG, _C_AMB_TX))
        fills_hr.append([None, None, _rag(score), lv_f])
    if not rows_hr:
        rows_hr  = [["\u2014",
                     "No Critical or High risk workloads \u2014 good standing",
                     "\u2014", "\u2014"]]
        fills_hr = [[(_C_GRN_BG, _C_GRN_TX)] + [None] * 3]
    _cw_hr = [2.5, 6.2, 1.5, 2.0]
    rows_hr  = rows_hr[:12];  fills_hr  = fills_hr[:12]
    # meth: rule + label + 5-row table (rh=0.22) + legend text
    _hr_meth_h = 0.015 + 0.05 + 0.22 + 0.27 + 5 * 0.22 + 0.08 + 0.22  # ≈ 1.98"
    _hr_y, _hr_my = _layout_mixed(len(rows_hr), _hr_meth_h)
    _table(slide, hdr_hr, rows_hr, _cw_hr,
           _cx(_cw_hr), _hr_y, fills=fills_hr)
    # ── Inline scoring methodology ──────────────────────────────────────────
    _rect(slide, 0.15, _hr_my,         W - 0.30, 0.015, "#CCCCCC")
    _txt(slide, "SCORING METHODOLOGY", 0.15, _hr_my + 0.05, W - 0.30, 0.22,
         size=7.5, bold=True, color="#888888")
    _risk_cols = [2.0, 1.0, 8.5]
    _table(slide,
           ["Factor", "Max Pts", "Detail"],
           [
               ["Last run status", "35 pts",
                "Success=35, Warning=25, Running=20, Skipped=10, Paused=5, Failed=0"],
               ["SLA compliance",  "25 pts",
                "25=no violation in the lookback window, 0=at least one SLA breach"],
               ["RPO gap",         "25 pts",
                "25=\u22644h, 20=\u22648h, 15=\u226424h, 8=\u226448h, 0=>48h since last success"],
               ["DataLock",        "15 pts",
                "15=Compliance or FortKnox-indelible, 10=Administrative, 0=none"],
           ],
           _risk_cols, _cx(_risk_cols), round(_hr_my + 0.32, 3), rh=0.22, body_sz=8)
    _txt(slide,
         "Low \u226575  \u00b7  Medium 50\u201374  \u00b7  High 25\u201349  \u00b7  Critical <25"
         "  \u00b7  Only Critical and High risk groups shown",
         0.15, round(_hr_my + 0.32 + 5 * 0.22 + 0.10, 3),
         W - 0.30, 0.22, size=7.5, color=_C_DGRY, align="center")
    _notes(slide,
           "Critical and High risk protection groups only, sorted worst-first by risk score. "
           "Only groups scoring below 50 are shown here — Medium and Low risk groups "
           "are omitted for brevity. "
           "Risk Score factors: last run status (35 pts — Success=35, Warning=25, "
           "Running=20, Skipped=10, Paused=5, Failed=0), SLA compliance (25 pts), "
           "RPO age from last successful run (25 pts), DataLock coverage (15 pts). "
           "Critical (< 25) = groups at highest recovery risk — likely failing backups, "
           "SLA violations, and no DataLock. "
           "High (25-49) = significant gaps needing prompt attention. "
           "See Excel Workload Risk Heatmap sheet for all groups including Medium and Low.")

    # ── Engineering Recommendations ───────────────────────────────────────────
    slide = _blank()
    _header(slide, "Engineering Recommendations",
            "Infrastructure, lifecycle, storage, agents, and protection findings")
    _footer(slide)
    _ENG_CATS = {"software lifecycle", "hardware", "infrastructure",
                 "storage", "agent", "protection", "source",
                 "policy", "alerts", "coverage", "backup"}
    eng_recs = []
    for cd in all_data:
        for rec in _build_recommendations(cd):
            if any(c in rec.get("category", "").lower() for c in _ENG_CATS):
                eng_recs.append(
                    (rec["priority"], rec["category"], cd["name"],
                     rec["finding"][:72], rec["action"][:72]))
    eng_recs.sort(key=lambda x:
                  ({"CRITICAL": 0, "HIGH": 1, "MEDIUM": 2,
                    "LOW": 3}.get(x[0], 4), x[1]))
    rows_er, fills_er = [], []
    for prio, cat, cname, finding, action in eng_recs[:15]:
        rows_er.append([prio, cat, cname, finding, action])
        fills_er.append([_P_COLORS.get(prio, (None, None))] + [None] * 4)
    if not rows_er:
        rows_er  = [["\u2014", "No engineering findings", "", "", ""]]
        fills_er = [[(_C_GRN_BG, _C_GRN_TX)] + [None] * 4]
    _cw_er = [1.2, 1.8, 2.0, 4.1, 3.8]
    _table(slide,
           ["Priority", "Category", "Cluster", "Finding", "Recommended Action"],
           rows_er, _cw_er, _cx(_cw_er), _tbl_y(len(rows_er)), fills=fills_er)
    _notes(slide,
           "Infrastructure and backup engineering findings, sorted by priority. "
           "Covers: software lifecycle (out-of-support versions), hardware EOL, "
           "storage capacity issues, protection group failures and SLA violations, "
           "agent health and upgrades, unprotected sources, and policy gaps. "
           "CRITICAL = immediate action needed (data at risk). "
           "HIGH = significant gap, address within days. "
           "MEDIUM = best-practice improvement, address within weeks. "
           "LOW = optimisation opportunity. "
           "See Excel Recommendations sheet for the full list with business impact "
           "descriptions and detailed remediation guidance.")

    # ── Populate Contents slide (now that all slide indices are finalised) ─────
    _rect(toc_slide, 0, 0, W, H, _C_BG)          # white background
    _rect(toc_slide, 0, 0, 0.06, H, _C_VACC)     # thin green left bar
    _header(toc_slide, "Contents", "Click a section to jump directly to it")
    _footer(toc_slide, show_page=False)

    _toc_items = [
        ("Executive Summary",
         "Environment snapshot  \u00b7  Health scorecard  \u00b7  "
         "Storage capacity  \u00b7  Protection  \u00b7  Ransomware readiness  \u00b7  Priority actions"),
        ("Environment Topology",
         "Cluster connections  \u00b7  Replication targets  \u00b7  "
         "Archival vaults  \u00b7  FortKnox / RPaaS"),
        ("Security Deep-Dive",
         "Security posture overview  \u00b7  User & MFA  \u00b7  "
         "Governance & audit activity  \u00b7  Security recommendations"),
        ("Backup Engineering",
         "Software & hardware lifecycle  \u00b7  Coverage gaps  \u00b7  "
         "Agent health  \u00b7  Source coverage  \u00b7  Workload risk  \u00b7  Recommendations"),
    ]
    _sec_dict    = dict(_sec_starts)
    _toc_row_h   = 1.18
    _toc_row_gap = 0.22
    _toc_total   = len(_toc_items) * _toc_row_h + (len(_toc_items) - 1) * _toc_row_gap
    _toc_y0      = (H - _toc_total) / 2 + 0.05   # vertically centred
    _toc_x0 = 0.18
    # Dynamic bar width — size to the widest text item (title at 15pt, desc at 8.5pt)
    # Estimate: char_count × pt_size × 0.58 / 72 (proportional font width in inches)
    def _est_w(text, pt): return len(text) * pt * 0.58 / 72
    _toc_max_content = max(
        max(_est_w(_ts, 15)  for _ts, _   in _toc_items),
        max(_est_w(_td, 8.5) for _,  _td  in _toc_items))
    # 1.60" fixed chrome: green block(0.30)+sep(0.38)+text-indent(0.50)+arrow(0.42)
    _toc_w = min(W - _toc_x0 - 0.18, round(_toc_max_content + 1.60, 2))

    for _ti, (_tsec, _tdesc) in enumerate(_toc_items):
        _ty = _toc_y0 + _ti * (_toc_row_h + _toc_row_gap)
        # Row card — light fill, subtle border on white background
        _row = toc_slide.shapes.add_shape(
            _MSO_SHAPE.ROUNDED_RECTANGLE,
            _PptxInches(_toc_x0), _PptxInches(_ty),
            _PptxInches(_toc_w), _PptxInches(_toc_row_h))
        _row.fill.solid()
        _row.fill.fore_color.rgb = _rgb(_C_BGROW)
        _row.line.color.rgb = _rgb("#D8E8D0")
        _row.line.width = _PptxInches(0.010)
        # Hyperlink to section slide
        _sidx = _sec_dict.get(_tsec)
        if _sidx is not None and _sidx < len(prs.slides):
            try:
                _row.click_action.target_slide = prs.slides[_sidx]
            except Exception:
                pass
        # Left green accent block on card
        _rect(toc_slide, _toc_x0, _ty, 0.30, _toc_row_h, _C_VACC)
        # Section number — white on green
        _txt(toc_slide, f"{_ti + 1}", _toc_x0 + 0.01, _ty + 0.32, 0.28, 0.52,
             size=15, bold=True, color=_C_WHT, align="center")
        # Vertical separator line
        _rect(toc_slide, _toc_x0 + 0.38, _ty + 0.14, 0.022, _toc_row_h - 0.28, "#C8D8C0")
        # Section name — dark on light
        _txt(toc_slide, _tsec, _toc_x0 + 0.50, _ty + 0.10, _toc_w - 1.15, 0.40,
             size=15, bold=True, color=_C_BLK)
        # Description — medium gray
        _txt(toc_slide, _tdesc, _toc_x0 + 0.50, _ty + 0.56, _toc_w - 1.15, 0.50,
             size=8.5, color="#555555")
        # Right arrow — green
        _txt(toc_slide, "\u203a", _toc_x0 + _toc_w - 0.48, _ty + 0.26, 0.40, 0.55,
             size=26, bold=True, color=_C_VACC, align="center")

    _notes(toc_slide,
           "Contents slide — click any section card to jump directly to that section. "
           "Executive Summary: overview KPIs, health scorecard, capacity, protection, "
           "ransomware readiness, and priority action items. "
           "Environment Topology: live diagram of cluster connections, replication, "
           "archival vaults, and FortKnox. "
           "Security Deep-Dive: security posture table, user & MFA status, "
           "governance and audit activity, security recommendations. "
           "Backup Engineering: software and hardware lifecycle, coverage gaps, "
           "agent health, source coverage, workload risk heatmap, recommendations.")

    # ── Thank You (end) slide ─────────────────────────────────────────────────
    slide_ty = _blank()
    _rect(slide_ty, 0, 0, W, H, _C_BG)
    _rect(slide_ty, 0, 0, 0.06, H, _C_VACC)
    _rect(slide_ty, 0.06, 2.10, 8.20, 3.10, _C_VACC)
    _txt(slide_ty, "Thank You", 0.28, 2.35, 7.80, 1.10,
         size=44, bold=True, color=_C_WHT, align="left")
    _rect(slide_ty, 0.28, 3.55, 3.80, 0.04, _C_WHT)
    _txt(slide_ty, "cohesity.com", 0.28, 3.68, 7.80, 0.42,
         size=13, color="#DDFFD0", align="left")
    _txt(slide_ty,
         "Questions? Contact your Cohesity SE or Customer Success Manager.",
         0.15, 5.15, 9.00, 0.38, size=10, color="#888888", align="left")
    _footer(slide_ty, show_page=False)
    _notes(slide_ty,
           "Thank you for your time. "
           "Please reach out to your Cohesity SE or Customer Success Manager "
           "with any questions following this review.")

    # ── Register native PowerPoint sections then save ─────────────────────────
    _register_sections(prs, _sec_starts)
    try:
        pptx_path = out_path + ".pptx"
        prs.save(pptx_path)
        print(f"  PowerPoint → {pptx_path}")
        return pptx_path
    except Exception as exc:
        print(f"  WARN: Could not write PowerPoint deck: {exc}")
        return None



# ── Native Word shapes topology ──────────────────────────────────────────────

def _render_topology_word_shapes(doc, all_data):
    """Render environment topology using native Word shapes (DrawingML).

    Creates editable rounded rectangles (source clusters in Cohesity green),
    cloud shapes (FortKnox), cylinders (archive), and connector arrows.
    All shapes are anchored inside a single paragraph with wrapNone positioning.
    Returns True on success, False if python-docx OxmlElement API is unavailable.
    """
    try:
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn, nsmap as _nsmap
        from docx.shared import Emu, Pt
    except ImportError:
        return False

    # python-docx ≤1.2 does not register the wps namespace; patch it so
    # OxmlElement('wps:wsp') resolves correctly.  setdefault is idempotent.
    _nsmap.setdefault(
        'wps', 'http://schemas.microsoft.com/office/word/2010/wordprocessingShape')

    try:
        return _render_topology_word_shapes_inner(doc, all_data,
                                                   OxmlElement, qn, Emu, Pt)
    except Exception as exc:
        print(f"  WARN: native Word shapes failed: {exc}")
        return False


def _render_topology_word_shapes_inner(doc, all_data, OxmlElement, qn, Emu, Pt):
    """Inner implementation — separated so the caller can catch all errors."""
    from docx.shared import RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    topo       = _topology_data(all_data)
    clusters   = topo["clusters"]
    repl_nodes = topo["repl_nodes"]
    arch_nodes = topo["arch_nodes"]
    fk_nodes   = topo["fk_nodes"]
    edges      = topo["edges"]

    if not clusters:
        return False

    # ── Constants (EMU: 914400 per inch) ─────────────────────────────────
    EMU = 914400

    CLR_CLUSTER = "96C73D";  CLR_CLUSTER_S = "6B8F2C"
    CLR_REPL    = "0062B1";  CLR_REPL_S    = "004A8A"
    CLR_ARCH    = "E07A00";  CLR_ARCH_S    = "B05A00"
    CLR_FK      = "1A5C3A";  CLR_FK_S      = "0E3B25"

    C_W  = int(2.05 * EMU);  C_H  = int(0.72 * EMU)
    T_W  = int(1.55 * EMU);  T_H  = int(0.58 * EMU)
    FK_W = int(1.70 * EMU);  FK_H = int(0.72 * EMU)
    V_GAP = int(0.20 * EMU)
    LABEL_H = int(0.28 * EMU)

    # ── Layout ───────────────────────────────────────────────────────────
    tcols = []  # (type, nodes, width, height, geom, fill, stroke, label)
    if repl_nodes:
        tcols.append(("repl", repl_nodes, T_W, T_H, "roundRect",
                       CLR_REPL, CLR_REPL_S, "Replication Targets"))
    if arch_nodes:
        tcols.append(("arch", arch_nodes, T_W, T_H, "can",
                       CLR_ARCH, CLR_ARCH_S, "Archival Vaults"))
    if fk_nodes:
        tcols.append(("fk", fk_nodes, FK_W, FK_H, "cloud",
                       CLR_FK, CLR_FK_S, "FortKnox"))

    # Column X positions (left edge)
    col_x = {"cluster": 0}
    x_cursor = C_W + int(0.55 * EMU)
    for ct, nodes, tw, *_ in tcols:
        col_x[ct] = x_cursor
        x_cursor += tw + int(0.30 * EMU)

    # Y positions
    SHAPES_TOP = LABEL_H + int(0.08 * EMU)

    def _ys(n, h):
        return [SHAPES_TOP + i * (h + V_GAP) for i in range(n)]

    cluster_ys = _ys(len(clusters), C_H)
    c_total_h  = len(clusters) * (C_H + V_GAP) - V_GAP if clusters else C_H

    tcol_ys = {}
    for ct, nodes, tw, th, *_ in tcols:
        ys = _ys(len(nodes), th)
        t_total = len(nodes) * (th + V_GAP) - V_GAP if nodes else th
        offset  = max((c_total_h - t_total) // 2, 0)
        tcol_ys[ct] = [y + offset for y in ys]

    max_right_h = max(
        (len(nodes) * (th + V_GAP) - V_GAP
         for _, nodes, _, th, *_ in tcols), default=0)
    total_h = SHAPES_TOP + max(c_total_h, max_right_h) + int(0.20 * EMU)

    # Position lookup: id → (left_x, top_y, width, height)
    pos = {}
    for i, cl in enumerate(clusters):
        pos[cl["id"]] = (col_x["cluster"], cluster_ys[i], C_W, C_H)
    for ct, nodes, tw, th, *_ in tcols:
        for i, n in enumerate(nodes):
            pos[n["id"]] = (col_x[ct], tcol_ys[ct][i], tw, th)

    # ── Shape ID counter ─────────────────────────────────────────────────
    _sid = [200]
    def _nid():
        _sid[0] += 1;  return _sid[0]

    # ── XML element helper ───────────────────────────────────────────────
    def _mk(tag, attrib=None, text=None):
        el = OxmlElement(tag)
        for k, v in (attrib or {}).items():
            el.set(k if ':' not in k else qn(k), str(v))
        if text is not None:
            el.text = str(text)
        return el

    def _anchor_frame(sid, x, y, w, h, z=0):
        """Build wp:anchor → a:graphic → a:graphicData; return (anchor, gd)."""
        anc = _mk('wp:anchor', {
            'distT': '0', 'distB': '0', 'distL': '0', 'distR': '0',
            'simplePos': '0',
            'relativeHeight': str(251660000 + z + sid),
            'behindDoc': '0', 'locked': '0',
            'layoutInCell': '1', 'allowOverlap': '1',
        })
        anc.append(_mk('wp:simplePos', {'x': '0', 'y': '0'}))
        ph = _mk('wp:positionH', {'relativeFrom': 'column'})
        ph.append(_mk('wp:posOffset', text=str(int(x))))
        anc.append(ph)
        pv = _mk('wp:positionV', {'relativeFrom': 'paragraph'})
        pv.append(_mk('wp:posOffset', text=str(int(y))))
        anc.append(pv)
        anc.append(_mk('wp:extent', {'cx': str(int(w)), 'cy': str(int(h))}))
        anc.append(_mk('wp:effectExtent',
                        {'l': '0', 't': '0', 'r': '19050', 'b': '19050'}))
        anc.append(_mk('wp:wrapNone'))
        anc.append(_mk('wp:docPr', {'id': str(sid), 'name': f'Shape{sid}'}))
        anc.append(_mk('wp:cNvGraphicFramePr'))
        gr  = _mk('a:graphic')
        gd  = _mk('a:graphicData', {
            'uri': 'http://schemas.microsoft.com/office/word/2010/wordprocessingShape'})
        gr.append(gd)
        anc.append(gr)
        return anc, gd

    def _append_to_para(para, anchor):
        r   = _mk('w:r')
        dwg = _mk('w:drawing')
        dwg.append(anchor)
        r.append(dwg)
        para._p.append(r)

    # ── Add a filled shape with text ─────────────────────────────────────
    def _shape(para, x, y, w, h, geom, fill, stroke,
               text_lines, fsz=10, fclr="FFFFFF", z=100):
        sid = _nid()
        anc, gd = _anchor_frame(sid, x, y, w, h, z)

        wsp = _mk('wps:wsp')
        wsp.append(_mk('wps:cNvSpPr'))

        spPr = _mk('wps:spPr')
        xfrm = _mk('a:xfrm')
        xfrm.append(_mk('a:off', {'x': '0', 'y': '0'}))
        xfrm.append(_mk('a:ext', {'cx': str(int(w)), 'cy': str(int(h))}))
        spPr.append(xfrm)
        pg = _mk('a:prstGeom', {'prst': geom})
        pg.append(_mk('a:avLst'))
        spPr.append(pg)
        if fill == "NONE":
            spPr.append(_mk('a:noFill'))
        else:
            sf = _mk('a:solidFill');  sf.append(_mk('a:srgbClr', {'val': fill}))
            spPr.append(sf)
        if stroke == "NONE":
            ln = _mk('a:ln');  ln.append(_mk('a:noFill'));  spPr.append(ln)
        else:
            ln = _mk('a:ln', {'w': '12700'})
            lf = _mk('a:solidFill');  lf.append(_mk('a:srgbClr', {'val': stroke}))
            ln.append(lf);  spPr.append(ln)
        wsp.append(spPr)

        # Text body
        txbx = _mk('wps:txbx')
        tc   = _mk('w:txbxContent')
        for idx, line in enumerate(text_lines):
            wp  = _mk('w:p')
            pPr = _mk('w:pPr')
            pPr.append(_mk('w:jc', {'w:val': 'center'}))
            sp = _mk('w:spacing')
            sp.set(qn('w:after'), '0')
            sp.set(qn('w:line'), '216')
            sp.set(qn('w:lineRule'), 'auto')
            pPr.append(sp)
            wp.append(pPr)
            wr  = _mk('w:r')
            rPr = _mk('w:rPr')
            rPr.append(_mk('w:color', {'w:val': fclr}))
            fs = fsz if idx == 0 else max(fsz - 1, 7)
            rPr.append(_mk('w:sz',   {'w:val': str(fs * 2)}))
            rPr.append(_mk('w:szCs', {'w:val': str(fs * 2)}))
            if idx == 0:
                rPr.append(_mk('w:b'));  rPr.append(_mk('w:bCs'))
            wr.append(rPr)
            t = _mk('w:t')
            t.set(qn('xml:space'), 'preserve')
            t.text = line
            wr.append(t)
            wp.append(wr)
            tc.append(wp)
        txbx.append(tc)
        wsp.append(txbx)

        bPr = _mk('wps:bodyPr', {
            'anchor': 'ctr', 'anchorCtr': '0', 'wrap': 'square',
            'lIns': '91440', 'tIns': '36576',
            'rIns': '91440', 'bIns': '36576',
        })
        wsp.append(bPr)
        gd.append(wsp)
        _append_to_para(para, anc)

    # ── Add a connector arrow ────────────────────────────────────────────
    def _connector(para, x1, y1, x2, y2, color, z=50):
        sid  = _nid()
        bx   = min(x1, x2)
        by   = min(y1, y2)
        cx   = abs(x2 - x1) or 1
        cy   = abs(y2 - y1) or 1
        anc, gd = _anchor_frame(sid, bx, by, cx, cy, z)

        wsp = _mk('wps:wsp')
        wsp.append(_mk('wps:cNvCnPr'))
        spPr = _mk('wps:spPr')
        xfrm = _mk('a:xfrm')
        if x2 < x1: xfrm.set('flipH', '1')
        if y2 < y1: xfrm.set('flipV', '1')
        xfrm.append(_mk('a:off', {'x': '0', 'y': '0'}))
        xfrm.append(_mk('a:ext', {'cx': str(int(cx)), 'cy': str(int(cy))}))
        spPr.append(xfrm)
        pg = _mk('a:prstGeom', {'prst': 'straightConnector1'})
        pg.append(_mk('a:avLst'))
        spPr.append(pg)
        ln = _mk('a:ln', {'w': '19050'})
        lf = _mk('a:solidFill');  lf.append(_mk('a:srgbClr', {'val': color}))
        ln.append(lf)
        ln.append(_mk('a:tailEnd', {'type': 'triangle', 'w': 'med', 'len': 'med'}))
        spPr.append(ln)
        wsp.append(spPr)
        wsp.append(_mk('wps:bodyPr'))
        gd.append(wsp)
        _append_to_para(para, anc)

    # ── Column label (no-fill text shape) ────────────────────────────────
    def _label(para, x, w, text, color, z=200):
        _shape(para, x, 0, w, LABEL_H, "rect", "NONE", "NONE",
               [text], fsz=9, fclr=color, z=z)

    # ── Build the diagram ────────────────────────────────────────────────
    p_diag = doc.add_paragraph()
    p_diag.paragraph_format.space_after = Emu(total_h + int(0.15 * EMU))

    # Column labels
    _label(p_diag, col_x["cluster"], C_W, "Source Clusters", CLR_CLUSTER)
    for ct, _, tw, *rest in tcols:
        _label(p_diag, col_x[ct], tw, rest[-1], rest[2])  # rest[-1]=label, rest[2]=fill

    # Connectors (draw first → lower z-order)
    drawn = set()
    for e in edges:
        key = (e["src"], e["tgt"])
        if key in drawn:
            continue
        drawn.add(key)
        sp = pos.get(e["src"])
        tp = pos.get(e["tgt"])
        if not sp or not tp:
            continue
        sx = sp[0] + sp[2]            # right edge of source
        sy = sp[1] + sp[3] // 2       # vertical center
        tx = tp[0]                     # left edge of target
        ty = tp[1] + tp[3] // 2       # vertical center
        clr = {"repl": CLR_REPL, "arch": CLR_ARCH, "fk": CLR_FK}[e["type"]]
        _connector(p_diag, sx, sy, tx, ty, clr)

    # Source cluster shapes
    for i, cl in enumerate(clusters):
        wls = ", ".join(cl["workloads"]) if cl["workloads"] else ""
        det = f"{cl['groups']} groups"
        if cl["sources"]:
            det += f" \u00b7 {cl['sources']} sources"
        if cl["cap_tb"] > 0.01:
            det += f" \u00b7 {cl['cap_tb']:.1f} TB"
        lines = [cl["name"]]
        if wls:
            lines.append(wls)
        lines.append(det)
        _shape(p_diag, col_x["cluster"], cluster_ys[i], C_W, C_H,
               "roundRect", CLR_CLUSTER, CLR_CLUSTER_S, lines, fsz=10)

    # Target shapes
    for ct, nodes, tw, th, geom, fill, stroke, label in tcols:
        sub = {"repl": "Replication Target", "arch": "Archival Vault",
               "fk": "FortKnox (Indelible)"}[ct]
        for i, n in enumerate(nodes):
            _shape(p_diag, col_x[ct], tcol_ys[ct][i], tw, th,
                   geom, fill, stroke, [n["name"], sub], fsz=9)

    # Caption
    cap = doc.add_paragraph(
        "Figure: Environment Topology \u2014 source clusters (green), "
        "replication targets (blue), archival vaults (amber), "
        "FortKnox (dark-green cloud). All shapes are editable in Word. "
        "A comprehensive PowerPoint deck (~31 slides) is also saved alongside this report."
    )
    cap.runs[0].font.size = Pt(8)
    cap.runs[0].font.color.rgb = RGBColor(0x80, 0x80, 0x80)
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    return True


# ── matplotlib PNG preview (fallback) ────────────────────────────────────────

def _render_topology_png(all_data):
    """Render the environment topology as a PNG byte-string using matplotlib.

    Cohesity brand colors: teal clusters, blue replication, amber archival,
    dark-green FortKnox. Returns PNG bytes, or None if matplotlib is unavailable.
    """
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        from matplotlib.patches import FancyBboxPatch
        import io as _io
    except ImportError:
        return None

    topo       = _topology_data(all_data)
    clusters   = topo["clusters"]
    repl_nodes = topo["repl_nodes"]
    arch_nodes = topo["arch_nodes"]
    fk_nodes   = topo["fk_nodes"]
    edges      = topo["edges"]

    # ── Figure setup ─────────────────────────────────────────────────────────
    FIG_W, FIG_H = 16.0, 9.0
    HEADER_H = 0.0
    FOOTER_H = 0.60

    CT = FIG_H - 0.30               # content top (no header)
    CB = FOOTER_H + 0.15            # content bottom
    CH = CT - CB                    # content height

    C_W, C_H = 2.80, 1.05          # cluster box
    T_W, T_H = 2.30, 0.78          # target box

    fig, ax = plt.subplots(figsize=(FIG_W, FIG_H), dpi=130)
    fig.patch.set_facecolor("#F8F8F8")
    ax.set_facecolor("#F8F8F8")
    ax.set_xlim(0, FIG_W)
    ax.set_ylim(0, FIG_H)
    ax.axis("off")

    # ── Determine active right columns ────────────────────────────────────────
    right_cols = []
    if repl_nodes: right_cols.append(("repl", repl_nodes, _DG_REPL, "#004A8A", "Replication Targets"))
    if arch_nodes: right_cols.append(("arch", arch_nodes, _DG_ARCH, "#B05A00", "Archival Vaults"))
    if fk_nodes:   right_cols.append(("fk",   fk_nodes,  _DG_FK,   "#0E3B25", "FortKnox / RPaaS"))

    # Compute column X centers
    X_CLUSTER = 1.85
    RIGHT_START = X_CLUSTER + C_W / 2 + 1.4
    RIGHT_END   = FIG_W - 0.35
    n_right     = len(right_cols)
    if n_right > 0:
        step = (RIGHT_END - RIGHT_START - T_W) / max(n_right - 1, 1) if n_right > 1 else 0
        col_xs = [RIGHT_START + T_W / 2 + i * step for i in range(n_right)]
        # if only 1 column centre it in available space
        if n_right == 1:
            col_xs = [(RIGHT_START + RIGHT_END) / 2]
    else:
        col_xs = []

    # Build (type → x_center) lookup
    col_x_by_type = {}
    col_info_by_type = {}
    for idx, (ctype, nodes, fc, ec, header) in enumerate(right_cols):
        col_x_by_type[ctype] = col_xs[idx]
        col_info_by_type[ctype] = (nodes, fc, ec, header)

    def _ys(n, box_h):
        """Y-centers for n items, distributed evenly in the content area."""
        if n == 0: return []
        if n == 1: return [CB + CH / 2]
        step = (CH - box_h) / (n - 1)
        return [CT - box_h / 2 - i * step for i in range(n)]

    cluster_ys = _ys(len(clusters), C_H)
    type_ys    = {}
    for ctype, nodes, *_ in right_cols:
        type_ys[ctype] = _ys(len(nodes), T_H)

    # Build id → (cx, cy)
    pos = {}
    for i, cl in enumerate(clusters):
        pos[cl["id"]] = (X_CLUSTER, cluster_ys[i])
    for ctype, nodes, *_ in right_cols:
        cx = col_x_by_type[ctype]
        for i, n in enumerate(nodes):
            pos[n["id"]] = (cx, type_ys[ctype][i])

    # ── Title ─────────────────────────────────────────────────────────────────
    ax.text(FIG_W / 2, FIG_H - 0.15, "Environment Topology",
            fontsize=14, fontweight="bold", color="#333333",
            va="center", ha="center", zorder=6)

    # ── Footer / legend ───────────────────────────────────────────────────────
    ax.add_patch(plt.Rectangle((0, 0), FIG_W, FOOTER_H, color="#EDEDED", zorder=1))
    leg_items = [
        (_DG_TEAL, "Source Cluster"),
    ] + [(c[2], c[4]) for c in right_cols]
    lx = 0.5
    for col, lbl in leg_items:
        ax.add_patch(plt.Rectangle((lx, 0.18), 0.35, 0.22,
                                    color=col, zorder=2, linewidth=0))
        ax.text(lx + 0.45, 0.29, lbl, va="center", fontsize=8.0, color="#444444", zorder=2)
        lx += 2.8 if len(leg_items) <= 4 else 2.2

    # ── Column headers ────────────────────────────────────────────────────────
    def _col_hdr(x, label, color, bw):
        y = CT + 0.12
        ax.plot([x - bw / 2, x + bw / 2], [y - 0.05, y - 0.05],
                color=color, lw=1.0, alpha=0.35, zorder=1)
        ax.text(x, y + 0.02, label, ha="center", va="bottom",
                fontsize=8.5, fontweight="bold", color=color, zorder=2)

    _col_hdr(X_CLUSTER, "Source Clusters", _DG_TEAL, C_W)
    for ctype, nodes, fc, ec, header in right_cols:
        _col_hdr(col_x_by_type[ctype], header, fc, T_W)

    # ── Vertical separator line ───────────────────────────────────────────────
    sep_x = X_CLUSTER + C_W / 2 + 0.65
    ax.plot([sep_x, sep_x], [CB, CT], color="#CCCCCC", lw=1.0,
            linestyle="--", zorder=1)

    # ── Edges (drawn before boxes so they appear underneath) ─────────────────
    drawn_pairs = set()
    for e in edges:
        key = (e["src"], e["tgt"])
        if key in drawn_pairs: continue
        drawn_pairs.add(key)
        sp, tp = pos.get(e["src"]), pos.get(e["tgt"])
        if sp is None or tp is None: continue
        if   e["type"] == "repl": col = _DG_REPL; bw = T_W
        elif e["type"] == "arch": col = _DG_ARCH; bw = T_W
        else:                     col = _DG_FK;   bw = T_W;
        sx, sy = sp[0] + C_W / 2, sp[1]
        tx, ty = tp[0] - bw / 2,  tp[1]
        ax.annotate("", xy=(tx, ty), xytext=(sx, sy),
                    arrowprops=dict(arrowstyle="-|>", color=col, lw=1.5,
                                   mutation_scale=14,
                                   connectionstyle="arc3,rad=0.0"),
                    zorder=2)

    # ── Box drawing helpers ───────────────────────────────────────────────────
    def _box(cx, cy, bw, bh, fc, ec, lines, bold_first=True):
        x0, y0 = cx - bw / 2, cy - bh / 2
        # Drop shadow
        ax.add_patch(FancyBboxPatch((x0 + 0.035, y0 - 0.035), bw, bh,
                                    boxstyle="round,pad=0.06",
                                    facecolor="#00000018", edgecolor="none", zorder=2))
        # Main box
        ax.add_patch(FancyBboxPatch((x0, y0), bw, bh,
                                    boxstyle="round,pad=0.06",
                                    facecolor=fc, edgecolor=ec,
                                    linewidth=1.4, zorder=3))
        # Text lines (top → bottom)
        n = len(lines)
        line_h = bh / (n + 0.3)
        for j, line in enumerate(lines):
            ty = cy + (n - 1) / 2 * line_h - j * line_h
            fw = "bold" if (bold_first and j == 0) else "normal"
            fs = 10.5 if j == 0 else 9.0
            ax.text(cx, ty, line, ha="center", va="center",
                    fontsize=fs, fontweight=fw, color="white", zorder=4,
                    clip_on=True)

    # ── Source cluster boxes ──────────────────────────────────────────────────
    for i, cl in enumerate(clusters):
        cx, cy = pos[cl["id"]]
        wls = ", ".join(cl["workloads"]) if cl["workloads"] else ""
        det = f"{cl['groups']} groups"
        if cl["sources"]: det += f"  ·  {cl['sources']} sources"
        if cl["cap_tb"] > 0.01: det += f"  ·  {cl['cap_tb']:.1f} TB"
        lines = [cl["name"]] + ([wls] if wls else []) + [det]
        _box(cx, cy, C_W, C_H, _DG_TEAL, _DG_TEAL_DK, lines)

    # ── Replication target boxes ──────────────────────────────────────────────
    if "repl" in col_x_by_type:
        for i, n in enumerate(repl_nodes):
            cx, cy = pos[n["id"]]
            _box(cx, cy, T_W, T_H, _DG_REPL, "#004A8A",
                 [n["name"], "Cluster"])

    # ── Archival vault boxes ──────────────────────────────────────────────────
    if "arch" in col_x_by_type:
        for i, n in enumerate(arch_nodes):
            cx, cy = pos[n["id"]]
            _box(cx, cy, T_W, T_H, _DG_ARCH, "#B05A00",
                 [n["name"], "Archival Vault"])

    # ── FortKnox boxes ────────────────────────────────────────────────────────
    if "fk" in col_x_by_type:
        for i, n in enumerate(fk_nodes):
            cx, cy = pos[n["id"]]
            _box(cx, cy, T_W, T_H, _DG_FK, "#0E3B25",
                 [n["name"], "FortKnox / RPaaS (indelible)"])

    plt.subplots_adjust(left=0, right=1, top=1, bottom=0)
    buf = _io.BytesIO()
    plt.savefig(buf, format="png", dpi=130, bbox_inches="tight",
                facecolor=fig.get_facecolor(), pad_inches=0.0)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


# ─────────────────────────────────────────────────────────────────────────────
# SHEET 20 — PER-WORKLOAD RISK HEATMAP
# ─────────────────────────────────────────────────────────────────────────────

def _sheet_risk_heatmap(wb, all_data):
    """Sheet 20: every protection group scored and sorted worst-first."""
    from openpyxl.styles import PatternFill, Font, Alignment

    RISK_COLORS = {
        "Critical": ("FF4C4C", "FFFFFF", True),
        "High":     ("FF8C42", "FFFFFF", True),
        "Medium":   ("FFF2CC", "000000", False),
        "Low":      ("E2EFDA", "000000", False),
    }

    ws = wb.create_sheet("Workload Risk Heatmap")
    ws.freeze_panes = "A4"
    _title(ws, "Per-Workload Risk Heatmap — Protection Groups Scored by Recovery Risk", "L")

    hdr_cols = ["Cluster", "Group Name", "Policy", "Risk Level", "Risk Score /100",
                "Last Run Status", "SLA Violated", "RPO Gap (h)",
                "DataLock", "Last Success", "Objects Failed", "Recommendation"]
    _hdr(ws, 2, hdr_cols)

    # Collect and score all groups across all clusters
    all_rows = []
    counts   = {"Critical": 0, "High": 0, "Medium": 0, "Low": 0}

    for cd in all_data:
        policy_by_id = cd.get("policy_by_id") or {p["id"]: p for p in (cd.get("policies") or []) if p.get("id")}

        for grp in (cd.get("groups") or []):
            score, level = _group_risk_score(grp, cd)
            counts[level] = counts.get(level, 0) + 1

            name      = grp.get("name", "")
            pol_id    = grp.get("policyId") or (grp.get("policy") or {}).get("id") or ""
            pol_name  = (policy_by_id.get(pol_id) or {}).get("name", pol_id or "—")
            policy    = policy_by_id.get(pol_id) or {}
            dl_str    = _policy_datalock_str(policy) if policy else "—"

            # Last run info
            last_run   = grp.get("lastRun") or {}
            run_status = ""
            sla_viol   = ""
            rpo_gap_h  = ""
            last_succ  = ""
            obj_fail   = ""
            run_sum    = last_run.get("backupRun") or last_run.get("protectionRun") or last_run
            if run_sum:
                run_status = (run_sum.get("status") or run_sum.get("runStatus")
                              or run_sum.get("localRun", {}).get("status", ""))
                sla_viol   = "Yes" if run_sum.get("isLocalSnapshotsDeleted") is False \
                              and run_sum.get("slaViolated") else \
                             "Yes" if run_sum.get("slaViolated") else "No"
                stats_local = (run_sum.get("localRun") or {}).get("stats") or \
                               run_sum.get("stats") or {}
                obj_fail    = stats_local.get("numFailedObjects") or \
                               stats_local.get("numFailed") or ""
                # Last success timestamp
                end_ts = (run_sum.get("endTimeUsecs")
                           or run_sum.get("localRun", {}).get("endTimeUsecs") or 0)
                if end_ts:
                    try:
                        import datetime as _dtt
                        last_succ = _dtt.datetime.fromtimestamp(
                            end_ts / 1_000_000).strftime("%Y-%m-%d %H:%M")
                    except Exception:
                        last_succ = ""

            # RPO gap from run history (hours since last success)
            rpo_hours = None
            _runs_list = cd.get("v1_runs") or {}
            grp_runs = _runs_list.get(name, []) if isinstance(_runs_list, dict) else []
            if not grp_runs:
                grp_runs = (cd.get("group_runs") or {}).get(name, [])
            now_usecs = int(time.time() * 1_000_000)
            for r in grp_runs:
                st = (r.get("status") or "").lower()
                if any(x in st for x in ("success", "warning")):
                    end_u = r.get("endTimeUsecs") or 0
                    if end_u:
                        rpo_hours = round((now_usecs - end_u) / 3_600_000_000, 1)
                        break
            rpo_gap_h = rpo_hours if rpo_hours is not None else ""

            # Recommendation
            if level == "Critical":
                rec = "Investigate immediately — backup protection severely degraded"
            elif level == "High":
                rec = "Review soon — recovery risk elevated"
            elif level == "Medium":
                rec = "Monitor — partial compliance issues detected"
            else:
                rec = "No immediate action required"

            all_rows.append({
                "cluster":    cd["name"],
                "name":       name,
                "policy":     pol_name,
                "level":      level,
                "score":      score,
                "status":     run_status,
                "sla":        sla_viol,
                "rpo":        rpo_gap_h,
                "datalock":   dl_str,
                "last_succ":  last_succ,
                "obj_fail":   obj_fail,
                "rec":        rec,
            })

    # Summary counts row (row 3)
    summary_str = "  |  ".join(
        f"{counts.get(lv, 0)} {lv}" for lv in ("Critical", "High", "Medium", "Low")
    )
    ws.merge_cells(f"A3:L3")
    sc = ws.cell(row=3, column=1, value=summary_str)
    sc.fill = PatternFill("solid", fgColor=LIGHT_GRAY)
    sc.font = Font(name="Calibri", size=10, bold=True)
    sc.alignment = Alignment(horizontal="center")

    # Sort worst first (ascending score)
    all_rows.sort(key=lambda x: x["score"])

    if not all_rows:
        ws.cell(row=4, column=1, value="No protection groups found").font = _FONT_NORMAL
    else:
        for row_data in all_rows:
            level = row_data["level"]
            bg, fg, bold = RISK_COLORS.get(level, (LIGHT_GRAY, "000000", False))
            rn = ws.max_row + 1
            vals = [row_data["cluster"], row_data["name"], row_data["policy"],
                    row_data["level"], row_data["score"], row_data["status"],
                    row_data["sla"], row_data["rpo"], row_data["datalock"],
                    row_data["last_succ"], row_data["obj_fail"], row_data["rec"]]
            for c, val in enumerate(vals, 1):
                cell = ws.cell(row=rn, column=c, value=_safe_cell(val))
                cell.fill = PatternFill("solid", fgColor=bg)
                cell.font = Font(name="Calibri", size=9, bold=bold,
                                 color=fg)

    auto_fit_columns(ws)
    ws.column_dimensions["B"].width = 40
    ws.column_dimensions["L"].width = 50


# ─────────────────────────────────────────────────────────────────────────────
# SHEET 21 — AUDIT LOG / CHANGE SUMMARY
# ─────────────────────────────────────────────────────────────────────────────

def _sheet_audit_log(wb, all_data):
    """Sheet 21: 30-day change summary and detail log from the audit trail."""
    from openpyxl.styles import PatternFill, Font, Alignment

    CAT_COLORS = {
        "Policy Changes":         ("BDD7EE", "000000"),
        "Group Changes":          ("E2EFDA", "000000"),
        "User Changes":           ("FFF2CC", "000000"),
        "Vault/Security Changes": ("FF4C4C", "FFFFFF"),
        "Cluster Config Changes": ("CCFFFF", "000000"),
        "Other":                  ("F2F2F2", "000000"),
    }
    HIGH_RISK_COLOR = ("FF4C4C", "FFFFFF")

    ws = wb.create_sheet("Audit Log")
    ws.freeze_panes = "A3"
    _title(ws, "Audit Log — Configuration & Change Activity (30 days)", "G")

    # ── Summary section ──────────────────────────────────────────────────────
    sum_hdr = ["Category", "Event Count", "Notable Events"]
    _hdr(ws, 2, sum_hdr)

    # Aggregate across all clusters
    cat_counts  = {}
    cat_notable = {}
    detail_rows = []

    denied_clusters = [cd["name"] for cd in all_data
                       if cd.get("audit_events") == "PERMISSION_DENIED"]

    for cd in all_data:
        cluster_name = cd["name"]
        events = cd.get("audit_events")
        if not isinstance(events, list):
            continue
        for event in events:
            cat = _audit_category(event)
            if cat is None:
                continue
            cat_counts[cat]  = cat_counts.get(cat, 0) + 1
            high_risk = _audit_high_risk(event)
            action    = event.get("action") or event.get("entityType") or ""
            entity    = event.get("entityName") or event.get("objectName") or ""
            if high_risk:
                notable = f"[HIGH RISK] {action}: {entity}"
                cat_notable.setdefault(cat, []).append(notable)

            ts_u  = (event.get("timestampUsecs")
                     or event.get("startTimestampUsecs") or 0)
            ts_dt = usecs_to_datetime(ts_u) if ts_u else ""
            detail_rows.append({
                "cluster":   cluster_name,
                "ts":        ts_dt,
                "user":      event.get("username") or event.get("user") or "",
                "domain":    event.get("domain") or "",
                "action":    action,
                "etype":     event.get("entityType") or "",
                "ename":     entity,
                "details":   event.get("details") or event.get("description") or "",
                "cat":       cat,
                "high_risk": high_risk,
            })

    for cat in ["Policy Changes", "Group Changes", "User Changes",
                "Vault/Security Changes", "Cluster Config Changes", "Other"]:
        if cat not in cat_counts:
            continue
        notable_str = "; ".join((cat_notable.get(cat) or [])[:3])
        rn = ws.max_row + 1
        bg, fg = CAT_COLORS.get(cat, ("F2F2F2", "000000"))
        for c, val in enumerate([cat, cat_counts[cat], notable_str or "—"], 1):
            cell = ws.cell(row=rn, column=c, value=_safe_cell(val))
            cell.fill = PatternFill("solid", fgColor=bg)
            cell.font = Font(name="Calibri", size=9, color=fg)

    # Divider
    divider_rn = ws.max_row + 1
    ws.merge_cells(f"A{divider_rn}:G{divider_rn}")
    dc = ws.cell(row=divider_rn, column=1, value="─── Detail Log ───")
    dc.font = Font(name="Calibri", size=9, bold=True, italic=True)
    dc.fill = PatternFill("solid", fgColor=LIGHT_GRAY)
    dc.alignment = Alignment(horizontal="center")

    # Detail header
    det_hdr = ["Timestamp", "Cluster", "User", "Domain",
               "Action", "Entity Type", "Entity Name"]
    _hdr(ws, ws.max_row + 1, det_hdr)

    # Detail rows — newest first
    detail_rows.sort(key=lambda x: x["ts"] or "", reverse=True)

    if denied_clusters:
        rn = ws.max_row + 1
        msg = (f"Access denied (HTTP 403) for: {', '.join(denied_clusters)}. "
               "Grant 'Audit Log View' privilege to the API key role: "
               "Cohesity UI → Settings → Access Management → API Keys → Edit Role.")
        c = ws.cell(row=rn, column=1, value=msg)
        c.font = _FONT_NORMAL
        c.fill = PatternFill("solid", fgColor="FFC7CE")
    if not detail_rows:
        ws.cell(row=ws.max_row + 1, column=1,
                value=("No audit log events found in the selected time window."
                       if not denied_clusters else
                       "No events returned — check API key role permissions above.")
                ).font = _FONT_NORMAL
    else:
        for ev in detail_rows:
            rn = ws.max_row + 1
            vals = [str(ev["ts"]) if ev["ts"] else "",
                    ev["cluster"], ev["user"], ev["domain"],
                    ev["action"], ev["etype"], ev["ename"]]
            bg, fg = HIGH_RISK_COLOR if ev["high_risk"] else CAT_COLORS.get(ev["cat"], ("FFFFFF", "000000"))
            bold = ev["high_risk"]
            for c, val in enumerate(vals, 1):
                cell = ws.cell(row=rn, column=c, value=_safe_cell(val))
                cell.fill = PatternFill("solid", fgColor=bg)
                cell.font = Font(name="Calibri", size=9, bold=bold, color=fg)

    auto_fit_columns(ws)
    ws.column_dimensions["A"].width = 20
    ws.column_dimensions["G"].width = 40


# ─────────────────────────────────────────────────────────────────────────────
# SHEET — AD / IDENTITY HEALTH
# ─────────────────────────────────────────────────────────────────────────────

def _sheet_ad_identity(wb, all_data):
    from openpyxl.styles import PatternFill, Font, Alignment

    ws = wb.create_sheet("AD & Identity Health")
    ws.freeze_panes = "A3"
    _title(ws, "Active Directory & LDAP Integration Health", "F")

    ad_cols = ["Cluster", "AD Domain", "Workgroup", "Connection Status",
               "Preferred DC", "Trusted Domains"]
    _hdr(ws, 2, ad_cols)

    has_ad = False
    for cd in all_data:
        cluster = cd["name"]
        for ad in (cd.get("active_directory") or []):
            has_ad = True
            raw_status   = ad.get("connectionStatus") or ad.get("status") or "Unknown"
            status_clean = raw_status.lstrip("k")
            pref_dcs     = ad.get("preferredDomainControllers") or []
            pref_dc_str  = ", ".join(
                (dc.get("domainName") or dc.get("ipAddress") or "")
                for dc in pref_dcs if isinstance(dc, dict)
            ) or "Auto"
            trusted = len(ad.get("trustedDomains") or [])
            rn = ws.max_row + 1
            vals = [cluster, ad.get("domainName") or "",
                    ad.get("workgroup") or "", status_clean,
                    pref_dc_str, trusted]
            for c, v in enumerate(vals, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(v)).font = _FONT_NORMAL
            sc = ws.cell(row=rn, column=4)
            if "connected" in status_clean.lower() and "dis" not in status_clean.lower():
                sc.fill = PatternFill("solid", fgColor="C6EFCE")
                sc.font = Font(name="Calibri", size=9, color="276221")
            elif any(w in status_clean.lower() for w in ("disconnect", "error", "fail")):
                sc.fill = PatternFill("solid", fgColor="FFC7CE")
                sc.font = Font(name="Calibri", size=9, color="9C0006", bold=True)

    if not has_ad:
        ws.cell(row=ws.max_row + 1, column=1,
                value="No Active Directory domains configured or data unavailable."
                ).font = _FONT_NORMAL

    # ── LDAP section ──────────────────────────────────────────────────────────
    ldap_start = ws.max_row + 2
    ws.merge_cells(f"A{ldap_start}:F{ldap_start}")
    hc = ws.cell(row=ldap_start, column=1, value="─── LDAP Providers ───")
    hc.font      = Font(name="Calibri", size=9, bold=True, italic=True)
    hc.fill      = PatternFill("solid", fgColor=LIGHT_GRAY)
    hc.alignment = Alignment(horizontal="center")

    _hdr(ws, ws.max_row + 1,
         ["Cluster", "Provider Name", "Server", "Port", "Base DN", "Status"])

    has_ldap = False
    for cd in all_data:
        cluster = cd["name"]
        for lp in (cd.get("ldap_providers") or []):
            has_ldap = True
            raw_status   = lp.get("connectionStatus") or lp.get("status") or "Unknown"
            status_clean = raw_status.lstrip("k")
            servers = lp.get("servers") or lp.get("server") or lp.get("serverIp") or ""
            if isinstance(servers, list):
                servers = ", ".join(str(s) for s in servers)
            rn = ws.max_row + 1
            vals = [cluster,
                    lp.get("name") or lp.get("id") or "",
                    servers,
                    lp.get("port") or "",
                    lp.get("baseDistinguishedName") or lp.get("baseDn") or "",
                    status_clean]
            for c, v in enumerate(vals, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(v)).font = _FONT_NORMAL
            sc = ws.cell(row=rn, column=6)
            if "connected" in status_clean.lower() and "dis" not in status_clean.lower():
                sc.fill = PatternFill("solid", fgColor="C6EFCE")
                sc.font = Font(name="Calibri", size=9, color="276221")
            elif any(w in status_clean.lower() for w in ("disconnect", "error", "fail")):
                sc.fill = PatternFill("solid", fgColor="FFC7CE")
                sc.font = Font(name="Calibri", size=9, color="9C0006", bold=True)

    if not has_ldap:
        ws.cell(row=ws.max_row + 1, column=1,
                value="No LDAP providers configured or data unavailable."
                ).font = _FONT_NORMAL

    auto_fit_columns(ws)


# ─────────────────────────────────────────────────────────────────────────────
# SHEET — CERTIFICATE INVENTORY
# ─────────────────────────────────────────────────────────────────────────────

def _sheet_cert_inventory(wb, all_data):
    from openpyxl.styles import PatternFill, Font

    ws = wb.create_sheet("Certificate Inventory")
    ws.freeze_panes = "A3"
    _title(ws, "TLS Certificate Inventory", "I")

    cols = ["Cluster", "Name", "Type", "Subject / CN", "Issuer",
            "Issued Date", "Expiry Date", "Days Left", "Status"]
    _hdr(ws, 2, cols)

    now_ms = int(time.time() * 1000)

    for cd in all_data:
        cluster   = cd["name"]
        cert_list = cd.get("all_certs") or cd.get("certs") or []
        if not cert_list:
            rn = ws.max_row + 1
            ws.cell(row=rn, column=1, value=cluster).font = _FONT_NORMAL
            ws.cell(row=rn, column=2,
                    value="No certificate data available").font = _FONT_NORMAL
            continue

        seen = set()
        for cert in cert_list:
            cname = cert.get("name") or cert.get("subject") or cert.get("commonName") or ""
            if cname in seen:
                continue
            seen.add(cname)

            ctype = cert.get("type") or cert.get("certificateType") or ""
            ctype_clean = ctype.lstrip("k") if ctype else "WebServer"
            subject = (cert.get("subject") or cert.get("commonName") or
                       cert.get("subjectCommonName") or "")
            issuer  = cert.get("issuer") or cert.get("issuerName") or ""

            expiry_ms = (cert.get("expiryDateMsecs") or cert.get("notAfterMsecs") or 0)
            if not expiry_ms:
                eu = cert.get("expiryDateUsecs") or cert.get("expirationDateUsecs") or 0
                if eu:
                    expiry_ms = eu // 1000

            issued_ms = cert.get("issuedDateMsecs") or cert.get("notBeforeMsecs") or 0
            if not issued_ms:
                iu = cert.get("issuedDateUsecs") or 0
                if iu:
                    issued_ms = iu // 1000

            expiry_str, issued_str, days_left = "", "", None
            if expiry_ms:
                expiry_str = datetime.datetime.utcfromtimestamp(
                    expiry_ms / 1000).strftime("%Y-%m-%d")
                days_left  = (expiry_ms - now_ms) / (1000 * 86400)
            if issued_ms:
                issued_str = datetime.datetime.utcfromtimestamp(
                    issued_ms / 1000).strftime("%Y-%m-%d")

            if days_left is None:
                status, sbg, sfg = "Unknown", "F2F2F2", "000000"
            elif days_left < 0:
                status, sbg, sfg = "EXPIRED",  "FF0000", "FFFFFF"
            elif days_left < 30:
                status, sbg, sfg = "CRITICAL", RED,     "FFFFFF"
            elif days_left < 90:
                status, sbg, sfg = "HIGH",     AMBER,   "000000"
            else:
                status, sbg, sfg = "OK",       "C6EFCE", "276221"

            rn   = ws.max_row + 1
            vals = [cluster, cname, ctype_clean, subject, issuer,
                    issued_str, expiry_str,
                    round(days_left) if days_left is not None else "",
                    status]
            for c, v in enumerate(vals, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(v)).font = _FONT_NORMAL
            sc = ws.cell(row=rn, column=9)
            sc.fill = PatternFill("solid", fgColor=sbg)
            sc.font = Font(name="Calibri", size=9, color=sfg,
                           bold=(status in ("EXPIRED", "CRITICAL")))

    auto_fit_columns(ws)
    ws.column_dimensions["D"].width = 30
    ws.column_dimensions["E"].width = 25


# ─────────────────────────────────────────────────────────────────────────────
# SHEET — CAPACITY & PERFORMANCE TRENDS
# ─────────────────────────────────────────────────────────────────────────────

def _sheet_perf_trends(wb, all_data):
    from openpyxl.chart import LineChart, Reference
    from openpyxl.chart.series import SeriesLabel

    ws = wb.create_sheet("Capacity & Perf Trends")
    ws.freeze_panes = "A3"
    _title(ws, "Capacity & Performance Trends (30 days) — Dedup Ratio & I/O Throughput", "G")

    cols = ["Cluster", "Date", "Logical (TB)", "Physical (TB)", "Dedup Ratio",
            "Read (GB/day)", "Write (GB/day)"]
    _hdr(ws, 2, cols)

    for cd in all_data:
        # ── Build daily physical bytes map (already fetched for cap trend) ────
        phys_map = {}
        for pt in (cd.get("cap_timeseries") or []):
            ts  = pt.get("timestampMsecs")
            val = (pt.get("data") or {}).get("int64Value") or 0
            if ts and val:
                day = datetime.datetime.utcfromtimestamp(ts / 1000).strftime("%Y-%m-%d")
                phys_map[day] = val

        # ── Logical bytes (new: kLogicalCapacityBytes) ────────────────────────
        log_map = {}
        for pt in (cd.get("dedup_ts") or []):
            ts  = pt.get("timestampMsecs")
            val = (pt.get("data") or {}).get("int64Value") or 0
            if ts and val:
                day = datetime.datetime.utcfromtimestamp(ts / 1000).strftime("%Y-%m-%d")
                log_map[day] = val

        # ── I/O maps ──────────────────────────────────────────────────────────
        io_ts    = cd.get("io_ts") or {}
        read_map = {}
        for pt in (io_ts.get("reads") or []):
            ts  = pt.get("timestampMsecs")
            val = (pt.get("data") or {}).get("int64Value") or 0
            if ts:
                day = datetime.datetime.utcfromtimestamp(ts / 1000).strftime("%Y-%m-%d")
                read_map[day] = round(val / 1e9, 2)
        write_map = {}
        for pt in (io_ts.get("writes") or []):
            ts  = pt.get("timestampMsecs")
            val = (pt.get("data") or {}).get("int64Value") or 0
            if ts:
                day = datetime.datetime.utcfromtimestamp(ts / 1000).strftime("%Y-%m-%d")
                write_map[day] = round(val / 1e9, 2)

        all_days = sorted(set(phys_map) | set(log_map) | set(read_map) | set(write_map))
        if not all_days:
            continue

        trend_start = ws.max_row + 1
        for day in all_days:
            phys  = phys_map.get(day, 0)
            log_b = log_map.get(day, 0)
            ratio = round(log_b / phys, 2) if phys and log_b else ""
            rn    = ws.max_row + 1
            row   = [cd["name"], day,
                     round(log_b / 1e12, 3) if log_b else "",
                     round(phys  / 1e12, 3) if phys  else "",
                     ratio,
                     read_map.get(day, ""),
                     write_map.get(day, "")]
            for c, val in enumerate(row, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL
        trend_end = ws.max_row
        n_days    = trend_end - trend_start + 1

        if n_days < 2:
            continue

        # ── Dedup ratio chart ─────────────────────────────────────────────────
        dc = LineChart()
        dc.style  = 10
        dc.height = 14
        dc.width  = 26
        try:
            from openpyxl.chart.title import Title as _CT
            from openpyxl.drawing.text import (RichTextBodyProperties as _BP,
                Paragraph as _Par, RegularTextRun as _Ru, RunProperties as _RP)
            from openpyxl.chart.text import RichText as _CRT, Text as _CTx
            _t = _CT(); _t.overlay = False
            _r2 = _CRT(); _r2.bodyPr = _BP()
            _p = _Par(); _ru = _Ru()
            _ru.t = f"{cd['name']} — Dedup Ratio Trend (Ransomware Indicator)"
            _ru.rPr = _RP(sz=1200, b=True)
            _p.r.append(_ru); _r2.p.append(_p); _t.tx = _CTx(rich=_r2)
            dc.title = _t
        except Exception:
            dc.title = f"{cd['name']} — Dedup Ratio Trend"

        dc.y_axis.title         = "Dedup Ratio (x)"
        dc.y_axis.numFmt        = "0.00"
        dc.y_axis.majorTickMark = "out"
        dc.y_axis.tickLblPos    = "nextTo"
        dc.y_axis.delete        = False
        dc.x_axis.title         = "Date"
        dc.x_axis.majorTickMark = "out"
        dc.x_axis.tickLblPos    = "low"
        dc.x_axis.delete        = False
        if n_days > 14:
            dc.x_axis.tickLblSkip = 2

        dc.add_data(Reference(ws, min_col=5, min_row=trend_start, max_row=trend_end))
        try:
            from openpyxl.chart.data_source import DataSource as _DS, StrRef as _SR
            dc.series[0].cat = _DS(strRef=_SR(
                f=f"'{ws.title}'!$B${trend_start}:$B${trend_end}"))
        except Exception:
            pass
        try:
            s = dc.series[0]
            s.title = SeriesLabel(v="Dedup Ratio")
            s.graphicalProperties.line.solidFill = "7030A0"   # purple
            s.graphicalProperties.line.width     = 18000
            s.marker.symbol  = "circle"
            s.marker.size    = 5
            s.marker.graphicalProperties.solidFill        = "7030A0"
            s.marker.graphicalProperties.line.solidFill   = "7030A0"
        except (IndexError, AttributeError):
            pass
        ws.add_chart(dc, f"I{trend_start}")

        # ── I/O throughput chart ───────────────────────────────────────────────
        ic = LineChart()
        ic.style  = 10
        ic.height = 14
        ic.width  = 26
        try:
            from openpyxl.chart.title import Title as _CT2
            from openpyxl.drawing.text import (RichTextBodyProperties as _BP2,
                Paragraph as _Par2, RegularTextRun as _Ru2, RunProperties as _RP2)
            from openpyxl.chart.text import RichText as _CRT2, Text as _CTx2
            _t2 = _CT2(); _t2.overlay = False
            _r3 = _CRT2(); _r3.bodyPr = _BP2()
            _p2 = _Par2(); _ru2 = _Ru2()
            _ru2.t = f"{cd['name']} — Daily I/O Throughput (GB)"
            _ru2.rPr = _RP2(sz=1200, b=True)
            _p2.r.append(_ru2); _r3.p.append(_p2); _t2.tx = _CTx2(rich=_r3)
            ic.title = _t2
        except Exception:
            ic.title = f"{cd['name']} — Daily I/O Throughput"

        ic.y_axis.title         = "GB / day"
        ic.y_axis.numFmt        = "0.0"
        ic.y_axis.majorTickMark = "out"
        ic.y_axis.tickLblPos    = "nextTo"
        ic.y_axis.delete        = False
        ic.x_axis.title         = "Date"
        ic.x_axis.majorTickMark = "out"
        ic.x_axis.tickLblPos    = "low"
        ic.x_axis.delete        = False
        if n_days > 14:
            ic.x_axis.tickLblSkip = 2

        ic.add_data(Reference(ws, min_col=6, min_row=trend_start, max_row=trend_end))
        ic.add_data(Reference(ws, min_col=7, min_row=trend_start, max_row=trend_end))
        try:
            from openpyxl.chart.data_source import DataSource as _DS3, StrRef as _SR3
            _cat_f = f"'{ws.title}'!$B${trend_start}:$B${trend_end}"
            for _si in range(2):
                ic.series[_si].cat = _DS3(strRef=_SR3(f=_cat_f))
        except Exception:
            pass
        try:
            ic.series[0].title = SeriesLabel(v="Read GB/day")
            ic.series[0].graphicalProperties.line.solidFill = "0070C0"
            ic.series[0].graphicalProperties.line.width     = 18000
            ic.series[0].marker.symbol = "circle"
            ic.series[0].marker.size   = 5
            ic.series[0].marker.graphicalProperties.solidFill      = "0070C0"
            ic.series[0].marker.graphicalProperties.line.solidFill = "0070C0"
            ic.series[1].title = SeriesLabel(v="Write GB/day")
            ic.series[1].graphicalProperties.line.solidFill = "FF4500"
            ic.series[1].graphicalProperties.line.width     = 18000
            ic.series[1].marker.symbol = "diamond"
            ic.series[1].marker.size   = 5
            ic.series[1].marker.graphicalProperties.solidFill      = "FF4500"
            ic.series[1].marker.graphicalProperties.line.solidFill = "FF4500"
        except (IndexError, AttributeError):
            pass
        try:
            from openpyxl.chart.legend import Legend as _CL2
            ic.legend = _CL2(); ic.legend.position = "b"
        except Exception:
            pass
        ws.add_chart(ic, f"I{trend_start + 28}")

    auto_fit_columns(ws)


# ─────────────────────────────────────────────────────────────────────────────
# GUIDE SHEET (first tab) — scoring reference and sheet descriptions
# ─────────────────────────────────────────────────────────────────────────────

def _sheet_guide(wb, all_data):
    """Tab 1: Static reference — what each sheet contains and how scoring works."""
    from openpyxl.styles import Alignment as _Aln, PatternFill as _PF, Font as _Fnt

    ws = wb.create_sheet("Guide")
    SPAN  = "F"   # merge across A:F
    NC    = 6     # number of meaningful columns

    ws.column_dimensions["A"].width = 8
    ws.column_dimensions["B"].width = 30
    ws.column_dimensions["C"].width = 58
    ws.column_dimensions["D"].width = 8
    ws.column_dimensions["E"].width = 8
    ws.column_dimensions["F"].width = 8

    _r = [0]

    def _nxt():
        _r[0] += 1
        return _r[0]

    def _blank():
        _r[0] += 1

    def _section(title):
        rn = _nxt()
        ws.merge_cells(f"A{rn}:{SPAN}{rn}")
        c = ws.cell(row=rn, column=1, value=title)
        c.fill = _PF("solid", fgColor=COH_GREEN)
        c.font = _Fnt(bold=True, color=WHITE, size=11, name="Calibri")
        c.alignment = _Aln(horizontal="left", vertical="center", indent=1)
        ws.row_dimensions[rn].height = 22

    def _sub_hdr(*cols):
        rn = _nxt()
        for ci, val in enumerate(cols, 1):
            cell = ws.cell(row=rn, column=ci, value=val)
            cell.fill = _PF("solid", fgColor=DARK_GREEN)
            cell.font = _Fnt(bold=True, color=WHITE, size=9, name="Calibri")
            cell.alignment = _Aln(horizontal="center", vertical="center")
        ws.row_dimensions[rn].height = 16

    def _row(a, b, c, bg=None, bold_a=False, bold_b=False, h=15, wrap_c=True):
        rn = _nxt()
        for ci, (val, bold) in enumerate([(a, bold_a), (b, bold_b), (c, False)], 1):
            cell = ws.cell(row=rn, column=ci, value=_safe_cell(val))
            cell.font = _Fnt(size=9, bold=bold, name="Calibri")
            cell.alignment = _Aln(wrap_text=(ci == 3 and wrap_c), vertical="top")
        if bg:
            for ci in range(1, NC + 1):
                ws.cell(row=rn, column=ci).fill = _PF("solid", fgColor=bg)
        ws.row_dimensions[rn].height = h

    def _color_row(a_val, a_bg, a_fg, b_val, c_val, h=16):
        rn = _nxt()
        c1 = ws.cell(row=rn, column=1, value=a_val)
        c1.fill = _PF("solid", fgColor=a_bg)
        c1.font = _Fnt(size=9, bold=True, color=a_fg, name="Calibri")
        c1.alignment = _Aln(horizontal="center", vertical="center")
        ws.cell(row=rn, column=2, value=b_val).font = _Fnt(size=9, bold=True, name="Calibri")
        c3 = ws.cell(row=rn, column=3, value=c_val)
        c3.font = _Fnt(size=9, name="Calibri")
        c3.alignment = _Aln(wrap_text=True, vertical="top")
        ws.row_dimensions[rn].height = h

    # ── Main title ────────────────────────────────────────────────────────────
    rn = _nxt()
    ws.merge_cells(f"A{rn}:{SPAN}{rn}")
    tc = ws.cell(row=rn, column=1,
                 value="Cohesity Health Check — Workbook Guide & Scoring Reference")
    tc.fill = _PF("solid", fgColor=DARK_GREEN)
    tc.font = _Fnt(bold=True, color=WHITE, size=14, name="Calibri")
    tc.alignment = _Aln(horizontal="center", vertical="center")
    ws.row_dimensions[rn].height = 32
    _blank()

    # ── About ─────────────────────────────────────────────────────────────────
    _section("About This Report")
    rn = _nxt()
    ws.merge_cells(f"A{rn}:{SPAN}{rn}")
    ac = ws.cell(row=rn, column=1,
                 value=(
                     "This workbook provides a comprehensive assessment of your Cohesity "
                     "environment, generated live from the Cohesity Helios API or directly "
                     "from a cluster. Each tab covers a specific area of infrastructure health, "
                     "protection group performance, storage capacity, security posture, and "
                     "recovery readiness. Alongside this workbook, the script also produces a "
                     "Word narrative report, and optionally a comprehensive ~31-slide PowerPoint "
                     "deck (.pptx, requires python-pptx) with RAG colour-coding and scoring "
                     "methodology reference slides. "
                     "Use this Guide tab as a quick reference for understanding what each sheet "
                     "contains and how the scores are calculated."
                 ))
    ac.font = _Fnt(size=9, name="Calibri")
    ac.alignment = _Aln(wrap_text=True, vertical="top")
    ws.row_dimensions[rn].height = 42
    _blank()

    # ── Workbook Contents ─────────────────────────────────────────────────────
    _section("Workbook Contents")
    _sub_hdr("#", "Sheet / Tab Name", "What It Contains")
    sheets_info = [
        ("Guide",                  "This tab — scoring methodology and sheet descriptions"),
        ("Executive Summary",      "Per-cluster health score, grade, success %, capacity used %, open criticals, "
                                   "ransomware readiness score, and capacity runway forecast"),
        ("Infrastructure",         "Software version, node count, healthy nodes, cluster/storage-domain encryption, "
                                   "DNS/NTP settings, software lifecycle status, end-of-support date, and fault tolerance margin"),
        ("Node Hardware",          "Per-node hardware model, serial, storage tiers, raw capacity, hardware EOL status, "
                                   "and software version match vs cluster version"),
        ("Disk Health",            "Per-disk status, SSD wear %, encryption, storage tier — CRITICAL on failed disks, "
                                   "HIGH on SSD wear ≥80%"),
        ("Protection Health",      "Per-group last-run status, SLA violations, RPO gap, object failure counts, "
                                   "and DataLock (WORM) coverage per group"),
        ("Storage & Capacity",     "Cluster and per-domain usable/used/free, data reduction/dedup/compression ratios, "
                                   "and predictive runway to 80% utilization"),
        ("Policy Audit",           "Policy retention schedules, replication targets, archival targets, "
                                   "and DataLock (WORM) mode and duration per policy"),
        ("Policy → Groups",        "Every protection group listed alongside the policy that governs it"),
        ("Alerts",                 "All open critical and warning alerts, sorted by severity, with age and description"),
        ("Security",               "21-column checklist: encryption, vault, FortKnox, replication, audit log, MFA, "
                                   "NTP auth, quorum, TLS cert expiry, ransomware score; plus KMS sub-section "
                                   "(KMS type, server, risk) and security & anomaly alert sub-section"),
        ("Agent Health",           "Per-host agent version, health status, upgradability, and last upgrade error"),
        ("Source Coverage",        "Registered protection sources with protected/unprotected object counts and coverage %"),
        ("Unprotected Objects",    "Sources with unprotected objects — lists object names where available; "
                                   "red if count >10, amber if count >0"),
        ("Recovery Audit",         "Recovery task history for the lookback window: status, duration, objects recovered, "
                                   "plus per-cluster summary of total recoveries, success rate, and days since last recovery"),
        ("Replication & Archive",  "Replication partner targets, archival vault names/types, FortKnox storage consumed, "
                                   "and replication lag in hours"),
        ("FortKnox Data Transfer", "Per-group data transfer to external vaults: logical TB, physical TB, "
                                   "storage consumed, snapshot count, last FK archival date and days since"),
        ("Data Services",          "NAS views with protocol, quota configuration, usage %, and near-quota warnings"),
        ("Data Exposure",          "NAS view access risk: SMB discovery, NFS open mount, S3 enabled, quota presence — "
                                   "HIGH/MEDIUM/LOW risk per view"),
        ("Coverage Gaps",          "Protection groups with a failed last run, paused state, or RPO gap exceeding threshold"),
        ("User Security",          "User accounts with domain, roles, MFA status, locked state, last login, and risk flags — "
                                   "Stale Admin (amber), Admin No MFA (red); plus Custom Role Definitions sub-section "
                                   "showing privilege analysis and risk level for each non-built-in role"),
        ("Trends (30d)",           "Daily backup success rate and average backup duration per cluster over the lookback window "
                                   "with embedded charts"),
        ("Recommendations",        "Prioritized action list (CRITICAL / HIGH / MEDIUM / LOW) with finding and suggested action"),
        ("Workload Risk Heatmap",  "All protection groups scored 0–100 on composite recovery risk, "
                                   "sorted worst-first with full-row color coding"),
        ("Audit Log",              "Last 30 days of configuration changes from the cluster audit trail, "
                                   "categorized by type with high-risk events highlighted"),
        ("DataLock Verification",  "Snapshot-level DataLock verification: checks actual worm/immutability status on recent "
                                   "snapshots for groups with DataLock policies — flags NOT LOCKED (red) vs VERIFIED (green)"),
        ("AD & Identity Health",   "Active Directory domain connection status (Connected/Disconnected) with preferred DC; "
                                   "LDAP provider status, server, port, base DN — flags disconnected identity sources in red"),
        ("Certificate Inventory",  "Full TLS certificate inventory: name, type, subject, issuer, issued/expiry dates, "
                                   "days remaining — CRITICAL (red) if <30 days, HIGH (amber) if <90 days, OK (green) otherwise"),
        ("Capacity & Perf Trends", "30-day dedup ratio trend (sudden drop = ransomware indicator — data becoming incompressible) "
                                   "and daily I/O throughput (read/write GB per day) with embedded line charts per cluster"),
    ]
    for i, (name, desc) in enumerate(sheets_info, 1):
        bg = LIGHT_GRAY if i % 2 == 0 else None
        rn = _nxt()
        ws.cell(row=rn, column=1, value=str(i)).font = _Fnt(size=9, name="Calibri")
        # Hyperlink to the target sheet — escape single quotes in sheet name
        safe = name.replace("'", "''")
        c2 = ws.cell(row=rn, column=2, value=name)
        c2.hyperlink = f"#'{safe}'!A1"
        c2.font      = _Fnt(size=9, bold=True, underline="single",
                            color="0563C1", name="Calibri")
        c3 = ws.cell(row=rn, column=3, value=desc)
        c3.font      = _Fnt(size=9, name="Calibri")
        c3.alignment = _Aln(wrap_text=True, vertical="top")
        if bg:
            for ci in range(1, NC + 1):
                ws.cell(row=rn, column=ci).fill = _PF("solid", fgColor=LIGHT_GRAY)
        ws.row_dimensions[rn].height = 28
    _blank()

    # ── Overall Health Score ──────────────────────────────────────────────────
    _section("Overall Health Score (0–100)")
    _sub_hdr("Weight", "Dimension", "What It Measures")
    score_rows = [
        ("25%", "Protection Success Rate",
         "Percentage of backup runs in the lookback window that succeeded"),
        ("20%", "SLA Compliance",
         "Percentage of runs that completed within the configured SLA window"),
        ("15%", "Infrastructure Health",
         "Percentage of cluster nodes in a healthy state"),
        ("15%", "Storage Capacity",
         "Penalized above 70% used; score reaches 0 at or above 85% utilization"),
        ("15%", "Security Posture",
         "Derived from the Security Score (see below) — weighted contribution to overall"),
        ("10%", "Alert Health",
         "Starts at 100; −10 per open critical alert, −2 per open warning alert (floor 0)"),
    ]
    for i, (wt, dim, desc) in enumerate(score_rows):
        bg = LIGHT_GRAY if i % 2 == 0 else None
        _row(wt, dim, desc, bg=bg, bold_b=True, h=18)

    _blank()
    _sub_hdr("Score Range", "Grade", "")
    grade_rows = [
        ("90–100", "Excellent", LT_GREEN,   "000000"),
        ("75–89",  "Good",      "C6EFCE",   "000000"),
        ("60–74",  "Fair",      YELLOW,     "000000"),
        ("40–59",  "At Risk",   ORANGE,     "FFFFFF"),
        ("<40",    "Critical",  RED,        "FFFFFF"),
    ]
    for score_r, lbl, bg_c, fg_c in grade_rows:
        _color_row(score_r, bg_c, fg_c, lbl, "", h=15)
    _blank()

    # ── Security Score ────────────────────────────────────────────────────────
    _section("Security Score (0–100)")
    rn = _nxt()
    ws.merge_cells(f"A{rn}:{SPAN}{rn}")
    nc = ws.cell(row=rn, column=1,
                 value="Score starts at 100. Points are deducted for each gap found:")
    nc.font      = _Fnt(size=9, italic=True, name="Calibri")
    nc.alignment = _Aln(vertical="center", indent=1)
    ws.row_dimensions[rn].height = 15
    _sub_hdr("Severity", "Deduction", "Finding")
    sec_rows = [
        ("HIGH",   "−20", "No cluster-wide encryption enabled"),
        ("HIGH",   "−15", "Storage domain encryption not configured"),
        ("HIGH",   "−10", "No vault / archival target configured"),
        ("HIGH",   "−10", "No DataLock (WORM) on any policy"),
        ("HIGH",   "−10", "Admin accounts without MFA enabled"),
        ("MEDIUM", "−5",  "No replication target configured"),
        ("MEDIUM", "−5",  "No FortKnox / RPaaS indelible vault"),
        ("MEDIUM", "−5",  "Quorum not enabled (Helios mode only)"),
        ("MEDIUM", "−5",  "Audit logging disabled"),
        ("MEDIUM", "−5",  "TLS certificate expired or expiring within 90 days"),
    ]
    SEV_BG = {"HIGH": "FFD0D0", "MEDIUM": "FFF2CC"}
    for i, (sev, ded, finding) in enumerate(sec_rows):
        bg = LIGHT_GRAY if i % 2 == 0 else None
        rn = _nxt()
        c1 = ws.cell(row=rn, column=1, value=sev)
        c1.fill      = _PF("solid", fgColor=SEV_BG.get(sev, LIGHT_GRAY))
        c1.font      = _Fnt(size=9, bold=True, name="Calibri")
        c1.alignment = _Aln(horizontal="center", vertical="center")
        ws.cell(row=rn, column=2, value=ded).font = _Fnt(size=9, bold=True, name="Calibri")
        c3 = ws.cell(row=rn, column=3, value=finding)
        c3.font = _Fnt(size=9, name="Calibri")
        if bg and i % 2 == 0:
            ws.cell(row=rn, column=3).fill = _PF("solid", fgColor=LIGHT_GRAY)
        ws.row_dimensions[rn].height = 15
    _blank()

    # ── Ransomware Readiness Score ────────────────────────────────────────────
    _section("Ransomware Readiness Score (0–100)")
    rn = _nxt()
    ws.merge_cells(f"A{rn}:{SPAN}{rn}")
    rrc = ws.cell(row=rn, column=1,
                  value='Answers: "If ransomware hit today, could we recover?" — '
                        'scored independently of the overall health score.')
    rrc.font      = _Fnt(size=9, italic=True, name="Calibri")
    rrc.alignment = _Aln(vertical="center", indent=1)
    ws.row_dimensions[rn].height = 15
    _sub_hdr("Max Pts", "Dimension", "Scoring Detail")
    rw_rows = [
        ("30", "DataLock (WORM) Coverage",
         "Full Compliance mode=30 | Full (any mode)=20 | "
         "Partial ≥50% of groups=15 | Partial <50%=8 | None=0"),
        ("25", "FortKnox / Indelible Vault",
         "Active (sending data)=25 | Configured but idle=10 | Not configured=0"),
        ("20", "Recovery Window Depth",
         "Oldest clean snapshot ≥30 days=20 | ≥14 days=15 | "
         "≥7 days=10 | ≥2 days=5 | No successful backup found=0"),
        ("15", "Admin MFA",
         "All admins have MFA enabled=15 | Partial=8 | No admins have MFA=0"),
        ("10", "Quorum",
         "Enabled (or N/A in direct-cluster mode)=10 | Disabled=0"),
    ]
    for i, (pts, dim, detail) in enumerate(rw_rows):
        bg = LIGHT_GRAY if i % 2 == 0 else None
        rn = _nxt()
        ws.cell(row=rn, column=1, value=pts).font = _Fnt(size=9, bold=True, name="Calibri")
        ws.cell(row=rn, column=2, value=dim).font = _Fnt(size=9, bold=True, name="Calibri")
        c3 = ws.cell(row=rn, column=3, value=detail)
        c3.font      = _Fnt(size=9, name="Calibri")
        c3.alignment = _Aln(wrap_text=True, vertical="top")
        if bg:
            for ci in range(1, NC + 1):
                ws.cell(row=rn, column=ci).fill = _PF("solid", fgColor=LIGHT_GRAY)
        ws.row_dimensions[rn].height = 30

    _blank()
    _sub_hdr("Score", "Label", "Interpretation")
    rw_labels = [
        ("≥80",  "Strong",    LT_GREEN,  "000000",
         "Well-positioned to recover — key protections are fully in place"),
        ("60–79","Moderate",  YELLOW,    "000000",
         "Most protections in place but notable gaps remain — review recommended"),
        ("40–59","High Risk", ORANGE,    "FFFFFF",
         "Multiple recovery gaps — ransomware recovery would be significantly difficult"),
        ("<40",  "Critical",  RED,       "FFFFFF",
         "Fundamental protections missing — recovery from ransomware is at severe risk"),
    ]
    for score_r, lbl, bg_c, fg_c, interp in rw_labels:
        _color_row(score_r, bg_c, fg_c, lbl, interp, h=20)
    _blank()

    # ── Workload Risk Score ───────────────────────────────────────────────────
    _section("Workload Risk Score (0–100) — Workload Risk Heatmap Tab")
    _sub_hdr("Max Pts", "Factor", "Scoring Detail")
    wrk_rows = [
        ("35", "Last Run Status",
         "Success=35 | Warning=25 | Running=20 | Skipped or Paused=10 | Failed or unknown=0"),
        ("25", "SLA Compliance",
         "No SLA violation=25 | SLA violated=0"),
        ("25", "RPO Gap (time since last success)",
         "≤4 hours=25 | ≤8h=20 | ≤24h=15 | ≤48h=8 | >48h or no successful run found=0"),
        ("15", "DataLock Coverage",
         "Compliance mode or FortKnox=15 | Administrative WORM=10 | No DataLock=0"),
    ]
    for i, (pts, factor, detail) in enumerate(wrk_rows):
        bg = LIGHT_GRAY if i % 2 == 0 else None
        rn = _nxt()
        ws.cell(row=rn, column=1, value=pts).font    = _Fnt(size=9, bold=True, name="Calibri")
        ws.cell(row=rn, column=2, value=factor).font = _Fnt(size=9, bold=True, name="Calibri")
        c3 = ws.cell(row=rn, column=3, value=detail)
        c3.font      = _Fnt(size=9, name="Calibri")
        c3.alignment = _Aln(wrap_text=True, vertical="top")
        if bg:
            for ci in range(1, NC + 1):
                ws.cell(row=rn, column=ci).fill = _PF("solid", fgColor=LIGHT_GRAY)
        ws.row_dimensions[rn].height = 18

    _blank()
    _sub_hdr("Score", "Risk Level", "Row Color in Heatmap")
    risk_levels = [
        ("≥75",  "Low",      LT_GREEN,  "000000"),
        ("50–74","Medium",   "FFF2CC",  "000000"),
        ("25–49","High",     "FF8C42",  "FFFFFF"),
        ("<25",  "Critical", RED,       "FFFFFF"),
    ]
    for score_r, lvl, bg_c, fg_c in risk_levels:
        _color_row(score_r, bg_c, fg_c, lvl, "Full row in the Workload Risk Heatmap tab is colored accordingly", h=15)
    _blank()

    # ── Recommendation Priorities ─────────────────────────────────────────────
    _section("Recommendation Priorities")
    _sub_hdr("Priority", "Expected Response", "Meaning")
    prio_rows = [
        ("CRITICAL", RED,      "FFFFFF", "Immediate",
         "Active risk — data loss, capacity exhaustion, node failure, or critical security exposure"),
        ("HIGH",     ORANGE,   "FFFFFF", "Within 30 days",
         "Significant gap — unprotected data, no replication, near-capacity, or security misconfiguration"),
        ("MEDIUM",   YELLOW,   "000000", "Within 90 days",
         "Best-practice gap — missing archival, quota governance, partial compliance"),
        ("LOW",      LT_GREEN, "000000", "Next review cycle",
         "Housekeeping — minor policy improvements, quota hygiene, optional enhancements"),
    ]
    for pri, bg_c, fg_c, resp, meaning in prio_rows:
        _color_row(pri, bg_c, fg_c, resp, meaning, h=20)
    _blank()

    # ── Color / RAG Legend ────────────────────────────────────────────────────
    _section("Color / RAG Legend")
    _sub_hdr("Color", "Name", "Meaning in This Workbook")
    legend_rows = [
        (LT_GREEN,  "000000", "Green",       "Healthy / within threshold / fully compliant"),
        ("C6EFCE",  "000000", "Light Green", "Good score (grade: Good) in score-based RAG cells"),
        (YELLOW,    "000000", "Yellow",      "Warning / approaching threshold / partial compliance"),
        (ORANGE,    "000000", "Orange",      "Elevated risk / near-critical / configured but inactive"),
        (RED,       "FFFFFF", "Red",         "Critical / failed / non-compliant / protection missing"),
        ("FFD0D0",  "000000", "Light Red",   "High-severity finding indicator (Security Score table)"),
        ("BDD7EE",  "000000", "Blue",        "Policy Changes category (Audit Log tab)"),
        ("CCFFFF",  "000000", "Teal",        "Cluster Config Changes category (Audit Log tab)"),
        ("FFF2CC",  "000000", "Amber",       "User Changes category / Medium-severity / Warning level"),
        (LIGHT_GRAY,"000000", "Gray",        "Alternating row shading for readability"),
    ]
    for bg_c, fg_c, name, meaning in legend_rows:
        rn = _nxt()
        c1 = ws.cell(row=rn, column=1, value=name)
        c1.fill      = _PF("solid", fgColor=bg_c)
        c1.font      = _Fnt(size=9, bold=True, color=fg_c, name="Calibri")
        c1.alignment = _Aln(horizontal="center", vertical="center")
        c2 = ws.cell(row=rn, column=2, value="")
        c2.fill = _PF("solid", fgColor=bg_c)
        c3 = ws.cell(row=rn, column=3, value=meaning)
        c3.font      = _Fnt(size=9, name="Calibri")
        c3.alignment = _Aln(vertical="center")
        ws.row_dimensions[rn].height = 15


# ─────────────────────────────────────────────────────────────────────────────
# NEW SHEETS — v1.64
# ─────────────────────────────────────────────────────────────────────────────

def _sheet_recovery_history(wb, all_data):
    """Recovery Audit — history of recovery tasks per cluster."""
    import datetime as _dt
    ws = wb.create_sheet("Recovery Audit")
    ws.freeze_panes = "A3"
    _title(ws, "Recovery Audit — Recovery Testing History & Recoverability Status", "L")

    cols = ["Cluster", "Recovery ID", "Name", "Start Time", "End Time",
            "Duration (min)", "Status", "Type", "Objects Recovered",
            "Source Cluster", "Target", "Notes"]
    _hdr(ws, 2, cols)

    # Per-cluster summary data for the bottom section
    cluster_summary = []

    for cd in all_data:
        recoveries = cd.get("recoveries") or []
        days       = cd.get("days", 30)

        total   = len(recoveries)
        success = 0
        last_succ_usecs = None

        for rec in recoveries:
            st = rec.get("status") or ""
            is_succ = "Succeeded" in st or "kSuccess" in st

            start_u = rec.get("startTimeUsecs") or rec.get("creationTimeUsecs") or 0
            end_u   = rec.get("endTimeUsecs") or 0

            if start_u:
                start_str = _dt.datetime.fromtimestamp(
                    start_u / 1_000_000, tz=_dt.timezone.utc).strftime("%Y-%m-%d %H:%M")
            else:
                start_str = ""

            if end_u:
                end_str = _dt.datetime.fromtimestamp(
                    end_u / 1_000_000, tz=_dt.timezone.utc).strftime("%Y-%m-%d %H:%M")
            else:
                end_str = ""

            if start_u and end_u:
                dur = round((end_u - start_u) / 60_000_000, 1)
            else:
                dur = "N/A"

            rec_type   = (rec.get("recoveryAction") or rec.get("type")
                          or rec.get("snapshotEnvironment") or "")
            obj_list   = rec.get("objects") or []
            obj_count  = len(obj_list) if isinstance(obj_list, list) else ""

            src_cluster = ""
            target      = ""
            params      = rec.get("vmwareParams") or rec.get("physicalParams") or {}
            if isinstance(params, dict):
                target = (params.get("targetHost") or params.get("recoveryTargetConfig")
                          or {})
                if isinstance(target, dict):
                    target = target.get("newSourceConfig", {}).get("host", {}).get("name", "") or \
                             target.get("originalSourceConfig", {}).get("host", {}).get("name", "")

            notes = ""
            if is_succ:
                success += 1
                if last_succ_usecs is None or start_u > last_succ_usecs:
                    last_succ_usecs = start_u
            elif "Failed" in st or "kFailed" in st:
                notes = "Failed"

            row = [cd["name"], _safe_cell(str(rec.get("id", ""))),
                   _safe_cell(rec.get("name") or ""),
                   start_str, end_str, dur,
                   _safe_cell(st),
                   _safe_cell(rec_type),
                   obj_count,
                   _safe_cell(str(src_cluster)),
                   _safe_cell(str(target)),
                   notes]
            rn = ws.max_row + 1
            for c, val in enumerate(row, 1):
                ws.cell(row=rn, column=c, value=val).font = _FONT_NORMAL

            # Status coloring (col 7)
            sc = ws.cell(row=rn, column=7)
            if is_succ:
                sc.fill = _fill(LT_GREEN)
            elif "Failed" in st or "kFailed" in st:
                sc.fill = _fill(RED); sc.font = _font(bold=True, color=WHITE)
            else:
                sc.fill = _fill(YELLOW)

        # Compute last-success days
        if last_succ_usecs:
            last_succ_dt  = _dt.datetime.fromtimestamp(
                last_succ_usecs / 1_000_000, tz=_dt.timezone.utc)
            last_succ_str = last_succ_dt.strftime("%Y-%m-%d %H:%M")
            days_since    = (cd["end_usecs"] - last_succ_usecs) // 86_400_000_000
        else:
            last_succ_str = "None"
            days_since    = days

        succ_rate = round(success / total * 100, 1) if total else "N/A"
        cluster_summary.append({
            "name":          cd["name"],
            "total":         total,
            "success":       success,
            "succ_rate":     succ_rate,
            "last_succ":     last_succ_str,
            "days_since":    days_since,
            "days":          days,
        })

    # ── Cluster Summary sub-section ───────────────────────────────────────────
    if cluster_summary:
        blank_rn = ws.max_row + 2
        ws.cell(row=blank_rn, column=1,
                value="Cluster Recovery Summary").font = _font(bold=True, size=11)
        sum_cols = ["Cluster", "Total Recoveries", "Successful",
                    "Success Rate %", "Last Successful Recovery", "Days Since Last Recovery"]
        _hdr(ws, blank_rn + 1, sum_cols)
        for s in cluster_summary:
            rn = ws.max_row + 1
            row = [s["name"], s["total"], s["success"],
                   s["succ_rate"], s["last_succ"], s["days_since"]]
            for c, val in enumerate(row, 1):
                ws.cell(row=rn, column=c, value=_safe_cell(val)).font = _FONT_NORMAL
            # Color days-since col
            dc = ws.cell(row=rn, column=6)
            if s["days_since"] >= s["days"] or s["total"] == 0:
                dc.fill = _fill(RED); dc.font = _font(bold=True, color=WHITE)
            elif s["days_since"] > 7:
                dc.fill = _fill(YELLOW); dc.font = _font(bold=True)
            else:
                dc.fill = _fill(LT_GREEN)

    auto_fit_columns(ws)


def _sheet_datalock_verification(wb, all_data):
    """DataLock Verification — checks actual DataLock status on recent snapshots."""
    ws = wb.create_sheet("DataLock Verification")
    ws.freeze_panes = "A3"
    _title(ws, "DataLock Verification — Snapshot Immutability Confirmed vs Policy", "J")

    cols = ["Cluster", "Protection Group", "Policy",
            "DataLock (Policy)", "Snapshots Checked",
            "DataLocked Snapshots", "Oldest DL Snapshot",
            "Newest DL Snapshot", "Verification Status", "Risk"]
    _hdr(ws, 2, cols)

    import datetime as _dt

    any_data = False
    for cd in all_data:
        entries = cd.get("snapshot_dl") or []
        if not entries:
            continue
        any_data = True

        policy_by_id = cd.get("policy_by_id") or {}
        groups_by_id = {g.get("id"): g for g in (cd.get("groups") or []) if g.get("id")}

        for entry in entries:
            gid       = entry.get("group_id", "")
            grp_name  = entry.get("group_name", gid)
            grp       = groups_by_id.get(gid) or {}
            pol_id    = grp.get("policyId") or (grp.get("policy") or {}).get("id") or ""
            policy    = policy_by_id.get(pol_id) or {}
            pol_name  = policy.get("name", pol_id or "—")
            dl_policy = _policy_datalock_str(policy) if policy else "—"

            snaps        = entry.get("snapshots") or []
            total_checked = len(snaps)
            locked_snaps  = []
            oldest_dl_u   = None
            newest_dl_u   = None

            for s in snaps:
                is_locked = bool(
                    s.get("dataLockConstraints") or
                    (isinstance(s.get("wormProperties"), dict) and
                     s["wormProperties"].get("isLocked")) or
                    s.get("retentionType") in ("kCompliance", "kAdministrative")
                )
                if is_locked:
                    locked_snaps.append(s)
                    exp_u = (s.get("expiryTimeUsecs") or
                             s.get("snapshotTimestampUsecs") or 0)
                    start_u = (s.get("snapshotTimestampUsecs") or
                               s.get("protectionRunStartTimeUsecs") or 0)
                    if start_u:
                        if oldest_dl_u is None or start_u < oldest_dl_u:
                            oldest_dl_u = start_u
                        if newest_dl_u is None or start_u > newest_dl_u:
                            newest_dl_u = start_u

            locked_count = len(locked_snaps)

            oldest_str = ""
            newest_str = ""
            if oldest_dl_u:
                oldest_str = _dt.datetime.fromtimestamp(
                    oldest_dl_u / 1_000_000, tz=_dt.timezone.utc).strftime("%Y-%m-%d")
            if newest_dl_u:
                newest_str = _dt.datetime.fromtimestamp(
                    newest_dl_u / 1_000_000, tz=_dt.timezone.utc).strftime("%Y-%m-%d")

            # Determine verification status
            has_dl_policy = dl_policy != "None" and dl_policy != "—"
            if not total_checked:
                status = "N/A"
                risk   = "N/A"
            elif not has_dl_policy:
                status = "N/A"
                risk   = "LOW"
            elif locked_count == 0:
                status = "NOT LOCKED"
                risk   = "HIGH"
            elif locked_count < total_checked:
                status = "PARTIAL"
                risk   = "MEDIUM"
            else:
                status = "VERIFIED"
                risk   = "LOW"

            row = [cd["name"], _safe_cell(grp_name), _safe_cell(pol_name),
                   _safe_cell(dl_policy),
                   total_checked, locked_count,
                   oldest_str, newest_str,
                   status, risk]
            rn = ws.max_row + 1
            for c, val in enumerate(row, 1):
                ws.cell(row=rn, column=c, value=val).font = _FONT_NORMAL

            # Status coloring (col 9)
            sc = ws.cell(row=rn, column=9)
            if status == "VERIFIED":
                sc.fill = _fill(LT_GREEN)
            elif status == "PARTIAL":
                sc.fill = _fill(YELLOW); sc.font = _font(bold=True)
            elif status == "NOT LOCKED":
                sc.fill = _fill(RED); sc.font = _font(bold=True, color=WHITE)

            # Risk coloring (col 10)
            rc = ws.cell(row=rn, column=10)
            if risk == "HIGH":
                rc.fill = _fill(RED); rc.font = _font(bold=True, color=WHITE)
            elif risk == "MEDIUM":
                rc.fill = _fill(YELLOW); rc.font = _font(bold=True)
            elif risk == "LOW":
                rc.fill = _fill(LT_GREEN)

    if not any_data:
        ws.cell(row=3, column=1,
                value="No DataLock groups found or snapshot verification skipped — "
                      "no DataLock policies detected on this cluster.") \
           .font = _font(color="595959")

    auto_fit_columns(ws)


# ─────────────────────────────────────────────────────────────────────────────
# EXCEL ORCHESTRATOR
# ─────────────────────────────────────────────────────────────────────────────

def write_excel(all_data, args):
    if not EXCEL_OK:
        print("  WARN: openpyxl not installed — skipping Excel output.")
        print("        Install: pip install openpyxl")
        return None

    wb = Workbook()
    wb.remove(wb.active)

    print("  Writing Excel sheets...")
    _sheet_guide(wb, all_data)             # Guide (first tab)
    _sheet_summary(wb, all_data)           # 1
    _sheet_infrastructure(wb, all_data)    # 2
    _sheet_hardware(wb, all_data)          # 3
    _sheet_disk_health(wb, all_data)       # 4
    _sheet_protection(wb, all_data)        # 5
    _sheet_storage(wb, all_data)           # 6
    _sheet_policies(wb, all_data)          # 7
    _sheet_policy_groups(wb, all_data)     # 8
    _sheet_alerts(wb, all_data)            # 9
    _sheet_security(wb, all_data)          # 10 (+ KMS sub-section)
    _sheet_agent_health(wb, all_data)      # 11
    _sheet_source_coverage(wb, all_data)   # 12
    _sheet_unprotected_objects(wb, all_data) # 13
    _sheet_recovery_history(wb, all_data)  # 14 NEW
    _sheet_replication(wb, all_data)       # 15
    _sheet_fortknox_detail(wb, all_data)   # 16
    _sheet_views(wb, all_data)             # 17
    _sheet_data_exposure(wb, all_data)     # 18
    _sheet_coverage(wb, all_data)          # 19
    _sheet_user_security(wb, all_data)     # 20 (+ Custom Roles sub-section)
    _sheet_trends(wb, all_data)            # 21
    _sheet_recommendations(wb, all_data)   # 22
    _sheet_risk_heatmap(wb, all_data)      # 23
    _sheet_audit_log(wb, all_data)         # 24
    _sheet_datalock_verification(wb, all_data)  # 25
    _sheet_ad_identity(wb, all_data)       # 26 NEW
    _sheet_cert_inventory(wb, all_data)    # 27 NEW
    _sheet_perf_trends(wb, all_data)       # 28 NEW

    out = f"{args.output}.xlsx"
    wb.save(out)
    print(f"  Excel  → {out}")
    return out


# ─────────────────────────────────────────────────────────────────────────────
# WORD DOCUMENT
# ─────────────────────────────────────────────────────────────────────────────

def write_word(all_data, args):
    if not DOCX_OK:
        print("  WARN: python-docx not installed — skipping Word output.")
        print("        Install: pip install python-docx")
        return None

    doc      = Document()
    customer = args.customer or "Customer"
    today    = datetime.date.today().strftime("%B %d, %Y")

    # ── base font ──────────────────────────────────────────────────────────
    for style_name in ("Normal", "Heading 1", "Heading 2"):
        try:
            doc.styles[style_name].font.name = "Calibri"
        except Exception:
            pass

    def _h1(text):
        p = doc.add_heading(text, level=1)
        for run in p.runs:
            run.font.color.rgb = RGBColor(0x70, 0xAD, 0x47)
        return p

    def _h2(text):
        p = doc.add_heading(text, level=2)
        for run in p.runs:
            run.font.color.rgb = RGBColor(0x70, 0xAD, 0x47)
        return p

    def _cell_shade(cell, hex_color):
        """Apply background color to a docx table cell."""
        tc   = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd  = OxmlElement("w:shd")
        shd.set(qn("w:val"),   "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"),  hex_color)
        tcPr.append(shd)

    def _word_autofit(table):
        """Mark a python-docx table as autofit so Word resizes columns to content."""
        tblPr = table._tbl.tblPr
        tblLayout = tblPr.find(qn("w:tblLayout"))
        if tblLayout is None:
            tblLayout = OxmlElement("w:tblLayout")
            tblPr.append(tblLayout)
        tblLayout.set(qn("w:type"), "autofit")

    def _word_fit_widths(table):
        """Redistribute column widths proportional to content. Call after rows are added."""
        _PAGE_W = 9360   # twips — 6.5 inch usable width at 1440 twips/inch
        _MIN_C  = 720    # 0.5 inch minimum per column
        _CAP_CH = 55     # cap contribution at 55 chars to prevent one column dominating

        nc = len(table.columns)
        col_lens = [1] * nc
        for row in table.rows:
            for ci, cell in enumerate(row.cells):
                for para in cell.paragraphs:
                    for line in para.text.split("\n"):
                        col_lens[ci] = max(col_lens[ci], min(len(line), _CAP_CH))

        total_len = sum(col_lens)
        raw = [max(_MIN_C, int(_PAGE_W * l / total_len)) for l in col_lens]

        # Scale to exactly _PAGE_W and fix any rounding residual
        diff = _PAGE_W - sum(raw)
        raw[-1] = max(_MIN_C, raw[-1] + diff)

        tblPr = table._tbl.tblPr
        tblW = tblPr.find(qn("w:tblW"))
        if tblW is None:
            tblW = OxmlElement("w:tblW")
            tblPr.insert(0, tblW)
        tblW.set(qn("w:w"), str(_PAGE_W))
        tblW.set(qn("w:type"), "dxa")

        col_els = table._tbl.tblGrid.findall(qn("w:gridCol"))
        for i, col_el in enumerate(col_els):
            if i < len(raw):
                col_el.set(qn("w:w"), str(raw[i]))

    def _tbl_header(table, cols, bg=COH_GREEN):
        row = table.rows[0].cells
        for i, col in enumerate(cols):
            row[i].text = col
            _cell_shade(row[i], bg)
            for para in row[i].paragraphs:
                for run in para.runs:
                    run.font.bold  = True
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    run.font.size  = Pt(9)

    # ── Cover page ─────────────────────────────────────────────────────────
    for _ in range(6):
        doc.add_paragraph("")

    cp = doc.add_paragraph()
    cp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = cp.add_run("Cohesity Health Check Report")
    run.font.size  = Pt(28)
    run.font.bold  = True
    run.font.color.rgb = RGBColor(0x70, 0xAD, 0x47)

    doc.add_paragraph("")
    cp2 = doc.add_paragraph()
    cp2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = cp2.add_run(f"Prepared for: {customer}\n{today}")
    r2.font.size = Pt(16)

    doc.add_paragraph("")
    cp3 = doc.add_paragraph()
    cp3.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r3 = cp3.add_run("CONFIDENTIAL — For authorized recipients only")
    r3.font.size   = Pt(10)
    r3.font.italic = True
    r3.font.color.rgb = RGBColor(0x59, 0x59, 0x59)

    doc.add_page_break()

    # ── 1. Executive Summary ───────────────────────────────────────────────
    _h1("1. Executive Summary")

    n_clusters = len(all_data)
    avg_score  = round(sum(cd["scores"]["overall"] for cd in all_data) / n_clusters, 1) \
                 if n_clusters else 0
    all_recs   = [(cd["name"], r) for cd in all_data for r in cd["recommendations"]]
    n_crit     = sum(1 for _, r in all_recs if r["priority"] == "CRITICAL")
    n_high     = sum(1 for _, r in all_recs if r["priority"] == "HIGH")

    summary = (
        f"This report covers a Cohesity environment comprising {n_clusters} "
        f"cluster(s) managed through Helios. The overall environment health score "
        f"is {avg_score}/100 ({grade(avg_score)}). "
    )
    if n_crit:
        summary += (f"There are {n_crit} critical finding(s) requiring immediate "
                    f"attention. ")
    if n_high:
        summary += f"{n_high} high-priority item(s) require action within 30 days. "
    if not n_crit and not n_high:
        summary += "No critical or high-priority gaps were identified. "
    summary += ("Detailed findings and prioritized recommendations are included "
                "in the sections below, with full data in the accompanying Excel workbook.")
    doc.add_paragraph(summary)

    t1 = doc.add_table(rows=1, cols=5)
    t1.style = "Table Grid"; _word_autofit(t1)
    _tbl_header(t1, ["Cluster", "Health Score", "Grade", "Open Criticals", "Top Finding"])
    for cd in all_data:
        row = t1.add_row().cells
        crits = sum(1 for a in cd["alerts"] if a.get("severity") == "kCritical")
        top   = cd["recommendations"][0]["finding"] if cd["recommendations"] else "None"
        row[0].text = cd["name"]
        row[1].text = str(cd["scores"]["overall"])
        row[2].text = cd["scores"]["grade"]
        row[3].text = str(crits)
        row[4].text = top
        for para in row[1].paragraphs:
            for run in para.runs:
                run.font.bold = True
    _word_fit_widths(t1)

    # ── Top At-Risk Workloads ──────────────────────────────────────────────
    _h2("Top At-Risk Workloads")
    doc.add_paragraph(
        "The following protection groups have the highest recovery risk across the "
        "environment, scored by last-run status, SLA compliance, RPO gap, and "
        "DataLock coverage. Groups are sorted by risk level (Critical and High "
        "shown here; full detail in the Workload Risk Heatmap sheet)."
    )

    # Collect and score all groups; take top 5 Critical/High
    _risk_rows = []
    for cd in all_data:
        pol_by_id = cd.get("policy_by_id") or {p["id"]: p for p in (cd.get("policies") or []) if p.get("id")}
        for grp in (cd.get("groups") or []):
            score, level = _group_risk_score(grp, cd)
            if level in ("Critical", "High"):
                pol_id   = grp.get("policyId") or (grp.get("policy") or {}).get("id") or ""
                pol_name = (pol_by_id.get(pol_id) or {}).get("name", pol_id or "—")
                last_run = grp.get("lastRun") or {}
                run_sum  = last_run.get("backupRun") or last_run.get("protectionRun") or last_run
                status   = (run_sum.get("status") or run_sum.get("runStatus") or "Unknown")
                # Primary issue
                if score < 10:
                    issue = "Multiple critical failures"
                elif level == "Critical":
                    issue = "Backup failing / no recent successful run"
                else:
                    issue = "SLA or RPO violation detected"
                _risk_rows.append((score, cd["name"], grp.get("name", ""), pol_name, level, issue))

    _risk_rows.sort(key=lambda x: x[0])  # worst first
    _risk_top = _risk_rows[:5]

    if _risk_top:
        t_risk = doc.add_table(rows=1, cols=5)
        t_risk.style = "Table Grid"; _word_autofit(t_risk)
        _tbl_header(t_risk, ["Cluster", "Group Name", "Policy", "Risk Level", "Primary Issue"])
        for _, cl_name, grp_name, pol_n, lvl, iss in _risk_top:
            row = t_risk.add_row().cells
            row[0].text = cl_name
            row[1].text = grp_name
            row[2].text = pol_n
            row[3].text = lvl
            row[4].text = iss
            if lvl == "Critical":
                _cell_shade(row[3], "FF4C4C")
                for para in row[3].paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                        run.font.bold = True
            elif lvl == "High":
                _cell_shade(row[3], "FF8C42")
                for para in row[3].paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                        run.font.bold = True
        _word_fit_widths(t_risk)
    else:
        doc.add_paragraph("No critical or high-risk protection groups identified.")

    doc.add_page_break()

    # ── 2. Environment Overview ────────────────────────────────────────────
    _h1("2. Environment Overview")
    doc.add_paragraph(
        f"The environment spans {n_clusters} cluster(s). "
        "The table below summarizes key infrastructure attributes for each cluster."
    )
    t2 = doc.add_table(rows=1, cols=7)
    t2.style = "Table Grid"; _word_autofit(t2)
    _tbl_header(t2, ["Cluster", "Version", "Nodes", "Healthy",
                     "Cluster Enc", "SD Enc", "Status"])
    for cd in all_data:
        info  = cd["info"]
        nodes = cd["nodes"]
        healthy = sum(1 for n in nodes
                      if ((n.get("nodeStatus") or {}).get("overallStatus") == "kNormal"
                          or n.get("removalState") in (None, "kDontRemove")))
        cluster_enc_w, sd_enc_w = _sd_enc_status(cd)
        sd_enc_str_w = ("Yes" if sd_enc_w is True else
                        "No"  if sd_enc_w is False else "N/A")
        status = "OK" if healthy == len(nodes) else f"WARN: {len(nodes)-healthy} node(s)"
        row = t2.add_row().cells
        for i, val in enumerate([cd["name"], info.get("clusterSoftwareVersion", ""),
                                  len(nodes), healthy,
                                  "Yes" if cluster_enc_w else "No",
                                  sd_enc_str_w, status]):
            row[i].text = str(val)
    _word_fit_widths(t2)

    _h2("Software & Hardware Lifecycle")
    doc.add_paragraph(
        "The table below summarizes the software lifecycle status for each cluster. "
        "Cohesity software versions follow a feature-release (9-month support) and "
        "LTS (12+ month support) model. All current in-support versions (6.8.2 LTS, "
        "7.1.2_u2 LTS, and 7.2.x feature releases) reach end of support on "
        "June 8, 2026. Version 7.3 and later follow a new lifecycle policy."
    )
    # SW lifecycle table
    tbl = doc.add_table(rows=1, cols=4)
    tbl.style = "Table Grid"; _word_autofit(tbl)
    hdr = tbl.rows[0].cells
    for i, h in enumerate(["Cluster", "Version", "Lifecycle Status", "EOS Date"]):
        hdr[i].text = h
        for run in hdr[i].paragraphs[0].runs:
            run.font.bold = True
    for cd in all_data:
        sw_ver = (cd["info"].get("clusterSoftwareVersion") or "").strip()
        _, sw_eos_date, sw_label = _sw_eos(sw_ver)
        row = tbl.add_row().cells
        row[0].text = cd["name"]
        row[1].text = sw_ver
        row[2].text = sw_label
        row[3].text = sw_eos_date.isoformat() if sw_eos_date else "—"
    _word_fit_widths(tbl)
    doc.add_paragraph("")

    doc.add_paragraph(
        "Hardware end-of-life (EOL) dates are provided by Cohesity's Platform EOL Terms. "
        "EOL hardware is no longer eligible for manufacturer support or standard replacement. "
        "Note: Cohesity C2000, C2500, and C3500 platform nodes reached EOL and are no longer "
        "eligible for support under the standard agreement. C4000/CN series hardware is in "
        "service until February 2027. See the Node Hardware sheet in the Excel workbook for "
        "per-node EOL status."
    )

    _h2("Environment Topology")
    doc.add_paragraph(
        "The table below shows each source cluster with its configured replication targets "
        "and archival/vault destinations, including FortKnox vaults."
    )
    t_topo = doc.add_table(rows=1, cols=4)
    t_topo.style = "Table Grid"; _word_autofit(t_topo)
    _tbl_header(t_topo, ["Source Cluster", "Replication Targets",
                          "Archival Targets", "FortKnox Vaults"])
    for cd in all_data:
        policies_t = cd["policies"]
        vaults_t   = cd["vaults"]

        # Collect unique replication target names across all policies
        repl_names = sorted({
            (t.get("targetName")
             or (t.get("remoteTargetConfig") or {}).get("clusterName")
             or (t.get("remoteTargetConfig") or {}).get("name")
             or "Unknown")
            for p in policies_t
            for t in ((p.get("remoteTargetPolicy") or {}).get("replicationTargets") or [])
        })
        # Collect unique archival target names (non-FortKnox)
        arch_names = sorted({
            (t.get("targetName")
             or (t.get("archivalTargetConfig") or {}).get("name")
             or (t.get("archivalTargetConfig") or {}).get("vaultName")
             or "Unknown")
            for p in policies_t
            for t in ((p.get("remoteTargetPolicy") or {}).get("archivalTargets") or [])
            if (t.get("archivalTargetConfig") or {}).get("targetType", "") != "kFortKnox"
        })
        # FortKnox: use cached fk_status for active/idle/none, then collect
        # vault names using the same comprehensive detection logic.
        fk_st_t   = cd.get("fk_status") or _fk_status(cd)
        _fkt_types = {"krpaas", "kfortknox", "kfort_knox", "krpaasarchival"}
        def _is_fk_t(tt, tname=""):
            tl = (tt or "").lower(); nl = (tname or "").lower()
            return any(x in tl for x in _fkt_types) or "fortknox" in tl or "fortknox" in nl
        fk_name_set = set()
        for v in vaults_t:
            vt = v.get("vaultType") or ""; vn = v.get("name") or ""
            if _is_fk_t(vt, vn):
                fk_name_set.add(vn or "FortKnox")
        for p in policies_t:
            for t in ((p.get("remoteTargetPolicy") or {}).get("archivalTargets") or []):
                cfg   = t.get("archivalTargetConfig") or {}
                tt    = cfg.get("targetType") or t.get("targetType") or ""
                tname = t.get("targetName") or cfg.get("name") or cfg.get("vaultName") or ""
                if tname and _is_fk_t(tt, tname):
                    fk_name_set.add(tname)
        if fk_st_t == "none":
            fk_cell_text = "No"
        elif fk_st_t == "active":
            fk_cell_text = "\n".join(sorted(fk_name_set)) if fk_name_set else "Active"
        else:  # idle — configured but no successful FK archival in lookback window
            fk_cell_text = (("\n".join(sorted(fk_name_set)) + "\n(Configured — not sending data)")
                            if fk_name_set else "Configured — not sending data")

        row_t = t_topo.add_row().cells
        row_t[0].text = cd["name"]
        row_t[1].text = "\n".join(repl_names) if repl_names else "None"
        row_t[2].text = "\n".join(arch_names) if arch_names else "None"
        row_t[3].text = fk_cell_text
        # Colour-code the FK cell: green=active, yellow=idle, no fill=none
        if fk_st_t == "active":
            _cell_shade(row_t[3], "E2EFDA")   # light green
        elif fk_st_t == "idle":
            _cell_shade(row_t[3], "FFF2CC")   # amber/yellow
    _word_fit_widths(t_topo)

    # ── Topology diagram ─────────────────────────────────────────────────────
    # Primary: native Word shapes (editable). Fallback: matplotlib PNG.
    if not _render_topology_word_shapes(doc, all_data):
        png_bytes = _render_topology_png(all_data)
        if png_bytes:
            from docx.shared import Inches
            import io as _io
            doc.add_paragraph("")
            p_img = doc.add_paragraph()
            p_img.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run_img = p_img.add_run()
            run_img.add_picture(_io.BytesIO(png_bytes), width=Inches(6.5))
            cap = doc.add_paragraph(
                "Figure: Environment Topology (PNG fallback)."
            )
            cap.runs[0].font.size = Pt(8)
            cap.runs[0].font.color.rgb = RGBColor(0x80, 0x80, 0x80)
            cap.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_page_break()

    # ── 3. Protection & Recovery ───────────────────────────────────────────
    _h1("3. Protection & Recovery Health")
    total_g   = sum(len(cd["groups"]) for cd in all_data)
    paused_g  = sum(sum(1 for g in cd["groups"] if g.get("isPaused"))
                    for cd in all_data)
    failed_g  = sum(sum(1 for g in cd["groups"]
                        if (g.get("lastRun") or {}).get("localBackupInfo", {})
                           .get("status", "") in ("kFailure", "Failed"))
                    for cd in all_data)
    doc.add_paragraph(
        f"The environment contains {total_g} active protection group(s) across all clusters. "
        f"{paused_g} group(s) are currently paused and {failed_g} have a failed last run. "
        "The table below summarizes protection health per cluster. "
        "Refer to the Protection Health and Coverage Gaps tabs in the Excel workbook "
        "for per-group detail."
    )
    t3 = doc.add_table(rows=1, cols=6)
    t3.style = "Table Grid"; _word_autofit(t3)
    _tbl_header(t3, ["Cluster", "Groups", "Paused", "Failed Last Run",
                      "30d Success %", "SLA Score"])
    for cd in all_data:
        succ_pct, sla_pct = _success_stats(cd)
        paused = sum(1 for g in cd["groups"] if g.get("isPaused"))
        failed = sum(1 for g in cd["groups"]
                     if (g.get("lastRun") or {}).get("localBackupInfo", {})
                        .get("status", "") in ("kFailure", "Failed"))
        row = t3.add_row().cells
        for i, val in enumerate([cd["name"], len(cd["groups"]),
                                  paused, failed, f"{succ_pct}%",
                                  f"{cd['scores']['sla']}/100"]):
            row[i].text = str(val)
    _word_fit_widths(t3)
    doc.add_page_break()

    # ── 4. Storage & Capacity ──────────────────────────────────────────────
    _h1("4. Storage & Capacity")
    doc.add_paragraph(
        "Cohesity applies inline deduplication and compression to maximize storage "
        "efficiency. The data reduction ratio shows how much raw data has been saved "
        "compared to the logical data ingested."
    )
    t4 = doc.add_table(rows=1, cols=7)
    t4.style = "Table Grid"; _word_autofit(t4)
    _tbl_header(t4, ["Cluster", "Usable (TB)", "Used (TB)", "Free (TB)",
                      "Used %", "Data Reduction", "Status"])
    for cd in all_data:
        info   = cd["info"]
        usable, used, cap_pct = _cap_stats(cd)
        free   = usable - used
        dr     = round((info.get("stats") or {}).get("dataReductionRatio") or 0, 2)
        pct_str = f"{cap_pct}%" if isinstance(cap_pct, float) else "N/A"
        status = ("N/A" if cap_pct == "N/A"
                  else "CRITICAL" if cap_pct >= CAPACITY_CRIT_PCT
                  else "WARNING" if cap_pct >= CAPACITY_WARN_PCT else "OK")
        row = t4.add_row().cells
        for i, val in enumerate([cd["name"],
                                  round(bytes_to_tb(usable), 2) or "N/A",
                                  round(bytes_to_tb(used), 2) or "N/A",
                                  round(bytes_to_tb(free), 2) or "N/A",
                                  pct_str, f"{dr}x", status]):
            row[i].text = str(val)
    _word_fit_widths(t4)

    # Capacity runway alerts
    _runway_alerts = []
    for cd in all_data:
        rw   = cd.get("capacity_runway") or {}
        d80  = rw.get("days_to_80")
        dq   = rw.get("data_quality", "no_data")
        proj = rw.get("projected_80")
        gr   = rw.get("daily_growth_tb")
        if dq in ("no_data", "insufficient"):
            continue
        if d80 is None:
            continue
        if isinstance(d80, int) and d80 < 90:
            proj_str = proj.isoformat() if proj else "soon"
            gr_str   = f"{round(gr, 3)} TB/d" if gr is not None else "unknown growth rate"
            if d80 == 0:
                _runway_alerts.append(
                    f"{cd['name']}: already above 80% utilization — immediate capacity "
                    f"review required (growth: {gr_str})."
                )
            elif d80 < 30:
                _runway_alerts.append(
                    f"{cd['name']}: projected to reach 80% capacity in {d80} days "
                    f"(est. {proj_str}, growth: {gr_str}) — urgent action required."
                )
            else:
                _runway_alerts.append(
                    f"{cd['name']}: projected to reach 80% capacity in {d80} days "
                    f"(est. {proj_str}, growth: {gr_str})."
                )

    if _runway_alerts:
        _h2("Capacity Runway Forecast")
        doc.add_paragraph(
            "Based on 30-day usage trends, the following cluster(s) are forecast to "
            "reach 80% storage utilization within 90 days. Plan capacity expansion or "
            "data management actions accordingly."
        )
        for alert in _runway_alerts:
            p = doc.add_paragraph(style="List Bullet")
            p.add_run(alert)

    doc.add_page_break()

    # ── 5. Security ────────────────────────────────────────────────────────
    _h1("5. Security Posture")
    doc.add_paragraph(
        "The following checklist evaluates each cluster against Cohesity security "
        "best practices. A fully hardened deployment enables cluster-wide and "
        "storage-domain encryption, at least one vault for archival, replication "
        "to a secondary cluster, and FortKnox or DataLock archival for indelible "
        "backups. Note: Cohesity's distributed filesystem is always immutable by "
        "design; DataLock (WORM) makes individual snapshots indelible — permanently "
        "retained even against administrator deletion."
    )

    # Pre-table: Ransomware readiness summary callout
    _rw_scores = [cd["scores"].get("ransomware") for cd in all_data
                  if cd["scores"].get("ransomware") is not None]
    if _rw_scores:
        _rw_avg  = round(sum(_rw_scores) / len(_rw_scores), 1)
        _rw_lbl  = ("Strong" if _rw_avg >= 80 else
                    "Moderate" if _rw_avg >= 60 else
                    "High Risk" if _rw_avg >= 40 else "Critical")
        _rw_para = doc.add_paragraph()
        _rw_run  = _rw_para.add_run(
            f"Ransomware Readiness (avg across {len(_rw_scores)} cluster(s)): "
            f"{_rw_avg}/100 — {_rw_lbl}. "
        )
        _rw_run.font.bold = True
        if _rw_avg < 60:
            _rw_run.font.color.rgb = RGBColor(0xFF, 0x4C, 0x4C)
        elif _rw_avg < 80:
            _rw_run.font.color.rgb = RGBColor(0xFF, 0x8C, 0x00)
        _rw_para.add_run(
            "This score measures how well-prepared the environment is to recover "
            "from a ransomware event, based on DataLock coverage, FortKnox vault "
            "status, recovery window depth, quorum, and admin MFA. "
            "Full detail is in the Security sheet of the Excel workbook."
        )

    t5 = doc.add_table(rows=1, cols=11)
    t5.style = "Table Grid"; _word_autofit(t5)
    _tbl_header(t5, ["Cluster", "Cluster Enc", "SD Enc", "Vault",
                      "Indelible (FK)", "DataLock (WORM)",
                      "Replication", "Quorum", "Admins w/o MFA",
                      "Score /100", "Ransomware Score"])
    _admin_roles_w = {"COHESITY_ADMIN", "Admin", "kAdmin", "kSuperAdmin"}
    for cd in all_data:
        vaults   = cd["vaults"]
        policies = cd["policies"]
        cluster_enc_s, sd_enc_s = _sd_enc_status(cd)
        sd_enc_str_s = ("Yes" if sd_enc_s is True else
                        "No"  if sd_enc_s is False else "N/A")
        has_vault   = bool(vaults)
        fk_st_s     = cd.get("fk_status") or _fk_status(cd)
        fk_label_s  = {"active": "Active", "idle": "Idle", "none": "No"}[fk_st_s]
        repl = any((p.get("remoteTargetPolicy") or {}).get("replicationTargets")
                   for p in policies)
        dl_status_s, dl_label_s, _, _ = _cluster_datalock(policies, cd["groups"])
        if cd.get("mode") == "direct":
            quorum_str_s = "N/A"
        else:
            quorum_str_s = "Yes" if cd.get("quorum_enabled") else "No"
        admins_no_mfa_w = sum(
            1 for u in (cd.get("users") or [])
            if isinstance(u, dict)
            and not (u.get("mfaEnabled") or u.get("isMfaEnabled"))
            and bool(set(u.get("roles") or []) & _admin_roles_w)
        )
        rw_score_w = cd["scores"].get("ransomware", "")
        rw_label_w = cd["scores"].get("ransomware_label", "")
        rw_disp    = f"{rw_score_w}/100 ({rw_label_w})" if rw_score_w != "" else "N/A"

        row = t5.add_row().cells
        yn  = lambda v: "Yes" if v else "No"
        for i, val in enumerate([cd["name"],
                                  yn(cluster_enc_s), sd_enc_str_s,
                                  yn(has_vault), fk_label_s, dl_label_s,
                                  yn(repl), quorum_str_s,
                                  admins_no_mfa_w,
                                  f"{cd['scores']['security']}/100",
                                  rw_disp]):
            row[i].text = str(val)

        # ── Color-code every security column ──────────────────────────────
        def _sec_red(cell):
            _cell_shade(cell, "FF4C4C")
            for _r in cell.paragraphs[0].runs:
                _r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                _r.font.bold = True

        def _sec_green(cell):
            _cell_shade(cell, "E2EFDA")

        def _sec_amber(cell):
            _cell_shade(cell, "FFF2CC")
            for _r in cell.paragraphs[0].runs:
                _r.font.bold = True

        # Col 1 — Cluster Encryption
        (_sec_green if cluster_enc_s else _sec_red)(row[1])
        # Col 2 — SD Encryption
        if sd_enc_str_s == "Yes":   _sec_green(row[2])
        elif sd_enc_str_s == "No":  _sec_red(row[2])
        # Col 3 — Vault
        (_sec_green if has_vault else _sec_red)(row[3])
        # Col 4 — FortKnox (Indelible)
        if fk_st_s == "active":   _sec_green(row[4])
        elif fk_st_s == "idle":   _sec_amber(row[4])
        else:                     _sec_red(row[4])
        # Col 5 — DataLock (WORM)
        if dl_status_s in ("full_compliance", "full_any"):
            _sec_green(row[5])
        elif dl_status_s == "partial":
            _sec_amber(row[5])
        else:
            _sec_red(row[5])
        # Col 6 — Replication
        (_sec_green if repl else _sec_red)(row[6])
        # Col 7 — Quorum
        if quorum_str_s == "Yes":   _sec_green(row[7])
        elif quorum_str_s == "No":  _sec_red(row[7])
        # Col 8 — Admins w/o MFA
        if admins_no_mfa_w > 0:     _sec_red(row[8])
        else:                       _sec_green(row[8])
        # Col 9 — Security Score
        sec_val = cd["scores"]["security"]
        if sec_val >= 75:      _sec_green(row[9])
        elif sec_val >= 50:    _sec_amber(row[9])
        else:                  _sec_red(row[9])
        # Col 10 — Ransomware Score
        if rw_score_w != "":
            if rw_score_w >= 80:   _sec_green(row[10])
            elif rw_score_w >= 60: _sec_amber(row[10])
            else:                  _sec_red(row[10])
    _word_fit_widths(t5)

    # ── Governance & Change Activity ───────────────────────────────────────
    _h2("Governance & Change Activity")

    # Aggregate audit data across all clusters
    _gov_cat_counts  = {}
    _gov_high_risk   = []
    _gov_available   = False
    for cd in all_data:
        events = cd.get("audit_events") or []
        if events:
            _gov_available = True
        for ev in events:
            cat = _audit_category(ev)
            if cat is None:
                continue
            _gov_cat_counts[cat] = _gov_cat_counts.get(cat, 0) + 1
            if _audit_high_risk(ev):
                action = ev.get("action") or ev.get("entityType") or ""
                entity = ev.get("entityName") or ev.get("objectName") or ""
                _gov_high_risk.append(f"{cd['name']}: {action} — {entity}")

    if not _gov_available:
        doc.add_paragraph(
            "Audit log data is not available. This may occur when the account used "
            "to run the health check does not have audit read permissions, or when "
            "audit logging is disabled on the cluster."
        )
    else:
        total_changes = sum(_gov_cat_counts.values())
        if total_changes == 0:
            doc.add_paragraph(
                "No configuration change events were recorded in the audit log during "
                "the reporting period. Read-only access events are excluded."
            )
        else:
            gov_summary = (
                f"The audit log recorded {total_changes} configuration change event(s) "
                f"over the past {args.days} day(s) across all clusters. "
            )
            cat_parts = ", ".join(
                f"{cnt} {cat.lower()}" for cat, cnt in sorted(
                    _gov_cat_counts.items(), key=lambda x: -x[1])
            )
            if cat_parts:
                gov_summary += f"By category: {cat_parts}. "
            doc.add_paragraph(gov_summary)

        if _gov_high_risk:
            _hr_para = doc.add_paragraph()
            _hr_run  = _hr_para.add_run(
                f"Warning: {len(_gov_high_risk)} high-risk event(s) detected "
                f"(vault deletion, user deletion, admin privilege grant, or encryption "
                f"configuration change). Review these events immediately: "
            )
            _hr_run.font.bold  = True
            _hr_run.font.color.rgb = RGBColor(0xFF, 0x8C, 0x00)
            for hr_item in _gov_high_risk[:5]:
                p = doc.add_paragraph(style="List Bullet")
                p.add_run(hr_item)

    doc.add_page_break()

    # ── 6. Agent Health & Source Coverage ────────────────────────────────────
    _h1("6. Agent Health & Source Coverage")

    total_agents    = sum(len(cd.get("agents") or []) for cd in all_data)
    def _ag_st(a):
        _hs = a.get("healthStatus") if isinstance(a, dict) else None
        return ((isinstance(_hs, dict) and _hs.get("status"))
                or (isinstance(_hs, str) and _hs)
                or (isinstance(a, dict) and (a.get("agentStatus") or a.get("status")))
                or "")
    unhealthy_ag    = sum(
        sum(1 for a in (cd.get("agents") or [])
            if _ag_st(a).lower()
            in ("kunhealthy", "unhealthy", "unreachable", "kfailed", "failed"))
        for cd in all_data)
    upgradable_ag   = sum(
        sum(1 for a in (cd.get("agents") or [])
            if "upgradable" in (a.get("upgradability") or "").lower())
        for cd in all_data)

    doc.add_paragraph(
        f"The environment has {total_agents} registered host agent(s) across all clusters. "
        f"{unhealthy_ag} agent(s) are unhealthy or unreachable. "
        f"{upgradable_ag} agent(s) are eligible for upgrade. "
        "Unhealthy agents should be remediated promptly as they will prevent backups "
        "from running. Agents running older versions MAY cause issues with newer cluster "
        "features; Cohesity supports a limited range of back-versions — review the agent "
        "compatibility matrix at docs.cohesity.com before upgrading. "
        "Refer to the Agent Health sheet in the Excel workbook for per-host detail."
    )

    total_sources = sum(len(cd.get("sources") or []) for cd in all_data)
    low_cov_src   = sum(
        sum(1 for s in (cd.get("sources") or [])
            if (lambda p, u: p + u > 0 and (u / (p + u)) >= 0.25)(
                (s.get("protectedObjectsCount") or s.get("numProtectedObjects")
                 or (s.get("stats") or {}).get("protectedObjectsCount") or 0),
                (s.get("unprotectedObjectsCount") or s.get("numUnprotectedObjects")
                 or (s.get("stats") or {}).get("unprotectedObjectsCount") or 0)))
        for cd in all_data)
    doc.add_paragraph(
        f"There are {total_sources} registered source(s). "
        f"{low_cov_src} source(s) have ≥25% unprotected objects — "
        "these objects will be unrecoverable in the event of data loss. "
        "Review the Source Coverage sheet in the Excel workbook for per-source breakdown."
    )
    doc.add_page_break()

    # ── 7. User Security ───────────────────────────────────────────────────────
    _h1("7. User Security")

    total_users  = sum(len(cd.get("users") or []) for cd in all_data)
    _ADMIN_SET   = {"cohesity_admin", "admin", "kadmin", "ksuperadmin", "super_admin"}
    admins_no_mfa = sum(
        sum(1 for u in (cd.get("users") or [])
            if not (u.get("mfaEnabled") or u.get("isMfaEnabled"))
            and bool({r.lower() for r in (u.get("roles") or [])} & _ADMIN_SET))
        for cd in all_data)
    locked_users = sum(
        sum(1 for u in (cd.get("users") or [])
            if u.get("locked") or u.get("isLocked"))
        for cd in all_data)
    has_sso_any  = any(cd.get("idps") for cd in all_data)

    doc.add_paragraph(
        f"The environment has {total_users} user account(s) across all clusters. "
        f"{admins_no_mfa} admin account(s) do not have MFA enabled — these are at "
        "elevated risk from credential theft. "
        f"{locked_users} account(s) are currently locked. "
        f"SSO / IDP: {'Configured' if has_sso_any else 'Not configured — local credentials only'}. "
        "Refer to the User Security sheet in the Excel workbook for per-user detail."
    )

    t_usr = doc.add_table(rows=1, cols=5)
    t_usr.style = "Table Grid"; _word_autofit(t_usr)
    _tbl_header(t_usr, ["Cluster", "Total Users", "Admins w/o MFA",
                         "Locked Accounts", "SSO Configured"])
    for cd in all_data:
        _A = {"cohesity_admin", "admin", "kadmin", "ksuperadmin", "super_admin"}
        ano_mfa = sum(1 for u in (cd.get("users") or [])
                      if not (u.get("mfaEnabled") or u.get("isMfaEnabled"))
                      and bool({r.lower() for r in (u.get("roles") or [])} & _A))
        locked  = sum(1 for u in (cd.get("users") or [])
                      if u.get("locked") or u.get("isLocked"))
        row = t_usr.add_row().cells
        for i, val in enumerate([cd["name"], len(cd.get("users") or []),
                                  ano_mfa, locked,
                                  "Yes" if cd.get("idps") else "No"]):
            row[i].text = str(val)
    _word_fit_widths(t_usr)
    doc.add_page_break()

    # ── 8. Recommendations ─────────────────────────────────────────────────
    _h1("8. Prioritized Recommendations")
    doc.add_paragraph(
        "The recommendations below are ranked by priority. Critical items should "
        "be addressed immediately; High-priority items within 30 days; Medium within "
        "90 days. Refer to the Recommendations tab in the Excel workbook for the "
        "full prioritized list."
    )
    sorted_recs = sorted(all_recs, key=lambda x: _PRIORITY_ORDER.get(x[1]["priority"], 99))
    if sorted_recs:
        t6 = doc.add_table(rows=1, cols=5)
        t6.style = "Table Grid"; _word_autofit(t6)
        _tbl_header(t6, ["Priority", "Cluster", "Category", "Finding", "Action"])
        for cname, rec in sorted_recs:
            row = t6.add_row().cells
            for i, val in enumerate([rec["priority"], cname, rec["category"],
                                      rec["finding"], rec["action"]]):
                row[i].text = str(val)
                if i == 0:
                    for para in row[i].paragraphs:
                        for run in para.runs:
                            run.font.bold = True
        _word_fit_widths(t6)
    else:
        doc.add_paragraph("No recommendations generated — environment appears healthy.")

    doc.add_page_break()

    # ── Appendix ───────────────────────────────────────────────────────────
    _h1("Appendix — Methodology & Thresholds")
    _h2("Scoring Model")
    doc.add_paragraph(
        "Health scores (0-100) are computed from five weighted dimensions: "
        "Protection success rate (25%), SLA compliance (20%), "
        "Infrastructure health (15%), Storage capacity (15%), "
        "Security posture (15%), Alert health (10%). "
        "Within the Security posture dimension, the 100-point security score is "
        "calculated across six areas: cluster encryption (15 pts), storage-domain "
        "encryption (15 pts), vault configured (15 pts), replication (15 pts), "
        "FortKnox/indelible copy (20 pts), and DataLock/WORM (20 pts). "
        "A penalty of −10 pts is applied when any administrator account lacks MFA. "
        "DataLock coverage is measured at the protection-group level — only groups "
        "whose governing policy has DataLock (WORM) explicitly configured count; "
        "FortKnox archival (indelible by design) scores separately in the FK dimension."
    )
    _h2("Thresholds")
    notes = [
        f"Capacity: Warning ≥{CAPACITY_WARN_PCT}%, Critical ≥{CAPACITY_CRIT_PCT}%",
        f"Backup success rate: Warning <{SUCCESS_WARN_PCT}%, Critical <{SUCCESS_CRIT_PCT}% "
        f"— runs with status 'Warning' (kWarning) are counted as successful in the "
        f"Success % calculation. They complete with data but may have minor issues "
        f"(e.g. some objects skipped) and are shown separately in the Warning column "
        f"for visibility.",
        f"RPO: Warning >{RPO_WARN_HOURS}h, Critical >{RPO_CRIT_HOURS}h since last success",
        ("Security: cluster enc (15) + SD enc (15) + vault (15) + replication (15) "
         "+ FortKnox/indelible (20) + DataLock (20) = 100 pts. "
         "Partial DataLock coverage earns proportional credit."),
        f"Alert penalty: −10 pts per open critical, −2 pts per open warning",
        f"View quota: Warning ≥{VIEW_QUOTA_WARN_PCT}%, Critical ≥{VIEW_QUOTA_CRIT_PCT}%",
        f"Lookback period used for this report: {args.days} days",
        f"Report generated: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')} (local time)",
    ]
    for note in notes:
        p = doc.add_paragraph(note, style="List Bullet")
        for run in p.runs:
            run.font.size = Pt(9)

    out = f"{args.output}.docx"
    doc.save(out)
    print(f"  Word   → {out}")
    return out


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────

def main():
    _t0 = time.time()
    parser = argparse.ArgumentParser(
        description=f"Cohesity Health Check Report v{__version__}",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    # ── Helios API key auth ───────────────────────────────────────────────────
    parser.add_argument("--apikey",           help="Helios API key (stored in keychain if omitted)")
    # ── Direct cluster username/password auth (requires --cluster-host) ───────
    parser.add_argument("--username",         help="Cluster username (use with --cluster-host)")
    parser.add_argument("--password",         help="Password (prompted if omitted; saved to keychain)")
    parser.add_argument("--domain",           default="LOCAL",
                        help="Auth domain for cluster login: LOCAL or AD name (default: LOCAL)")
    parser.add_argument("--mfa-code",         dest="mfa_code",
                        help="TOTP/OTP code for MFA-enabled accounts")
    # ── Direct cluster mode ───────────────────────────────────────────────────
    parser.add_argument("--cluster-host",     dest="cluster_host",
                        help="Bypass Helios — connect directly to this cluster hostname/IP")
    # ── Credential management ─────────────────────────────────────────────────
    parser.add_argument("--clear-credentials", action="store_true",
                        help="Remove stored credentials from keychain and exit")
    # ── Helios cluster filter ─────────────────────────────────────────────────
    parser.add_argument("--cluster",          help="(Helios) Limit to one cluster (partial name match)")
    # ── Report options ────────────────────────────────────────────────────────
    parser.add_argument("--days",             type=int, default=30,
                        help="Lookback window in days (default 30)")
    parser.add_argument("--customer",         default="",
                        help="Customer name for report cover page")
    parser.add_argument("--output",           default="cohesity_health_check",
                        help="Output base path without extension "
                             "(default: cohesity_health_check)")
    parser.add_argument("--quick",            action="store_true",
                        help="Skip per-group run history; use last-run only (faster)")
    parser.add_argument("--excel-only",       action="store_true",
                        help="Generate Excel only, skip Word document")
    parser.add_argument("--word-only",        action="store_true",
                        help="Generate Word only, skip Excel")
    parser.add_argument("--debug",            action="store_true",
                        help="Print HTTP status for each API call")
    # ── TLS options ───────────────────────────────────────────────────────────
    parser.add_argument("--ca-bundle",        dest="ca_bundle", default=None,
                        metavar="PATH",
                        help="Path to CA bundle (.pem) for TLS verification "
                             "(use for self-signed or corporate proxy certs)")
    parser.add_argument("--insecure",         action="store_true",
                        help="Disable TLS certificate validation (NOT recommended)")
    args = parser.parse_args()

    # ── Input validation ───────────────────────────────────────────────────────
    # Credentials passed on the command line are visible to every local user
    # via /proc/<pid>/cmdline (Linux) or process-monitoring tools elsewhere.
    # Warn — but don't block — so automated runners can still use them.
    if args.password:
        print("WARNING: --password on the command line is visible to other users "
              "via the process list. Prefer the interactive prompt or system keychain.")
    if args.mfa_code:
        print("WARNING: --mfa-code on the command line is visible to other users "
              "via the process list.")

    if not (1 <= args.days <= 3650):
        parser.error("--days must be between 1 and 3650")

    if args.cluster_host:
        _valid_host = re.compile(
            r'^(?:[a-zA-Z0-9](?:[a-zA-Z0-9\-]{0,61}[a-zA-Z0-9])?\.)*'
            r'[a-zA-Z0-9](?:[a-zA-Z0-9\-]{0,61}[a-zA-Z0-9])?$'
            r'|^(?:\d{1,3}\.){3}\d{1,3}$'
        )
        if not _valid_host.match(args.cluster_host):
            parser.error(f"--cluster-host: '{args.cluster_host}' is not a valid "
                         "hostname or IPv4 address")

    if ".." in os.path.normpath(args.output).split(os.sep):
        parser.error("--output: path must not contain '..' directory traversal sequences")

    if args.ca_bundle and not os.path.isfile(args.ca_bundle):
        parser.error(f"--ca-bundle: file not found: {args.ca_bundle}")

    # ── Clear stored credentials and exit ─────────────────────────────────────
    if args.clear_credentials:
        clear_stored_credentials(
            cluster=args.cluster_host,
            username=args.username or "admin",
            domain=args.domain,
        )
        sys.exit(0)

    # ── Pre-requisite check ─────────────────────────────────────────────────
    # Build a full status table covering all packages (required + optional).
    # Always printed so the user can see at a glance what is and isn't available.
    def _pkg_ok(name):
        try:
            __import__(name)
            return True
        except ImportError:
            return False

    # requests is guaranteed present here (checked at module import level).
    # openpyxl / python-docx flags come from the module-level try/except.
    _PKG_TABLE = [
        # (name,          installed,    role,       description)
        ("requests",      True,         "required", "HTTP client for Cohesity REST API"),
        ("openpyxl",      EXCEL_OK,     "required", "Excel workbook output (skip with --word-only)"),
        ("python-docx",   DOCX_OK,      "required", "Word document generation (skip with --excel-only)"),
        ("python-pptx",   PPTX_OK,      "required", "PowerPoint deck output (~31 slides, always generated)"),
        ("matplotlib",    _pkg_ok("matplotlib"), "optional", "Topology PNG fallback for Word diagram"),
        ("keyring",       _pkg_ok("keyring"),    "optional", "Stores API key in OS secure vault (macOS Keychain / Windows Credential Manager) — prevents credentials appearing in shell history and process lists"),
    ]

    # Apply mode exemptions: some required packages are not needed in certain modes
    _exempt = set()
    if args.word_only:   _exempt.add("openpyxl")
    if args.excel_only:  _exempt.add("python-docx")

    _need_required = [n for n, ok, r, d in _PKG_TABLE
                      if r == "required" and not ok and n not in _exempt]
    _missing_opt   = [n for n, ok, r, d in _PKG_TABLE if r == "optional" and not ok]
    _all_missing   = _need_required + _missing_opt

    # Always print the full status table
    print()
    print("  Prerequisites")
    print("  " + "─" * 68)
    for name, ok, role, desc in _PKG_TABLE:
        if name in _exempt:
            mark   = "–"
            status = "skipped  "
            note   = f"  (not needed with {'--word-only' if name == 'openpyxl' else '--excel-only'})"
        else:
            mark   = "✓" if ok else "✗"
            status = "installed" if ok else "NOT FOUND"
            note   = ""
        role_tag = f"[{role}]"
        print(f"  {mark}  {name:16s}  {status}  {role_tag:11s}  {desc}{note}")
    print("  " + "─" * 68)

    if _all_missing:
        print()
        if _need_required:
            print("  ERROR: Missing required package(s) — script cannot continue.")
        else:
            print("  Optional package(s) not installed — affected features will be skipped.")
        print()
        print("  Install everything missing with one command:")
        print(f"    pip install {' '.join(_all_missing)}")
        print()
        if _need_required:
            sys.exit(1)
    else:
        print(f"  All {len(_PKG_TABLE)} packages OK")

    if not _pkg_ok("keyring"):
        print()
        print("  SECURITY NOTE: keyring is not installed. API keys passed via --apikey")
        print("  will appear in shell history and process lists (ps aux). Installing")
        print("  keyring stores credentials in the OS secure vault instead.")
        print("  Not applicable on headless Linux servers (no D-Bus/Secret Service).")
    print()

    # Build output filename: include customer name (if given) + local timestamp
    if args.output == "cohesity_health_check":
        ts    = datetime.datetime.now().strftime("%Y%m%d_%H%M")
        parts = ["cohesity_health_check", f"v{__version__}"]
        if args.customer:
            safe = "".join(c if c.isalnum() or c in "-_" else "_"
                           for c in args.customer).strip("_")
            parts.append(safe)
        parts.append(ts)
        args.output = "_".join(parts)

    # ── TLS verification ──────────────────────────────────────────────────────
    if args.insecure:
        _verify = False
        print("=" * 65)
        print("  WARN: --insecure — TLS certificate validation DISABLED.")
        print("        Credentials and tokens may be intercepted (MITM).")
        print("=" * 65)
        try:
            import urllib3 as _u3
            _u3.disable_warnings(_u3.exceptions.InsecureRequestWarning)
        except ImportError:
            requests.packages.urllib3.disable_warnings()
    elif args.ca_bundle:
        _verify = args.ca_bundle
    else:
        _verify = True

    # ── Authenticate and discover clusters ────────────────────────────────────
    session = requests.Session()
    session.verify = _verify
    api_key = None

    if args.cluster_host:
        # ── Direct cluster mode ───────────────────────────────────────────────
        if not args.username:
            print("ERROR: --username is required with --cluster-host.")
            sys.exit(1)
        pwd   = get_cluster_password(args.cluster_host, args.username, args.domain,
                                     cli_password=args.password)
        token = get_auth_token(args.cluster_host, args.username, pwd, args.domain,
                               mfa_code=args.mfa_code, verify=_verify)
        session.base_url = f"https://{args.cluster_host}"
        # Fetch cluster name from the cluster itself
        try:
            _r = session.get(
                f"https://{args.cluster_host}/irisservices/api/v1/public/cluster",
                headers=make_headers(token), timeout=30)
            _r.raise_for_status()
            cname = _r.json().get("name", args.cluster_host)
        except Exception:
            cname = args.cluster_host
        clusters = [{"name": cname, "clusterId": None,
                     "mode": "direct", "host": args.cluster_host, "token": token}]
        auth_desc = f"Direct ({args.cluster_host})"

    elif args.username and not args.cluster_host:
        # --username without --cluster-host makes no sense: Helios does not
        # expose a scriptable username/password login endpoint (it uses
        # OIDC/OAuth2 web-based auth). Direct the user to create an API key.
        print(
            "ERROR: --username is only supported with --cluster-host.\n"
            "\n"
            "  Helios does not provide a scriptable username/password login\n"
            "  endpoint. Its web login uses OIDC/OAuth2 and cannot be driven\n"
            "  from a script.\n"
            "\n"
            "  To use Helios, generate an API key and pass it with --apikey:\n"
            "    helios.cohesity.com → Settings → Access Management → API Keys\n"
            "\n"
            "  To run against a single cluster directly (no Helios required):\n"
            "    --cluster-host <ip/hostname> --username <user>"
        )
        sys.exit(1)

    else:
        # ── Helios API key (default) ──────────────────────────────────────────
        api_key = get_api_key(args.apikey)
        session.base_url = f"https://{HELIOS_HOST}"
        clusters = get_helios_clusters(api_key, verify=_verify)
        if not clusters:
            print("ERROR: No clusters returned from Helios.")
            sys.exit(1)
        auth_desc = "Helios API key"
        # Quorum is a Helios control-plane feature configured in Security Center →
        # Security Posture → Quorum. When any enabled quorum group exists, the
        # quorum requirement applies to all clusters managed by that Helios tenant.
        # We therefore set quorum_enabled=True for every cluster when at least one
        # enabled group is returned by the API (rather than matching per clusterId,
        # which is unreliable when clusterIdentifiers is unpopulated).
        _quorum_enabled = False
        try:
            for _qg in _helios_quorum_groups(session, api_key, args.debug):
                if _qg.get("isEnabled"):
                    _quorum_enabled = True
                    break
        except Exception:
            pass
        for _c in clusters:
            _c["quorum_enabled"] = _quorum_enabled

    # Apply --cluster name filter (Helios modes only)
    if args.cluster and not args.cluster_host:
        clusters = [c for c in clusters
                    if args.cluster.lower() in c.get("name", "").lower()]
        if not clusters:
            print(f"ERROR: No cluster matching '{args.cluster}'")
            sys.exit(1)

    print(f"\nCohesity Health Check Report  v{__version__}")
    print(f"  Auth     : {auth_desc}")
    print(f"  Customer : {args.customer or '(not specified)'}")
    print(f"  Lookback : {args.days} days")
    print(f"  Mode     : {'Quick (last-run only)' if args.quick else 'Full (run history)'}")
    print(f"  Clusters : {', '.join(c.get('name', '?') for c in clusters)}")
    print()

    all_data = []
    for cluster in clusters:
        try:
            cd = collect_cluster(session, api_key, cluster, args)
            s  = cd["scores"]
            print(f"  [{cd['name']}] score={s['overall']}/100 ({s['grade']}) "
                  f"| {len(cd['recommendations'])} recommendation(s)")
            all_data.append(cd)
        except Exception as exc:
            import traceback
            print(f"  ERROR collecting [{cluster.get('name','?')}]: {exc}")
            if args.debug:
                traceback.print_exc()
            continue

    if not all_data:
        print("ERROR: No data collected from any cluster.")
        sys.exit(1)

    print()
    if not args.word_only:
        write_excel(all_data, args)
    if not args.excel_only:
        write_word(all_data, args)

    pptx_path = write_pptx(all_data, args, args.output)

    out_parts = [f"{args.output}.xlsx", f"{args.output}.docx"]
    if pptx_path:
        out_parts.append(pptx_path)
    elapsed = time.time() - _t0
    print(f"\nComplete.  Output: {' / '.join(out_parts)}")
    print(f"  Runtime  : {elapsed:.1f}s ({elapsed/60:.1f} min)")


if __name__ == "__main__":
    main()
